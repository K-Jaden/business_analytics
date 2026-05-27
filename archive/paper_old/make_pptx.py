# -*- coding: utf-8 -*-
"""
논문 발표용 PPT 생성 스크립트.
출력: paper/presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

import sys
OUT = r"C:\Users\이승현\Finance\paper\presentation.pptx"
# 파일이 잠겨 있으면 _v2 로 자동 우회
def _resolve_out(p):
    try:
        if not os.path.exists(p):
            return p
        with open(p, "ab") as _f:
            pass
        return p
    except PermissionError:
        base, ext = os.path.splitext(p)
        return base + "_v2" + ext
OUT = _resolve_out(OUT)
FIG_DIR = r"C:\Users\이승현\Finance\results\figures"

# ---------- 색상 팔레트 ----------
NAVY   = RGBColor(0x0B, 0x2E, 0x4F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)   # 강조 빨강
SUB    = RGBColor(0x3B, 0x5B, 0x7E)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x1E, 0x82, 0x4C)

# ---------- 슬라이드 크기 16:9 ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]  # blank


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="맑은 고딕"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=NAVY,
                bullet="• ", spacing=4, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        # 줄 안에 *볼드* 표시: '**...**' 구문 지원
        text = bullet + it
        chunks = []
        rem = text
        while "**" in rem:
            before, _, rest = rem.partition("**")
            bold_part, _, after = rest.partition("**")
            chunks.append((before, False))
            chunks.append((bold_part, True))
            rem = after
        chunks.append((rem, False))
        for s, b in chunks:
            if not s:
                continue
            r = p.add_run()
            r.text = s
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.color.rgb = color
    return tb


def slide_header(slide, title, subtitle=None):
    # 상단 띠
    add_rect(slide, 0, 0, SW, Inches(0.85), NAVY)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.6),
             title, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.4),
                 subtitle, size=14, color=SUB)


def add_picture_fit(slide, path, x, y, max_w, max_h):
    """원본 비율을 유지하며 영역 안에 중앙 정렬."""
    from PIL import Image as _Image
    im = _Image.open(path)
    iw, ih = im.size
    sw = max_w / iw
    sh = max_h / ih
    sc = min(sw, sh)
    w = int(iw * sc)
    h = int(ih * sc)
    cx = int(x + (max_w - w) / 2)
    cy = int(y + (max_h - h) / 2)
    return slide.shapes.add_picture(path, cx, cy, w, h)


def add_footer(slide, page, total):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.35),
             "리셀 시장 미디어 선행 지표 연구", size=10, color=GRAY)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35),
             f"{page} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ====================================================================
# 슬라이드 1 — 표지
# ====================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
# 액센트 띠
add_rect(s, 0, Inches(3.05), SW, Inches(0.06), ACCENT)

add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "Empirical Study", size=16, color=RGBColor(0xC9, 0xD6, 0xE3))
add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.4),
         "미디어 채널은 리셀 시장 가격을\n선행하는가?",
         size=40, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.35), Inches(12), Inches(0.7),
         "자산 유형별 Granger 인과성 분석과 머신러닝 검증",
         size=22, color=RGBColor(0xE0, 0xE8, 0xF0))

add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
         "스니커즈 · 트레이딩 카드 · 레고  |  2022.01 – 2025.12  |  N=48 months × 15 items",
         size=14, color=RGBColor(0xB7, 0xC4, 0xD2))
add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         "발표자: [저자명]    |    2026.05.19",
         size=14, color=RGBColor(0xB7, 0xC4, 0xD2))


# ====================================================================
# 슬라이드 2 — 연구 동기
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "연구 동기")

add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
         "리셀 시장은 정보 비대칭과 ‘관심(attention)’에 민감한 대안 자산",
         size=18, bold=True, color=ACCENT)

add_bullets(s, Inches(0.6), Inches(2.1), Inches(12.0), Inches(2.0), [
    "스니커즈(StockX), PSA 카드, 단종 LEGO — 표준화된 디지털 거래 인프라로 빠른 성장",
    "주식 시장 대비 **거래량은 작지만**, 정보·관심에 의한 가격 결정 메커니즘은 유사",
    "기존 금융 연구: 검색·뉴스·SNS 무드가 자산 가격을 선행 (Da et al. 2011; Tetlock 2007; Bollen et al. 2011)",
], size=16)

add_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0), LIGHT)
add_text(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.4),
         "❓ 연구 공백", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.3), [
    "리셀 시장에 동일 직관이 작동하는지에 대한 **실증 연구가 거의 부재**",
    "자산 유형별로 **어떤 미디어 채널이 선행하는지** 비교 가능한 설계가 없음",
], size=15)
add_footer(s, 2, 19)


# ====================================================================
# 슬라이드 3 — 연구 문제 (RQ1·RQ2·RQ3)
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "연구 문제")

bx_y = Inches(1.5)
bx_h = Inches(5.2)
bx_w = Inches(4.05)
gap  = Inches(0.2)
left = Inches(0.55)

titles = ["RQ1", "RQ2", "RQ3"]
heads = [
    "채널의 선행성",
    "채널 구성의 이질성",
    "예측력 등가성",
]
descs = [
    "다섯 가지 미디어 채널은\n각 리셀 자산의 가격을\nGranger 의미에서 선행하는가?",
    "유의한 채널 구성이\n자산 유형(스니커즈·카드·레고)\n마다 다른가?",
    "Granger로 선별한 채널만으로\n학습한 모형(B)이 전 채널 모형(A)과\n동등 또는 더 우수한가?",
]
methods = [
    "자산별 단독 Granger\n3 자산 × 5 채널 = 15회",
    "유의 채널 집합 간\nJaccard 유사도",
    "XGBoost + SHAP\nDiebold–Mariano 검정",
]
colors_top = [NAVY, SUB, ACCENT]

for i in range(3):
    x = left + (bx_w + gap) * i
    add_rect(s, x, bx_y, bx_w, Inches(0.7), colors_top[i])
    add_text(s, x, bx_y, bx_w, Inches(0.7), titles[i],
             size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, bx_y + Inches(0.7), bx_w, bx_h - Inches(0.7), LIGHT)

    add_text(s, x + Inches(0.2), bx_y + Inches(0.9), bx_w - Inches(0.4), Inches(0.5),
             heads[i], size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), bx_y + Inches(1.55), bx_w - Inches(0.4), Inches(2.0),
             descs[i], size=14, color=GRAY)

    add_rect(s, x + Inches(0.2), bx_y + Inches(3.7), bx_w - Inches(0.4), Inches(0.05), NAVY)
    add_text(s, x + Inches(0.2), bx_y + Inches(3.85), bx_w - Inches(0.4), Inches(0.4),
             "검증 방법", size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), bx_y + Inches(4.25), bx_w - Inches(0.4), Inches(1.0),
             methods[i], size=14, color=ACCENT, bold=True)

add_footer(s, 3, 19)


# ====================================================================
# 슬라이드 4 — 데이터 패널
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "데이터 패널", "2022.01 – 2025.12 (48개월) | 3 자산군 × 5 아이템 = 15 아이템")

# 좌측: 자산별 아이템 표
def asset_block(x, y, w, h, color, title, src, items):
    add_rect(s, x, y, w, Inches(0.6), color)
    add_text(s, x, y, w, Inches(0.6), title,
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, y + Inches(0.6), w, h - Inches(0.6), LIGHT)
    add_text(s, x + Inches(0.15), y + Inches(0.65), w, Inches(0.35),
             src, size=11, color=GRAY)
    add_bullets(s, x + Inches(0.15), y + Inches(1.05),
                w - Inches(0.3), h - Inches(1.2),
                items, size=12, color=NAVY, bullet="– ", spacing=2)

y0 = Inches(1.55)
hh = Inches(3.4)
ww = Inches(4.1)
gp = Inches(0.15)
x0 = Inches(0.5)

asset_block(x0,                     y0, ww, hh, NAVY,
            "Sneakers (StockX)", "US Size 10, 월별 평균 가격",
            ["Jordan 1 Bordeaux",
             "Nike Dunk Low Panda",
             "Yeezy 350 V2 Zebra",
             "Travis Scott Jordan 1",
             "New Balance 550 White/Green"])
asset_block(x0 + ww + gp,           y0, ww, hh, SUB,
            "Cards (PriceCharting)", "PSA 10, 월별 평균 가격",
            ["Charizard VMAX SV107",
             "Charizard GX Hidden Fates",
             "Umbreon VMAX Alt Art",
             "Rayquaza VMAX Alt Art",
             "Pikachu VMAX"])
asset_block(x0 + (ww + gp) * 2,     y0, ww, hh, ACCENT,
            "LEGO (BrickRanker)", "New Sealed, 월별 가격",
            ["Millennium Falcon 75192",
             "Hogwarts Castle 71043",
             "Titanic 10294",
             "Porsche 911 42096",
             "Bugatti Chiron 42083"])

# 하단: 핵심 수치
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), LIGHT)
add_text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.4),
         "패널 규모 및 한계", size=15, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.0), [
    "총 관측치: **15 아이템 × 48 개월 = 720 아이템-월**",
    "tx_count<3 월은 선형보간 / 연속 3개월 이상 결측은 한계 명시 (예: yeezy 4개월 보간)",
    "LEGO porsche·bugatti: CH1 zero값 비율 다소 높음 → 한계 절 기재",
], size=13)
add_footer(s, 4, 19)


# ====================================================================
# 슬라이드 5 — 5개 미디어 채널
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "미디어 채널 (5개)", "원시 데이터 → 점수 변환 후 분석 투입")

cols = [
    ("CH1", "Google Trends", "pytrends 월평균", "0 – 100", NAVY),
    ("CH2", "뉴스 보도량",   "GDELT timelinevolnorm 월평균", "0 – 1", SUB),
    ("CH3", "뉴스 감성",     "GDELT timelinetone ÷ 10", "-1 ~ +1", ACCENT),
    ("CH4", "YouTube 조회수","월별 합산 → log1p → z-score", "z-score", GREEN),
    ("CH5", "YouTube 댓글 감성", "FinBERT  P(pos) − P(neg)", "-1 ~ +1", GRAY),
]

cw = Inches(2.45)
ch = Inches(3.6)
cy = Inches(1.55)
cx0 = Inches(0.45)
cgap = Inches(0.12)

for i, (cid, name, conv, rng, col) in enumerate(cols):
    x = cx0 + (cw + cgap) * i
    # 헤더
    add_rect(s, x, cy, cw, Inches(0.8), col)
    add_text(s, x, cy, cw, Inches(0.45), cid,
             size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, cy + Inches(0.4), cw, Inches(0.4), name,
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 본문
    add_rect(s, x, cy + Inches(0.8), cw, ch - Inches(0.8), LIGHT)
    add_text(s, x + Inches(0.15), cy + Inches(1.0), cw - Inches(0.3), Inches(0.35),
             "변환 방식", size=11, bold=True, color=GRAY)
    add_text(s, x + Inches(0.15), cy + Inches(1.35), cw - Inches(0.3), Inches(1.6),
             conv, size=13, color=NAVY)
    add_text(s, x + Inches(0.15), cy + Inches(2.7), cw - Inches(0.3), Inches(0.35),
             "범위", size=11, bold=True, color=GRAY)
    add_text(s, x + Inches(0.15), cy + Inches(3.05), cw - Inches(0.3), Inches(0.4),
             rng, size=14, color=ACCENT, bold=True)

# 한계 박스
add_rect(s, Inches(0.45), Inches(5.4), Inches(12.45), Inches(1.4),
         RGBColor(0xFD, 0xF3, 0xE6))
add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
         "⚠ CH3 한계", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.9),
         "당초 NewsAPI + FinBERT 설계였으나, 무료 플랜의 30일 역사 자료 제한으로 GDELT timelinetone 대체.\n"
         "GDELT 톤은 일반 어휘 자동 산출 점수 → 도메인 특화도 낮음 (논문 한계 명시).",
         size=12, color=GRAY)
add_footer(s, 5, 19)


# ====================================================================
# 슬라이드 6 — 연구 설계 (분석 파이프라인)
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "분석 파이프라인")

steps = [
    ("01", "수집",       "StockX · PriceCharting\nBrickRanker · pytrends\nGDELT · YouTube API"),
    ("02", "점수 변환",   "CH1~CH5 표준화\nFinBERT(영어 댓글)"),
    ("03", "패널 통합",   "panel_monthly.csv\nasset_series.csv"),
    ("04", "ADF",        "정상성 검정\n비정상 → 1차 차분"),
    ("05", "Granger",    "15회 단독 검정\nLag = BIC 선택"),
    ("06", "Jaccard",    "자산쌍 채널 집합\n유사도"),
    ("07", "XGBoost+SHAP", "Model A vs B\nTimeSeriesSplit"),
    ("08", "DM 검정",     "RMSE 등가성\n(T=24/아이템)"),
]

# 8단계를 2줄로 배치
n = len(steps)
per_row = 4
bw = Inches(2.95)
bh = Inches(2.05)
gx = Inches(0.15)
gy = Inches(0.35)
ox = Inches(0.45)
oy = Inches(1.55)

for i, (no, name, desc) in enumerate(steps):
    r = i // per_row
    c = i % per_row
    x = ox + (bw + gx) * c
    y = oy + (bh + gy) * r
    # 배경
    add_rect(s, x, y, bw, bh, LIGHT)
    # 번호 원
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              x + Inches(0.15), y + Inches(0.15),
                              Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT if r == 0 else NAVY
    circ.line.fill.background()
    tf = circ.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rn = p.add_run(); rn.text = no
    rn.font.size = Pt(14); rn.font.bold = True
    rn.font.color.rgb = WHITE; rn.font.name = "맑은 고딕"

    add_text(s, x + Inches(0.85), y + Inches(0.2),
             bw - Inches(0.95), Inches(0.5),
             name, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(0.85),
             bw - Inches(0.3), bh - Inches(0.95),
             desc, size=12, color=GRAY)
    # 화살표 (행 끝 제외)
    if c < per_row - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + bw, y + bh / 2 - Inches(0.1),
                                 gx, Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = SUB
        arr.line.fill.background()

add_footer(s, 6, 19)


# ====================================================================
# 슬라이드 7 — Granger 검정 결과 (RQ1)
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "RQ1 — Granger 인과성 검정", "15회 검정: 유의 4 · 경계 1 · 비유의 10")

# 표 직접 구성
rows = [
    ("자산", "채널", "Lag", "F", "p", "판정"),
    ("sneakers", "CH3 뉴스 감성",     "1", "13.38", "0.001", "SIGNIFICANT"),
    ("sneakers", "CH4 YT 조회수",     "1", "4.73",  "0.035", "SIGNIFICANT"),
    ("sneakers", "CH1 Google Trends", "2", "3.47",  "0.041", "SIGNIFICANT"),
    ("cards",    "CH1 Google Trends", "1", "8.42",  "0.006", "SIGNIFICANT"),
    ("lego",     "CH1 Google Trends", "1", "3.92",  "0.054", "MARGINAL"),
    ("(나머지 10개 조합)", "—", "—", "<2", ">0.18", "NOT_SIG"),
]
rows_n = len(rows)
cols_w = [Inches(1.6), Inches(2.6), Inches(0.9), Inches(1.1), Inches(1.1), Inches(1.9)]
tbl_x = Inches(0.5)
tbl_y = Inches(1.55)
row_h = Inches(0.4)

# 헤더
cx = tbl_x
for j, w in enumerate(cols_w):
    add_rect(s, cx, tbl_y, w, row_h, NAVY)
    add_text(s, cx, tbl_y, w, row_h, rows[0][j],
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w

for i in range(1, rows_n):
    cx = tbl_x
    y = tbl_y + row_h * i
    bg = LIGHT if i % 2 == 1 else WHITE
    is_sig = rows[i][5] == "SIGNIFICANT"
    is_marg = rows[i][5] == "MARGINAL"
    for j, w in enumerate(cols_w):
        add_rect(s, cx, y, w, row_h, bg,
                 line=RGBColor(0xDD, 0xDD, 0xDD))
        if j == 5 and is_sig:
            color = ACCENT; bold = True
        elif j == 5 and is_marg:
            color = RGBColor(0xC2, 0x86, 0x12); bold = True
        else:
            color = NAVY; bold = (j == 3 and is_sig)
        add_text(s, cx, y, w, row_h, rows[i][j],
                 size=12, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# 우측 요약
sx = tbl_x + sum(cols_w, Emu(0)) + Inches(0.3)
sw = SW - sx - Inches(0.4)
add_rect(s, sx, tbl_y, sw, Inches(4.5), LIGHT)
add_text(s, sx + Inches(0.2), tbl_y + Inches(0.15),
         sw - Inches(0.4), Inches(0.4),
         "핵심 발견", size=15, bold=True, color=NAVY)
add_bullets(s, sx + Inches(0.2), tbl_y + Inches(0.6),
            sw - Inches(0.4), Inches(3.8), [
    "**sneakers**: 3개 채널 유의\n(CH1·CH3·CH4)",
    "**cards**: CH1만 유의",
    "**lego**: 모두 비유의 (CH1 경계)",
    "CH3 sneakers F=13.38 — 15회 검정 최대 효과 크기",
    "BH 보정 후: 4 → 2개 유의\n(sneakers CH3, cards CH1)",
], size=12, color=NAVY)

# 결론 박스
add_rect(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.65),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.65),
         "✔ 일부 채널이 자산별로 가격을 선행하는 ‘잠정적’ 증거 — 탐색적(exploratory) 해석",
         size=14, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 7, 19)


# ====================================================================
# 슬라이드 8 — Lag 견고성
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "RQ1 견고성 — Lag 1·2·3 고정 재검정",
             "BIC 단일 lag 의존성을 보완: 다중 lag 일관성 확인")

rows = [
    ("자산", "채널", "L=1", "L=2", "L=3", "견고성"),
    ("sneakers", "CH3 뉴스 감성", "13.38 [<.001]", "7.02 [.002]", "6.65 [.001]", "✓"),
    ("sneakers", "CH1 Google",    "0.50 [.483]",   "3.47 [.041]", "2.30 [.093]", "△"),
    ("sneakers", "CH4 YT 조회수", "4.73 [.035]",   "2.79 [.073]", "1.45 [.245]", "△"),
    ("sneakers", "CH5 댓글감성",  "0.09 [.764]",   "4.58 [.016]", "4.47 [.009]", "✓ (지연 효과)"),
    ("cards",    "CH1 Google",    "8.42 [.006]",   "4.56 [.016]", "2.44 [.080]", "✓"),
    ("lego",     "(모든 채널)",    "비유의",         "비유의",       "비유의",      "—"),
]

cols_w = [Inches(1.5), Inches(2.4), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
tbl_x = Inches(0.45)
tbl_y = Inches(1.55)
row_h = Inches(0.42)

cx = tbl_x
for j, w in enumerate(cols_w):
    add_rect(s, cx, tbl_y, w, row_h, NAVY)
    add_text(s, cx, tbl_y, w, row_h, rows[0][j],
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w

for i in range(1, len(rows)):
    cx = tbl_x
    y = tbl_y + row_h * i
    bg = LIGHT if i % 2 == 1 else WHITE
    rob = rows[i][5]
    for j, w in enumerate(cols_w):
        add_rect(s, cx, y, w, row_h, bg, line=RGBColor(0xDD, 0xDD, 0xDD))
        if j == 5:
            if rob.startswith("✓"):
                color = GREEN; bold = True
            elif rob == "△":
                color = RGBColor(0xC2, 0x86, 0x12); bold = True
            else:
                color = GRAY; bold = False
        else:
            color = NAVY; bold = False
        add_text(s, cx, y, w, row_h, rows[i][j],
                 size=11, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# 하단 요약
add_rect(s, Inches(0.45), Inches(5.5), Inches(12.45), Inches(1.4),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
         "관찰점", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(5.95), Inches(12), Inches(1.0), [
    "**sneakers CH3** 모든 lag에서 p<0.01 — 가장 견고한 발견",
    "**sneakers CH5** BIC=1에선 비유의나 L=2,3에서 유의 — **댓글 감성의 2~3개월 지연 효과** 가능성 (신규 발견)",
    "sneakers CH1·CH4: 단일 lag에서만 유의 → lag 민감, 신뢰도 낮음",
], size=12)
add_footer(s, 8, 19)


# ====================================================================
# 슬라이드 9 — RQ2 Jaccard
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "RQ2 — 자산별 유의 채널 구성 비교", "Jaccard 유사도")

# 자산별 집합
sets = [
    ("Sneakers", "{CH1, CH3, CH4}", NAVY),
    ("Cards",    "{CH1}",            SUB),
    ("LEGO",     "∅ (없음)",         GRAY),
]
ssx = Inches(0.55)
ssy = Inches(1.65)
ssw = Inches(4.0)
ssh = Inches(1.3)
for i, (a, st, col) in enumerate(sets):
    y = ssy + (ssh + Inches(0.15)) * i
    add_rect(s, ssx, y, ssw, ssh, LIGHT)
    add_rect(s, ssx, y, Inches(0.18), ssh, col)
    add_text(s, ssx + Inches(0.35), y + Inches(0.15), Inches(3.5), Inches(0.4),
             a, size=16, bold=True, color=NAVY)
    add_text(s, ssx + Inches(0.35), y + Inches(0.6), Inches(3.5), Inches(0.6),
             st, size=18, color=ACCENT, bold=True)

# 우측: Jaccard 표
jx = Inches(5.0)
jy = Inches(1.65)
jw = Inches(7.8)
jh = Inches(2.6)
rows = [
    ("자산쌍", "비교 집합", "J"),
    ("sneakers – cards", "{CH1,CH3,CH4} vs {CH1}", "0.333"),
    ("sneakers – lego",  "{CH1,CH3,CH4} vs ∅",      "0.000"),
    ("cards – lego",     "{CH1} vs ∅",              "0.000"),
]
cw = [Inches(2.3), Inches(3.7), Inches(1.8)]
rh = Inches(0.5)
cx = jx
for j, w in enumerate(cw):
    add_rect(s, cx, jy, w, rh, NAVY)
    add_text(s, cx, jy, w, rh, rows[0][j],
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w
for i in range(1, len(rows)):
    cx = jx
    y = jy + rh * i
    bg = LIGHT if i % 2 == 1 else WHITE
    for j, w in enumerate(cw):
        add_rect(s, cx, y, w, rh, bg, line=RGBColor(0xDD, 0xDD, 0xDD))
        color = ACCENT if j == 2 else NAVY
        bold  = (j == 2)
        add_text(s, cx, y, w, rh, rows[i][j],
                 size=13, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# 임계 안내
add_rect(s, jx, jy + Inches(2.5), jw, Inches(1.1), RGBColor(0xF5, 0xF7, 0xFA))
add_text(s, jx + Inches(0.2), jy + Inches(2.55), jw, Inches(0.4),
         "임계 기준", size=12, bold=True, color=NAVY)
add_text(s, jx + Inches(0.2), jy + Inches(2.9), jw, Inches(0.7),
         "J<0.3 : 구성 다름 (RQ2 지지)    |   0.3≤J≤0.6 : 부분적 차이    |   J>0.6 : 유사",
         size=12, color=GRAY)

# 정의 캡션
add_rect(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.45),
         RGBColor(0xFD, 0xF3, 0xE6))
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.45),
         "📌 Jaccard 유사도 = |A ∩ B| / |A ∪ B|  ·  두 집합의 겹침을 0~1로 표현 "
         "(1=완전히 같음, 0=완전히 다름)",
         size=12, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# 결론
add_rect(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.55),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.55),
         "✔ RQ2 대체로 지지: 자산 유형별 유의 채널 구성 상이. "
         "다만 sneakers ↔ cards는 CH1을 공유 — 검색 트렌드는 ‘공통 선행 신호’ 후보",
         size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 9, 19)


# ====================================================================
# 슬라이드 10 — RQ3 DM 결과
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "RQ3 — XGBoost + Diebold–Mariano 검정",
             "Model A (전채널) vs Model B (Granger 선별) | TimeSeriesSplit, T=24/아이템")

# 좌측 표
rows = [
    ("자산", "아이템", "DM", "p", "판정"),
    ("sneakers", "jordan1", "−4.29", "<.001", "Model A 우수"),
    ("sneakers", "panda",   "−1.47", ".155",  "차이 없음"),
    ("sneakers", "yeezy",   "−1.09", ".288",  "차이 없음"),
    ("sneakers", "travis",  "+0.38", ".705",  "차이 없음"),
    ("sneakers", "nb550",   "+1.96", ".062",  "B 우세(경계)"),
    ("cards", "(5개 모두)", "—",     "≥.150", "차이 없음"),
    ("lego",  "falcon",     "−2.28", ".032",  "Model A 우수"),
    ("lego",  "hogwarts",   "+1.27", ".217",  "차이 없음"),
    ("lego",  "titanic",    "−0.56", ".579",  "차이 없음"),
    ("lego",  "porsche",    "−1.75", ".094",  "A 우세(경계)"),
    ("lego",  "bugatti",    "+0.42", ".677",  "차이 없음"),
]
cw = [Inches(1.4), Inches(1.7), Inches(1.0), Inches(1.0), Inches(2.4)]
tx = Inches(0.45); ty = Inches(1.55); rh = Inches(0.36)

cx = tx
for j, w in enumerate(cw):
    add_rect(s, cx, ty, w, rh, NAVY)
    add_text(s, cx, ty, w, rh, rows[0][j],
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w
for i in range(1, len(rows)):
    cx = tx
    y = ty + rh * i
    bg = LIGHT if i % 2 == 1 else WHITE
    judge = rows[i][4]
    for j, w in enumerate(cw):
        add_rect(s, cx, y, w, rh, bg, line=RGBColor(0xDD, 0xDD, 0xDD))
        if j == 4:
            if "A 우수" in judge:
                color = ACCENT; bold = True
            elif "B 우세" in judge:
                color = GREEN; bold = True
            elif "경계" in judge:
                color = RGBColor(0xC2, 0x86, 0x12); bold = True
            else:
                color = GRAY; bold = False
        else:
            color = NAVY; bold = False
        add_text(s, cx, y, w, rh, rows[i][j],
                 size=11, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# 우측: 핵심 수치 박스
rx = tx + sum(cw, Emu(0)) + Inches(0.25)
rw = SW - rx - Inches(0.4)

add_rect(s, rx, ty, rw, Inches(1.8), LIGHT)
add_text(s, rx + Inches(0.2), ty + Inches(0.15), rw - Inches(0.4), Inches(0.4),
         "DM 등가성 비율", size=14, bold=True, color=NAVY)
add_text(s, rx + Inches(0.2), ty + Inches(0.6), rw - Inches(0.4), Inches(1.2),
         "11 / 15  (73%)\n차이 없음",
         size=28, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# RMSE 비교
add_rect(s, rx, ty + Inches(1.9), rw, Inches(2.6), LIGHT)
add_text(s, rx + Inches(0.2), ty + Inches(2.0), rw - Inches(0.4), Inches(0.4),
         "자산별 평균 RMSE  (A vs B)", size=14, bold=True, color=NAVY)
rmse_lines = [
    ("sneakers", "452", "569"),
    ("cards",    "334", "329"),
    ("lego",     "176", "199"),
]
for k, (a, ra, rb) in enumerate(rmse_lines):
    yy = ty + Inches(2.5) + Inches(0.55) * k
    add_text(s, rx + Inches(0.2), yy, Inches(1.4), Inches(0.45),
             a, size=13, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, rx + Inches(1.6), yy, Inches(0.9), Inches(0.45),
             ra, size=14, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, rx + Inches(2.5), yy, Inches(0.3), Inches(0.45),
             "vs", size=12, color=GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, rx + Inches(2.8), yy, Inches(0.9), Inches(0.45),
             rb, size=14, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 정의 캡션
add_rect(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
         RGBColor(0xFD, 0xF3, 0xE6))
add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.4),
         "📌 Diebold–Mariano 검정  ·  두 모형(A·B)의 예측 오차 시계열이 "
         "통계적으로 다른지 평가 (p<0.05 → 차이 있음)",
         size=12, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# 결론
add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.45),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.45),
         "절제된 해석: ‘B가 A보다 우수’ 아님 → ‘대다수 아이템에서 큰 손해 없이 채널 수 축소 가능’",
         size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 10, 19)


# ====================================================================
# 슬라이드 11 — RMSE 비교 그래프 (논문 Figure 1)
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "RQ3 — 자산별 RMSE 비교 (Figure 1)",
             "Model A (전채널) vs Model B (Granger 선별) · 낮을수록 정확")

# 좌측 그림
fig_path = os.path.join(FIG_DIR, "reg_rmse_bar_00.png")
add_rect(s, Inches(0.45), Inches(1.55), Inches(8.0), Inches(5.2), WHITE,
         line=RGBColor(0xDD, 0xDD, 0xDD))
add_picture_fit(s, fig_path,
                Inches(0.55), Inches(1.65), Inches(7.8), Inches(5.0))

# 우측 해설
rx = Inches(8.65); rw = Inches(4.3)
add_rect(s, rx, Inches(1.55), rw, Inches(5.2), LIGHT)
add_text(s, rx + Inches(0.2), Inches(1.65), rw - Inches(0.4), Inches(0.4),
         "관찰", size=15, bold=True, color=NAVY)
add_bullets(s, rx + Inches(0.2), Inches(2.05), rw - Inches(0.4), Inches(2.6), [
    "**sneakers** : A 452  <  B 569",
    "**cards**    : A 334  ≈  B 329",
    "**lego**     : A 176  <  B 199",
], size=13)

add_rect(s, rx + Inches(0.2), Inches(4.0), rw - Inches(0.4), Inches(0.05), NAVY)
add_text(s, rx + Inches(0.2), Inches(4.15), rw - Inches(0.4), Inches(0.4),
         "해석", size=15, bold=True, color=NAVY)
add_bullets(s, rx + Inches(0.2), Inches(4.55), rw - Inches(0.4), Inches(2.1), [
    "sneakers · lego는 Model A 평균이 일관되게 낮음",
    "cards는 양 모형 사실상 동등",
    "DM ‘차이 없음’ 판정과 RMSE 평균은 측정 대상이 달라 함께 보고 필요",
], size=12)

add_footer(s, 11, 19)


# ====================================================================
# 슬라이드 12 — SHAP vs Granger
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "SHAP vs Granger — 측정 대상의 차이",
             "Model A 평균 |SHAP|  &  Granger 판정")

rows = [
    ("자산", "지표", "CH1", "CH2", "CH3", "CH4", "CH5"),
    ("sneakers", "|SHAP|",      "0.125", "0.000", "0.011", "0.080", "0.124"),
    ("sneakers", "Granger",     "SIG",   "NOT",   "SIG",   "SIG",   "NOT"),
    ("cards",    "|SHAP|",      "0.048", "0.031", "0.025", "0.091", "0.004"),
    ("cards",    "Granger",     "SIG",   "NOT",   "NOT",   "NOT",   "NOT"),
    ("lego",     "|SHAP|",      "0.036", "0.003", "0.000", "0.070", "0.021"),
    ("lego",     "Granger",     "MARG",  "NOT",   "NOT",   "NOT",   "NOT"),
]
cw = [Inches(1.3), Inches(1.4), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1)]
tx = Inches(1.05); ty = Inches(1.55); rh = Inches(0.42)

cx = tx
for j, w in enumerate(cw):
    add_rect(s, cx, ty, w, rh, NAVY)
    add_text(s, cx, ty, w, rh, rows[0][j],
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w
for i in range(1, len(rows)):
    cx = tx
    y = ty + rh * i
    bg = LIGHT if i % 2 == 1 else WHITE
    for j, w in enumerate(cw):
        add_rect(s, cx, y, w, rh, bg, line=RGBColor(0xDD, 0xDD, 0xDD))
        cell = rows[i][j]
        if cell == "SIG":
            color = ACCENT; bold = True
        elif cell == "MARG":
            color = RGBColor(0xC2, 0x86, 0x12); bold = True
        elif cell == "NOT":
            color = GRAY; bold = False
        else:
            color = NAVY; bold = (j >= 2)
        add_text(s, cx, y, w, rh, cell,
                 size=12, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# 해설 박스
add_rect(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(1.85), LIGHT)
add_text(s, Inches(0.7), Inches(5.15), Inches(12), Inches(0.4),
         "왜 불일치하는가?", size=15, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(5.55), Inches(12), Inches(1.4), [
    "**Granger**: 채널 단독 · 가격 라그만 통제 · **선형 선행성**",
    "**SHAP**: 모든 변수 동시 투입 · 비선형 트리 모형 · **한계 기여도**",
    "가격 모멘텀(price_vs_ma3, lag1)이 분산을 흡수 → 선형 선행 채널의 SHAP 잠식",
    "특히 sneakers CH3: Granger F 1위 ↔ SHAP 최하위 — 두 지표 차이가 가장 큰 사례",
], size=13)
add_footer(s, 12, 19)


# ====================================================================
# 슬라이드 13 — SHAP Summary (Sneakers)
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "SHAP Summary — Sneakers",
             "Model A (분류) · 좌: Gain Importance · 우: SHAP Summary Plot")

fig_path = os.path.join(FIG_DIR, "clf_shap_asset_00.png")
add_rect(s, Inches(0.45), Inches(1.55), Inches(12.45), Inches(4.4), WHITE,
         line=RGBColor(0xDD, 0xDD, 0xDD))
add_picture_fit(s, fig_path,
                Inches(0.55), Inches(1.65), Inches(12.25), Inches(4.2))

add_rect(s, Inches(0.45), Inches(6.05), Inches(12.45), Inches(0.85),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "관찰: CH1(0.125) ≈ CH5(0.124) > CH4(0.080) > CH3(0.011) > CH2(0.000)",
         size=13, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         "Granger F 1위 CH3(뉴스 감성)가 SHAP에선 최하위 — 가격 모멘텀이 선형 신호를 흡수",
         size=12, color=GRAY)
add_footer(s, 13, 19)


# ====================================================================
# 슬라이드 14 — SHAP Summary (Cards)
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "SHAP Summary — Cards",
             "Model A (분류) · 좌: Gain Importance · 우: SHAP Summary Plot")

fig_path = os.path.join(FIG_DIR, "clf_shap_asset_02.png")
add_rect(s, Inches(0.45), Inches(1.55), Inches(12.45), Inches(4.4), WHITE,
         line=RGBColor(0xDD, 0xDD, 0xDD))
add_picture_fit(s, fig_path,
                Inches(0.55), Inches(1.65), Inches(12.25), Inches(4.2))

add_rect(s, Inches(0.45), Inches(6.05), Inches(12.45), Inches(0.85),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "관찰: CH4(0.091) > CH1(0.048) > CH2(0.031) > CH3(0.025) > CH5(0.004)",
         size=13, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         "Granger 유일 유의 CH1은 SHAP 2위 — 선형 인과·비선형 기여가 비교적 일관",
         size=12, color=GRAY)
add_footer(s, 14, 19)


# ====================================================================
# 슬라이드 15 — SHAP Summary (LEGO)
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "SHAP Summary — LEGO",
             "Model A (분류) · 좌: Gain Importance · 우: SHAP Summary Plot")

fig_path = os.path.join(FIG_DIR, "clf_shap_asset_04.png")
add_rect(s, Inches(0.45), Inches(1.55), Inches(12.45), Inches(4.4), WHITE,
         line=RGBColor(0xDD, 0xDD, 0xDD))
add_picture_fit(s, fig_path,
                Inches(0.55), Inches(1.65), Inches(12.25), Inches(4.2))

add_rect(s, Inches(0.45), Inches(6.05), Inches(12.45), Inches(0.85),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "관찰: CH4(0.070) > CH1(0.036) > CH5(0.021) > CH2(0.003) > CH3(0.000)",
         size=13, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         "Granger 비유의 CH4가 SHAP 1위 — falcon Star Wars 이벤트 등 비선형 관심 충격 시사",
         size=12, color=GRAY)
add_footer(s, 15, 19)


# ====================================================================
# 슬라이드 16 — 4개 태스크 종합
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "4개 태스크 종합 — 견고한 채널 후보",
             "Granger · Model B · SHAP · DM 4개 신호의 일관성")

note = ("‘다수 태스크에서 일관된 양의 신호’를 보이는 조합이 실용적으로 유용한 선행 지표")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
         note, size=14, color=ACCENT)

rows = [
    ("자산", "채널", "Granger", "Model B", "SHAP", "DM 단서"),
    ("sneakers", "CH1 Google Trends", "✓ SIG",    "포함", "1위",  "—"),
    ("sneakers", "CH3 뉴스 감성",      "✓ SIG (F최대)","포함","4위", "—"),
    ("sneakers", "CH4 YT 조회수",      "✓ SIG",    "포함", "3위",  "—"),
    ("sneakers", "CH5 댓글 감성",      "지연 효과","제외", "2위",  "—"),
    ("cards",    "CH1 Google Trends", "✓ SIG",    "포함", "2위",  "차이 없음"),
    ("lego",     "CH1 Google Trends", "△ MARG",   "포함", "2위",  "—"),
    ("lego",     "CH4 YT 조회수",      "비유의",    "제외", "1위",  "—"),
]
cw = [Inches(1.4), Inches(2.6), Inches(2.0), Inches(1.5), Inches(1.2), Inches(2.0)]
tx = Inches(0.55); ty = Inches(1.95); rh = Inches(0.4)

cx = tx
for j, w in enumerate(cw):
    add_rect(s, cx, ty, w, rh, NAVY)
    add_text(s, cx, ty, w, rh, rows[0][j],
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w
for i in range(1, len(rows)):
    cx = tx
    y = ty + rh * i
    bg = LIGHT if i % 2 == 1 else WHITE
    for j, w in enumerate(cw):
        add_rect(s, cx, y, w, rh, bg, line=RGBColor(0xDD, 0xDD, 0xDD))
        cell = rows[i][j]
        if "SIG" in cell or "포함" in cell or "1위" in cell or "2위" in cell:
            color = ACCENT; bold = True
        elif "지연" in cell or "MARG" in cell:
            color = RGBColor(0xC2, 0x86, 0x12); bold = True
        elif "제외" in cell or "비유의" in cell:
            color = GRAY; bold = False
        else:
            color = NAVY; bold = False
        add_text(s, cx, y, w, rh, cell,
                 size=12, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# 견고 후보
add_rect(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.2),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.4),
         "★ 가장 견고한 선행 지표 후보", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.7),
         "① sneakers의 CH1 (Granger 유의 + Model B + SHAP 1위)\n"
         "② cards의 CH1     (Granger 유의 + Model B + SHAP 2위)  →  검색 트렌드(CH1)가 공통 후보",
         size=12, color=NAVY)
add_footer(s, 16, 19)


# ====================================================================
# 슬라이드 17 — 한계 & 향후 연구
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "한계 및 향후 연구")

# 좌측: 한계
add_rect(s, Inches(0.45), Inches(1.5), Inches(6.3), Inches(5.3),
         RGBColor(0xFD, 0xF0, 0xEE))
add_text(s, Inches(0.65), Inches(1.6), Inches(6), Inches(0.4),
         "한계 (Limitations)", size=17, bold=True, color=ACCENT)
add_bullets(s, Inches(0.65), Inches(2.05), Inches(6), Inches(4.7), [
    "**소표본** N=48 — 다중 검정 검정력 제약",
    "ADF: sneakers CH1 차분 후 p=0.120 잔존 비정상",
    "**분석 단위 불일치**: 자산 Granger vs 아이템 XGBoost",
    "**CH3 대리변수** GDELT tone (FinBERT 대체) — 도메인 특화도 낮음",
    "yeezy 4개월 선형보간 / lego porsche·bugatti CH1 zero값 다수",
    "**z-score 전역 적합** — 미미한 정보 누설 가능성",
    "DM 검정 태스크 불일치 (회귀 RMSE vs 분류 모형)",
    "RMSE 평균 ↔ DM 등가의 괴리",
    "탐색적 연구 — 직접 비교 가능 기준선 부재",
], size=12, color=GRAY)

# 우측: 향후 연구
add_rect(s, Inches(7.0), Inches(1.5), Inches(5.85), Inches(5.3),
         RGBColor(0xE9, 0xF1, 0xFA))
add_text(s, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.4),
         "향후 연구 (Future Work)", size=17, bold=True, color=NAVY)
add_bullets(s, Inches(7.2), Inches(2.05), Inches(5.5), Inches(4.7), [
    "더 긴 패널(5년+) · 아이템 확대 → 검정력 강화",
    "CH3를 **도메인 특화 FinBERT**로 대체",
    "**패널 Granger** (Dumitrescu–Hurlin) 또는 아이템별 단독 Granger",
    "**SPA / MCS** 등 다중 모형 비교 절차 도입",
    "**사전 등록 설계**로 확증적(confirmatory) 연구 확장",
    "lag 범위 확장 — 댓글 감성 지연 효과 정밀 검증",
    "명품·시계·예술품 등 타 리셀 자산군 적용",
], size=12, color=NAVY)
add_footer(s, 17, 19)


# ====================================================================
# 슬라이드 18 — 시사점
# ====================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "실무 시사점", "세 집단별 함의")

groups = [
    ("리셀 투자자 · 수집가", NAVY,
     ["sneakers · cards에서 **CH1 (Google Trends)** 는 lag 견고성이 확인된 선행 채널",
      "검색 관심도를 **1–2개월 선행 지표**로 활용",
      "기계적 매매 신호 ✗ — **보조 지표**로 종합 판단"]),
    ("리셀 플랫폼 운영자", SUB,
     ["가격 이상 탐지 · 수요 예측에 채널 신호 통합",
      "자산별 **유효 채널이 다름** → 맞춤형 피처 집합",
      "Granger 선별 채널로 파이프라인 유지비 절감 (73% 등가)"]),
    ("학술 연구자", ACCENT,
     ["공개 무료 데이터 → Granger → XGBoost+SHAP+DM 파이프라인",
      "수집 비용 낮음 → 명품·시계·예술품 등 타 자산군 확장 용이",
      "탐색적 비교 연구의 **재현 가능한 설계 틀**"]),
]
bw = Inches(4.05)
bh = Inches(5.0)
by = Inches(1.55)
bgap = Inches(0.2)
bx0 = Inches(0.5)
for i, (h, col, items) in enumerate(groups):
    x = bx0 + (bw + bgap) * i
    add_rect(s, x, by, bw, Inches(0.7), col)
    add_text(s, x, by, bw, Inches(0.7), h,
             size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, by + Inches(0.7), bw, bh - Inches(0.7), LIGHT)
    add_bullets(s, x + Inches(0.2), by + Inches(0.95),
                bw - Inches(0.4), bh - Inches(1.0),
                items, size=13, color=NAVY, spacing=8)

add_footer(s, 18, 19)


# ====================================================================
# 슬라이드 19 — 결론
# ====================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.0), SW, Inches(0.06), ACCENT)

add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.5),
         "Conclusion", size=16, color=RGBColor(0xC9, 0xD6, 0xE3))
add_text(s, Inches(0.6), Inches(1.05), Inches(12), Inches(0.9),
         "검색 트렌드(CH1)가 가장 견고한 공통 선행 신호 후보",
         size=28, bold=True, color=WHITE)

# 3개 핵심 메시지
items = [
    ("RQ1", "15회 중 4 SIG · 1 MARG\n— 탐색적 발견"),
    ("RQ2", "Jaccard < 0.34\n— 자산별 채널 구성 상이"),
    ("RQ3", "73% 아이템에서 DM 등가\n— ‘큰 손해 없이 간결화’"),
]
bw = Inches(4.0); bh = Inches(2.7); by = Inches(3.55)
bgap = Inches(0.2); bx0 = Inches(0.5)
for i, (rq, txt) in enumerate(items):
    x = bx0 + (bw + bgap) * i
    add_rect(s, x, by, bw, bh, RGBColor(0x12, 0x3C, 0x66))
    add_text(s, x, by + Inches(0.2), bw, Inches(0.6), rq,
             size=26, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), by + Inches(0.95),
             bw - Inches(0.4), bh - Inches(1.1), txt,
             size=15, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         "Thank you.   Q & A",
         size=18, bold=True, color=RGBColor(0xE0, 0xE8, 0xF0))

# ---------- 저장 ----------
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
