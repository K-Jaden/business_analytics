# -*- coding: utf-8 -*-
"""
최종 발표자료 — 원본 중간발표 PPT의 앞부분(서론·관련연구·방법론, 그림 포함)을
그대로 재사용하고, 실험 결과(정량/정성, figure 중심)·결론만 새로 추가해 병합.
저자: 이승현(1) · 윤재은(2) · 최형호(3)
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

INK="1A1A1A"; WHITE="FFFFFF"; SUB="AAAAAA"; SUB2="888888"
BODY="555555"; MUTE="999999"; LINE="E5E5E5"; BOX="F2F2F2"; BOX2="F7F7F7"
RED="C00000"; GREEN="2A7A2A"; GOLD="B8860B"; BLUE="185FA5"
LAT="Calibri"; EA="Malgun Gothic"
EMU=914400
def IN(v): return Emu(int(v*EMU))
HERE=os.path.dirname(os.path.abspath(__file__)); ROOT=os.path.dirname(HERE)
FIGD=os.path.join(HERE,"fig")
TEMPLATE=os.path.join(ROOT,"비지니스어낼리틱스_중간발표 최종.pptx")

prs=Presentation(TEMPLATE)
LAYOUT=prs.slide_layouts[0]

def add_slide(bg=None):
    s=prs.slides.add_slide(LAYOUT)
    if bg:
        f=s.background.fill; f.solid(); f.fore_color.rgb=RGBColor.from_string(bg)
    return s

def style_run(run,size=None,bold=None,color=None,latin=LAT,ea=EA,italic=None):
    f=run.font
    if size is not None: f.size=Pt(size)
    if bold is not None: f.bold=bold
    if italic is not None: f.italic=italic
    if color is not None: f.color.rgb=RGBColor.from_string(color)
    rPr=run._r.get_or_add_rPr()
    for t_,face in [('a:latin',latin),('a:ea',ea),('a:cs',latin)]:
        el=rPr.find(qn(t_))
        if el is None: el=rPr.makeelement(qn(t_),{}); rPr.append(el)
        el.set('typeface',face)

def rect(slide,x,y,w,h,fill=None,line=None,line_w=0.75,shape=MSO_SHAPE.RECTANGLE):
    sp=slide.shapes.add_shape(shape,IN(x),IN(y),IN(w),IN(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=RGBColor.from_string(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=RGBColor.from_string(line); sp.line.width=Pt(line_w)
    sp.shadow.inherit=False
    return sp

def tb(slide,x,y,w,h,paras,anchor=MSO_ANCHOR.TOP,wrap=True):
    box=slide.shapes.add_textbox(IN(x),IN(y),IN(w),IN(h)); tf=box.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    first=True
    for p in paras:
        text,st=p[0],(p[1] if len(p)>1 else {})
        para=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        para.alignment=st.get('align',PP_ALIGN.LEFT)
        if 'space_after' in st: para.space_after=Pt(st['space_after'])
        if 'space_before' in st: para.space_before=Pt(st['space_before'])
        if 'line_spacing' in st: para.line_spacing=st['line_spacing']
        if len(p)>2 and isinstance(p[2],list):
            for rtext,rst in p[2]:
                r=para.add_run(); r.text=rtext
                style_run(r,rst.get('size',st.get('size')),rst.get('bold',st.get('bold')),
                          rst.get('color',st.get('color')),rst.get('latin',LAT),rst.get('ea',EA),rst.get('italic'))
        else:
            r=para.add_run(); r.text=text
            style_run(r,st.get('size'),st.get('bold'),st.get('color'),st.get('latin',LAT),st.get('ea',EA),st.get('italic'))
    return box

def header(slide,title,sub=None):
    rect(slide,0,0,10,0.045,fill=INK)
    tb(slide,0.55,0.18,9.3,0.52,[(title,{'size':26,'bold':True,'color':INK})],wrap=False)
    if sub: tb(slide,0.55,0.72,9.0,0.28,[(sub,{'size':12.5,'color':SUB})])
    rect(slide,0.55,1.04,9.0,0.014,fill=LINE)

def tpl_divider(num,title,sub):
    s=add_slide(bg=INK); rect(s,0,0,10,0.04,fill=WHITE)
    tb(s,0.6,1.5,2.0,1.2,[(num,{'size':64,'bold':True,'color':WHITE})])
    tb(s,2.4,1.8,7.2,0.7,[(title,{'size':30,'bold':True,'color':WHITE})])
    tb(s,2.4,2.55,7.2,0.45,[(sub,{'size':15,'color':SUB})])
    return s

def pic(slide,name,x,y,w=None,h=None):
    path=os.path.join(FIGD,name); iw,ih=Image.open(path).size; ar=iw/ih
    if w and not h: h=w/ar
    elif h and not w: w=h*ar
    slide.shapes.add_picture(path,IN(x),IN(y),IN(w),IN(h)); return w,h

def tagbox(slide,x,y,text,w=1.0,h=0.34,fsz=12.5):
    sp=rect(slide,x,y,w,h,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf=sp.text_frame; tf.word_wrap=False
    tf.margin_left=Pt(4); tf.margin_right=Pt(4); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=text
    style_run(r,fsz,True,WHITE)

def explain(slide,x,y,w,h,paras,accent=INK):
    rect(slide,x,y,w,h,fill=WHITE,line=LINE,line_w=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide,x,y+0.06,0.05,h-0.12,fill=accent)
    tb(slide,x+0.22,y+0.12,w-0.4,h-0.2,paras,anchor=MSO_ANCHOR.MIDDLE)

def takeaway(slide,y,runs,x=0.55,w=8.9,h=0.6):
    rect(slide,x,y,w,h,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(slide,x+0.22,y,w-0.44,h,[("",{},runs)],anchor=MSO_ANCHOR.MIDDLE)

def caption(slide,x,y,w,text):
    tb(slide,x,y,w,0.3,[(text,{'size':9.5,'italic':True,'color':SUB2,'align':PP_ALIGN.CENTER})])

def notebox(slide,x,y,w,h,items,anchor=MSO_ANCHOR.MIDDLE):
    """예시 양식: 둥근 흰 박스 + 불릿. items=[(text,color,bold),...] (핵심은 RED)"""
    rect(slide,x,y,w,h,fill=WHITE,line=LINE,line_w=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    paras=[("• "+t,{'size':12,'color':c,'bold':b,'space_after':5,'line_spacing':1.08}) for (t,c,b) in items]
    tb(slide,x+0.26,y+0.12,w-0.5,h-0.22,paras,anchor=anchor)

def setnotes(slide,text):
    try:
        ns=slide.notes_slide
        tf=ns.notes_text_frame
        if tf is None:
            ph=ns.notes_placeholder
            tf=ph.text_frame if ph is not None else None
        if tf is not None:
            tf.text=text; return True
    except Exception:
        pass
    return False

# ---------- 원본 슬라이드의 정적 페이지번호 제거(병합 후 번호 불일치 방지) ----------
for s in prs.slides:
    for sh in list(s.shapes):
        if sh.has_text_frame and sh.left is not None and sh.left>IN(9.0) and sh.top is not None and sh.top<IN(1.0):
            t=sh.text_frame.text.strip()
            if re.fullmatch(r"\d{1,2}",t):
                sh._element.getparent().remove(sh._element)

# ---------- 원본 슬라이드 인용 오류 교정 (references.bib 대조) ----------
def fix_citations(slide,repls):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for para in sh.text_frame.paragraphs:
            if not para.runs: continue
            full=''.join(r.text for r in para.runs); new=full
            for a,b in repls:
                if a in new: new=new.replace(a,b)
            if new!=full:
                para.runs[0].text=new
                for r in para.runs[1:]: r.text=''
fix_citations(prs.slides[7], [("Jang et al. (2024, JFE)","Jang & Jun (2025, IRFA)")])
fix_citations(prs.slides[12],[("Jang et al. 2024","Jang & Jun 2025"),
                              ("Jang 2024; Huang 2023","Jang & Jun 2025; Huang 2023")])
fix_citations(prs.slides[13],[("Da et al. (2011, RFS)","Da et al. (2011, JF)"),
                              ("Jang et al. (2024, JFE)","Jang & Jun (2025, IRFA)")])
fix_citations(prs.slides[22],[("In search of attention. Review of Financial Studies, 24(6), 1461-1499.",
                               "In search of attention. Journal of Finance, 66(5), 1461-1499."),
                              ("59, 101547.","59, 101539.")])

# 원본 24개 슬라이드 순서 보존 (삭제는 새 슬라이드 추가 후 맨 마지막에 → partname 충돌 방지)
orig=list(prs.slides._sldIdLst)

# =====================================================================
# 새 슬라이드 (append 후 reorder)
# =====================================================================
# --- 목차 (6항목, 원본 양식) ---
s=add_slide(); header(s,"발표 목차","Contents")
toc=[("01","Introduction · 연구 배경 & 개요","리셀 시장 성장 · 연구 질문 (RQ1–RQ3)"),
     ("02","Related Works · 기존 연구 & 포지션","선행연구 · 공통 한계 · 차별점"),
     ("03","Proposed Method · 분석 대상 & 채널","3종 자산 · 5개 미디어 채널"),
     ("04","Proposed Method · 분석 방법론","FinBERT · Granger · 견고성 검정 · XGBoost"),
     ("05","Experiment Results · 실험 결과","정량평가(RQ1–3) · 정성평가(SHAP·IRF·한계)"),
     ("06","Conclusion · 결론","종합 결론 · 시사점 · 향후 연구")]
yy=1.16
for i,(num,en,kr) in enumerate(toc):
    rect(s,0.55,yy,8.9,0.62,fill=(BOX if i%2==0 else WHITE))
    rect(s,0.55,yy,0.06,0.62,fill=INK)
    tb(s,0.78,yy,0.7,0.62,[(num,{'size':18,'bold':True,'color':INK,'align':PP_ALIGN.CENTER})],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,1.6,yy+0.12,0.012,0.38,fill=LINE)
    tb(s,1.76,yy+0.08,7.5,0.3,[(en,{'size':14,'bold':True,'color':INK})])
    tb(s,1.76,yy+0.36,7.6,0.24,[(kr,{'size':11.5,'color':SUB2})])
    yy+=0.70

# --- 추가 분석 방법: 견고성 + Jaccard ---
s=add_slide(); header(s,"추가 분석 방법 — 견고성 검정 & Jaccard","단독 Granger의 약점 보완 + 채널 구성 차이의 정량화")
rect(s,0.55,1.14,8.9,0.62,fill=BOX2); rect(s,0.55,1.14,0.06,0.62,fill=INK)
tb(s,0.75,1.16,8.6,0.6,[("",{},[("왜 추가 검정?  ",{'size':12.5,'bold':True,'color':INK}),
    ("단독 Granger는 ①자산 평균 ②단일 채널 ③15회 반복 → 평균 아티팩트·다중공선성·다중검정 문제. 두 보완 도구를 추가.",{'size':12,'color':BODY})])],anchor=MSO_ANCHOR.MIDDLE)
# 왼쪽: 견고성 삼각검증
rect(s,0.55,1.92,4.35,3.3,fill=BOX2); rect(s,0.55,1.92,4.35,0.42,fill=INK)
tb(s,0.55,1.92,4.35,0.42,[("견고성 삼각검증 (Triangulation)",{'size':13,'bold':True,'color':WHITE,'align':PP_ALIGN.CENTER})],anchor=MSO_ANCHOR.MIDDLE)
robs=[("① 단독 Granger","자산 평균 시계열에 채널 1개씩 투입 (주 분석, 15회)"),
      ("② DH 패널","아이템 5개 횡단면 → 평균이 만든 거짓신호 적발"),
      ("③ VAR 블록","5채널 동시 투입 → 다중공선성 통제 후 순효과")]
yy=2.5
for t_,d in robs:
    tb(s,0.74,yy,4.0,0.3,[(t_,{'size':12.5,'bold':True,'color':RED})])
    tb(s,0.74,yy+0.30,4.0,0.5,[(d,{'size':11.5,'color':BODY,'line_spacing':1.05})]); yy+=0.92
# 오른쪽: Jaccard
rect(s,5.1,1.92,4.35,3.3,fill=BOX2); rect(s,5.1,1.92,4.35,0.42,fill=INK)
tb(s,5.1,1.92,4.35,0.42,[("Jaccard 유사도",{'size':13,'bold':True,'color':WHITE,'align':PP_ALIGN.CENTER})],anchor=MSO_ANCHOR.MIDDLE)
tb(s,5.3,2.5,4.0,2.7,[
    ("",{},[("정의  ",{'size':12.5,'bold':True,'color':INK}),("J = |A ∩ B| / |A ∪ B|  ∈ [0, 1]",{'size':12,'color':BODY})]),
    ("",{'space_before':6},[("무엇을  ",{'size':12.5,'bold':True,'color':INK}),("두 자산의 ‘유의 채널 집합’이 얼마나 겹치는지",{'size':12,'color':BODY})]),
    ("",{'space_before':6},[("왜 쓰나  ",{'size':12.5,'bold':True,'color':RED}),("자산별 선행 채널 구성이 다른지(RQ2/H2)를 한 숫자로 정량화",{'size':12,'color':BODY})]),
    ("",{'space_before':8},[("판정  ",{'size':12.5,'bold':True,'color':INK}),("<0.3 다름 · 0.3~0.6 부분 · >0.6 유사",{'size':11.5,'color':BODY})])])

# --- 05 실험 결과 divider ---
tpl_divider("05","실험 결과","정량평가 (RQ1–RQ3) · 정성평가 (SHAP · IRF · 한계)")

# R1 정량① Granger
s=add_slide(); header(s,"실험 결과 · 정량평가 ① — RQ1 Granger 인과검정","15회 중 유의 4 · 경계 1 · 비유의 10 (raw p, 래그 BIC 선택)")
pic(s,"fig_granger.png",0.45,1.25,h=3.35)
caption(s,0.3,4.62,5.2,"그림 1. 자산×채널 15쌍의 Granger F-통계량 (초록 유의 · 노랑 경계 · 회색 비유의)")
notebox(s,5.75,1.45,3.7,3.4,[
    ("스니커즈는 뉴스 감성(F=13.38)·조회수·검색이, 카드는 검색량(F=8.42)이 가격을 앞서 움직였다.",INK,False),
    ("레고에서는 검색이 경계 수준일 뿐, 뚜렷한 선행 채널이 관찰되지 않았다.",INK,False),
    ("선행 채널의 구성이 자산 유형마다 다르게 나타난다.",RED,True)])

# R2 정량② 삼각검증
s=add_slide(); header(s,"실험 결과 · 정량평가 ② — 견고성 삼각검증","단독·DH·VAR 세 검정 + BH 보정을 모두 통과하는 채널만 신뢰")
pic(s,"fig_triangulation.png",0.5,1.20,h=3.15)
caption(s,0.4,4.5,5.9,"그림 2. 단독·DH·VAR·BH 네 검정의 채널별 유의성 (초록 통과 · 회색 탈락)")
notebox(s,6.6,1.42,2.85,3.45,[
    ("카드 검색량(CH1)은 네 검정을 모두 통과한 유일한 채널이다.",INK,False),
    ("단독 검정에서 가장 강했던 스니커즈 뉴스감성은 아이템별로 나누면 사라진다(5개 모두 F<0.5).",INK,False),
    ("카드 검색량만이 평균화에 흔들리지 않는 진짜 선행 신호다.",RED,True)])

# R3 정량③ RQ2
s=add_slide(); header(s,"실험 결과 · 정량평가 ③ — RQ2 채널 구성 차이","자산별 유의 채널이 다른가? + 선행 신호 시각 확인")
pic(s,"fig_jaccard.png",0.6,1.22,h=2.35)
pic(s,"fig_leadlag.png",5.2,1.22,h=2.35)
caption(s,0.4,3.6,4.6,"그림 3. 자산별 유의 채널 집합")
caption(s,5.0,3.6,4.6,"그림 4. 카드: 검색량(초록)과 가격(검정), 정규화")
notebox(s,0.55,3.98,8.9,1.2,[
    ("유의 채널 집합은 스니커즈 {CH1,CH3,CH4}, 카드 {CH1}, 레고 {∅}로 나타났다.",INK,False),
    ("자산쌍 Jaccard 유사도는 모두 0.6 미만이다(스니커즈–카드 0.33, 그 외 0).",INK,False),
    ("자산 유형에 따라 선행 채널의 구성이 뚜렷이 갈린다 (H2 지지).",RED,True)])

# R4 정량④ RQ3
s=add_slide(); header(s,"실험 결과 · 정량평가 ④ — RQ3 예측 검증","Model A(전채널) vs Model B(Granger 선별) — H3 부분 지지")
pic(s,"fig_rmse.png",0.55,1.2,h=2.7)
pic(s,"fig_dm.png",5.2,1.16,h=2.74)
caption(s,0.4,3.95,4.6,"그림 5. 자산별 예측 오차(RMSE): Model A vs B")
caption(s,5.0,3.95,4.6,"그림 6. 아이템별 Diebold–Mariano 통계량")
notebox(s,0.55,4.26,8.9,0.94,[
    ("15개 아이템 중 11개에서 두 모형의 예측 오차에 유의한 차이가 없었다.",INK,False),
    ("스니커즈·레고에서는 전채널 모형(A)이 오히려 약간 우세했다.",INK,False),
    ("선별 모형이 더 낫다기보다, 큰 손해 없이 채널 수를 줄일 수 있다는 의미다 (H3 부분 지지).",RED,True)])

# R5 정성① SHAP
s=add_slide(); header(s,"실험 결과 · 정성평가 ① — SHAP 해석","인과성(Granger) ≠ 예측 기여(SHAP) — 체계적 불일치")
fw,fh=pic(s,"fig_shap.png",1.2,1.12,w=7.6); yb=1.12+fh
caption(s,0.55,yb+0.04,8.9,"그림 7. 채널별 평균 |SHAP| 기여도(좌)와 가격 모멘텀의 압도적 비중(우)")
notebox(s,0.55,yb+0.34,8.9,5.22-(yb+0.34),[
    ("Granger에서 가장 강했던 뉴스 감성(CH3)이 SHAP 기여도에서는 최하위로 나타났다.",INK,False),
    ("가격 모멘텀(price_vs_ma3)이 예측 분산의 대부분을 흡수하기 때문이다.",INK,False),
    ("통계적 인과성과 예측 기여도는 다른 개념 — 두 지표를 함께 볼 때 채널의 역할이 온전히 드러난다.",RED,True)])

# R6 정성② IRF
s=add_slide(); header(s,"실험 결과 · 정성평가 ② — 충격반응 분석(IRF)","두 BH-견고 채널의 동학: 부호·시차·감쇄")
fw,fh=pic(s,"fig_irf.png",1.2,1.12,w=7.6); yb=1.12+fh
caption(s,0.55,yb+0.04,8.9,"그림 8. 채널 충격에 대한 가격의 반응 함수 (음영 = 95% 신뢰구간)")
notebox(s,0.55,yb+0.34,8.9,5.22-(yb+0.34),[
    ("두 채널 모두 충격 후 1~2개월에 반응이 정점에 이른 뒤 4개월 안에 사라진다.",INK,False),
    ("부호는 정반대다 — 카드 검색은 가격을 끌어올리고(+), 스니커즈 뉴스 감성은 끌어내린다(−).",INK,False),
    ("같은 ‘선행’이라도 작동 방향이 자산마다 다르다.",RED,True)])

# R7 정성③ 경제학 2x2
s=add_slide(); header(s,"실험 결과 · 정성평가 ③ — 경제학적 해석","‘정보 비대칭 강도 × 공급 탄력성’ 2×2 프레임")
pic(s,"fig_2x2.png",0.4,1.3,h=3.3)
caption(s,0.3,4.68,5.3,"그림 9. 정보 비대칭 × 공급 탄력성 좌표 위의 세 자산")
econ=[("스니커즈","정보 비대칭이 높아 뉴스 감성(CH3)이 우세","출시·셀럽·가품 정보가 많아 보도의 방향이 가격을 움직인다"),
      ("트레이딩 카드","정보 비대칭이 낮아 검색량(CH1)이 우세","PSA 등급제로 품질이 표준화돼 가격이 수요(검색)에 직접 반응한다"),
      ("레고","공급이 극히 비탄력적이라 미디어가 매개되지 않음","단종 이후엔 잔존 재고와 거래 빈도가 가격을 지배한다")]
yy=1.5
for name,head,desc in econ:
    rect(s,5.75,yy,3.7,1.02,fill=BOX2); rect(s,5.75,yy,0.05,1.02,fill=INK)
    tb(s,5.9,yy+0.08,3.5,0.28,[(name,{'size':12,'bold':True,'color':INK})])
    tb(s,5.9,yy+0.34,3.5,0.3,[(head,{'size':10.5,'bold':True,'color':RED})])
    tb(s,5.9,yy+0.66,3.5,0.32,[(desc,{'size':9.5,'color':BODY,'line_spacing':1.0})]); yy+=1.06

# R8 정성④ 한계 (까닭 포함)
s=add_slide(); header(s,"실험 결과 · 정성평가 ④ — 한계와 그 배경","다수 채널이 비유의였고 H1·H3가 부분 지지에 그친 데에는 구조적 이유가 있다")
lims=[("소표본과 다중검정","신뢰할 수 있는 월별 체결가가 2022년부터 확보되고 레고 단종 세트는 그 이전 데이터가 없어, 패널이 48개월로 제한된다. 표본이 작고 채널을 반복 검정하므로 결과는 탐색적으로 해석한다."),
      ("평균화의 부작용","자산 대표 시계열은 아이템 5개를 평균해 만들기에 개별 잡음이 줄지만, 그 과정에서 거짓 공통신호가 증폭될 수 있다. DH 검정으로 이를 가려냈다."),
      ("분석 단위의 불일치","Granger는 자산 단위 비교를, XGBoost는 아이템 단위 예측을 목적으로 한다. 단위가 달라 자산 선별 채널을 모든 아이템에 적용할 때 정보 손실이 생긴다."),
      ("가격 모멘텀의 우위","리셀 가격의 자기상관이 강해 가격 라그가 예측 분산을 먼저 차지한다. 그 결과 SHAP에서 채널의 한계 기여가 작게 측정된다."),
      ("감성 채널의 대리변수","계획했던 NewsAPI는 무료 플랜의 역사 데이터가 30일로 제한돼 4년치를 모을 수 없었다. 대신 GDELT 논조를 사용했고, 도메인 특화도가 낮은 점은 한계로 남는다."),
      ("직접 비교 기준의 부재","리셀 다자산·다채널 선행성 비교는 사실상 첫 시도여서 동일 정의·기간의 기준선 연구가 없다. 효과 크기의 상대적 위치를 가늠하기 어렵다.")]
positions=[(0.55,1.18),(5.1,1.18),(0.55,2.52),(5.1,2.52),(0.55,3.86),(5.1,3.86)]
for (px,py),(title,desc) in zip(positions,lims):
    rect(s,px,py,4.35,1.24,fill=BOX2); rect(s,px,py,0.05,1.24,fill=RED)
    tb(s,px+0.18,py+0.1,4.0,0.3,[(title,{'size':12.5,'bold':True,'color':INK})])
    tb(s,px+0.18,py+0.4,4.05,0.8,[(desc,{'size':9.8,'color':BODY,'line_spacing':1.05})])
tb(s,0.55,5.22,9.0,0.3,[("이러한 한계를 투명하게 밝히고, BH 보정·DH·VAR로 결과를 스스로 검증해 탐색적 발견으로 한정했다.",{'size':11,'bold':True,'color':INK})])

# --- 06 결론 divider ---
tpl_divider("06","결론","종합 결론 · 시사점 · 향후 연구")

# C1 결론
s=add_slide(); header(s,"결론 — 가설 판정과 핵심 발견","48개월 패널 · 3자산 × 5채널 · Granger + 견고성 3종 + 머신러닝 검증")
concl=[("RQ1 / H1","부분 지지",GOLD,"15회 중 유의 4·경계 1. BH·DH·VAR 후 일관 생존은 카드 CH1 1개. 확증보다 ‘검증할 가설’."),
       ("RQ2 / H2","지지",GREEN,"Jaccard 모두 <0.6 (sn–cards=0.333, 그외 0). 자산 유형별 선행 채널 구성이 대체로 상이."),
       ("RQ3 / H3","부분 지지",GOLD,"73% 아이템서 A·B 예측 동등. ‘B가 우수’는 아니나 큰 손해 없이 채널 축소 가능.")]
yy=1.20
for t_,vtxt,vcol,desc in concl:
    rect(s,0.55,yy,8.9,0.92,fill=BOX2); rect(s,0.55,yy,0.06,0.92,fill=INK)
    tb(s,0.74,yy,1.5,0.92,[(t_,{'size':14,'bold':True,'color':INK})],anchor=MSO_ANCHOR.MIDDLE)
    tb(s,2.25,yy,1.4,0.92,[(vtxt,{'size':14,'bold':True,'color':vcol})],anchor=MSO_ANCHOR.MIDDLE)
    tb(s,3.7,yy,5.6,0.92,[(desc,{'size':11.8,'color':BODY,'line_spacing':1.08})],anchor=MSO_ANCHOR.MIDDLE); yy+=1.0
rect(s,0.55,4.25,8.9,0.92,fill=INK)
tb(s,0.74,4.25,8.55,0.92,[("",{},[("핵심 발견  ",{'size':14,'bold':True,'color':"FFE08A"}),
    ("‘카드의 검색량(CH1)’이 단독·DH·VAR·BH 모든 검정을 통과한 ",{'size':12.5,'color':WHITE}),
    ("유일한 견고한 자산-공통 선행 채널",{'size':12.5,'bold':True,'color':"FFE08A"}),
    (". 미디어 선행성은 자산 유형별로 분화된다.",{'size':12.5,'color':WHITE})])],anchor=MSO_ANCHOR.MIDDLE)

# C2 시사점 & 향후 연구
s=add_slide(); header(s,"시사점 & 향후 연구","탐색적 발견의 실무적 함의와 확증 연구로의 발전 경로")
tb(s,0.55,1.16,9.0,0.3,[("실무 시사점 — 3개 집단",{'size':13.5,'bold':True,'color':INK})])
impl=[("리셀 투자자·수집가","검색 관심도(CH1)를 1~2개월 선행 보조지표로 모니터링. 기계적 매매신호가 아닌 수급 정보와 함께 종합 판단."),
      ("리셀 플랫폼 운영자","자산 유형별 맞춤 피처 권장. Granger 선별 채널 중심 구성 시 73% 아이템서 예측력 유지하며 비용 절감."),
      ("학술 연구자","‘공개 무료데이터 → Granger 선별 → XGBoost+SHAP+DM’ 파이프라인은 명품·시계·예술품 등으로 확장 가능한 재현 틀.")]
bx=0.55
for name,desc in impl:
    rect(s,bx,1.5,2.92,1.9,fill=BOX2); rect(s,bx,1.5,2.92,0.46,fill=INK)
    tb(s,bx,1.5,2.92,0.46,[(name,{'size':12.5,'bold':True,'color':WHITE,'align':PP_ALIGN.CENTER})],anchor=MSO_ANCHOR.MIDDLE)
    tb(s,bx+0.16,2.04,2.6,1.3,[(desc,{'size':11.3,'color':BODY,'line_spacing':1.1})]); bx+=3.0
tb(s,0.55,3.60,9.0,0.3,[("향후 연구",{'size':13.5,'bold':True,'color':INK})])
rect(s,0.55,3.94,8.9,0.014,fill=LINE)
tb(s,0.6,4.04,9.0,1.2,[
    ("① 더 긴 패널(5년+)·아이템 확대로 N·T 양 차원 검정력 강화",{'size':12,'color':BODY,'space_after':4}),
    ("② CH3를 도메인 특화 FinBERT로 대체해 감성 신호 품질 개선",{'size':12,'color':BODY,'space_after':4}),
    ("③ DH 잔차 부트스트랩 · ④ DM 한계 보완을 위한 SPA / MCS 다중모형 비교 절차 적용",{'size':12,'color':BODY,'space_after':4}),
    ("⑤ 사전 등록(pre-registration) 설계를 통한 확증(confirmatory) 연구로 발전",{'size':12,'bold':True,'color':INK})])

# =====================================================================
# reorder: 앞부분(원본) + 새 슬라이드 + refs + thank
# =====================================================================
allids=list(prs.slides._sldIdLst)
new=allids[len(orig):]   # 14개 새 슬라이드(생성 순)
# 원본 삭제 대상(관계 제거): 1 목차, 19 기대결과div, 20 예상결과, 21 투고전략
for i in (1,19,20,21):
    prs.part.drop_rel(orig[i].get(qn('r:id')))
# orig: 0 title,2 §01,3 연구개요,4 배경,5 리셀,6 §02,7 기존,8 한계,9 차이,10 §03,
#       11 분석대상,12 5채널,13 채널선정,14 §04,15 제안모델,16 FinBERT,17 Granger,18 XGBoost,22 refs,23 thank
# new : 0 목차,1 bridge,2 §05div,3 R1,4 R2,5 R3,6 R4,7 R5,8 R6,9 R7,10 R8,11 §06div,12 C1,13 C2
order=[orig[0], new[0]] + orig[2:6] + orig[6:10] + orig[10:14] + orig[14:19] \
      + [new[1], new[2]] + new[3:11] + [new[11], new[12], new[13]] + [orig[22], orig[23]]
lst=prs.slides._sldIdLst
for el in list(lst): lst.remove(el)
for el in order: lst.append(el)

# ---------- 슬라이드별 발표 노트(스피커 노트) ----------
NOTES=[
"안녕하세요. 저희 연구는 ‘미디어 채널이 리셀 시장 가격을 미리 움직이는가’를 묻습니다. 스니커즈·트레이딩 카드·레고 세 자산에 대해 Granger 인과성과 머신러닝으로 검증했습니다. 발표는 이승현·윤재은·최형호가 함께 준비했습니다.",
"발표는 서론, 관련 연구, 제안 방법, 실험 결과, 결론 순으로 진행합니다. 실험 결과는 정량평가와 정성평가로 나누어 보여드립니다.",
"먼저 연구 배경과 문제의식, 세 가지 연구 질문을 말씀드립니다.",
"리셀 가격이 오르기 전에 어떤 미디어 채널이 먼저 반응하는지를 자산별로 비교하는 것이 목적입니다. 질문은 셋입니다 — 채널이 가격을 선행하는가(H1), 선행 채널이 자산마다 다른가(H2), 선별 채널 모형이 전채널과 동등 이상인가(H3).",
"리셀 시장은 최근 5년간 급성장해 2024년 약 490억 달러, 2029년 740억 달러로 전망됩니다. 가격이 관심과 정보에 민감하다는 점에서 금융시장과 본질을 공유하지만, 어떤 정보 채널이 가격을 움직이는지에 대한 실증은 부족했습니다.",
"리셀을 택한 이유는 진입 용이성, 큰 미디어 영향, 그리고 투자 자산으로서의 시장 성장 세 가지입니다.",
"다음은 관련 연구, 그리고 기존 연구가 리셀에 그대로 적용되지 않는 지점입니다.",
"검색량은 Da와 동료들이, 트위터 감성은 Bollen이, 유튜브 조회수는 Jang과 Jun이 주가를 선행한다고 보였습니다. 저희는 검증된 이 채널들을 차용하되 리셀로 확장합니다.",
"다만 기존 연구는 주식 단일 자산, 단일 채널이고 볼륨과 감성을 분리하지 않았습니다. 리셀은 드롭·단종처럼 이벤트로 움직이는 다른 시장입니다.",
"그래서 저희는 리셀 세 자산, 다섯 채널, 볼륨·감성 분리, Granger와 머신러닝 결합이라는 차별점을 둡니다.",
"이제 분석 대상과 채널 구성입니다.",
"세 자산 각 5개, 총 15개 아이템을 48개월 패널로 구축했습니다. 세 자산은 정보 생성 경로가 서로 달라 H2 검증에 이상적입니다.",
"채널은 검색·뉴스·유튜브 세 정보원을 볼륨과 감성으로 나눠 다섯 개로 구성했습니다. 같은 보도라도 긍정·부정에 따라 가격 방향이 달라지므로 분리가 핵심입니다.",
"각 채널은 선행 연구에서 가격 선행성이 검증된 것만 채택했습니다.",
"분석 방법론입니다.",
"전체 흐름은 데이터 수집, FinBERT 감성 점수화, Granger 선별, XGBoost 예측 검증의 4단계입니다.",
"뉴스와 댓글은 금융 특화 모델 FinBERT로 긍정에서 부정을 뺀 감성 점수로 변환했습니다.",
"Granger 검정은 가격 과거값만 쓴 모형과 채널을 더한 모형을 비교해, 채널을 넣었을 때 예측 오차가 유의하게 줄면 선행한다고 판정합니다.",
"예측 검증은 전채널 모형 A와 Granger 선별 모형 B를 XGBoost로 비교하고, DM 검정과 SHAP으로 해석합니다.",
"단독 Granger는 자산 평균·단일 채널·반복 검정이라는 약점이 있어 세 검정으로 삼각 검증했습니다. 자산별 유의 채널 집합이 얼마나 겹치는지는 Jaccard로 정량화하며, 0.6 미만이면 구성이 다르다고 봅니다.",
"이제 핵심인 실험 결과입니다. 정량평가 넷, 정성평가 넷으로 나눠 보여드립니다.",
"15회 Granger 검정에서 스니커즈는 뉴스 감성·조회수·검색이, 카드는 검색량이 가격을 선행했고 레고는 뚜렷한 선행 채널이 없었습니다. 핵심은 선행 채널 구성이 자산마다 다르다는 점입니다.",
"이 표가 가장 중요합니다. 네 검정을 모두 통과한 채널은 카드 검색량 하나뿐입니다. 단독에서 가장 강했던 스니커즈 뉴스 감성은 아이템별로 보면 사라집니다 — 평균이 만든 신호였던 겁니다. 카드 검색량만이 진짜 견고한 신호입니다.",
"유의 채널 집합은 자산마다 다르고 모든 자산쌍의 Jaccard가 0.6 미만입니다. 오른쪽 그림처럼 카드는 검색이 가격보다 먼저 움직입니다. H2는 지지됩니다.",
"예측에서는 15개 중 11개가 두 모형 간 차이가 없었고 스니커즈·레고는 전채널 모형이 약간 우세했습니다. 선별이 더 낫다기보다 큰 손해 없이 채널을 줄일 수 있다는 절제된 결론이며 H3는 부분 지지입니다.",
"흥미롭게도 Granger에서 가장 강했던 뉴스 감성이 SHAP 기여도에서는 최하위입니다. 가격 모멘텀이 분산을 대부분 흡수하기 때문입니다. 인과성과 예측 기여는 다른 개념이라 두 지표를 함께 봐야 합니다.",
"충격반응을 보면 두 채널 모두 1~2개월에 정점을 찍고 4개월 안에 사라집니다. 그런데 부호가 반대입니다 — 검색은 가격을 올리고 부정 뉴스 감성은 내립니다. 같은 선행이라도 방향이 다릅니다.",
"이 분화를 경제학적으로 보면, 정보 비대칭이 높은 스니커즈는 뉴스 감성이, 낮은 카드는 검색이 우세하고, 공급이 극히 비탄력적인 레고는 미디어가 가격에 매개되지 않습니다.",
"결과가 약하게 나온 데에는 구조적 이유가 있습니다. 데이터가 48개월로 제한되고 평균화·단위 불일치·가격 모멘텀·대리변수가 작용했습니다. 이를 투명하게 밝히고 탐색적 발견으로 한정했습니다.",
"마지막으로 결론입니다.",
"정리하면 H1은 부분 지지, H2는 지지, H3는 부분 지지입니다. 모든 검정을 통과한 카드 검색량이 가장 견고한 자산-공통 선행 채널이며, 미디어 선행성은 자산 유형별로 분화됩니다.",
"투자자에게는 검색량이 보조 선행지표가, 플랫폼에는 자산별 맞춤 피처가, 연구자에게는 무료 데이터 기반 재현 파이프라인이 의미가 있습니다. 향후 더 긴 패널, 도메인 특화 모델, 사전 등록 설계로 확증 연구로 발전시키겠습니다.",
"참고 문헌입니다.",
"이상입니다. 감사합니다. 질문 받겠습니다.",
]
embedded=0
for sl,note in zip(prs.slides,NOTES):
    if setnotes(sl,note): embedded+=1
# 발표 대본을 별도 파일로도 저장(임베딩 실패 대비 + 인쇄용)
titles=["표지","목차","[01] 서론 표지","연구 개요·질문","연구 배경","리셀 선정 이유",
        "[02] 관련연구 표지","기존 연구","공통 한계","차별점","[03] 방법 표지(대상·채널)",
        "분석 대상","5개 채널","채널 선정","[04] 방법론 표지","제안 모델","FinBERT",
        "Granger 검정","XGBoost","견고성 검정·Jaccard","[05] 실험결과 표지",
        "정량① Granger","정량② 견고성 삼각검증","정량③ RQ2","정량④ RQ3",
        "정성① SHAP","정성② IRF","정성③ 경제학 해석","정성④ 한계",
        "[06] 결론 표지","결론","시사점·향후연구","참고문헌","Thank you"]
lines=["발표 대본 (슬라이드별) — 미디어 채널의 리셀 가격 선행성\n"]
for i,(t,n) in enumerate(zip(titles,NOTES),1):
    lines.append("[슬라이드 %d] %s\n%s\n"%(i,t,n))
open(os.path.join(ROOT,"발표대본.txt"),"w",encoding="utf-8").write("\n".join(lines))
print("notes embedded:",embedded,"/",len(NOTES))

out=os.path.join(ROOT,"발표자료_최종.pptx")
try:
    prs.save(out)
except PermissionError:
    out=os.path.join(ROOT,"발표자료_최종_new.pptx"); prs.save(out)
    print("WARN: 원본이 열려있어 새 파일로 저장")
print("SAVED:",out,"| slides:",len(prs.slides._sldIdLst))
