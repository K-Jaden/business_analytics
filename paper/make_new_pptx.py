# -*- coding: utf-8 -*-
"""
수업 발표자료 — XGBoost 중심 재구성 (2026-06-02)
구조: 0막~4막, 13슬라이드, 25분
출력: Finance/수업발표자료.pptx
"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ── Palette ──────────────────────────────────────────────────────────
INK  = "1A1A1A"; WHITE = "FFFFFF"; SUB  = "AAAAAA"; SUB2  = "888888"
BODY = "555555"; LINE  = "E5E5E5"; BOX  = "F2F2F2"; BOX2  = "F7F7F7"
RED  = "C00000"; GREEN = "2A7A2A"; GOLD = "B8860B"; BLUE  = "185FA5"
LAT  = "Calibri"; EA = "Malgun Gothic"
EMU  = 914400
def IN(v): return Emu(int(v * EMU))

HERE     = os.path.dirname(os.path.abspath(__file__))
ROOT     = os.path.dirname(HERE)
FIGD     = os.path.join(HERE, "fig")
RES_FIG  = os.path.join(ROOT, "results", "figures")
TEMPLATE = os.path.join(ROOT, "발표자료_최종.pptx")

prs    = Presentation(TEMPLATE)
LAYOUT = prs.slide_layouts[0]

# 원본 슬라이드 목록 저장 (나중에 제거)
_orig_ids = list(prs.slides._sldIdLst)

# ── Helpers ──────────────────────────────────────────────────────────
def add_slide(bg=None):
    s = prs.slides.add_slide(LAYOUT)
    if bg:
        f = s.background.fill; f.solid()
        f.fore_color.rgb = RGBColor.from_string(bg)
    return s

def style_run(run, size=None, bold=None, color=None,
              latin=LAT, ea=EA, italic=None):
    f = run.font
    if size   is not None: f.size   = Pt(size)
    if bold   is not None: f.bold   = bold
    if italic is not None: f.italic = italic
    if color  is not None: f.color.rgb = RGBColor.from_string(color)
    rPr = run._r.get_or_add_rPr()
    for tag, face in [("a:latin", latin), ("a:ea", ea), ("a:cs", latin)]:
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", face)

def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75,
         shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = RGBColor.from_string(line); sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def tb(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf  = box.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for p in paras:
        text, st = p[0], (p[1] if len(p) > 1 else {})
        para = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        para.alignment = st.get("align", PP_ALIGN.LEFT)
        if "space_after"  in st: para.space_after  = Pt(st["space_after"])
        if "space_before" in st: para.space_before = Pt(st["space_before"])
        if "line_spacing" in st: para.line_spacing  = st["line_spacing"]
        if len(p) > 2 and isinstance(p[2], list):
            for rtext, rst in p[2]:
                r = para.add_run(); r.text = rtext
                style_run(r, rst.get("size", st.get("size")),
                          rst.get("bold", st.get("bold")),
                          rst.get("color", st.get("color")),
                          rst.get("latin", LAT), rst.get("ea", EA),
                          rst.get("italic"))
        else:
            r = para.add_run(); r.text = text
            style_run(r, st.get("size"), st.get("bold"), st.get("color"),
                      st.get("latin", LAT), st.get("ea", EA), st.get("italic"))
    return box

def header(slide, title, sub=None):
    rect(slide, 0, 0, 10, 0.042, fill=INK)
    tb(slide, 0.45, 0.14, 9.3, 0.52,
       [(title, {"size": 22, "bold": True, "color": INK})], wrap=False)
    if sub:
        tb(slide, 0.45, 0.68, 9.0, 0.26,
           [(sub, {"size": 11, "color": SUB2})])
    rect(slide, 0.45, 0.98, 9.1, 0.013, fill=LINE)

def pic(slide, path, x, y, w=None, h=None):
    iw, ih = Image.open(path).size; ar = iw / ih
    if w and not h: h = w / ar
    elif h and not w: w = h * ar
    slide.shapes.add_picture(path, IN(x), IN(y), IN(w), IN(h))
    return w, h

def notebox(slide, x, y, w, h, items, anchor=MSO_ANCHOR.MIDDLE):
    rect(slide, x, y, w, h, fill=WHITE, line=LINE, lw=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    paras = [("• " + t, {"size": 11.5, "color": c, "bold": b,
                          "space_after": 4, "line_spacing": 1.08})
             for t, c, b in items]
    tb(slide, x + 0.22, y + 0.1, w - 0.44, h - 0.18, paras, anchor=anchor)

def takeaway(slide, y, runs, x=0.45, w=9.1, h=0.55):
    rect(slide, x, y, w, h, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(slide, x + 0.2, y, w - 0.4, h, [("", {}, runs)],
       anchor=MSO_ANCHOR.MIDDLE)

def caption(slide, x, y, w, text):
    tb(slide, x, y, w, 0.28,
       [(text, {"size": 9, "italic": True, "color": SUB2,
                "align": PP_ALIGN.CENTER})])

def setnotes(slide, text):
    try:
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        if tf is None:
            ph = ns.notes_placeholder
            tf = ph.text_frame if ph else None
        if tf is not None: tf.text = text
    except Exception:
        pass

# 셀 단위 표 그리기 헬퍼
def table_cell(slide, x, y, w, h, text, bg=BOX2, tc=BODY,
               bold=False, fsz=10.5, align=PP_ALIGN.CENTER, line=LINE):
    rect(slide, x, y, w, h, fill=bg, line=line, lw=0.5)
    tb(slide, x + 0.06, y, w - 0.1, h,
       [(text, {"size": fsz, "bold": bold, "color": tc, "align": align})],
       anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# Slide 1 — 표지
# =====================================================================
s = add_slide(bg=INK)
rect(s, 0, 0, 10, 0.04, fill=WHITE)
rect(s, 0, 5.585, 10, 0.04, fill=WHITE)
tb(s, 0.8, 0.9, 8.4, 1.1,
   [("미디어 채널의 리셀 가격 선행성",
     {"size": 30, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER})],
   anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.8, 2.1, 8.4, 0.7,
   [("스니커즈 · 트레이딩 카드 · 레고  |  5개 미디어 채널  |  XGBoost 예측 분석",
     {"size": 14, "color": SUB, "align": PP_ALIGN.CENTER})])
rect(s, 3.0, 2.98, 4.0, 0.013, fill=SUB)
tb(s, 0.8, 3.2, 8.4, 0.4,
   [("이승현 · 윤재은 · 최형호",
     {"size": 13, "color": SUB, "align": PP_ALIGN.CENTER})])
tb(s, 0.8, 3.7, 8.4, 0.38,
   [("2026", {"size": 12, "color": SUB2, "align": PP_ALIGN.CENTER})])
setnotes(s, "안녕하세요. 리셀 시장에서 미디어 채널이 가격을 선행하는지를 Granger + XGBoost로 분석한 연구입니다.")

# =====================================================================
# Slide 2 — [0막-1] 선행연구 (2×2 그리드, 4번째: Granger→ML)
# =====================================================================
s = add_slide()
header(s, "관련 연구 — 4개 연구 흐름과 공통 한계",
       "미디어·가격·멀티채널·Granger→예측 파이프라인")

pillars = [
    # (title, research, limit, accent_color)
    ("미디어 → 가격\n선행성 연구",
     "Tetlock (2007): 뉴스 감성 → 주가 하락\nDa et al. (2011, JF): Google Trends → 주가 공포\nPreis et al. (2013): 검색량 → 시장 전략",
     "전통 금융자산(주식)만\n리셀 자산 적용 없음",
     INK),
    ("리셀/대체자산\n가격 연구",
     "StockX 스니커즈 프리미엄 연구\nDobrynskaya & Kishilova (2018): LEGO 투자\n거시 경제 변수·공급 분석 중심",
     "미디어 채널 효과 미포함\n채널 선행성 미검증",
     INK),
    ("멀티채널\n미디어 분석",
     "Bollen et al. (2011): Twitter → DJIA 예측\nJang & Jun (2025): YouTube → 주가\n볼륨·감성 분리 활용",
     "단일 자산(주식) 위주\n리셀 3종 자산 비교 없음",
     INK),
    ("Granger → ML\n예측 파이프라인",
     "Bollen (2011): Granger → SVM 예측\nPeng et al. (2020): Granger + LSTM 결합\n공통 전제: 선행 채널을 ML 피처로 그냥 사용",
     "선행 채널이 ML 예측에 실제로\n유효한지 검증한 연구 없음\n★ 이 가정 자체를 검증한 것이 본 연구",
     RED),
]

positions = [(0.45, 1.1), (5.25, 1.1), (0.45, 2.95), (5.25, 2.95)]
w_box, h_box = 4.3, 1.72

for (bx, by), (title, research, limit, col) in zip(positions, pillars):
    rect(s, bx, by, w_box, h_box, fill=BOX2,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, bx, by, w_box, 0.44, fill=col,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, bx, by, w_box, 0.44,
       [(title, {"size": 11.5, "bold": True, "color": WHITE,
                 "align": PP_ALIGN.CENTER})],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, bx + 0.16, by + 0.5, w_box - 0.3, 0.72,
       [(research, {"size": 9.5, "color": BODY, "line_spacing": 1.1})])
    rect(s, bx + 0.16, by + 1.26, w_box - 0.3, 0.013, fill=LINE)
    tb(s, bx + 0.16, by + 1.32, w_box - 0.3, 0.32,
       [("▶ 한계: " + limit,
         {"size": 9, "color": RED if col == RED else BODY,
          "bold": col == RED, "line_spacing": 1.1})])

takeaway(s, 4.78,
         [("본 연구  ",
           {"size": 12, "bold": True, "color": "FFE08A"}),
          ("① 리셀 × 멀티채널 선행성 비교 (최초)  "
           "② '선행 채널 = ML 유효 피처'라는 기존 가정을 직접 검증 → 틀렸음을 발견",
           {"size": 12, "color": WHITE})])
setnotes(s, "기존 연구는 4개 흐름입니다. Bollen(2011)·Peng(2020) 등은 Granger로 선행성을 확인한 뒤 ML 예측에 그 채널을 씁니다. 그러나 누구도 '과연 Granger 선별 채널이 ML 예측에 실제로 유효한가'를 검증하지 않았습니다. 그냥 유효하다고 가정합니다. 본 연구는 이 가정 자체를 리셀 시장에서 검증했고 — 결과는 그 가정이 틀렸습니다.")

# =====================================================================
# Slide 3 — [0막-2] 갭 매트릭스
# =====================================================================
s = add_slide()
header(s, "연구 배경 — 우리가 채우는 갭",
       "리셀 × 멀티채널 선행성 = 최초 시도")
gm_path = os.path.join(FIGD, "fig_gap_matrix.png")
if os.path.exists(gm_path):
    pic(s, gm_path, 0.5, 1.1, h=3.22)
else:
    tb(s, 0.5, 1.1, 9.0, 1.0,
       [("[fig_gap_matrix.png — scripts/make_presentation_figs.py 먼저 실행]",
         {"size": 12, "color": RED})])

# 연구 목적 명시 박스
rect(s, 0.45, 4.42, 9.1, 0.64, fill="FFF3CD", line=GOLD, lw=1.5,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.45, 4.42, 0.07, 0.64, fill=GOLD)
tb(s, 0.65, 4.46, 8.75, 0.56, [
    ("",{},[
        ("연구 목적  ",{"size":12,"bold":True,"color":GOLD}),
        ("예측 성능 향상이 목적이 아닙니다. ",{"size":12,"color":INK}),
        ("어떤 미디어 채널이 가격보다 먼저 움직이는가",{"size":12,"bold":True,"color":INK}),
        ("를 자산별로 찾는 것입니다. "
         "XGBoost는 그 채널이 실제로도 유용한지 검증하는 도구로 씁니다.",
         {"size":12,"color":BODY}),
    ])
], anchor=MSO_ANCHOR.MIDDLE)

takeaway(s, 5.14,
         [("★ 본 연구  ",
           {"size": 12, "bold": True, "color": "FFE08A"}),
          ("리셀 3종 × 5채널 선행성 비교 (최초) + Granger 선행 채널이 예측에도 유효한지 검증",
           {"size": 12, "color": WHITE})])
setnotes(s, "이 매트릭스가 우리 연구의 위치를 보여줍니다. 중요한 것은 연구 목적입니다. 저희는 예측 성능을 높이려는 게 아닙니다. 어떤 미디어가 가격보다 먼저 반응하는지 찾는 것이 목적이고, XGBoost는 그 채널이 실제 예측에도 쓸모 있는지 확인하는 검증 도구입니다.")

# =====================================================================
# Slide 4 — [1막-1] Granger 헤드라인
# =====================================================================
s = add_slide()
header(s, "정당화 — 미디어가 정말 가격을 선행하는가?",
       "Granger 인과검정 15회 · 역인과·우연이 아님을 확인")
pic(s, os.path.join(FIGD, "fig_granger.png"), 0.4, 1.15, h=3.28)
caption(s, 0.3, 4.48, 5.2,
        "그림 1. 자산×채널 15쌍 Granger F-통계량 (초록=유의 · 노랑=경계 · 회색=비유의)")
notebox(s, 5.65, 1.2, 3.9, 3.5, [
    ("Sneakers: 뉴스감성(F=13.4)·YT조회수(F=4.7)·검색(F=3.5) 선행", INK, False),
    ("Cards: 검색량(F=8.4) 선행",                                      INK, False),
    ("LEGO: 유의 채널 없음 (경계: 검색 F=3.9)",                        INK, False),
    ("BH 보정 후 sneakers CH3 · cards CH1 두 채널만 견고하게 생존",    BODY, False),
    ("역인과·우연이 아님을 확인 — XGBoost 피처 채택의 통계적 근거",     RED, True),
])
tb(s, 5.65, 4.78, 3.9, 0.44,
   [("⚠ Granger는 '정당화 근거'이지 XGBoost의 필수조건이 아닙니다.",
     {"size": 9.5, "italic": True, "color": BODY})])
setnotes(s, "Granger는 채널을 추가했을 때 예측 오차가 유의하게 줄면 선행한다고 판정합니다. 스니커즈 3채널, 카드 1채널이 유의했고 레고는 없었습니다. 이 결과를 바탕으로 Model A(전채널) vs Model B(Granger 선별) 비교로 넘어갑니다.")

# =====================================================================
# Slide 5 — [1막-2] Jaccard
# =====================================================================
s = add_slide()
header(s, "정당화 — 자산마다 유효 채널이 다르다 (H2 지지)",
       "유의 채널 집합 비교 · Jaccard 유사도")
pic(s, os.path.join(FIGD, "fig_jaccard.png"), 0.4, 1.12, h=3.3)
caption(s, 0.3, 4.45, 5.3,
        "그림 2. 자산별 유의 채널 집합 (초록=유의 · 회색=비유의)")
notebox(s, 5.65, 1.2, 3.9, 3.6, [
    ("Sneakers: {CH1, CH3, CH4}",             INK, False),
    ("Cards: {CH1}  — 검색량 1개",             INK, False),
    ("LEGO: {∅}  — 유의 채널 없음",           INK, False),
    ("Jaccard  sn–cards=0.333 / 그 외 0.000", BODY, False),
    ("모든 자산쌍 Jaccard < 0.6 → 채널 구성이 자산마다 뚜렷이 다름 (H2 지지)",
     RED, True),
])
setnotes(s, "자산별 유의 채널 집합을 비교하면 Jaccard가 모두 0.6 미만으로 채널 구성이 뚜렷이 다릅니다. H2 가설이 지지됩니다.")

# =====================================================================
# Slide 6 — [2막-1] XGBoost 설계도
# =====================================================================
s = add_slide()
header(s, "2단계 검증 — Granger 선행 채널이 예측에도 유용한가?",
       "Model A (전채널) vs Model B (Granger 선별) · XGBoost + DM 검정 · TimeSeriesSplit(5)")

rect(s, 0.45, 1.15, 4.3, 3.95, fill=BOX2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.45, 1.15, 4.3, 0.48, fill=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tb(s, 0.45, 1.15, 4.3, 0.48,
   [("Model A — 전채널", {"size": 13, "bold": True, "color": WHITE,
                          "align": PP_ALIGN.CENTER})],
   anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.65, 1.72, 3.9, 3.25, [
    ("피처 (9개):", {"size": 11, "bold": True, "color": INK}),
    ("price_vs_ma3, price_chg_lag1~3\n(가격 모멘텀 통제 변수)",
     {"size": 10.5, "color": BODY, "space_after": 6}),
    ("CH1 Google Trends 검색량",       {"size": 10.5, "color": BODY}),
    ("CH2 뉴스 보도량 (GDELT)",         {"size": 10.5, "color": BODY}),
    ("CH3 뉴스 감성 (GDELT tone)",      {"size": 10.5, "color": BODY}),
    ("CH4 YouTube 조회수",              {"size": 10.5, "color": BODY}),
    ("CH5 YouTube 댓글 감성 (FinBERT)", {"size": 10.5, "color": BODY}),
])

tb(s, 4.82, 2.8, 0.6, 0.5,
   [("→", {"size": 22, "bold": True, "color": INK, "align": PP_ALIGN.CENTER})],
   anchor=MSO_ANCHOR.MIDDLE)

rect(s, 5.45, 1.15, 4.1, 3.95, fill=BOX2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 5.45, 1.15, 4.1, 0.48, fill=RED,  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tb(s, 5.45, 1.15, 4.1, 0.48,
   [("Model B — Granger 선별", {"size": 13, "bold": True, "color": WHITE,
                                 "align": PP_ALIGN.CENTER})],
   anchor=MSO_ANCHOR.MIDDLE)
tb(s, 5.65, 1.72, 3.7, 3.25, [
    ("피처 (4~5개):", {"size": 11, "bold": True, "color": INK}),
    ("price_vs_ma3, price_chg_lag1~3\n(동일 통제 변수)",
     {"size": 10.5, "color": BODY, "space_after": 6}),
    ("Granger 유의 채널만:", {"size": 10.5, "bold": True, "color": RED}),
    ("Sneakers → CH1, CH3, CH4",  {"size": 10.5, "color": BODY}),
    ("Cards → CH1",               {"size": 10.5, "color": BODY}),
    ("LEGO → (없음, 전채널 폴백)", {"size": 10.5, "color": BODY}),
])
takeaway(s, 5.18,
         [("질문: Granger에서 선행한다고 나온 채널을 ML 예측 피처로 쓰면 — 실제로 더 잘 맞히나?",
           {"size": 12, "bold": True, "color": "FFE08A"})])
setnotes(s, "XGBoost는 예측 성능 향상이 목적이 아니라 검증 도구입니다. Granger로 선행성이 확인된 채널을 피처로 썼을 때 예측이 나아지는지 확인합니다. 이게 핵심 질문입니다.")

# =====================================================================
# Slide 7 — [2막-2] DM 검정 결과
# =====================================================================
s = add_slide()
header(s, "DM 검정 — A·B 모형 예측력 비교 결과",
       "Diebold-Mariano · 15개 아이템 × 2모형 · 24개월 테스트 구간")

rect(s, 0.45, 1.15, 2.8, 2.25, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tb(s, 0.45, 1.15, 2.8, 1.4,
   [("11 / 15", {"size": 36, "bold": True, "color": WHITE,
                 "align": PP_ALIGN.CENTER})],
   anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.45, 2.55, 2.8, 0.75,
   [("아이템에서 A·B 차이 없음\n(73% — 간결성 확보)",
     {"size": 11.5, "color": SUB, "align": PP_ALIGN.CENTER, "line_spacing": 1.2})],
   anchor=MSO_ANCHOR.TOP)

details = [
    ("Sneakers", "Jordan1: Model A 유의 우세 (p<0.001)\n나머지 4개: 차이 없음", RED,   "A 우세"),
    ("Cards",    "5개 전부 차이 없음\nGranger 선별(CH1)로 간결성 확보",           GREEN, "동등"),
    ("LEGO",     "Falcon: Model A 약간 우세 (p=0.032)\n나머지 4개: 차이 없음",   GOLD,  "A 약간 우세"),
]
yy = 1.15
for asset, desc, col, verdict in details:
    rect(s, 3.45, yy, 6.1, 1.28, fill=BOX2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 3.45, yy, 0.06, 1.28, fill=col)
    tb(s, 3.6,  yy + 0.1,  1.5, 0.36,
       [(asset, {"size": 13, "bold": True, "color": INK})])
    tb(s, 3.6,  yy + 0.44, 4.2, 0.78,
       [(desc,  {"size": 10.5, "color": BODY, "line_spacing": 1.15})])
    tb(s, 8.25, yy + 0.36, 1.2, 0.56,
       [(verdict, {"size": 11.5, "bold": True, "color": col,
                   "align": PP_ALIGN.CENTER})],
       anchor=MSO_ANCHOR.MIDDLE)
    yy += 1.36

takeaway(s, 5.18,
         [("답: Granger 선별 채널이 전채널보다 낫지 않다 — ",
           {"size": 12, "color": WHITE}),
          ("Granger 선행성 ≠ ML 예측 유효성",
           {"size": 12, "bold": True, "color": "FFE08A"}),
          ("  (H3 부분 지지, 다음 슬라이드에서 더 명확히 확인)",
           {"size": 11, "color": SUB})])
setnotes(s, "검증 결과입니다. Granger 유의 채널을 선별해서 써도 전채널 대비 예측이 나아지지 않습니다. 즉 Granger 선행성이 ML 예측 유효성을 보장하지 않습니다. 이것이 이 연구의 두 번째 기여입니다.")

# =====================================================================
# Slide 8 — [2막-3] SHAP 표 (채널별 기여도 수치표)
# =====================================================================
s = add_slide()
header(s, "SHAP — 채널별 예측 기여도",
       "mean |SHAP| (Model A, XGBoost) · 인과성(Granger) ≠ 예측 기여(SHAP)")

# ── 표 레이아웃 ──────────────────────────────────────────────────────
# 열: 채널 | 설명 | Sneakers | Cards | LEGO | 핵심 포인트
col_w  = [1.05, 2.15, 1.45, 1.45, 1.45, 2.06]
col_x  = [0.45]
for cw in col_w[:-1]:
    col_x.append(col_x[-1] + cw)
hdr_h  = 0.44
row_h  = 0.55
y0     = 1.12

# 헤더행
hdr_labels = ["채널", "설명", "Sneakers", "Cards", "LEGO", "핵심 포인트"]
hdr_colors = [INK, INK, BLUE, BLUE, BLUE, RED]
for ci, (label, col_color) in enumerate(zip(hdr_labels, hdr_colors)):
    rect(s, col_x[ci], y0, col_w[ci], hdr_h, fill=col_color, line=WHITE, lw=1)
    tb(s, col_x[ci] + 0.05, y0, col_w[ci] - 0.08, hdr_h,
       [(label, {"size": 11, "bold": True, "color": WHITE,
                 "align": PP_ALIGN.CENTER})],
       anchor=MSO_ANCHOR.MIDDLE)

# 데이터 행
# (채널ID, 채널명, 설명, sn_val, ca_val, le_val, 핵심포인트, best_col 0-based in [sn,ca,le])
rows = [
    ("CH1", "Google Trends\n검색량",   "0.125", "0.048", "0.036", "Sneakers 최고 기여 · Cards Granger 유일 채널", 0),
    ("CH2", "뉴스 보도량\n(GDELT)",    "0.000", "0.031", "0.003", "전체적으로 기여 미미",                         1),
    ("CH3", "뉴스 감성\n(GDELT tone)", "0.011", "0.025", "0.000", "Granger 1위(F=13.4) → SHAP 최하위\n▶ 인과성 ≠ 예측 기여", -1),
    ("CH4", "YT 조회수",               "0.080", "0.091", "0.070", "Cards·LEGO 최고 기여\n3자산 공통 중요 채널",   1),
    ("CH5", "YT 댓글감성\n(FinBERT)",  "0.124", "0.004", "0.021", "Sneakers CH1과 동급 기여",                   0),
]

# best_col: 0=sneakers, 1=cards, 2=lego, -1=없음(강조)
vals_cols = [2, 3, 4]  # col_x indices for sn/ca/le
row_bgs   = [BOX2, BOX, BOX2, BOX, BOX2]

for ri, (ch, desc, sn, ca, le, note, best) in enumerate(rows):
    ry = y0 + hdr_h + ri * row_h
    bg = row_bgs[ri]
    vals = [sn, ca, le]

    # 채널 ID
    rect(s, col_x[0], ry, col_w[0], row_h, fill=INK if ri % 2 == 0 else "333333",
         line=WHITE, lw=1)
    tb(s, col_x[0] + 0.04, ry, col_w[0] - 0.06, row_h,
       [(ch, {"size": 11.5, "bold": True, "color": WHITE,
              "align": PP_ALIGN.CENTER})],
       anchor=MSO_ANCHOR.MIDDLE)

    # 설명
    rect(s, col_x[1], ry, col_w[1], row_h, fill=bg, line=LINE, lw=0.5)
    tb(s, col_x[1] + 0.08, ry, col_w[1] - 0.14, row_h,
       [(desc, {"size": 9.5, "color": BODY, "line_spacing": 1.1})],
       anchor=MSO_ANCHOR.MIDDLE)

    # 수치 3열
    for vi, (val, vci) in enumerate(zip(vals, [2, 3, 4])):
        is_best  = (vi == best)
        is_alert = (ch == "CH3")  # Granger 역설
        cell_bg  = ("D4EDDA" if is_best else
                    ("FADBD8" if is_alert else bg))
        cell_tc  = (GREEN   if is_best else
                    (RED     if is_alert and val == "0.011" else BODY))
        rect(s, col_x[vci], ry, col_w[vci], row_h,
             fill=cell_bg, line=LINE, lw=0.5)
        tb(s, col_x[vci] + 0.04, ry, col_w[vci] - 0.06, row_h,
           [(val, {"size": 12, "bold": is_best, "color": cell_tc,
                   "align": PP_ALIGN.CENTER})],
           anchor=MSO_ANCHOR.MIDDLE)

    # 핵심 포인트
    is_paradox = (ch == "CH3")
    rect(s, col_x[5], ry, col_w[5], row_h,
         fill="FADBD8" if is_paradox else bg, line=LINE, lw=0.5)
    tb(s, col_x[5] + 0.08, ry, col_w[5] - 0.14, row_h,
       [(note, {"size": 8.5,
                "color": RED if is_paradox else BODY,
                "bold": is_paradox, "line_spacing": 1.1})],
       anchor=MSO_ANCHOR.MIDDLE)

# 가격 래그 참고행
ry_ref = y0 + hdr_h + 5 * row_h + 0.06
tb(s, col_x[0], ry_ref, 9.1, 0.3,
   [("",{},[("참고  ",{"size":9.5,"bold":True,"color":INK}),
            ("price_vs_ma3 = 2.93 / price_chg_lag1 = 0.97  →  가격 모멘텀이 예측 분산의 대부분 흡수 (채널 대비 20~30배)",
             {"size":9.5,"color":BODY,"italic":True})])])

takeaway(s, 5.28,
         [("CH3(Granger 최강) → SHAP 최하위 · CH4(Granger 비유의) → 3자산 공통 최상위  ▶  ",
           {"size": 11.5, "color": WHITE}),
          ("선행성 ≠ 예측 기여",
           {"size": 11.5, "bold": True, "color": "FFE08A"})])
setnotes(s, "SHAP 표입니다. 가격 모멘텀(price_vs_ma3)이 2.93으로 압도적이고, 채널들의 기여는 0.1 수준입니다. 특히 Granger에서 F=13.4로 최강이었던 CH3(뉴스감성)이 SHAP에서는 0.011로 최하위입니다. 반면 Granger 비유의였던 CH4가 3자산 모두에서 중요 채널로 나타납니다. 인과성과 예측 기여는 다른 개념입니다.")

# =====================================================================
# Slide 9 — [2막-4] Ablation (상대 스케일 2패널)
# =====================================================================
s = add_slide()
header(s, "추가 실험 — 가격 래그의 역할과 Granger 채널 제거 효과",
       "채널만으로는 예측 불가 · Granger 유의 채널을 빼면 어떻게 되나?")

abl2_path = os.path.join(FIGD, "fig_ablation2.png")
pic(s, abl2_path, 0.35, 1.1, h=3.55)
caption(s, 0.3, 4.7, 9.4,
        "그림 3(좌) 4모형 AUC 전체 스케일 — 채널만(회색 빗금)은 무작위 수준  |  "
        "그림 3(우) Model A 대비 차이 — Cards에서 Granger 채널 제거 시 유의한 향상(★)")

takeaway(s, 5.18,
         [("① 채널만으로는 AUC≈0.5(무작위)  |  ② Granger 채널 제거 시 Cards AUC ",
           {"size": 11.5, "color": WHITE}),
          ("오히려 향상",
           {"size": 11.5, "bold": True, "color": "FFE08A"}),
          ("  (0.859→0.867, p=.004)  →  인과성 ≠ 예측력",
           {"size": 11.5, "color": WHITE})])
setnotes(s, "두 가지 핵심 실험입니다. 첫째, 채널만으로는 AUC 0.47~0.53으로 무작위 수준입니다. 가격 래그 없이 미디어 채널만으로는 예측이 불가능합니다. 둘째, Granger 유의 채널을 제거(A-dropGranger)하면 Cards에서 오히려 AUC가 향상됩니다. 이것이 핵심 역설입니다.")

# =====================================================================
# Slide 10 — 핵심 발견: Granger 선행성 ≠ ML 예측력
# =====================================================================
s = add_slide()
header(s, "핵심 발견 — Granger 선행 채널이 예측에는 '독'이 될 수 있다",
       "A-dropGranger 실험 · 선행성과 예측력은 서로 다른 차원의 개념")

# 중앙 역설 강조 박스
rect(s, 0.45, 1.1, 9.1, 1.3, fill="FFF3CD", line=GOLD, lw=2.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.45, 1.1, 0.08, 1.3, fill=GOLD)
tb(s, 0.65, 1.18, 8.8, 0.52,
   [("",{},[("역설  ",{"size":13,"bold":True,"color":GOLD}),
            ("Granger 검정에서 유의하게 선행하는 채널을 제거했더니 — 예측 성능이 오히려 향상됐다",
             {"size":13,"bold":True,"color":INK})])])
tb(s, 0.65, 1.7, 8.8, 0.62,
   [("",{},[("Cards:  Granger 유의 채널 CH1 포함(Model A) AUC=0.859  →  CH1 제거(A-dropGranger) AUC=",
             {"size":12,"color":BODY}),
            ("0.867  (DM p=0.004, 통계적 유의)",
             {"size":12,"bold":True,"color":RED})])])

# 3자산 결과 카드
cards_data = [
    ("스니커즈",
     "Granger 유의: CH1·CH3·CH4\nA-dropGranger AUC: 0.847\nModel A AUC:      0.841",
     "p=0.916 → 차이 없음\n채널 제거해도 성능 유지",
     GOLD, "무관"),
    ("트레이딩 카드",
     "Granger 유의: CH1\nA-dropGranger AUC: 0.867 ★\nModel A AUC:      0.859",
     "p=0.004 → 유의 향상!\nGranger 채널 제거 → 더 좋음",
     RED, "향상"),
    ("레고",
     "Granger 유의: 없음\nA-dropGranger = Model A\n(제거 채널 없음)",
     "비교 불가\n(Granger 유의채널 없음)",
     BLUE, "동일"),
]
yy3 = 2.54
for asset, nums, verdict, col, tag in cards_data:
    bx3 = 0.45 + cards_data.index((asset, nums, verdict, col, tag)) * 3.07
    rect(s, bx3, yy3, 2.97, 2.54, fill=BOX2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, bx3, yy3, 2.97, 0.44, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, bx3, yy3, 2.97, 0.44,
       [(asset, {"size": 12, "bold": True, "color": WHITE,
                 "align": PP_ALIGN.CENTER})],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, bx3 + 0.14, yy3 + 0.52, 2.7, 1.0,
       [(nums, {"size": 10, "color": BODY, "line_spacing": 1.2})])
    rect(s, bx3 + 0.14, yy3 + 1.56, 2.7, 0.013, fill=LINE)
    tb(s, bx3 + 0.14, yy3 + 1.64, 2.7, 0.78,
       [(verdict, {"size": 10, "bold": (col == RED), "color": col,
                   "line_spacing": 1.15})])

takeaway(s, 5.18,
         [("Granger 선행성 ≠ ML 예측력  — 선행성이 확인된 채널이 오히려 예측에 노이즈로 작용할 수 있음",
           {"size": 12, "bold": True, "color": "FFE08A"})])
setnotes(s, "핵심 발견입니다. Granger 검정에서 유의하게 선행하는 채널을 제거했더니 Cards에서는 오히려 AUC가 향상됐습니다(0.859→0.867, p=0.004). 즉 Granger 선행성이 확인된 채널이 ML 예측에서는 노이즈로 작용할 수 있습니다. Welch & Goyal(2008)의 'in-sample 유의성 ≠ out-of-sample 예측력' 개념과 일맥상통하는 발견입니다.")

# =====================================================================
# Slide 11 — [4막-1] 실무 함의
# =====================================================================
s = add_slide()
header(s, "결론 — 자산마다 미디어 신호가 다르게 작동한다",
       "3개 집단별 실무 시사점 · 일괄 적용 불가")

impls = [
    ("리셀 투자자·수집가",
     "• 카드: Google Trends를 1~2개월 선행 보조지표로 활용\n"
     "• 스니커즈: 뉴스 감성·YT 댓글 모니터링 가치 높음\n"
     "• 레고: 미디어보다 단종·재고 정보 우선"),
    ("리셀 플랫폼 운영자",
     "• 자산 유형별 맞춤 피처 구성 권장\n"
     "• Granger 선별 채널로 73% 아이템 예측력 유지\n"
     "• 채널 축소해도 성능 손실 없음 → 비용 효율"),
    ("학술 연구자",
     "• 공개 데이터 → Granger 선별 → XGBoost+SHAP+DM\n  파이프라인은 명품·시계 등으로 확장 가능\n"
     "• Granger 선행성 ≠ ML 예측력 — 두 개념 구분 필요"),
]
bx = 0.45
for title, desc in impls:
    rect(s, bx, 1.15, 2.97, 3.38, fill=BOX2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, bx, 1.15, 2.97, 0.5, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, bx, 1.15, 2.97, 0.5,
       [(title, {"size": 11.5, "bold": True, "color": WHITE,
                 "align": PP_ALIGN.CENTER})],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, bx + 0.16, 1.74, 2.65, 2.75,
       [(desc, {"size": 10.5, "color": BODY, "line_spacing": 1.18})])
    bx += 3.07

takeaway(s, 5.18,
         [("핵심: 미디어 선행성은 자산 유형마다 다름 — ",
           {"size": 12, "color": WHITE}),
          ("cards CH1이 3종 검정 모두 통과한 유일하고 견고한 신호",
           {"size": 12, "bold": True, "color": "FFE08A"})])
setnotes(s, "실무 시사점입니다. 자산별로 다른 채널 전략이 필요합니다.")

# =====================================================================
# Slide 13 — [4막-2] 한계 & 후속연구
# =====================================================================
s = add_slide()
header(s, "한계 & 후속 연구",
       "탐색적 발견으로 한정 · 확증 연구를 위한 발전 경로")

limits = [
    ("소표본 (N=48)",
     "2022년부터 확보 가능한 월별 체결가\n→ 결과는 탐색적, 확증 해석 지양"),
    ("Google Trends 결측\n(레고 2종)",
     "porsche·bugatti: 21~25/48 zero\n→ 레고 CH1 신호 품질 제한"),
    ("GDELT 감성 (CH3)",
     "NewsAPI 30일 제한으로 대체 사용\n→ FinBERT 대비 도메인 특화도 낮음"),
    ("가격 모멘텀 압도",
     "SHAP에서 price_vs_ma3가 채널의\n20~30배 → 채널 기여 과소 측정 위험"),
    ("단위 불일치",
     "Granger=자산, XGBoost=아이템\n→ 선별 채널 일반화 시 정보 손실"),
    ("sneakers_yeezy 보간",
     "4개월 선형보간 처리\n→ 논문 limitation 명시됨"),
]
px = [0.45, 3.45, 6.45, 0.45, 3.45, 6.45]
py = [1.12, 1.12, 1.12, 2.88, 2.88, 2.88]
for (title, desc), x0, y0 in zip(limits, px, py):
    rect(s, x0, y0, 2.85, 1.66, fill=BOX2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x0, y0, 0.05, 1.66, fill=RED)
    tb(s, x0 + 0.18, y0 + 0.1, 2.5, 0.36,
       [(title, {"size": 10.5, "bold": True, "color": INK, "line_spacing": 1.1})])
    tb(s, x0 + 0.18, y0 + 0.48, 2.5, 1.12,
       [(desc, {"size": 9.5, "color": BODY, "line_spacing": 1.15})])

tb(s, 0.45, 4.66, 9.1, 0.38, [
    ("향후:  ① 5년+ 패널  ② FinBERT CH3 교체  ③ CatBoost-XGBoost DM 검정  "
     "④ 사전 등록(pre-registration) 확증 연구",
     {"size": 10.5, "bold": True, "color": INK})])
takeaway(s, 5.18,
         [("한계를 투명하게 밝히고 BH·DH·VAR 삼각검증으로 스스로 검증 — 탐색적 발견으로 한정",
           {"size": 12, "color": WHITE})])
setnotes(s, "이상입니다. 감사합니다.")

# =====================================================================
# Reorder: 원본 슬라이드 제거, 새 슬라이드만 유지
# =====================================================================
all_ids = list(prs.slides._sldIdLst)
new_ids = [el for el in all_ids if el not in _orig_ids]

for el in _orig_ids:
    try:
        prs.part.drop_rel(el.get(qn("r:id")))
    except Exception:
        pass

lst = prs.slides._sldIdLst
for el in list(lst):
    lst.remove(el)
for el in new_ids:
    lst.append(el)

# =====================================================================
# Save
# =====================================================================
out = os.path.join(ROOT, "수업발표자료.pptx")
try:
    prs.save(out)
    print(f"SAVED: {out} | slides: {len(prs.slides)}")
except PermissionError:
    alt = os.path.join(ROOT, "수업발표자료_new.pptx")
    prs.save(alt)
    print(f"WARN: 파일 열려 있어 → {alt}")
