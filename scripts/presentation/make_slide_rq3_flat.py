#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RQ3 모델 비교 — 단일 슬라이드 (차이 작아 보이는 버전)
make_slide11.py와 동일 데이터, y축만 공통 0.5~1.0(무작위 기준)으로 펴서
실제 미미한 AUC 차이가 시각적으로도 작게 보이도록 재작도.
Usage : python scripts/make_slide_rq3_flat.py
Output: slide_rq3_model_comparison_flat.pptx
"""
import os, io
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RES  = os.path.join(BASE, "results")
OUT  = os.path.join(BASE, "slide_rq3_model_comparison_flat.pptx")

# ── 데이터 로드 (make_slide11.py와 동일) ──────────────────────────────────
abl = pd.read_csv(os.path.join(RES, "ablation_results.csv"))
mc  = pd.read_csv(os.path.join(RES, "model_c_results.csv"))
abl_auc = abl.pivot_table(index="model", columns="asset_type", values="auc")

def get(m, a):
    try:   return float(abl_auc.loc[m, a])
    except: return None

ASSETS = ["sneakers", "cards", "lego"]
A_auc  = {a: get("A", a)              for a in ASSETS}
C_auc  = {r["asset_type"]: r["auc_C"] for _, r in mc.iterrows()}
NA_auc = {a: get("A-dropGranger", a)  for a in ASSETS}
B_auc = {
    "sneakers": round((get("CH1-only","sneakers") +
                       get("CH3-only","sneakers") +
                       get("CH4-only","sneakers")) / 3, 4),
    "cards":    get("CH1-only","cards"),
    "lego":     get("A","lego"),
}
CH_models = ["CH1-only","CH2-only","CH3-only","CH4-only","CH5-only"]

# ── 색상 ──────────────────────────────────────────────────────────────────
C_A   = "#1F2D4E"
C_B   = "#E87722"
C_C   = "#2E6DA4"
C_NA  = "#7B3FA0"
C_CH  = ["#60B8E8","#70C87A","#F0A050","#F8C840","#F07878"]
ASSET_TITLE = {"sneakers": "스니커즈", "cards": "카드", "lego": "레고"}

# ★ 핵심 변경: 세 패널 공통 y축 (무작위 0.5 ~ 완벽 1.0)
Y_LO, Y_HI = 0.5, 1.0

plt.rcParams["font.family"] = ["Malgun Gothic", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False

fig, axes = plt.subplots(1, 3, figsize=(15.5, 5.2), sharey=True)
fig.patch.set_facecolor("white")

def draw_bar(ax, x, h, style, color, width=0.62):
    if style == "hollow":
        bars = ax.bar(x, h, width, color="none", edgecolor=color,
                      linewidth=2.0, zorder=3)
        for b in bars: b.set_linestyle("--")
    elif style == "ch":
        ax.bar(x, h, width, color=color, alpha=0.62,
               edgecolor="white", linewidth=0.5, zorder=3)
    else:
        ax.bar(x, h, width, color=color, edgecolor="white",
               linewidth=0.5, zorder=3)

def annotate(ax, x, y, txt, color, fs=8.0):
    ax.text(x, y + 0.008, txt, ha="center", va="bottom",
            fontsize=fs, color=color, fontweight="bold")

for ax, asset in zip(axes, ASSETS):
    ax.set_facecolor("#F9F9F9")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4, zorder=0)
    ax.set_ylim(Y_LO, Y_HI)

    if asset == "lego":
        main = [(0, A_auc[asset], "solid",  C_A), (1, C_auc[asset], "hollow", C_C)]
        labels = ["A=B=N-A","C","CH1","CH2","CH3","CH4","CH5"]
        for xi, h, st, col in main:
            draw_bar(ax, xi, h, st, col); annotate(ax, xi, h, f"{h:.3f}", col)
        for j, (ch, col) in enumerate(zip(CH_models, C_CH)):
            h = get(ch, asset); draw_bar(ax, 2+j, h, "ch", col)
            annotate(ax, 2+j, h, f"{h:.3f}", "#555555", fs=7)
        ax.set_xticks(range(7)); ax.set_xticklabels(labels, fontsize=9, fontweight="bold")
        ax.annotate("레고: Granger 유의 채널 없음 → A=B=N-A",
                    xy=(0.5, 0.04), xycoords="axes fraction", ha="center",
                    fontsize=8, color="#888888", style="italic")
    else:
        main = [(0, A_auc[asset], "solid",  C_A),
                (1, B_auc[asset], "hollow", C_B),
                (2, C_auc[asset], "hollow", C_C),
                (3, NA_auc[asset],"hollow", C_NA)]
        bl = "B≈" if asset == "sneakers" else "B"
        labels = ["A", bl, "C", "N-A", "CH1","CH2","CH3","CH4","CH5"]
        for xi, h, st, col in main:
            draw_bar(ax, xi, h, st, col); annotate(ax, xi, h, f"{h:.3f}", col)
        for j, (ch, col) in enumerate(zip(CH_models, C_CH)):
            h = get(ch, asset); draw_bar(ax, 4+j, h, "ch", col)
            annotate(ax, 4+j, h, f"{h:.3f}", "#555555", fs=7)
        ax.set_xticks(range(9)); ax.set_xticklabels(labels, fontsize=9, fontweight="bold")
        ax.axvline(3.5, color="#BBBBBB", linewidth=0.9, linestyle=":", zorder=4)

    # A 기준선 + 무작위 수준선
    ax.axhline(A_auc[asset], color=C_A, linestyle="--", linewidth=0.9, alpha=0.5, zorder=4)
    ax.axhline(0.5, color="#C0392B", linestyle="--", linewidth=1.0, alpha=0.7, zorder=2)
    ax.set_title(ASSET_TITLE[asset], fontsize=14, fontweight="bold", pad=6)
    if asset == "sneakers":
        ax.set_ylabel("AUC-ROC", fontsize=10)
        ax.annotate("무작위 수준 0.5", xy=(0.05, 0.505), fontsize=7.5,
                    color="#C0392B", va="bottom")

legend_handles = [
    mpatches.Patch(facecolor=C_A, edgecolor="white", label="A  전채널"),
    mpatches.Patch(facecolor="none", edgecolor=C_B, linestyle="--", linewidth=1.5, label="B  Granger 선별"),
    mpatches.Patch(facecolor="none", edgecolor=C_C, linestyle="--", linewidth=1.5, label="C  SHAP 선별"),
    mpatches.Patch(facecolor="none", edgecolor=C_NA, linestyle="--", linewidth=1.5, label="N-A  유의채널 제거"),
    mpatches.Patch(facecolor=C_CH[0], alpha=0.7, label="CH1~5 단독"),
]
fig.legend(handles=legend_handles, loc="lower center", ncol=5,
           fontsize=9.5, bbox_to_anchor=(0.5, -0.02),
           framealpha=0.97, edgecolor="#cccccc")
fig.suptitle("모델별 AUC-ROC 비교  (공통 축 0.5~1.0 · 0.5=무작위 기준)",
             fontsize=13, fontweight="bold", y=1.0)

plt.tight_layout(rect=[0, 0.07, 1, 0.97])
fig_buf = io.BytesIO()
fig.savefig(fig_buf, dpi=160, bbox_inches="tight", facecolor="white")
plt.close(fig)

# ── PPT (단일 슬라이드) ────────────────────────────────────────────────────
NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1E, 0x7E, 0x34)

def add_box(sld, l, t, w, h, fill=None):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else:    tb.fill.background()
    tb.line.fill.background()
    return tb

def add_text(sld, text, l, t, w, h, bold=False, italic=False,
             size=11, color=RGBColor(0x0A,0x0A,0x0A), align=PP_ALIGN.LEFT):
    tb = add_box(sld, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.bold=bold; r.font.italic=italic
    r.font.size=Pt(size); r.font.color.rgb=color; r.font.name="Calibri"
    return tb

prs  = Presentation()
prs.slide_width  = Inches(13.33); prs.slide_height = Inches(7.5)
sld  = prs.slides.add_slide(prs.slide_layouts[6])

hbar = add_box(sld, 0, 0, 13.33, 0.78, fill=NAVY)
tf = hbar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  RQ3 모델별 예측 성능 — 채널을 어떻게 바꿔도 거의 동일"
r.font.bold=True; r.font.size=Pt(20); r.font.color.rgb=WHITE; r.font.name="Calibri"
add_box(sld, 0, 0.78, 13.33, 0.025, fill=BLUE)
add_text(sld,
    "A=전채널 | B=Granger 선별 | C=SHAP 선별 | N-A=유의채널 제거 | CH1~5=단독  "
    "—  무작위(0.5) 기준 공통 축에서 모든 구성의 막대 높이가 사실상 동일",
    0.2, 0.82, 13.0, 0.32, size=9.5, color=DGRAY, italic=True)

fig_buf.seek(0)
sld.shapes.add_picture(fig_buf, Inches(0.13), Inches(1.30), Inches(13.07), Inches(5.05))

add_text(sld,
    "핵심:  모든 모델 AUC 0.84~0.96 — 최대 차이 0.01 내외(DM p>0.5).  "
    "채널을 더하든 빼든, 선별 기준을 Granger→SHAP로 바꾸든 예측력은 통계적으로 구분되지 않음.",
    0.2, 6.55, 13.0, 0.75, size=12, color=GREEN, bold=True)

prs.save(OUT)
print(f"[OK] {OUT}")
