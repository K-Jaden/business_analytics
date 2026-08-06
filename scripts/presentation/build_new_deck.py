# -*- coding: utf-8 -*-
"""
처음부터 새로 만드는 발표자료 (교수님 피드백 전면 반영)
출력: 발표자료_최종_피드백반영.pptx

피드백 반영 원칙
  - 슬라이드당 글씨 최소화, 핵심만, 인포그래픽/도형 위주
  - 모든 분석 도구마다 '왜 이 도구/검정을 썼는가' 명시 (ADF·FinBERT·Granger·BH·XGBoost·SHAP·DM·IRF·Jaccard)
  - 채널별 스케일이 다른 이유 + z-score 통일(전처리 결과) 표기
  - 분석 대상(아이템) 슬라이드는 자산별 대표·출처 위주로 간소화
  - 모든 그래프에 '그래프 해석' 캡션
  - 모델설계는 ㅁ→ㅁ→ㅁ 플로우 도형
내용 정본: 최종_버전_수정.pptx (48개월·Model A/B/C·DM 확정)
구조 참고: 발표자료_최종 (중간).pptx (왜 이 도구·그래프 해석 흐름)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

import os
OUT = os.environ.get("DECK_OUT", "발표자료_최종_피드백반영.pptx")

# ---- 팔레트 ----
NAVY="1F2D4E"; BLUE="2E6DA4"; GREEN="277A4A"; PURPLE="6C3D91"
ORANGE="C8772E"; GRAY="555555"; INK="111111"
LBLUE="EAF1F8"; LGREEN="E7F2EC"; LPUR="EFEAF5"; LORANGE="FBEEDD"
LRED="F8D7DA"; DRED="B23B3B"; LGRAY="F2F3F5"; MGRAY="9AA5B1"; WHITE="FFFFFF"
FONT="맑은 고딕"
EMU=914400

def C(s): return RGBColor.from_string(s)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)

def _font(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', name)

def no_shadow(shp):
    try: shp.shadow.inherit=False
    except Exception: pass

def box(s, l,t,w,h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=s.shapes.add_shape(shape, Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
    no_shadow(sp)
    return sp

def tb(s, l,t,w,h):
    return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))

def settext(shp, paras, anchor=MSO_ANCHOR.TOP, wrap=True, ml=5,mr=5,mt=3,mb=3):
    """paras: list of (runs, align, space_after); runs: list of (text,size,color,bold)"""
    tf=shp.text_frame; tf.clear(); tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=Pt(ml); tf.margin_right=Pt(mr); tf.margin_top=Pt(mt); tf.margin_bottom=Pt(mb)
    first=True
    for runs,align,space in paras:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.space_before=Pt(0)
        if space is not None: p.space_after=Pt(space)
        for text,size,color,bold in runs:
            r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold
            r.font.color.rgb=C(color); _font(r)
    return shp

def arrow(s,l,t,w,h,fill=MGRAY,shape=MSO_SHAPE.RIGHT_ARROW):
    a=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb=C(fill); a.line.fill.background(); no_shadow(a)
    return a

def title_bar(s, title, subtitle=None, pageno=None):
    bar=box(s,0,0,13.333,0.62, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    settext(bar,[([(title,21,WHITE,True)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE,ml=18)
    box(s,0,0.62,13.333,0.04, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
    if subtitle:
        st=tb(s,0.2,0.70,13.0,0.34)
        settext(st,[([(subtitle,11,GRAY,False)],PP_ALIGN.LEFT,0)])
    if pageno:
        pn=tb(s,12.5,7.16,0.75,0.3)
        settext(pn,[([(pageno,9,GRAY,False)],PP_ALIGN.RIGHT,0)])

def divider(s, num, name, sub):
    box(s,0,0,13.333,7.5, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    box(s,0.9,3.0,0.18,1.5, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
    n=tb(s,1.3,2.55,3.0,1.4); settext(n,[([(num,72,BLUE,True)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE)
    nm=tb(s,3.4,2.95,8.8,0.9); settext(nm,[([(name,34,WHITE,True)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE)
    sb=tb(s,3.45,3.95,8.8,0.5); settext(sb,[([(sub,14,"BBC4D0",False)],PP_ALIGN.LEFT,0)])

def fit_image(s, path, l,t,bw,bh, halign='center', valign='middle'):
    pic=s.shapes.add_picture(path, Inches(l),Inches(t))
    w,h=pic.width,pic.height
    sc=min(Inches(bw)/w, Inches(bh)/h)
    pic.width=int(w*sc); pic.height=int(h*sc)
    if halign=='center': pic.left=Inches(l)+(Inches(bw)-pic.width)//2
    elif halign=='left': pic.left=Inches(l)
    if valign=='middle': pic.top=Inches(t)+(Inches(bh)-pic.height)//2
    elif valign=='top': pic.top=Inches(t)
    return pic

def caption(s, l,t,w, text, color=GRAY):
    c=tb(s,l,t,w,0.5)
    settext(c,[([("그래프 해석  ",10,BLUE,True),(text,10,color,False)],PP_ALIGN.LEFT,0)])
    return c

def chip(s,l,t,w,h,label,fill,tcolor=WHITE,size=11,bold=True):
    b=box(s,l,t,w,h,fill=fill)
    settext(b,[([(label,size,tcolor,bold)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
    return b

# “왜 이 도구?” 공통 박스
def why_box(s,l,t,w,h, head, lines, accent=BLUE, lfill=LBLUE):
    b=box(s,l,t,w,h,fill=lfill,line=accent,lw=1.25)
    paras=[([("왜 "+head+"인가?",12.5,accent,True)],PP_ALIGN.LEFT,4)]
    for ln in lines:
        paras.append(([(ln,11.5,INK,False)],PP_ALIGN.LEFT,2))
    settext(b,paras,anchor=MSO_ANCHOR.TOP,ml=10,mt=8)
    return b

PG=[0]
def pg():
    PG[0]+=1
    return "%d"%PG[0]

# ============================================================
# 1. Title
# ============================================================
def s_title():
    s=slide()
    box(s,0,0,13.333,7.5,fill=NAVY,shape=MSO_SHAPE.RECTANGLE)
    box(s,0,5.0,13.333,0.06,fill=BLUE,shape=MSO_SHAPE.RECTANGLE)
    t=tb(s,0.9,2.45,11.5,1.6)
    settext(t,[
        ([("리셀 시장에서 미디어 채널은 가격을 선행하는가?",33,WHITE,True)],PP_ALIGN.LEFT,6),
    ])
    sub=tb(s,0.95,3.9,11.5,0.5)
    settext(sub,[([("스니커즈 · 트레이딩 카드 · 레고 — 다중자산 Granger 인과분석 & XGBoost 검증",16,"BBC4D0",False)],PP_ALIGN.LEFT,0)])
    meta=tb(s,0.95,5.2,11.5,0.4)
    settext(meta,[([("2022.01 – 2025.12  |  48개월 × 15개 아이템 × 5개 미디어 채널",13,"8FA0B5",False)],PP_ALIGN.LEFT,0)])
    au=tb(s,0.95,6.4,11.5,0.4)
    settext(au,[([("이승현 · 윤재은 · 최형호",14,WHITE,False)],PP_ALIGN.LEFT,0)])

# ============================================================
# 2. 목차
# ============================================================
def s_toc():
    s=slide()
    title_bar(s,"목차","발표 흐름 한눈에 보기 — 6개 섹션")
    items=[
        ("01","연구 배경 & 질문","리셀 시장 · 연구문제 RQ1–RQ3",BLUE,LBLUE),
        ("02","기존 연구 & 한계","선행연구의 공통 한계 · 본 연구의 차별점",GREEN,LGREEN),
        ("03","분석 대상 & 채널","3종 자산 · 5개 미디어 채널 · 전처리",PURPLE,LPUR),
        ("04","분석 방법론 (왜 이 도구?)","FinBERT · Granger · XGBoost · SHAP · DM",ORANGE,LORANGE),
        ("05","실험 결과","RQ1–RQ3 · SHAP · IRF · 한계",BLUE,LBLUE),
        ("06","결론 & 시사점","연구문제별 결론 · 실무 함의 · 향후 연구",GREEN,LGREEN),
    ]
    y=1.35; h=0.86; gap=0.07
    for num,name,desc,acc,lf in items:
        b=box(s,0.7,y,11.9,h,fill=lf,line="DDDDDD",lw=0.75)
        box(s,0.7,y,0.12,h,fill=acc,shape=MSO_SHAPE.RECTANGLE)
        nn=tb(s,0.95,y,1.0,h); settext(nn,[([(num,26,acc,True)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE)
        nm=tb(s,2.0,y+0.08,9.5,0.45); settext(nm,[([(name,16,INK,True)],PP_ALIGN.LEFT,0)])
        dd=tb(s,2.0,y+0.46,9.5,0.35); settext(dd,[([(desc,11,GRAY,False)],PP_ALIGN.LEFT,0)])
        y+=h+gap

# ============================================================
# 4. 연구 배경 + RQ 인포그래픽
# ============================================================
def s_rq():
    s=slide()
    title_bar(s,"연구 배경 & 연구 문제 (RQ1–RQ3)",
              "리셀 시장은 미디어에 민감하나 다채널 선행성 비교는 미답 — 세 단계로 해부",pg())
    # 배경 한 줄 띠
    bgt=box(s,0.3,1.15,12.7,0.55,fill=NAVY)
    settext(bgt,[([("리셀 시장(’24 美 중고 약 490억$, 전년 +14%)은 SNS·뉴스·유튜브에 민감 → ",11.5,WHITE,False),
                   ("어떤 미디어가 ‘가격보다 먼저’ 움직이는가?",11.5,"FFD9A0",True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)
    cards=[
        dict(x=0.30,acc=BLUE,lf=LBLUE,badge="RQ1",stage="시간적 선행성",
             q="어떤 미디어 채널이\n가격을 선행하는가?",
             test="Granger 인과검정",
             why="회귀는 ‘관계’만 봄.\n방향·시차를 함께 보려면\nGranger가 필요",
             metric="F 통계량 · p · 래그"),
        dict(x=4.62,acc=GREEN,lf=LGREEN,badge="RQ2",stage="채널 이질성",
             q="선행 채널 구성이\n자산마다 다른가?",
             test="유의채널 집합 · Jaccard",
             why="자산별 채널 차이를\n한 숫자로 정량화하려고\nJaccard 사용",
             metric="Jaccard 유사도"),
        dict(x=8.94,acc=PURPLE,lf=LPUR,badge="RQ3",stage="예측 간결성",
             q="선별 채널 모델이\n전채널과 동등한가?",
             test="XGBoost A/B/C · DM검정",
             why="AUC 숫자차가 우연인지\n실제인지 통계로 보려면\nDM 검정이 필요",
             metric="AUC · ΔAUC · DM p"),
    ]
    cw=3.55; ctop=1.95; chh=4.7
    for c in cards:
        x=c["x"]
        box(s,x,ctop,cw,chh,fill=WHITE,line="DDDDDD",lw=1.0)
        hd=box(s,x,ctop,cw,0.9,fill=c["acc"])
        settext(hd,[([(c["badge"]+"   ",20,WHITE,True),(c["stage"],14,WHITE,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        q=tb(s,x+0.15,ctop+1.0,cw-0.3,0.85)
        settext(q,[([(ln,13,INK,True)],PP_ALIGN.CENTER,2) for ln in c["q"].split("\n")],anchor=MSO_ANCHOR.MIDDLE)
        tbx=box(s,x+0.15,ctop+1.9,cw-0.3,0.5,fill=c["lf"],line=c["acc"],lw=1.0)
        settext(tbx,[([("검정  ",10,c["acc"],True),(c["test"],11,INK,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        wy=tb(s,x+0.15,ctop+2.5,cw-0.3,1.2)
        rr=[([("왜 이 검정?",10.5,c["acc"],True)],PP_ALIGN.LEFT,2)]
        for ln in c["why"].split("\n"): rr.append(([(ln,10.5,GRAY,False)],PP_ALIGN.LEFT,0))
        settext(wy,rr)
        mt=box(s,x,ctop+chh-0.55,cw,0.55,fill=LGRAY,shape=MSO_SHAPE.RECTANGLE)
        settext(mt,[([("측정  ",10,c["acc"],True),(c["metric"],10.5,INK,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
    arrow(s,4.12,ctop+1.9,0.48,0.8,fill=BLUE)
    arrow(s,8.44,ctop+1.9,0.48,0.8,fill=GREEN)

# ============================================================
# 6. 기존 연구 & 공통 한계
# ============================================================
def s_related():
    s=slide()
    title_bar(s,"기존 연구 & 공통 한계",
              "세 연구 흐름 모두 ‘리셀 × 다채널 × 선행성 비교’의 교점은 비워 둠 — 한계가 핵심",pg())
    cols=[
        dict(x=0.30,acc=BLUE,title="미디어 → 주식 선행성",
             body="Tetlock (2007) · 뉴스 비관론 → S&P500\nDa et al. (2011) · Google 검색량 → 주가 선행",
             lim="전통 금융자산(주식)에만 적용"),
        dict(x=4.62,acc=GREEN,title="리셀 상품 가격 예측",
             body="Campello et al. (2021) · StockX 스니커즈 예측\nDobrynskaya & Kishilova (2023) · 레고 대체자산",
             lim="미디어 채널 신호 미포함 · 단일 자산"),
        dict(x=8.94,acc=PURPLE,title="멀티채널 예측 방법론",
             body="Bollen et al. (2011) · Twitter 감성 → DJIA 선행\nZhang et al. (2018) · 뉴스+감성+가격 결합",
             lim="주식 단일 · 채널 비교 없음 · Granger 미적용"),
    ]
    cw=3.55; ctop=1.2
    for c in cols:
        x=c["x"]
        box(s,x,ctop,cw,5.05,fill=WHITE,line="DDDDDD",lw=1.0)
        hd=box(s,x,ctop,cw,0.55,fill=c["acc"])
        settext(hd,[([(c["title"],13.5,WHITE,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        bd=tb(s,x+0.15,ctop+0.7,cw-0.3,2.1)
        settext(bd,[([(ln,11.5,INK,False)],PP_ALIGN.LEFT,6) for ln in c["body"].split("\n")])
        lim=box(s,x+0.15,ctop+2.95,cw-0.3,1.85,fill=LRED,line=DRED,lw=1.25)
        settext(lim,[([("⚠ 한계",13,DRED,True)],PP_ALIGN.LEFT,5),
                     ([(c["lim"],12.5,INK,True)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE,ml=8)
    bottom=box(s,0.3,6.45,12.7,0.7,fill=NAVY)
    settext(bottom,[([("본 연구  ",13,"FFD9A0",True),
                      ("스니커즈·카드·레고 3개 자산 × 5채널 × Granger × SHAP+DM → 세 흐름의 빈 교점을 메움",12.5,WHITE,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=14)

# ============================================================
# 7. 차별점
# ============================================================
def s_diff():
    s=slide()
    title_bar(s,"기존 연구와의 차별점","무엇이 새로운가 — 4가지 축",pg())
    rows=[
        ("분석 자산","주식 단일 자산","리셀 수집품 3종 동시 비교"),
        ("채널 수","단일 채널 (트위터 or 구글)","다중 채널 5개 (볼륨+감성 분리)"),
        ("채널 선정","사전 고정","Granger 검정으로 자산별 탐색"),
        ("분석 방향","‘선행하는가 Yes/No’","XGBoost·SHAP로 자산별 기여 정량화"),
    ]
    # 헤더
    y=1.35
    hd1=box(s,0.5,y,3.0,0.6,fill=NAVY); settext(hd1,[([("구분",13,WHITE,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
    hd2=box(s,3.6,y,4.6,0.6,fill="8A99AD"); settext(hd2,[([("기존 연구",13,WHITE,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
    hd3=box(s,8.3,y,4.5,0.6,fill=GREEN); settext(hd3,[([("제안 연구 (본 연구)",13,WHITE,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.72
    for k,old,new in rows:
        c1=box(s,0.5,y,3.0,0.95,fill=LGRAY,line="DDDDDD",lw=0.75)
        settext(c1,[([(k,12.5,INK,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        c2=box(s,3.6,y,4.6,0.95,fill=WHITE,line="DDDDDD",lw=0.75)
        settext(c2,[([(old,12,GRAY,False)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        c3=box(s,8.3,y,4.5,0.95,fill=LGREEN,line=GREEN,lw=1.0)
        settext(c3,[([(new,12,INK,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        arrow(s,8.0,y+0.27,0.35,0.4,fill=GREEN)
        y+=1.05

def find_last(s):
    return s.shapes[-1]

# ============================================================
# 9. 분석 대상 (간소화)
# ============================================================
def s_items():
    s=slide()
    title_bar(s,"분석 대상: 3종 리셀 자산",
              "자산별 대표 5개 아이템 · 2022.01–2025.12 (48개월) · 공개 실거래가",pg())
    assets=[
        dict(x=0.30,acc=BLUE,lf=LBLUE,name="스니커즈",src="StockX · US Sz10 평균가",
             items=["Jordan 1 Bordeaux","Nike Dunk Low Panda","Yeezy 350 V2 Zebra","Travis Scott Jordan 1","New Balance 550"]),
        dict(x=4.62,acc=GREEN,lf=LGREEN,name="트레이딩 카드",src="PriceCharting · PSA 10",
             items=["Charizard VMAX (SV107)","Umbreon VMAX Alt","Rayquaza VMAX Alt","Pikachu VMAX","Charizard GX Hidden"]),
        dict(x=8.94,acc=PURPLE,lf=LPUR,name="레고",src="BrickRanker · 새 봉인가",
             items=["Millennium Falcon","Hogwarts Castle","Titanic 10294","Porsche 911","Bugatti Chiron"]),
    ]
    cw=3.55; ctop=1.3; chh=4.9
    for a in assets:
        x=a["x"]
        box(s,x,ctop,cw,chh,fill=WHITE,line="DDDDDD",lw=1.0)
        hd=box(s,x,ctop,cw,0.95,fill=a["acc"])
        settext(hd,[([(a["name"],18,WHITE,True)],PP_ALIGN.CENTER,3),
                    ([(a["src"],10.5,"E6ECF3",False)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        yy=ctop+1.15
        for it in a["items"]:
            dot=box(s,x+0.25,yy+0.12,0.12,0.12,fill=a["acc"],shape=MSO_SHAPE.OVAL)
            li=tb(s,x+0.5,yy,cw-0.7,0.55)
            settext(li,[([(it,13,INK,False)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE)
            yy+=0.66
    note=tb(s,0.3,6.45,12.7,0.4)
    settext(note,[([("모든 아이템 48/48 커버리지 · 호가가 아닌 체결가 사용 (선행연구 다수 채택)",10.5,GRAY,False)],PP_ALIGN.LEFT,0)])

# ============================================================
# 10. 5채널 + 전처리·정상성
# ============================================================
def s_channels():
    s=slide()
    title_bar(s,"5개 미디어 채널 & 전처리",
              "볼륨 3 + 감성 2 · 채널마다 단위가 달라 스케일이 다름 → 전처리로 통일",pg())
    # 채널 표 (좌)
    rows=[
        ("ID","채널","유형","전처리 → 최종 스케일"),
        ("CH1","Google Trends","볼륨","월평균 → 0–100"),
        ("CH2","뉴스 보도량","볼륨","GDELT → 0–1 정규화"),
        ("CH3","뉴스 감성","감성","GDELT tone → −1~+1"),
        ("CH4","유튜브 조회수","볼륨","log1p → z-score"),
        ("CH5","유튜브 댓글 감성","감성","FinBERT → −1~+1"),
    ]
    gt=s.shapes.add_table(6,4,Inches(0.3),Inches(1.3),Inches(7.7),Inches(3.5)).table
    widths=[0.7,2.3,1.0,3.7]
    for i,w in enumerate(widths): gt.columns[i].width=Inches(w)
    for ri in range(6):
        for ci in range(4):
            cell=gt.cell(ri,ci); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Pt(5);cell.margin_right=Pt(3);cell.margin_top=Pt(2);cell.margin_bottom=Pt(2)
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER if ci in (0,2) else PP_ALIGN.LEFT
            r=p.add_run(); r.text=rows[ri][ci]; _font(r)
            if ri==0:
                r.font.size=Pt(11);r.font.bold=True;r.font.color.rgb=C(WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb=C(NAVY)
            else:
                r.font.size=Pt(11.5);r.font.bold=(ci==0);r.font.color.rgb=C(INK)
                cell.fill.solid(); cell.fill.fore_color.rgb=C(WHITE if ri%2 else LGRAY)
    # 왜 스케일 다른가 (우상)
    why_box(s,8.2,1.3,4.8,1.55,"채널마다 스케일이 다른가? → 통일",
            ["검색량·보도량·감성은 서로 다른 단위.",
             "→ XGBoost 투입 직전 StandardScaler로",
             "    z-score 통일 (전처리한 수치 사용)"],accent=BLUE,lfill=LBLUE)
    # 왜 ADF (우하)
    why_box(s,8.2,3.0,4.8,1.85,"정상성(ADF) 검정",
            ["비정상 시계열은 ‘허구 인과’ 위험",
             "→ Granger의 전제조건 (반드시 선검정)",
             "결과: 정상 9개 그대로 · 비정상 9개 차분",
             "예외: 스니커즈 CH1 차분 후 p=0.120(경계)"],accent=GREEN,lfill=LGREEN)
    # 하단: 볼륨/감성 분리 이유
    sep=box(s,0.3,5.05,7.7,1.05,fill=LORANGE,line=ORANGE,lw=1.0)
    settext(sep,[([("왜 볼륨과 감성을 분리하나?",12,ORANGE,True)],PP_ALIGN.LEFT,3),
                 ([("같은 보도량도 긍정/부정에 따라 가격 방향이 다름 — 합치면 신호가 상쇄됨 (Mao 2011 표준)",11.5,INK,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=10)
    note=tb(s,0.3,6.25,12.7,0.5)
    settext(note,[([("제한점  CH3는 GDELT tone(자동추출) 사용 — FinBERT 대비 도메인 특화도 낮음 · CH4는 ‘resell+review’ 키워드로 일반 hype 과소 가능",10,GRAY,False)],PP_ALIGN.LEFT,0)])

# ============================================================
# 12. 분석 파이프라인 플로우
# ============================================================
def s_pipeline():
    s=slide()
    title_bar(s,"분석 파이프라인",
              "무엇을 모으고 → 어떻게 수치화하고 → 선행 채널을 찾고 → 예측력을 검증하는가",pg())
    steps=[
        dict(x=0.30,fill=NAVY,t="STEP 1\n데이터 수집",sub="StockX·PriceCharting\nGoogle·GDELT·YouTube"),
        dict(x=3.55,fill=ORANGE,t="STEP 2\n감성 수치화",sub="FinBERT로 댓글 감성\n→ −1~+1 점수"),
        dict(x=6.80,fill=BLUE,t="STEP 3\nGranger 인과",sub="3자산×5채널\n선행 채널 식별 (RQ1·2)"),
        dict(x=10.05,fill=GREEN,t="STEP 4\nXGBoost·DM",sub="A/B/C 성능 비교\n+ SHAP·DM 검증 (RQ3)"),
    ]
    ft=1.5; fh=1.7
    for st in steps:
        b=box(s,st["x"],ft,2.75 if st["x"]<10 else 2.95,fh,fill=st["fill"])
        settext(b,[([(ln,15,WHITE,True)],PP_ALIGN.CENTER,2) for ln in st["t"].split("\n")]+
                  [([(ln,10.5,WHITE,False)],PP_ALIGN.CENTER,0) for ln in st["sub"].split("\n")],anchor=MSO_ANCHOR.MIDDLE)
    for ax in [3.07,6.32,9.57]:
        arrow(s,ax,ft+0.55,0.46,0.6,fill=MGRAY)
    # 입력 패널 결과 띠
    panel=box(s,0.3,3.7,12.7,0.7,fill=LGRAY,line="DDDDDD",lw=0.75)
    settext(panel,[([("입력 패널 완성  ",12,NAVY,True),
                     ("볼륨 3(검색·뉴스량·조회수) + 감성 2(뉴스·댓글) × 15아이템 × 48개월",12,INK,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)
    # 다음 3장 안내(왜 이 도구)
    guide=box(s,0.3,4.65,12.7,1.9,fill=WHITE,line=ORANGE,lw=1.25)
    settext(guide,[([("다음 3장 — ‘왜 이 도구인가?’",13,ORANGE,True)],PP_ALIGN.LEFT,6),
                   ([("①  FinBERT  ",12,INK,True),("금융 맥락 감성 방향 포착에 특화 (일반 BERT 대비 정확)",12,GRAY,False)],PP_ALIGN.LEFT,4),
                   ([("②  Granger  ",12,INK,True),("회귀는 ‘관계’만 — 방향·시차를 함께 보려면 Granger",12,GRAY,False)],PP_ALIGN.LEFT,4),
                   ([("③  XGBoost·SHAP·DM  ",12,INK,True),("소표본에 강함 · 채널 기여 해석 · 성능차의 통계 검정",12,GRAY,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.TOP,ml=14,mt=10)

# ============================================================
# 13. FinBERT
# ============================================================
def s_finbert():
    s=slide()
    title_bar(s,"FinBERT — 감성 점수화 (왜 이 도구인가?)",
              "금융 텍스트 특화 BERT (Araci 2019) · 유튜브 댓글 감성(CH5) 산출에 사용",pg())
    # 플로우 4단계
    steps=[("①  문장 입력","뉴스·유튜브 댓글"),
           ("②  토큰화·벡터화","BERT Transformer"),
           ("③  감성 확률","P(pos)·P(neu)·P(neg)"),
           ("④  감성 점수","P(pos)−P(neg) ∈ [−1,+1]")]
    x=0.3; w=3.05; t=1.35; h=1.25
    for i,(a,b) in enumerate(steps):
        bx=box(s,x,t,w,h,fill=LBLUE,line=BLUE,lw=1.0)
        settext(bx,[([(a,13,BLUE,True)],PP_ALIGN.CENTER,3),([(b,11,INK,False)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        if i<3: arrow(s,x+w+0.02,t+0.42,0.26,0.4,fill=MGRAY)
        x+=w+0.28
    # 왜 FinBERT
    why_box(s,0.3,2.9,6.2,1.9,"FinBERT",
            ["“가격 급등 예상”·“공급 부족 우려” 등",
             "금융 맥락의 감성 방향 포착에 특화",
             "→ 일반 BERT보다 금융 텍스트 분류 정확",
             "    (Huang et al. 2023, CAR)"],accent=BLUE,lfill=LBLUE)
    # 예시
    ex=box(s,6.7,2.9,6.3,1.9,fill=WHITE,line="DDDDDD",lw=1.0)
    settext(ex,[([("예시",12,NAVY,True)],PP_ALIGN.LEFT,5),
                ([("“limited release drives huge demand”   ",12,INK,False),("+0.85",13,GREEN,True)],PP_ALIGN.LEFT,4),
                ([("“resale prices collapse”   ",12,INK,False),("−0.72",13,DRED,True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)
    # 적용/제한
    ap=box(s,0.3,5.0,12.7,1.4,fill=LORANGE,line=ORANGE,lw=1.0)
    settext(ap,[([("적용 대상  ",12,ORANGE,True),("유튜브 댓글 감성(CH5) — 아이템당 상위 15영상 × 30댓글, 영어만",12,INK,False)],PP_ALIGN.LEFT,4),
                ([("제한점  ",12,ORANGE,True),("뉴스 감성(CH3)은 NewsAPI 역사데이터 제약으로 GDELT tone 대체 — 도메인 특화도가 낮은 점은 한계로 명시",11.5,GRAY,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)

# ============================================================
# 14. Granger
# ============================================================
def s_granger():
    s=slide()
    title_bar(s,"Granger 인과검정 — 선행성 판별 (왜 이 도구인가?)",
              "‘미디어가 먼저 움직이면 가격도 따라오는가’를 두 모델 비교로 검정",pg())
    # 핵심 아이디어
    idea=box(s,0.3,1.2,12.7,0.85,fill=NAVY)
    settext(idea,[([("핵심  ",12.5,"FFD9A0",True),
                    ("미디어의 과거값을 알면 가격의 미래를 더 잘 맞히는가? → YES면 그 채널이 가격을 선행한다",13,WHITE,True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=14)
    # 두 모델 비교 (ㅁ vs ㅁ → 비교)
    A=box(s,0.3,2.3,5.0,1.7,fill=LGRAY,line="DDDDDD",lw=1.0)
    settext(A,[([("모델 A · 기준",13,NAVY,True)],PP_ALIGN.LEFT,4),
               ([("이번달 가격 = 지난달·지지난달 가격",12,INK,False)],PP_ALIGN.LEFT,3),
               ([("미디어 없이 가격 관성만 이용",11,GRAY,False)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE,ml=10)
    B=box(s,5.6,2.3,5.0,1.7,fill=LBLUE,line=BLUE,lw=1.25)
    settext(B,[([("모델 B · 비교",13,BLUE,True)],PP_ALIGN.LEFT,4),
               ([("이번달 가격 = 가격 + 지난달 미디어",12,INK,False)],PP_ALIGN.LEFT,3),
               ([("미디어 추가 시 오차가 줄면 선행",11,GRAY,False)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE,ml=10)
    arrow(s,10.75,2.85,0.5,0.6,fill=BLUE)
    res=box(s,11.4,2.3,1.6,1.7,fill=GREEN)
    settext(res,[([("B≪A 오차\n& p<0.05\n→ 선행",12,WHITE,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
    # 왜 Granger + 설정
    why_box(s,0.3,4.25,6.2,2.0,"Granger 검정",
            ["연구 질문이 ‘미디어가 가격에 선행하는가’.",
             "회귀분석은 동시적 ‘관계’만 봄.",
             "→ 방향과 시차를 동시에 보는 Granger 선택"],accent=BLUE,lfill=LBLUE)
    setbox=box(s,6.7,4.25,6.3,2.0,fill=WHITE,line="DDDDDD",lw=1.0)
    settext(setbox,[([("검정 설정",12,NAVY,True)],PP_ALIGN.LEFT,5),
                    ([("• 분석단위 자산(아이템 평균 대표 시계열)",11.5,INK,False)],PP_ALIGN.LEFT,3),
                    ([("• 이변량 VAR · 래그 AIC로 1–4 선택",11.5,INK,False)],PP_ALIGN.LEFT,3),
                    ([("• 유의 기준 ",11.5,INK,False),("p < 0.05",11.5,DRED,True),(" → 자산별 선행 채널 선별",11.5,INK,False)],PP_ALIGN.LEFT,3),
                    ([("   (F 통계량으로 채널 간 선행 강도 비교)",10.5,GRAY,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)

# ============================================================
# 15. XGBoost·SHAP·DM + A/B/C
# ============================================================
def s_xgb():
    s=slide()
    title_bar(s,"XGBoost · SHAP · DM — 예측 검증 (왜 이 도구인가?)",
              "Granger 선행성을 실제 예측력으로 검증 · 모델 A/B/C 설계",pg())
    # 왜 3박스
    why_box(s,0.3,1.2,4.15,2.0,"XGBoost",
            ["N=48 소표본 — 딥러닝은 과적합 위험.",
             "트리 앙상블이 소표본·비선형에 강함."],accent=GREEN,lfill=LGREEN)
    why_box(s,4.6,1.2,4.15,2.0,"SHAP",
            ["채널별 ‘예측 기여도’를 수치화.",
             "Granger(선행성)와 교차검증하는 핵심."],accent=BLUE,lfill=LBLUE)
    why_box(s,8.9,1.2,4.1,2.0,"DM 검정",
            ["AUC 숫자차가 우연인지 실제인지",
             "월별 예측오차 시계열로 통계 판정."],accent=PURPLE,lfill=LPUR)
    # A/B/C 설계 플로우 ㅁ→ㅁ→ㅁ
    fl=tb(s,0.3,3.35,12.7,0.35)
    settext(fl,[([("모델 설계 흐름",13,NAVY,True),("   전채널(A) → SHAP 기여분석 → 상위 선별(B·C) → DM 동등성 검정",12,GRAY,False)],PP_ALIGN.LEFT,0)])
    steps=[dict(x=0.3,w=2.75,fill=NAVY,t="① Model A\n전채널",sub="CH1~CH5 전부"),
           dict(x=3.55,w=2.75,fill=BLUE,t="② SHAP 분석",sub="실제 기여 채널"),
           dict(x=6.80,w=2.75,fill=GREEN,t="③ Model B·C\n선별",sub="B:Granger / C:SHAP2"),
           dict(x=10.05,w=2.95,fill=PURPLE,t="④ DM 검정",sub="A vs 선별 동등?")]
    ft=3.8; fh=1.3
    for st in steps:
        b=box(s,st["x"],ft,st["w"],fh,fill=st["fill"])
        settext(b,[([(ln,13.5,WHITE,True)],PP_ALIGN.CENTER,1) for ln in st["t"].split("\n")]+
                  [([(st["sub"],10.5,WHITE,False)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
    for ax in [3.07,6.32,9.57]: arrow(s,ax,ft+0.42,0.46,0.45,fill=MGRAY)
    # A/B/C 정의 미니
    minis=[("Model A · 전채널",NAVY,LGRAY,"CH1+CH2+CH3+CH4+CH5 (기준)"),
           ("Model B · Granger 선별",GREEN,LGREEN,"스니커즈 CH3 · 카드 CH1 · 레고 —"),
           ("Model C · SHAP 상위2",PURPLE,LPUR,"스니커즈 CH1+CH5 · 카드 CH4+CH1 · 레고 CH5+CH4")]
    mt=5.35; x=0.3
    for name,acc,lf,body in minis:
        b=box(s,x,mt,4.05,1.0,fill=lf,line="DDDDDD",lw=0.75)
        settext(b,[([(name,12,acc,True)],PP_ALIGN.LEFT,3),([(body,11,INK,False)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE,ml=10)
        x+=4.32
    cm=tb(s,0.3,6.5,12.7,0.35)
    settext(cm,[([("공통  ",11,BLUE,True),("통제피처 price_vs_ma3·lag1~3 | TimeSeriesSplit(5) | 평가 AUC-ROC",11,GRAY,False)],PP_ALIGN.LEFT,0)])

# ============================================================
# 17. RQ1 Granger 결과
# ============================================================
def s_rq1():
    s=slide()
    title_bar(s,"RQ1 · Granger 인과검정 결과",
              "15회 중 유의 4 · 경계 1 · 비유의 10 (유의 기준 p<0.05)",pg())
    fit_image(s,"_imgs/granger.png",0.3,1.2,5.5,4.6,halign='center',valign='top')
    caption(s,0.3,5.95,5.5,"자산×채널 15쌍의 Granger F값 — 높을수록(스니커즈 CH3·카드 CH1) 선행성 강함.")
    # 표
    data=[("자산","채널","래그","F","p","판정"),
          ("스니커즈","CH3 뉴스감성","1","13.39","0.001","유의"),
          ("카드","CH1 Google","1","8.42","0.006","유의"),
          ("스니커즈","CH4 유튜브뷰","1","4.73","0.035","유의"),
          ("스니커즈","CH1 Google","2","3.47","0.041","유의"),
          ("레고","CH1 Google","1","3.92","0.054","경계"),
          ("나머지 10쌍","—","—","<2",">0.18","비유의")]
    gt=s.shapes.add_table(7,6,Inches(6.1),Inches(1.25),Inches(6.9),Inches(3.6)).table
    for i,w in enumerate([1.4,1.9,0.7,0.9,1.0,1.0]): gt.columns[i].width=Inches(w)
    for ri in range(7):
        for ci in range(6):
            cell=gt.cell(ri,ci); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Pt(4);cell.margin_right=Pt(2);cell.margin_top=Pt(1);cell.margin_bottom=Pt(1)
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER if ci>=2 else PP_ALIGN.LEFT
            r=p.add_run(); r.text=data[ri][ci]; _font(r); r.font.size=Pt(11)
            if ri==0:
                r.font.bold=True;r.font.color.rgb=C(WHITE);cell.fill.solid();cell.fill.fore_color.rgb=C(NAVY)
            else:
                sig = data[ri][5]=="유의"
                r.font.bold=sig; r.font.color.rgb=C(GREEN if sig else INK)
                cell.fill.solid(); cell.fill.fore_color.rgb=C(LGREEN if sig else (WHITE if ri%2 else LGRAY))
    concl=box(s,6.1,5.05,6.9,1.4,fill=LBLUE,line=BLUE,lw=1.0)
    settext(concl,[([("결론 RQ1",12,BLUE,True)],PP_ALIGN.LEFT,3),
                   ([("미디어의 가격 선행성 ",12,INK,False),("부분 지지",12,GREEN,True),
                     (" — 스니커즈 3개(CH3·CH4·CH1)·카드 1개(CH1) 채널이 유의(p<0.05)",12,INK,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)

# ============================================================
# 18. RQ2 Jaccard
# ============================================================
def s_rq2():
    s=slide()
    title_bar(s,"RQ2 · 자산별 채널 구성 차이 (Jaccard)",
              "유의 채널 집합이 자산마다 다른가? — Jaccard로 정량화",pg())
    fit_image(s,"paper/fig/fig_jaccard.png",0.3,1.25,6.2,4.6,halign='center',valign='top')
    caption(s,0.3,5.95,6.2,"자산별 유의 채널 집합. 겹치는 채널이 적을수록 채널 구성이 이질적.")
    # 우측 요약
    sets=box(s,6.8,1.3,6.2,2.0,fill=WHITE,line="DDDDDD",lw=1.0)
    settext(sets,[([("유의 채널 집합",13,NAVY,True)],PP_ALIGN.LEFT,5),
                  ([("스니커즈  {CH1, CH3, CH4}",12.5,BLUE,True)],PP_ALIGN.LEFT,3),
                  ([("카드        {CH1}",12.5,GREEN,True)],PP_ALIGN.LEFT,3),
                  ([("레고        {∅}  (유의 채널 없음)",12.5,PURPLE,True)],PP_ALIGN.LEFT,0)],
             anchor=MSO_ANCHOR.MIDDLE,ml=12)
    jac=box(s,6.8,3.45,6.2,1.5,fill=LGRAY,line="DDDDDD",lw=0.75)
    settext(jac,[([("Jaccard 유사도 (모두 < 0.6)",12,NAVY,True)],PP_ALIGN.LEFT,3),
                 ([("스니커즈–카드 ",12,INK,False),("0.333",12,INK,True),("   ·   그 외 쌍 ",12,INK,False),("0.000",12,INK,True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)
    concl=box(s,6.8,5.1,6.2,1.35,fill=LGREEN,line=GREEN,lw=1.0)
    settext(concl,[([("결론 RQ2",12,GREEN,True)],PP_ALIGN.LEFT,3),
                   ([("자산 유형마다 선행 채널 구성이 다름 → ",12,INK,False),("RQ2 지지",12,GREEN,True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=12)

# ============================================================
# 19. RQ3 성능 + DM
# ============================================================
def s_rq3():
    s=slide()
    title_bar(s,"RQ3 · 예측 성능 비교 & DM 검정",
              "Model A(전채널) vs B/C(선별) — 채널을 줄여도 성능이 유지되는가?",pg())
    fit_image(s,"_imgs/modelcmp.png",0.3,1.2,7.4,3.5,halign='center',valign='top')
    caption(s,0.3,4.55,7.4,"A/B/C/단독채널 AUC. 어느 구성이든 차이 최대 0.01 내외 → 채널 줄여도 성능 유지.")
    # AUC/DM 표 (우)
    data=[("","스니커즈","카드","레고"),
          ("Model A AUC","0.841","0.859","0.960"),
          ("Model C AUC","0.840","0.859","0.963"),
          ("DM p값","0.678","0.504","0.646"),
          ("판정","동등","동등","동등")]
    gt=s.shapes.add_table(5,4,Inches(7.9),Inches(1.25),Inches(5.1),Inches(2.7)).table
    for i,w in enumerate([1.8,1.1,1.1,1.1]): gt.columns[i].width=Inches(w)
    for ri in range(5):
        for ci in range(4):
            cell=gt.cell(ri,ci);cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Pt(4);cell.margin_right=Pt(2);cell.margin_top=Pt(1);cell.margin_bottom=Pt(1)
            p=cell.text_frame.paragraphs[0];p.alignment=PP_ALIGN.CENTER if ci>=1 else PP_ALIGN.LEFT
            r=p.add_run();r.text=data[ri][ci];_font(r);r.font.size=Pt(11.5)
            if ri==0 or ci==0:
                r.font.bold=True;r.font.color.rgb=C(WHITE if ri==0 else INK)
                cell.fill.solid();cell.fill.fore_color.rgb=C(NAVY if ri==0 else LGRAY)
            else:
                eq = data[ri][0]=="판정"
                r.font.bold=eq; r.font.color.rgb=C(GREEN if eq else INK)
                cell.fill.solid();cell.fill.fore_color.rgb=C(LGREEN if eq else WHITE)
    why_box(s,7.9,4.2,5.1,1.05,"DM 검정",
            ["AUC 숫자차가 우연인지 실제인지를",
             "월별 예측오차 시계열로 통계 판정"],accent=PURPLE,lfill=LPUR)
    concl=box(s,0.3,5.25,12.7,1.2,fill=LBLUE,line=BLUE,lw=1.0)
    settext(concl,[([("결론 RQ3",12.5,BLUE,True)],PP_ALIGN.LEFT,3),
                   ([("3자산 모두 DM p>0.50 → SHAP 상위 2채널만으로도 전채널과 ",12,INK,False),
                     ("통계적으로 동등",12,GREEN,True),
                     ("한 예측 성능 → 간결성(Parsimony) 확보",12,INK,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=14)

# ============================================================
# 20. SHAP & IRF 정성
# ============================================================
def s_qual():
    s=slide()
    title_bar(s,"정성 분석 · SHAP & IRF",
              "선행성(Granger) ≠ 예측 기여(SHAP) · 충격의 방향·지속을 IRF로 추적",pg())
    # SHAP (좌)
    fit_image(s,"_imgs/shap.png",0.3,1.25,6.3,2.7,halign='center',valign='top')
    caption(s,0.3,3.95,6.3,"채널별 평균|SHAP|. 스니커즈 Granger 1위였던 CH3가 SHAP에선 최하위.")
    sb=box(s,0.3,4.55,6.3,1.95,fill=LBLUE,line=BLUE,lw=1.0)
    settext(sb,[([("SHAP 핵심",12,BLUE,True)],PP_ALIGN.LEFT,3),
                ([("Granger 최강 채널(스니커즈 CH3)이 SHAP 기여는 4위",11.5,INK,False)],PP_ALIGN.LEFT,2),
                ([("→ 가격 모멘텀이 분산을 흡수 · ",11.5,INK,False),("선행성 ≠ 예측력",11.5,DRED,True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=10)
    # IRF (우)
    fit_image(s,"_imgs/irf_sneakers.png",6.8,1.25,3.0,2.5,halign='center',valign='top')
    fit_image(s,"_imgs/irf_cards.png",9.9,1.25,3.0,2.5,halign='center',valign='top')
    caption(s,6.8,3.85,6.2,"채널 충격→가격 반응(6개월). 음영=95% CI. 부호가 자산마다 반대.")
    ib=box(s,6.8,4.55,6.2,1.95,fill=LGREEN,line=GREEN,lw=1.0)
    settext(ib,[([("IRF 핵심  (왜 IRF? Granger는 유무만, IRF는 방향·크기·지속)",11.5,GREEN,True)],PP_ALIGN.LEFT,3),
                ([("카드 검색(CH1)↑ → 가격 ▲ (래그1 정점, 래그3 감쇠)",11.5,INK,False)],PP_ALIGN.LEFT,2),
                ([("스니커즈 뉴스감성(CH3)↓ → 가격 ▼ (래그1–2, 래그4 소멸)",11.5,INK,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=10)

# ============================================================
# 21. 제한점
# ============================================================
def s_limit():
    s=slide()
    title_bar(s,"연구 제한점",
              "과도한 일반화를 막기 위해 한계를 투명하게 명시 — 탐색적 발견으로 한정",pg())
    lims=[("소표본 (N=48)","자산별 48개월 — 검정력 낮음","탐색적 해석 · 검정력 명시",DRED,LRED),
          ("스니커즈 CH1 약정상성","차분 후도 p=0.120 경계","해석 시 주의 명시",ORANGE,LORANGE),
          ("CH3 도메인 특화도","GDELT tone(자동추출) 대체","민감도 분석·한계 명시",ORANGE,LORANGE),
          ("레고 GT 0주 비율","Porsche·Bugatti 0값 많음","강한 결론 배제",ORANGE,LORANGE),
          ("단일 플랫폼","YouTube만 — Tiktok·X 제외","향후 연구 기회",BLUE,LBLUE),
          ("z-score 데이터 누출","전체 패널 fit (fold별 X)","N=48 영향 미미·명시",BLUE,LBLUE)]
    x0=0.3; y0=1.3; w=4.1; h=1.55; gx=0.32; gy=0.2
    for i,(t,d,r2,acc,lf) in enumerate(lims):
        col=i%3; row=i//3
        x=x0+col*(w+gx); y=y0+row*(h+gy)
        b=box(s,x,y,w,h,fill=lf,line=acc,lw=1.0)
        settext(b,[([(t,12.5,acc,True)],PP_ALIGN.LEFT,3),
                   ([(d,11,INK,False)],PP_ALIGN.LEFT,3),
                   ([("대응  ",10.5,GRAY,True),(r2,10.5,GRAY,False)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE,ml=10)
    foot=box(s,0.3,4.95,12.7,0.7,fill=NAVY)
    settext(foot,[([("DH 패널 · VAR 블록 교차검증으로 결과를 스스로 검증 → 확증이 아닌 ‘검증할 가설’로 제시",12,WHITE,True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=14)

# ============================================================
# 23. 결론
# ============================================================
def s_concl():
    s=slide()
    title_bar(s,"결론 · 연구문제별 답변",
              "Granger → Jaccard → XGBoost+DM, 세 단계 분석의 핵심 결과",pg())
    rows=[("RQ1","미디어가 가격을 선행하는가?","부분 지지",
           "스니커즈 CH3(F=13.4)·CH4·CH1 · 카드 CH1(F=8.4) 유의 (p<0.05, 4개)",BLUE,LBLUE),
          ("RQ2","선행 채널 구성이 자산마다 다른가?","지지",
           "Jaccard 모두 <0.6 (스–카 0.33, 그 외 0) — 자산별 채널 이질성",GREEN,LGREEN),
          ("RQ3","선별 모델이 전채널과 동등한가?","동등(간결성)",
           "DM p>0.50 (스 0.678·카 0.504·레 0.646) — 2채널로 5채널 동등",PURPLE,LPUR)]
    y=1.35
    for rq,q,verd,ev,acc,lf in rows:
        b=box(s,0.4,y,12.5,1.45,fill=WHITE,line="DDDDDD",lw=1.0)
        box(s,0.4,y,0.14,1.45,fill=acc,shape=MSO_SHAPE.RECTANGLE)
        rl=tb(s,0.65,y,1.4,1.45);settext(rl,[([(rq,24,acc,True)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.MIDDLE)
        qq=tb(s,2.1,y+0.12,6.5,0.5);settext(qq,[([(q,13,INK,True)],PP_ALIGN.LEFT,0)])
        ee=tb(s,2.1,y+0.62,6.6,0.7);settext(ee,[([(ev,11.5,GRAY,False)],PP_ALIGN.LEFT,0)])
        vb=box(s,8.9,y+0.28,3.8,0.9,fill=lf,line=acc,lw=1.25)
        settext(vb,[([(verd,16,acc,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        y+=1.6
    key=box(s,0.4,y,12.5,0.75,fill=NAVY)
    settext(key,[([("핵심  ",12.5,"FFD9A0",True),
                   ("미디어 선행성은 자산별로 분화 — 스니커즈는 뉴스감성, 카드는 검색량이 핵심 선행 채널",12.5,WHITE,True)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.MIDDLE,ml=14)

# ============================================================
# 24. 시사점 & 향후
# ============================================================
def s_impl():
    s=slide()
    title_bar(s,"시사점 & 향후 연구","탐색적 발견의 실무 함의와 확증 연구로의 발전 경로",pg())
    cards=[("리셀 투자자·수집가",BLUE,LBLUE,"검색 관심도(CH1)를 1–2개월 선행 보조지표로 모니터링 — 기계적 신호가 아닌 종합 판단"),
           ("리셀 플랫폼 운영자",GREEN,LGREEN,"자산별 맞춤 피처 · Granger 선별 채널 중심 구성 시 예측력 유지하며 비용 절감"),
           ("학술 연구자",PURPLE,LPUR,"‘공개데이터 → Granger 선별 → XGBoost+SHAP+DM’ 파이프라인은 명품·시계 등으로 확장 가능")]
    x=0.3; w=4.1
    for name,acc,lf,body in cards:
        b=box(s,x,1.35,w,2.6,fill=lf,line=acc,lw=1.0)
        hd=box(s,x,1.35,w,0.6,fill=acc)
        settext(hd,[([(name,13.5,WHITE,True)],PP_ALIGN.CENTER,0)],anchor=MSO_ANCHOR.MIDDLE)
        bd=tb(s,x+0.2,2.1,w-0.4,1.7);settext(bd,[([(body,12.5,INK,False)],PP_ALIGN.LEFT,0)],anchor=MSO_ANCHOR.TOP)
        x+=4.32
    fut=box(s,0.3,4.25,12.7,2.3,fill=WHITE,line=ORANGE,lw=1.25)
    settext(fut,[([("향후 연구",14,ORANGE,True)],PP_ALIGN.LEFT,6),
                 ([("①  더 긴 패널(5년+)·아이템 확대로 검정력 강화",12.5,INK,False)],PP_ALIGN.LEFT,4),
                 ([("②  CH3를 도메인 특화 FinBERT로 대체해 감성 신호 품질 개선",12.5,INK,False)],PP_ALIGN.LEFT,4),
                 ([("③  월별 → 주별 고빈도 분석 · 드롭/한정판 이벤트 효과 통제",12.5,INK,False)],PP_ALIGN.LEFT,4),
                 ([("④  SPA/MCS 다중모형 비교 · 사전등록(pre-registration)으로 확증 연구화",12.5,INK,False)],PP_ALIGN.LEFT,0)],
            anchor=MSO_ANCHOR.TOP,ml=14,mt=10)

# ============================================================
# 25. 참고문헌
# ============================================================
def s_refs():
    s=slide()
    title_bar(s,"참고 문헌 (References)","주요 인용 · 데이터 출처",pg())
    col1=["미디어·감성 → 자산",
          "Tetlock (2007), JF 62(3)",
          "Da, Engelberg & Gao (2011), JF 66(5)",
          "Bollen, Mao & Zeng (2011), J.Comp.Sci 2(1)",
          "Antweiler & Frank (2004), JF 59(3)",
          "Loughran & McDonald (2011), JF 66(1)",
          "",
          "계량·통계",
          "Granger (1969), Econometrica 37(3)",
          "Dickey & Fuller (1979), JASA 74",
          "Diebold & Mariano (1995), JBES 13(3)",
          "Akaike (1974), IEEE TAC 19(6)"]
    col2=["머신러닝·NLP",
          "Chen & Guestrin (2016), KDD",
          "Lundberg & Lee (2017), NeurIPS 30",
          "Araci (2019), FinBERT arXiv:1908.10063",
          "Huang et al. (2023), CAR",
          "",
          "대안자산·수집품",
          "Dobrynskaya & Kishilova (2022), RIBF 59",
          "Campello/Raditya et al. (2021), Procedia CS",
          "",
          "데이터 출처",
          "StockX · PriceCharting · BrickRanker",
          "Google Trends · GDELT · YouTube Data API"]
    for cx,col in [(0.6,col1),(6.9,col2)]:
        y=1.25
        for line in col:
            t=tb(s,cx,y,6.0,0.32)
            if line=="":
                y+=0.16; continue
            head = not line[0].isalpha() or ('→' in line) or line in ("계량·통계","머신러닝·NLP","대안자산·수집품","데이터 출처","미디어·감성 → 자산")
            head = line in ("미디어·감성 → 자산","계량·통계","머신러닝·NLP","대안자산·수집품","데이터 출처")
            if head:
                settext(t,[([(line,12.5,BLUE,True)],PP_ALIGN.LEFT,0)])
            else:
                settext(t,[([(line,11,INK,False)],PP_ALIGN.LEFT,0)])
            y+=0.42

# ============================================================
# 26. Thank you
# ============================================================
def s_thanks():
    s=slide()
    box(s,0,0,13.333,7.5,fill=NAVY,shape=MSO_SHAPE.RECTANGLE)
    box(s,0,4.35,13.333,0.06,fill=BLUE,shape=MSO_SHAPE.RECTANGLE)
    t=tb(s,0.9,2.7,11.5,1.2);settext(t,[([("Thank you",44,WHITE,True)],PP_ALIGN.LEFT,0)])
    sub=tb(s,0.95,4.6,11.5,0.6)
    settext(sub,[([("리셀 시장 미디어 선행 지표 연구 — 스니커즈 · 트레이딩 카드 · 레고",15,"BBC4D0",False)],PP_ALIGN.LEFT,0)])
    au=tb(s,0.95,5.4,11.5,0.4);settext(au,[([("이승현 · 윤재은 · 최형호",13,"8FA0B5",False)],PP_ALIGN.LEFT,0)])

# ---------------- build ----------------
s_title()
s_toc()
divider(slide(),"01","연구 배경 & 질문","리셀 시장 · 연구문제 RQ1–RQ3")
s_rq()
divider(slide(),"02","기존 연구 & 한계","선행연구의 공통 한계 · 본 연구의 차별점")
s_related()
s_diff()
divider(slide(),"03","분석 대상 & 채널","3종 자산 · 5개 미디어 채널 · 전처리")
s_items()
s_channels()
divider(slide(),"04","분석 방법론 — 왜 이 도구인가?","FinBERT · Granger · XGBoost · SHAP · DM")
s_pipeline()
s_finbert()
s_granger()
s_xgb()
divider(slide(),"05","실험 결과","RQ1–RQ3 · SHAP · IRF · 한계")
s_rq1()
s_rq2()
s_rq3()
s_qual()
s_limit()
divider(slide(),"06","결론 & 시사점","연구문제별 결론 · 실무 함의 · 향후 연구")
s_concl()
s_impl()
s_refs()
s_thanks()

prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
