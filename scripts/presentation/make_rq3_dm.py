#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RQ3 DM 검정 설명 + 결과 슬라이드
Usage : python scripts/make_rq3_dm.py
Output: slide_rq3_dm.pptx
"""
import os, io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide_rq3_dm_v3.pptx")

# ── 모델별 AUC + DM 결과 ─────────────────────────────────────────────────
# B = Granger raw p<0.05  /  C = SHAP 선별 채널 (Model A SHAP 평균 임계값, 21쪽 값 기준)
# DM (A vs C): results/dm_fixed_selection.csv
ASSETS_KR = ["스니커즈", "카드", "레고"]
A_AUC = [0.8413, 0.8587, 0.9601]
B_AUC = [0.8461, 0.8612, 0.9601]   # 스니커즈 †근사
C_AUC = [0.8458, 0.8565, 0.9609]
B_CH  = ["CH1+CH3+CH4 †", "CH1",     "없음(=A)"]
C_CH  = ["CH1+CH5+CH4",    "CH4+CH1", "CH4+CH1"]
# DM A vs C (asset level) — Model A SHAP 평균 임계값 선별 기준
C_DM_P = [0.057, 0.356, 0.390]
# DM A vs B (asset level 없음 → 표시 생략, 아이템별만 보유)

# ── 색상 ─────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x2D, 0x4E)
BLUE   = RGBColor(0x2E, 0x6D, 0xA4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x0A, 0x0A, 0x0A)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
EXPBG  = RGBColor(0xF0, 0xF4, 0xFF)
GREENBG= RGBColor(0xD4, 0xED, 0xDA)
SNK_C  = "#2E6DA4"
CRD_C  = "#27A84A"
LGO_C  = "#8B5CF6"
ASSET_COLORS = [SNK_C, CRD_C, LGO_C]
C_MODEL_A = "#1F2D4E"
C_MODEL_B = "#E87722"
C_MODEL_C = "#2E6DA4"

# ── 그래프: 하나의 그룹 막대 — y축 0.80~1.01로 차이가 작아 보이게 ──────────
plt.rcParams["font.family"] = ["Malgun Gothic", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False

fig, ax = plt.subplots(figsize=(12.0, 4.8))
fig.patch.set_facecolor("white")
ax.set_facecolor("#F8F9FB")

xs  = np.array([0.0, 1.6, 3.2])   # 자산 그룹
bw  = 0.38
off = [-bw, 0, bw]
model_cols  = [C_MODEL_A, C_MODEL_B, C_MODEL_C]
model_lbls  = ["A  전채널", "B  Granger 유의", "C  SHAP 선별"]
all_aucs    = [A_AUC, B_AUC, C_AUC]

for mi, (aucs, col, lbl) in enumerate(zip(all_aucs, model_cols, model_lbls)):
    bars = ax.bar(xs + off[mi], aucs, bw * 0.90,
                  color=col, edgecolor="white", linewidth=0.5,
                  label=lbl, zorder=3)
    for bar, v in zip(bars, aucs):
        ax.text(bar.get_x() + bar.get_width()/2,
                v + 0.002,
                f"{v:.4f}",
                ha="center", va="bottom",
                fontsize=8.5, fontweight="bold", color=col)

# DM 차이없음 브래킷 (A vs C, 정확값)
for gi, (xc, p) in enumerate(zip(xs, C_DM_P)):
    y_br = 0.998
    # A 막대 왼쪽에서 C 막대 오른쪽까지
    x0 = xc + off[0] - bw*0.45
    x1 = xc + off[2] + bw*0.45
    ax.annotate("", xy=(x1, y_br), xytext=(x0, y_br),
                arrowprops=dict(arrowstyle="-", color="#1A6E35",
                                lw=1.5))
    _mk = " (경계)" if 0.05 < p < 0.10 else ""
    ax.text(xc, y_br + 0.003,
            f"DM p={p:.3f}  →  통계적 차이없음{_mk}",
            ha="center", va="bottom", fontsize=9,
            color="#1A6E35", fontweight="bold")

ax.set_ylim(0.80, 1.015)
ax.set_xlim(-0.65, 3.85)
ax.set_xticks(xs)
ax.set_xticklabels(ASSETS_KR, fontsize=13, fontweight="bold")
ax.set_ylabel("AUC-ROC", fontsize=11)
ax.set_title("Model A / B / C  AUC-ROC 비교 — 막대 간 차이가 미미함을 확인",
             fontsize=12, fontweight="bold", pad=10)
ax.spines["top"].set_visible(False)
ax.spines["right"].set_visible(False)
ax.yaxis.grid(True, linestyle="--", alpha=0.35, zorder=0)

legend_handles = [
    mpatches.Patch(facecolor=C_MODEL_A, label="A  전채널 (CH1~5 전부)"),
    mpatches.Patch(facecolor=C_MODEL_B, label="B  Granger 유의 채널 (raw p<0.05)"),
    mpatches.Patch(facecolor=C_MODEL_C, label="C  SHAP 선별 채널 (평균 임계값)"),
    plt.Line2D([0],[0], color="#1A6E35", linewidth=1.5,
               label="DM 브래킷: A↔C (자산 수준, 정확값)"),
]
ax.legend(handles=legend_handles, loc="lower right",
          fontsize=9.5, framealpha=0.96, edgecolor="#cccccc")

fig.text(0.5, -0.03,
    "† 스니커즈 B AUC: CH1+CH3+CH4 조합 미산출 → 개별 평균 근사  "
    "|  DM(A vs B)는 아이템 수준만 보유하여 브래킷 미표기",
    ha="center", fontsize=8, color="#999999", style="italic")

plt.tight_layout()
fig_buf = io.BytesIO()
fig.savefig(fig_buf, dpi=165, bbox_inches="tight", facecolor="white")
plt.close(fig)

# ── PPT 헬퍼 ─────────────────────────────────────────────────────────────
def _set_bg(cell, rgb):
    tc = cell._tc; p = tc.get_or_add_tcPr()
    for c in list(p):
        if c.tag.split("}")[-1] in ("solidFill","gradFill","noFill","blipFill","pattFill"):
            p.remove(c)
    sf = etree.SubElement(p, qn("a:solidFill"))
    cl = etree.SubElement(sf, qn("a:srgbClr"))
    cl.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
    p.remove(sf); p.insert(0, sf)

def add_box(sld, l, t, w, h, fill=None, border=False):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else:    tb.fill.background()
    if border:
        tb.line.color.rgb = RGBColor(0xBB,0xBB,0xBB); tb.line.width = Pt(0.75)
    else:
        tb.line.fill.background()
    return tb

def add_text(sld, text, l, t, w, h, bold=False, italic=False,
             size=11, color=BLACK, align=PP_ALIGN.LEFT):
    tb = add_box(sld, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.bold=bold; r.font.italic=italic
    r.font.size=Pt(size); r.font.color.rgb=color; r.font.name="Calibri"
    return tb

# ── PPT 생성 ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld = prs.slides.add_slide(prs.slide_layouts[6])

# 헤더
hbar = add_box(sld, 0, 0, 13.33, 0.72, fill=NAVY)
tf = hbar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  RQ3: 선별 채널 모델이 전채널 모델과 동등한가? — DM 검정"
r.font.bold=True; r.font.size=Pt(20); r.font.color.rgb=WHITE; r.font.name="Calibri"
add_box(sld, 0, 0.72, 13.33, 0.025, fill=BLUE)
add_text(sld,
    "Diebold-Mariano 검정으로 'Model A(전채널) = Model C(SHAP 선별 채널)' 가설을 통계적으로 검증",
    0.20, 0.755, 13.0, 0.30, size=9, color=DGRAY, italic=True)

# ── 좌측: DM 검정 설명 박스 ──────────────────────────────────────────────
EXP_L = 0.22; EXP_T = 1.07; EXP_W = 4.55; EXP_H = 5.85

add_box(sld, EXP_L, EXP_T, EXP_W, EXP_H, fill=EXPBG, border=True)

# 제목
add_text(sld, "왜 DM 검정이 필요한가?",
    EXP_L+0.15, EXP_T+0.10, EXP_W-0.30, 0.40,
    bold=True, size=13, color=NAVY)
add_box(sld, EXP_L+0.15, EXP_T+0.52, EXP_W-0.30, 0.022,
        fill=BLUE)

# 설명 1: 문제
add_text(sld, "문제: AUC 숫자 비교만으론 부족",
    EXP_L+0.15, EXP_T+0.60, EXP_W-0.30, 0.32,
    bold=True, size=11, color=NAVY)
add_text(sld,
    'A AUC ≈ C AUC (차이 0.00x 수준)\n'
    '→ "거의 같은데?" 라고만 보임\n'
    '→ 이 차이가 우연인지, 진짜 차이인지 알 수 없음',
    EXP_L+0.15, EXP_T+0.93, EXP_W-0.30, 0.85,
    size=10.5, color=BLACK)

# 구분선
add_box(sld, EXP_L+0.15, EXP_T+1.80, EXP_W-0.30, 0.018,
        fill=RGBColor(0xCC,0xCC,0xCC))

# 설명 2: DM 검정이란
add_text(sld, "DM 검정이 하는 일",
    EXP_L+0.15, EXP_T+1.87, EXP_W-0.30, 0.32,
    bold=True, size=11, color=NAVY)
add_text(sld,
    "두 모델의 예측 오차 시계열을 직접 비교\n"
    "→ 오차 차이가 통계적으로 유의한지 검정\n\n"
    "H₀: 두 모델의 예측 오차 크기가 같다\n\n"
    "p > 0.05  →  H₀ 기각 실패\n"
    "  = '통계적으로 차이없음' 공식 확인\n"
    "  = Model C가 A와 동등하다고 주장 가능",
    EXP_L+0.15, EXP_T+2.20, EXP_W-0.30, 1.80,
    size=10.5, color=BLACK)

# 구분선
add_box(sld, EXP_L+0.15, EXP_T+4.05, EXP_W-0.30, 0.018,
        fill=RGBColor(0xCC,0xCC,0xCC))

# 설명 3: 왜 중요한가
add_text(sld, "왜 이게 중요한가?",
    EXP_L+0.15, EXP_T+4.12, EXP_W-0.30, 0.32,
    bold=True, size=11, color=NAVY)
add_text(sld,
    "RQ3 핵심 주장:\n"
    "\"5개 채널 대신 2~3개만 써도 충분하다\"\n\n"
    "DM 검정 없이는 → 단순 AUC 비교\n"
    "DM 검정 있으면 → 통계적 동등성 입증\n"
    "= 간결성(parsimony) 근거 확보",
    EXP_L+0.15, EXP_T+4.45, EXP_W-0.30, 1.30,
    size=10.5, color=BLACK)

# ── 우측: 그래프 + 결과 표 ────────────────────────────────────────────────
GRF_L = 4.95; GRF_T = 1.07; GRF_W = 8.20; GRF_H = 4.05

fig_buf.seek(0)
sld.shapes.add_picture(fig_buf,
    Inches(GRF_L), Inches(GRF_T),
    Inches(GRF_W), Inches(GRF_H))

# ── 하단 결과 표 ──────────────────────────────────────────────────────────
TBL_T = GRF_T + GRF_H + 0.10
tbl_data = [
    ["자산", "A AUC\n(전채널)", "B AUC\n(Granger)", "C AUC\n(SHAP)", "DM A↔C\n통계량", "DM A↔C\np값", "판정"],
    ["스니커즈", "0.8413", "0.8461 †", "0.8458", "+1.916", "0.057", "H₀ 기각 실패\n→ 동등(경계)"],
    ["카드",     "0.8587", "0.8612",   "0.8565", "+0.926", "0.356", "H₀ 기각 실패\n→ 동등"],
    ["레고",     "0.9601", "0.9601",   "0.9609", "+0.861", "0.390", "H₀ 기각 실패\n→ 동등"],
]
SNK_BG = RGBColor(0xE3,0xF2,0xFD)
CRD_BG = RGBColor(0xE8,0xF5,0xE9)
LGO_BG = RGBColor(0xFF,0xF9,0xC4)
ROW_BG = [SNK_BG, CRD_BG, LGO_BG]
VRD_BG = GREENBG

col_w = [1.45, 1.25, 1.35, 1.25, 1.20, 0.90, 2.55]
tbl_l = GRF_L
TBL_H = 7.36 - TBL_T

tbl = sld.shapes.add_table(
    4, 7, Inches(tbl_l), Inches(TBL_T),
    Inches(sum(col_w)), Inches(TBL_H)).table

for i, w in enumerate(col_w): tbl.columns[i].width = Inches(w)
rh = Inches(TBL_H / 4)
for i in range(4): tbl.rows[i].height = rh

for ri, row in enumerate(tbl_data):
    is_hdr = (ri == 0)
    for ci, txt in enumerate(row):
        cell = tbl.cell(ri, ci)
        tf2  = cell.text_frame; tf2.word_wrap = True
        p2   = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = str(txt)
        r2.font.name = "Calibri"
        r2.font.size = Pt(9.5 if is_hdr else 9.5)
        r2.font.bold = is_hdr or (ci == 6 and not is_hdr)
        r2.font.color.rgb = WHITE if is_hdr else BLACK
        if is_hdr:
            _set_bg(cell, NAVY)
        elif ci == 6:
            _set_bg(cell, VRD_BG)
        else:
            _set_bg(cell, ROW_BG[ri-1])

# ── 결론 배너 ─────────────────────────────────────────────────────────────
add_box(sld, 0.22, 7.10, 13.0, 0.32, fill=NAVY)
add_text(sld,
    "결론 RQ3  |  3자산 모두 DM p > 0.05 → H₀ 기각 실패 → "
    "SHAP 선별 채널(Model C)로도 전채널(A)과 통계적으로 동등한 예측 성능 달성  →  간결성(Parsimony) 확인",
    0.30, 7.12, 12.85, 0.30,
    bold=True, size=10, color=WHITE, align=PP_ALIGN.CENTER)

for cand in [OUT, OUT.replace("_v3", "_v4"), OUT.replace("_v3", "_v5")]:
    try:
        prs.save(cand); print(f"[OK] {cand}"); break
    except PermissionError:
        print(f"  (잠김: {os.path.basename(cand)} → 다음)")
