# -*- coding: utf-8 -*-
"""
비교군 슬라이드 1장 (개정판 v3)
'기존 논문 방법(선형회귀·랜덤포레스트) vs 우리 모델 — 동일 데이터'

반영: 선형회귀 복귀 / 출처(Raditya 2021) / 슬라이드에 쉬운 설명 직접 기재 /
      가격만(혼란 원인) 제외 / 우리 모델 강조 / 기대효과 / 무작위(0.50) 기준선

숫자 출처(실제 계산값, 새로 학습 X):
  - results/model_comparison_auc.csv : ElasticNet(선형회귀), RandomForest, XGBoost
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

for fname in ['Malgun Gothic', 'NanumGothic', 'AppleGothic', 'DejaVu Sans']:
    try:
        fm.findfont(fm.FontProperties(family=fname), fallback_to_default=False)
        plt.rcParams['font.family'] = fname
        break
    except Exception:
        continue
plt.rcParams['axes.unicode_minus'] = False

ASSETS = ['sneakers', 'cards', 'lego']
ASSET_KR = {'sneakers': '스니커즈', 'cards': '카드', 'lego': '레고'}

comp = pd.read_csv('results/model_comparison_auc.csv')
def gc(asset, model):
    return float(comp[(comp.asset_type == asset) & (comp.model == model)]['auc'].iloc[0])

METHODS = [
    ('선형회귀',   [gc(a, 'ElasticNet')   for a in ASSETS], '#9E9E9E'),
    ('랜덤포레스트', [gc(a, 'RandomForest') for a in ASSETS], '#66BB6A'),
    ('우리 모델',   [gc(a, 'XGBoost')      for a in ASSETS], '#EF6C00'),
]
lin  = [gc(a, 'ElasticNet') for a in ASSETS]
ours = [gc(a, 'XGBoost')    for a in ASSETS]
print('=== 숫자 ===')
for label, vals, _ in METHODS:
    print(f"  {label:10s} | " + " | ".join(f"{ASSET_KR[a]} {v:.3f}" for a, v in zip(ASSETS, vals)))

# ── 그래프 ───────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(8.6, 5.3))
x = np.arange(len(ASSETS)); n = len(METHODS); w = 0.25
for i, (label, vals, color) in enumerate(METHODS):
    offset = (i - (n - 1) / 2) * w
    is_ours = (i == n - 1)
    bars = ax.bar(x + offset, vals, w, label=label, color=color,
                  edgecolor='black' if is_ours else 'white',
                  linewidth=2.2 if is_ours else 0.6, zorder=3)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.013, f'{v:.2f}',
                ha='center', va='bottom', fontsize=10 if is_ours else 9,
                fontweight='bold' if is_ours else 'normal')
ax.axhline(0.5, color='#D32F2F', linestyle='--', linewidth=1.3, zorder=2)
ax.text(len(ASSETS) - 0.5, 0.515, '찍기(동전) 수준 = 0.50', color='#D32F2F',
        fontsize=9, va='bottom', fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels([ASSET_KR[a] for a in ASSETS], fontsize=13, fontweight='bold')
ax.set_ylabel('방향 예측 정확도 (AUC)', fontsize=12)
ax.set_ylim(0.4, 1.07)
ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.07), ncol=3, fontsize=11, frameon=False)
ax.spines[['top', 'right']].set_visible(False)
ax.grid(axis='y', alpha=0.25, zorder=0)
plt.tight_layout()
png = 'results/figures/benchmark_comparison.png'
fig.savefig(png, dpi=180, bbox_inches='tight')
print(f'그림 → {png}')

# ── PPT ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

NAVY = RGBColor(0x1A, 0x23, 0x7E); GREY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x22, 0x22, 0x22); ORANGE = RGBColor(0xE6, 0x5C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TAG_RESELL = RGBColor(0x15, 0x65, 0xC0); FAINT = RGBColor(0x88, 0x88, 0x88)

def add_box(left, top, width, height, paragraphs, align=PP_ALIGN.LEFT, line_spacing=1.05):
    """paragraphs: [[(text,bold,color,size), ...], ...]  (문단 리스트)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for pi, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        for t, b, c, s in runs:
            r = p.add_run(); r.text = t
            r.font.size = Pt(s); r.font.bold = b; r.font.color.rgb = c
    return tb

def add_rect(left, top, width, height, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top), Inches(width), Inches(height))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

# 제목 + 부제
add_box(0.55, 0.28, 12.3, 0.7, [[('기존 방법과의 비교 — 동일 데이터로 정면 평가', True, NAVY, 27)]])
add_box(0.58, 1.02, 12.3, 0.5, [[
    ('기존 리셀 예측 방법(선형회귀·랜덤포레스트)과 주식의 대체데이터 기법을 같은 데이터로 비교', False, GREY, 13)]])

# 그래프 (왼쪽)
slide.shapes.add_picture(png, Inches(0.45), Inches(1.62), width=Inches(7.7))

# 한눈에 보기 배너 (그래프 아래)
add_rect(0.45, 6.05, 7.7, 0.95, RGBColor(0x1A, 0x23, 0x7E))
add_box(0.65, 6.12, 7.35, 0.85, [
    [('한눈에 보기   ', True, RGBColor(0xFF, 0xCC, 0x80), 12.5),
     ('옛 방식(선형회귀)은 잘 못 맞혔지만(0.60~0.71),', False, WHITE, 12.5)],
    [('우리 ML 모델은 훨씬 정확(0.84~0.96) → 머신러닝 + 미디어 접근이 리셀 예측에 효과적', False, WHITE, 12.5)],
], line_spacing=1.1)

# 오른쪽 패널
bx = 8.55; bw = 4.45
add_box(bx, 1.65, bw, 0.4, [[('■ 비교한 방법 (출처)', True, NAVY, 14)]])
add_box(bx, 2.08, bw, 2.0, [
    [('선형회귀 ', True, DARK, 11.5), ('[리셀·Raditya 2021] ', True, TAG_RESELL, 9),
     ('— 직선으로 관계를 긋는 전통 통계', False, GREY, 10)],
    [('랜덤포레스트 ', True, DARK, 11.5), ('[리셀·Raditya 2021] ', True, TAG_RESELL, 9),
     ('— 여러 결정트리를 합친 머신러닝', False, GREY, 10)],
    [('우리 모델 ', True, ORANGE, 11.5), ('[주식→리셀] ', True, ORANGE, 9),
     ('— 주식에서 검증된 미디어 예측법(검색·감성)을 리셀에 적용', False, GREY, 10)],
    [('※ Raditya 2021은 2개 모델(선형·RF)만 비교 — 우리는 XGBoost·미디어를 추가 확장', False, FAINT, 8.5)],
], line_spacing=1.13)

add_box(bx, 4.05, bw, 0.4, [[('■ AUC가 뭔가요?', True, NAVY, 14)]])
add_box(bx, 4.48, bw, 0.85, [
    [('가격이 ', False, DARK, 11), ('"오를지·내릴지" 방향', True, DARK, 11), ('을 맞히는 정확도.', False, DARK, 11)],
    [('0.5 = 찍기(동전) · 1.0 = 완벽 · 0.9↑ = 매우 잘 맞힘', False, DARK, 11)],
], line_spacing=1.1)

# 기대효과 박스
add_rect(bx, 5.45, bw, 1.55, RGBColor(0xFF, 0xF3, 0xE0))
add_box(bx + 0.2, 5.55, bw - 0.4, 0.4, [[('■ 기대효과', True, ORANGE, 14)]])
add_box(bx + 0.2, 5.98, bw - 0.4, 1.0, [
    [('· 투자자: 가격 방향을 미리 파악', False, DARK, 10.5)],
    [('· 플랫폼: 수요 급증을 사전에 감지', False, DARK, 10.5)],
    [('· 학술: 주식의 대체데이터 예측법을 리셀 시장에 처음 적용·검증', False, DARK, 10.5)],
], line_spacing=1.08)

# 하단 주석(정직)
add_box(0.55, 7.05, 12.3, 0.4, [[
    ('※ 동일 데이터·동일 검증(시간순 교차검증 5-fold) 평균 AUC. 우리 모델은 랜덤포레스트와 대등하며, '
     '어떤 미디어가 가격을 선행하는지 해석(SHAP) 가능한 점이 차별점.', False, GREY, 9)]])

# 저장 (잠김 대비 fallback)
for out in ['slide_benchmark_comparison_v3.pptx', 'slide_benchmark_comparison_v4.pptx']:
    try:
        prs.save(out); print(f'PPT → {out}'); break
    except PermissionError:
        print(f'  (열려있어 잠김: {out} → 다음 이름 시도)')
