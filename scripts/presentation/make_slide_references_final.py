#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
참고 문헌 (References) — 발표용 출처 슬라이드 (1장)
slides_prior_work_v4에서 새로 추가된 인용(유재필·김문선 2023, Olaniyan 외 2024)을
기존 발표자료 슬라이드29 출처 목록과 통합 + 미검증/오류 인용 교정:
  - "Campello et al. (2021)" -> Raditya et al. (2021) (Procedia CS 179, 동일 논문)
  - "Dobrynskaya & Kishilova (2023)" -> (2022) (연도 오타 수정, RIBF 59)
  - "Zhang et al. (2018)" (미검증) -> Joseph, Wintoki & Zhang (2011), IJF 27(4)
  - "Huang et al. (2023), CAR" (미검증) -> Devlin et al. (2019), NAACL-HLT (BERT)
모두 paper/references.bib 또는 직접 확인한 PDF 기준 검증된 항목만 사용.
Usage : python scripts/make_slide_references_final.py
Output: slide_references_final.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide_references_final.pptx")

NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE= RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK   = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 21, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "참고 문헌 — 본 발표 인용 논문 & 데이터 출처",
       [("계량·통계 방법론 · 머신러닝/NLP · 미디어 선행성 · 대안자산 · 본 연구 직접 비교 대상", 13, False, DGRAY)])


def ref_box(x, t, w, h, tag, items, tag_fill=BLUE):
    box(s1, x, t, w, h, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s1, x, t, w, 0.40, fill=tag_fill, round_=True)
    text(s1, x, t + 0.02, w, 0.36, [[(tag, 11.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    runs = []
    for (cite, note) in items:
        runs.append([("•  ", 10, True, BLUE), (cite, 10, True, NAVY)])
        runs.append([("     " + note, 9.5, False, INK)])
    text(s1, x + 0.10, t + 0.48, w - 0.20, h - 0.56, runs, space=4)


GRID_T = 1.20
GAP_X, GAP_Y = 0.20, 0.18
MARGIN = 0.30
COL_W = (13.33 - 2 * MARGIN - 2 * GAP_X) / 3
ROW_H = 2.85
xs = [MARGIN, MARGIN + COL_W + GAP_X, MARGIN + 2 * (COL_W + GAP_X)]
ys = [GRID_T, GRID_T + ROW_H + GAP_Y]

# ── Row 1 ────────────────────────────────────────────────────────────────
ref_box(xs[0], ys[0], COL_W, ROW_H, "미디어·감성 → 주식 선행성", [
    ("Tetlock (2007), JF 62(3)", "뉴스 비관론 → S&P500 하락"),
    ("Da, Engelberg & Gao (2011), JF 66(5)", "Google 검색량 → 주가 선행"),
    ("Bollen, Mao & Zeng (2011), J.Comp.Sci 2(1)", "Twitter 감성 → DJIA 선행"),
    ("Antweiler & Frank (2004), JF 59(3)", "온라인 게시판 정보량 분석"),
    ("Loughran & McDonald (2011), JF 66(1)", "금융 텍스트 감성사전 구축"),
])

ref_box(xs[1], ys[0], COL_W, ROW_H, "Granger 선별 + 예측 (본 연구 직접 비교)", [
    ("유재필·김문선 (2023), 정보화연구", "유튜브 데이터 → ETF 선행 (Granger)"),
    ("Olaniyan 외 (2024), Electronics·MDPI", "Granger+상관분석 선별 → LSTM 예측"),
    ("Joseph, Wintoki & Zhang (2011), IJF 27(4)", "검색기반 투자자감성 → 수익률·거래량"),
], tag_fill=ORANGE)

ref_box(xs[2], ys[0], COL_W, ROW_H, "계량·통계 방법론", [
    ("Granger (1969), Econometrica 37(3)", "Granger 인과검정"),
    ("Dickey & Fuller (1979), JASA 74", "ADF 단위근(정상성) 검정"),
    ("Diebold & Mariano (1995), JBES 13(3)", "예측 정확도 비교 (DM 검정)"),
    ("Akaike (1974), IEEE TAC 19(6)", "AIC 기준 VAR 래그 선택"),
])

# ── Row 2 ────────────────────────────────────────────────────────────────
ref_box(xs[0], ys[1], COL_W, ROW_H, "머신러닝 · NLP", [
    ("Chen & Guestrin (2016), KDD", "XGBoost"),
    ("Lundberg & Lee (2017), NeurIPS 30", "SHAP (TreeExplainer)"),
    ("Araci (2019), arXiv:1908.10063", "FinBERT (금융 감성분석)"),
    ("Devlin et al. (2019), NAACL-HLT", "BERT (FinBERT의 기반 모델)"),
])

ref_box(xs[1], ys[1], COL_W, ROW_H, "대안자산 · 수집품 시장", [
    ("Dobrynskaya & Kishilova (2022), RIBF 59", "LEGO를 대체투자자산으로 분석"),
    ("Raditya et al. (2021), Procedia CS 179", "스니커즈 리셀가 ML 예측"),
])

box(s1, xs[2], ys[1], COL_W, ROW_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
box(s1, xs[2], ys[1], COL_W, 0.40, fill=GREEN, round_=True)
text(s1, xs[2], ys[1] + 0.02, COL_W, 0.36, [[("데이터 출처", 11.5, True, WHITE)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s1, xs[2] + 0.10, ys[1] + 0.48, COL_W - 0.20, ROW_H - 0.56,
     [[("가격", 10, True, NAVY)],
      [("   StockX · PriceCharting · BrickRanker", 9.5, False, INK)],
      [("미디어", 10, True, NAVY)],
      [("   Google Trends · GDELT · YouTube Data API v3", 9.5, False, INK)],
      [("기간", 10, True, NAVY)],
      [("   2022.01 ~ 2025.12 (48개월), 15개 아이템 (3자산×5)", 9.5, False, INK)]],
     space=4)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
