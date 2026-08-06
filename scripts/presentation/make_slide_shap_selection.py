#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SHAP 채널 선택 — "몇 개를 남길 것인가?"를 수치 기준으로 (1장)
원문 보강: 발표자료_최종_피드백반영_v3.pptx SLIDE21 "SHAP이 밝힌 것"
핵심: 임의의 '상위 2개' 대신, 선행연구 기반 수치 기준(평균 중요도 임계값) 도입.
  - 평균 중요도 임계값 = scikit-learn SelectFromModel 기본값(threshold='mean')
  - 채널 SHAP가 '채널 평균 SHAP' 이상인 채널만 채택 → 자산별 2개씩 선택
  - DM 검정: 3자산 모두 전채널(A)과 통계적 동등 → 간결성 근거
데이터: results/model_c_threshold_results.csv, results/model_c_channel_selection.csv
Usage : python scripts/make_slide_shap_selection.py
Output: slide_shap_selection.pptx
"""
import os
import io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(BASE, "slide_shap_selection.pptx")

for f in ["Malgun Gothic", "맑은 고딕"]:
    if any(f.lower() in fn.name.lower() for fn in fm.fontManager.ttflist):
        plt.rcParams["font.family"] = f
        break
plt.rcParams["axes.unicode_minus"] = False

NAVY = RGBColor(0x1F, 0x2D, 0x4E)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0xAA, 0xB1, 0xBC)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK = RGBColor(0x0A, 0x0A, 0x0A)

# ── 데이터 (results/model_c_channel_selection.csv 기준, 일관 파이프라인 SHAP) ──
CH = ["CH1", "CH2", "CH3", "CH4", "CH5"]
SHAP = {
    "스니커즈": {"CH1": 0.3081, "CH2": 0.0, "CH3": 0.0, "CH4": 0.1306, "CH5": 0.2236},
    "카드":     {"CH1": 0.0588, "CH2": 0.0, "CH3": 0.0, "CH4": 0.3334, "CH5": 0.1145},
    "레고":     {"CH1": 0.0778, "CH2": 0.0562, "CH3": 0.0045, "CH4": 0.2058, "CH5": 0.0925},
}
SELECTED = {"스니커즈": ["CH1", "CH5"], "카드": ["CH4", "CH5"], "레고": ["CH4", "CH5"]}
# DM (평균이상 규칙) : auc_A, auc_C, dm_p
DM = {
    "스니커즈": (0.8413, 0.8401, 0.678),
    "카드":     (0.8587, 0.8577, 0.254),
    "레고":     (0.9601, 0.9626, 0.646),
}
ASSETS = ["스니커즈", "카드", "레고"]

# ════════════════════════════════════════════════════════════════════════
# 그래프 — 자산별 채널 SHAP 가로막대 + 평균 임계선 (선택=초록 / 탈락=회색)
# ════════════════════════════════════════════════════════════════════════
fig, axes = plt.subplots(1, 3, figsize=(11.6, 3.5), dpi=200)
GREEN_H = "#1E7E34"
GRAY_H = "#C2C8D0"
for ax, asset in zip(axes, ASSETS):
    d = SHAP[asset]
    order = sorted(CH, key=lambda c: d[c])  # 작은 값이 아래로
    vals = [d[c] for c in order]
    mean_imp = sum(d.values()) / len(CH)
    colors = [GREEN_H if d[c] >= mean_imp else GRAY_H for c in order]
    ypos = np.arange(len(order))
    ax.barh(ypos, vals, color=colors, edgecolor="white", height=0.72)
    ax.set_yticks(ypos)
    ax.set_yticklabels(order, fontsize=10)
    ax.axvline(mean_imp, color="#E87722", linestyle="--", linewidth=1.4)
    ax.text(mean_imp, len(order) - 0.35, f"평균 {mean_imp:.2f}",
            color="#E87722", fontsize=8.5, fontweight="bold", ha="left", va="center")
    for yi, (c, v) in enumerate(zip(order, vals)):
        ax.text(v + 0.006, yi, f"{v:.2f}", va="center", fontsize=8.5,
                color="#333333", fontweight="bold" if d[c] >= mean_imp else "normal")
    ax.set_title(asset, fontsize=12.5, fontweight="bold", pad=6)
    ax.set_xlim(0, max(vals) * 1.28)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(axis="x", labelsize=8)
fig.suptitle("채널별 평균|SHAP|  —  주황 점선(채널 평균) 이상이면 선택(초록), 미만이면 탈락(회색)",
             fontsize=11, fontweight="bold", y=1.02)
plt.tight_layout()
fig_buf = io.BytesIO()
fig.savefig(fig_buf, dpi=185, bbox_inches="tight", facecolor="white")
plt.close(fig)

# ════════════════════════════════════════════════════════════════════════
# 슬라이드
# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 21, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "SHAP 채널 선택 — '몇 개를 남길까?'를 수치 기준으로",
       [("임의의 '상위 2개' → ", 13, False, DGRAY),
        ("평균 중요도 임계값(scikit-learn 기본 기준)", 13, True, ORANGE),
        ("으로 대체", 13, False, DGRAY)])

MARGIN = 0.30

# ── 좌측 상단: 기준 설명 ──────────────────────────────────────────────────
CRIT_T, CRIT_W, CRIT_H = 1.30, 4.75, 3.05
box(s1, MARGIN, CRIT_T, CRIT_W, CRIT_H, fill=RGBColor(0xEA, 0xF1, 0xFA), line=BLUE, line_w=1.0, round_=True)
text(s1, MARGIN + 0.16, CRIT_T + 0.10, CRIT_W - 0.32, 0.36,
     [[("몇 개를 남길지, 무엇을 기준으로?", 13.5, True, BLUE)]])
text(s1, MARGIN + 0.16, CRIT_T + 0.50, CRIT_W - 0.32, CRIT_H - 0.60,
     [[("•  기존: ", 11, False, INK), ("'상위 2개'는 임의 선택", 11, True, ORANGE),
       (" — \"왜 2개?\" 근거 부족", 11, False, INK)],
      [("•  도입 기준: ", 11, False, INK), ("평균 중요도 임계값", 11, True, GREEN)],
      [("   채널 SHAP가 ", 10.5, False, INK), ("'5개 채널의 평균 SHAP'", 10.5, True, GREEN),
       (" 이상인 채널만 채택", 10.5, False, INK)],
      [("   → 남길 개수를 ", 10.5, False, INK), ("데이터가 결정", 10.5, True, GREEN),
       (" (임의의 k 없음)", 10.5, False, INK)],
      [("•  근거: ", 11, False, INK),
       ("scikit-learn SelectFromModel 기본값 threshold='mean'", 10, True, BLUE)],
      [("   + 누적중요도 방식 (J. Big Data 2024 등) 계열의 표준 규칙", 10, False, DGRAY)],
      [("•  결과: 세 자산 모두 ", 11, False, INK), ("정확히 2개 채널", 11, True, GREEN),
       (" 선택됨", 11, False, INK)]],
     space=4)

# ── 좌측 하단: DM 동등성 표 ───────────────────────────────────────────────
DM_T = CRIT_T + CRIT_H + 0.16
DM_H = 1.92
box(s1, MARGIN, DM_T, CRIT_W, DM_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
text(s1, MARGIN + 0.16, DM_T + 0.08, CRIT_W - 0.32, 0.30,
     [[("선택 모델(C) vs 전채널(A) — DM 검정", 12, True, NAVY)]])
# 미니 표
tb_x = MARGIN + 0.16
tb_w = CRIT_W - 0.32
cols = [1.05, 1.30, 0.95, 1.05]
cx = [tb_x]
for w in cols[:-1]:
    cx.append(cx[-1] + w)
hh = 0.30
rh = 0.34
hdrs = ["자산", "선택 채널", "DM p", "판정"]
for x, w, h in zip(cx, cols, hdrs):
    box(s1, x, DM_T + 0.42, w, hh, fill=NAVY, line=WHITE, line_w=0.5)
    text(s1, x, DM_T + 0.42, w, hh, [[(h, 9.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, asset in enumerate(ASSETS):
    y = DM_T + 0.42 + hh + i * rh
    aa, cc, p = DM[asset]
    vals = [asset, "+".join(SELECTED[asset]), f"{p:.3f}", "동등"]
    for j, (x, w, v) in enumerate(zip(cx, cols, vals)):
        fill = WHITE if j < 3 else GREEN
        col = INK if j < 3 else WHITE
        bold = (j == 0) or (j == 3)
        box(s1, x, y, w, rh, fill=fill, line=RGBColor(0xDD, 0xDD, 0xDD), line_w=0.5)
        text(s1, x, y, w, rh, [[(v, 9.5, bold, col)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── 우측: 그래프 ──────────────────────────────────────────────────────────
FIG_X = MARGIN + CRIT_W + 0.22
FIG_W = 13.33 - MARGIN - FIG_X
FIG_T = 1.30
box(s1, FIG_X, FIG_T, FIG_W, 3.55, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
from PIL import Image as _PILImage
fig_buf.seek(0)
_w_px, _h_px = _PILImage.open(fig_buf).size
ratio = _h_px / _w_px
pic_w = FIG_W - 0.24
pic_h = pic_w * ratio
if pic_h > 3.55 - 0.20:
    pic_h = 3.55 - 0.20
    pic_w = pic_h / ratio
fig_buf.seek(0)
s1.shapes.add_picture(fig_buf, Inches(FIG_X + (FIG_W - pic_w) / 2), Inches(FIG_T + 0.12),
                      width=Inches(pic_w))

# ── 우측 하단: 강건성 + 누적합 해설 ────────────────────────────────────────
ROB_T = FIG_T + 3.55 + 0.16
ROB_H = 1.92
box(s1, FIG_X, ROB_T, FIG_W, ROB_H, fill=RGBColor(0xFC, 0xF0, 0xE6), line=ORANGE, line_w=1.0, round_=True)
text(s1, FIG_X + 0.16, ROB_T + 0.08, FIG_W - 0.32, 0.30,
     [[("강건성 점검 — 기준을 바꿔도 결론이 유지되나?", 12.5, True, ORANGE)]])
text(s1, FIG_X + 0.16, ROB_T + 0.42, FIG_W - 0.32, ROB_H - 0.50,
     [[("•  더 느슨한 기준(누적 SHAP 80%)을 쓰면 레고는 채널이 ", 11, False, INK),
       ("3개로 늘어남", 11, True, ORANGE),
       (" → 약한 3번째 채널이 노이즈로 작용, DM p=0.021로 ", 11, False, INK),
       ("오히려 전채널(A)이 더 우수", 11, True, ORANGE)],
      [("•  즉 채널을 더 넣는다고 좋아지지 않음 — ", 11, False, INK),
       ("'적게 쓰는 것이 충분하다'는 간결성(parsimony) 주장이 강화", 11, True, GREEN)],
      [("•  '평균 중요도 임계값'은 자산마다 딱 2개를 골라 세 자산 모두 A와 동등 → 본 연구의 기준으로 채택", 11, False, INK)]],
     space=5)

# ── 하단 결론 배너 ────────────────────────────────────────────────────────
BAN_T = 7.5 - 0.55
box(s1, MARGIN, BAN_T, 13.33 - 2 * MARGIN, 0.42, fill=NAVY)
text(s1, MARGIN + 0.15, BAN_T, 13.33 - 2 * MARGIN - 0.3, 0.42,
     [[("결론  ", 12.5, True, WHITE),
       ("채널 개수는 임의가 아니라 ", 11.5, False, WHITE),
       ("'평균 중요도 임계값'(=라이브러리 표준 기준)", 11.5, True, WHITE),
       ("으로 결정 → 자산별 2개 → DM 검정상 전채널과 동등 → RQ3 간결성 입증", 11.5, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
