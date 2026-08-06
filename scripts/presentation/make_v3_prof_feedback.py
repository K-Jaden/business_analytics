# -*- coding: utf-8 -*-
"""
교수님 피드백 반영 v2 -> v3  (발표자료_최종_피드백반영_v2.pptx -> _v3.pptx)
원본은 손대지 않고 항상 새 파일로 저장 (idempotent).

피드백 매핑:
 1. Jaccard·BH 제거            -> S4 RQ2 카드, S25 RQ2 근거·부제 수정 (BH는 v2에 없음 확인)
 2. 비교군 "무조건"             -> 벤치마크 비교 슬라이드 신설 (S22 뒤) — 동일 데이터로
                                  기존 리셀 연구 방법(선형회귀·RF, Raditya 2021) vs 우리 모델
 3. FinBERT 연결고리 어색       -> '데이터 수집' 슬라이드 신설 (파이프라인 S12 뒤, FinBERT 앞)
                                  텍스트 데이터 -> 점수화 필요 -> FinBERT 브리지
 4. 본 연구 좀 더 크게          -> S6 하단 '본 연구' 배너 확대 + 폰트 키움
 5. 용어 풀이 (GDELT·z-score)   -> S10 용어 박스 + 수집 슬라이드에 GDELT 설명
 6. 수치화/점수화 통일          -> 전체 '수치화' -> '점수화' 치환
 7. 수치 설명 부족              -> S17 F·p값 풀이, S18 AUC·ΔAUC 풀이, S21 DM p 풀이
 8. 빈 공간에 기대효과          -> S12 파이프라인 하단 기대효과 배너
 9. IRF 잔재 제거               -> TOC·섹션 부제 'IRF' -> '기존 방법 비교'
10. 페이지 번호 재계산
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")
import copy
import re
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "발표자료_최종_피드백반영_v2.pptx"
DST = "발표자료_최종_피드백반영_v3.pptx"

NAVY = "1F2D4E"; BLUE = "2E6DA4"; GREEN = "277A4A"; PURPLE = "6C3D91"
ORANGE = "C8772E"; GRAY = "555555"; INK = "111111"
LBLUE = "EAF1F8"; LGREEN = "E7F2EC"; LPUR = "EFEAF5"; LORANGE = "FBEEDD"
LRED = "F8D7DA"; DRED = "B23B3B"; LGRAY = "F2F3F5"; MGRAY = "9AA5B1"
WHITE = "FFFFFF"; AMBER = "FFD27F"
FONT = "맑은 고딕"

BENCH_PNG = "results/figures/benchmark_comparison.png"

def C(s):
    return RGBColor.from_string(s)

def find(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

def remove_shape(sh):
    sh._element.getparent().remove(sh._element)

def no_shadow(shp):
    try:
        shp.shadow.inherit = False
    except Exception:
        pass

def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = C(line); shp.line.width = Pt(line_w)
    no_shadow(shp)
    return shp

def set_text(shp, paras, anchor=MSO_ANCHOR.TOP, wrap=True, ml=5, mr=5, mt=3, mb=3):
    """paras: [(runs, align, space_after)], runs = [(text, size, color, bold)]"""
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(ml); tf.margin_right = Pt(mr)
    tf.margin_top = Pt(mt); tf.margin_bottom = Pt(mb)
    first = True
    for runs, align, space in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if space is not None:
            p.space_after = Pt(space)
        p.space_before = Pt(0)
        for text, size, color, bold in runs:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = FONT; r.font.color.rgb = C(color)
    return shp

def textbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def clone_shape(src_sh, dst_slide):
    el = copy.deepcopy(src_sh._element)
    dst_slide.shapes._spTree.append(el)
    for sh in dst_slide.shapes:
        if sh._element is el:
            return sh
    return None

def make_table(slide, l, t, w, h, rows_data, col_widths, font_sizes=(11, 10.5)):
    rows, cols = len(rows_data), len(rows_data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for ci, cw in enumerate(col_widths):
        gt.columns[ci].width = Inches(cw)
    for ri in range(rows):
        for ci in range(cols):
            cell = gt.cell(ri, ci)
            cell.margin_left = Pt(5); cell.margin_right = Pt(5)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            first = True
            for line in str(rows_data[ri][ci]).split("\n"):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
                r = p.add_run(); r.text = line
                r.font.name = FONT
                if ri == 0:
                    r.font.size = Pt(font_sizes[0]); r.font.bold = True
                    r.font.color.rgb = C(WHITE)
                else:
                    r.font.size = Pt(font_sizes[1]); r.font.bold = (ci == 0)
                    r.font.color.rgb = C(INK)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C(NAVY)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C(WHITE if ri % 2 else LGRAY)
    return gt

# ================= main =================
shutil.copyfile(SRC, DST)
prs = Presentation(DST)
S = list(prs.slides)

s2, s4, s6, s10, s12 = S[1], S[3], S[5], S[9], S[11]
s16, s17, s18, s21, s25 = S[15], S[16], S[17], S[20], S[24]

# ------------------------------------------------------------
# 1. S4 — RQ2 카드에서 Jaccard 제거
# ------------------------------------------------------------
set_text(find(s4, "Rounded Rectangle 15"), [
    ([("검정   ", 10.5, GREEN, True), ("자산별 유의 채널 집합 비교", 11.5, INK, True)],
     PP_ALIGN.CENTER, 0),
], anchor=MSO_ANCHOR.MIDDLE)
set_text(find(s4, "TextBox 16"), [
    ([("왜 이 검정?", 10.5, GREEN, True)], PP_ALIGN.LEFT, 2),
    ([("자산마다 '유의한 채널'이", 11, GRAY, False)], PP_ALIGN.LEFT, 0),
    ([("다른지 직접 맞대어 비교해", 11, GRAY, False)], PP_ALIGN.LEFT, 0),
    ([("자산별 특이성 확인", 11, GRAY, False)], PP_ALIGN.LEFT, 0),
], anchor=MSO_ANCHOR.TOP)
set_text(find(s4, "Rectangle 17"), [
    ([("측정   ", 10, GREEN, True), ("유의 채널 집합의 일치 여부", 11, INK, True)],
     PP_ALIGN.CENTER, 0),
], anchor=MSO_ANCHOR.MIDDLE)

# ------------------------------------------------------------
# 2. S6 — '본 연구' 배너 확대
# ------------------------------------------------------------
banner = find(s6, "Rounded Rectangle 17")
banner.top = Inches(6.10); banner.height = Inches(1.10)
banner.left = Inches(0.30); banner.width = Inches(12.70)
banner.fill.solid(); banner.fill.fore_color.rgb = C(NAVY)
banner.line.fill.background(); no_shadow(banner)
set_text(banner, [
    ([("본 연구   ", 17, AMBER, True),
      ("리셀 3개 자산(스니커즈·카드·레고) × 5개 미디어 채널 × Granger 선행성 + XGBoost·SHAP·DM 예측 검증",
       15, WHITE, True)], PP_ALIGN.CENTER, 3),
    ([("→ 주식에서 검증된 '미디어 선행성' 분석을 리셀 시장에 다채널로 처음 확장 — 세 흐름의 빈 교점을 메움",
       12.5, "D9E2F3", False)], PP_ALIGN.CENTER, 0),
], anchor=MSO_ANCHOR.MIDDLE)

# ------------------------------------------------------------
# 3. S10 — 용어 풀이 박스 (GDELT · z-score)
# ------------------------------------------------------------
b = box(s10, 8.20, 4.95, 4.80, 1.15, fill=LBLUE, line=BLUE, line_w=1.0)
set_text(b, [
    ([("용어 풀이", 11.5, BLUE, True)], PP_ALIGN.LEFT, 2),
    ([("GDELT", 10.5, INK, True),
      (" — 전 세계 뉴스 기사를 자동 수집·분석하는 공개 뉴스 빅데이터 (Google 지원)",
       10.5, INK, False)], PP_ALIGN.LEFT, 2),
    ([("z-score", 10.5, INK, True),
      (" — 평균 0·표준편차 1로 맞춘 표준 점수 → 단위가 다른 채널을 같은 잣대로 비교",
       10.5, INK, False)], PP_ALIGN.LEFT, 0),
], anchor=MSO_ANCHOR.MIDDLE, ml=8)

# ------------------------------------------------------------
# 4. S12 — 파이프라인 보강 (단계별 상세 + 기대효과 배너)
# ------------------------------------------------------------
step_names = ["Rounded Rectangle 5", "Rounded Rectangle 6",
              "Rounded Rectangle 7", "Rounded Rectangle 8"]
for nm in step_names:
    sh = find(s12, nm)
    sh.top = Inches(1.45); sh.height = Inches(1.30)
for nm in ["Right Arrow 9", "Right Arrow 10", "Right Arrow 11"]:
    sh = find(s12, nm)
    sh.top = Inches(1.80)

detail_cards = [
    (0.30, 2.75, LBLUE, BLUE,
     ["가격 — StockX·PriceCharting·BrickRanker",
      "미디어 — Trends·GDELT·YouTube",
      "15아이템 × 48개월 (2022–2025)"]),
    (3.55, 2.75, LGREEN, GREEN,
     ["뉴스 논조·유튜브 댓글은 '텍스트'",
      "→ FinBERT로 감성 점수화",
      "5개 채널 점수 완성"]),
    (6.80, 2.75, LPUR, PURPLE,
     ["어떤 채널이 가격보다",
      "먼저 움직이는가? (RQ1·RQ2)",
      "자산별 유의 채널 집합 도출"]),
    (10.05, 2.95, LORANGE, ORANGE,
     ["선별 채널 예측력 검증 (RQ3)",
      "XGBoost A/B/C + DM 검정",
      "+ 기존 방법과 성능 비교"]),
]
for x, w, fill, line, lines in detail_cards:
    card = box(s12, x, 2.95, w, 1.85, fill=fill, line=line, line_w=1.0)
    set_text(card, [([(t, 11, INK, False)], PP_ALIGN.CENTER, 3) for t in lines],
             anchor=MSO_ANCHOR.MIDDLE)

eff = box(s12, 0.30, 5.30, 12.70, 1.45, fill=LORANGE, line=ORANGE, line_w=1.25)
set_text(eff, [
    ([("기대효과", 13.5, ORANGE, True)], PP_ALIGN.LEFT, 4),
    ([("· 투자자 — 미디어 신호로 가격 방향을 1–2개월 먼저 파악하는 보조지표 확보",
       11.5, INK, False)], PP_ALIGN.LEFT, 2),
    ([("· 플랫폼 — 수요 급증을 사전 감지 → 재고·프로모션 전략에 활용",
       11.5, INK, False)], PP_ALIGN.LEFT, 2),
    ([("· 학술 — 주식에서 검증된 대체데이터 예측법을 리셀 시장에 처음 적용·검증",
       11.5, INK, False)], PP_ALIGN.LEFT, 0),
], anchor=MSO_ANCHOR.MIDDLE, ml=12)

# ------------------------------------------------------------
# 5. 신규 슬라이드 A — 데이터 수집 (S12 뒤 = FinBERT 앞)
# ------------------------------------------------------------
blank = prs.slide_layouts[6]
sA = prs.slides.add_slide(blank)
for ph in list(sA.placeholders):
    remove_shape(ph)
# 제목 바·부제·페이지번호는 S10에서 스타일 복제
for nm in ["Rectangle 1", "Rectangle 2", "TextBox 3", "TextBox 4"]:
    clone_shape(find(s10, nm), sA)
tsh = find(sA, "Rectangle 1")
p0 = tsh.text_frame.paragraphs[0]
for r in list(p0.runs):
    r._r.getparent().remove(r._r)
r = p0.add_run(); r.text = "  데이터 수집 — 무엇을, 어디서, 어떻게"
r.font.size = Pt(20); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = C(WHITE)
sub = find(sA, "TextBox 3")
p0 = sub.text_frame.paragraphs[0]
for r in list(p0.runs):
    r._r.getparent().remove(r._r)
r = p0.add_run(); r.text = "가격 + 5개 채널, 모두 공개 데이터 — 텍스트 채널(뉴스·댓글)은 점수화가 필요"
r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = C(GRAY)

make_table(sA, 0.30, 1.25, 12.70, 4.40, [
    ["데이터", "출처", "수집 방법", "형태·규모"],
    ["가격 (타깃)", "StockX · PriceCharting · BrickRanker",
     "아이템별 월 실거래 평균가", "숫자 · 15아이템×48개월"],
    ["CH1 검색 관심도", "Google Trends",
     "아이템별 월간 검색지수 (0~100)", "숫자 · 720 아이템-월"],
    ["CH2·CH3 뉴스\n보도량 · 논조", "GDELT — 전 세계 뉴스 기사를 자동\n수집·분석하는 공개 빅데이터 (Google 지원)",
     "기사량·논조(tone)를 일별 수집 → 월평균", "숫자 · 720 아이템-월"],
    ["CH4 유튜브 조회수", "YouTube Data API",
     "해당 월에 게시된 관련 영상 상위 50개 조회수 합산", "숫자 · 720 아이템-월"],
    ["CH5 유튜브 댓글", "YouTube Data API",
     "조회수 상위 15개 영상 × 댓글 30개 (영어만)", "텍스트 → 점수화 필요"],
], [1.85, 4.05, 4.55, 2.25])

bridge = box(sA, 0.30, 5.95, 12.70, 0.95, fill=NAVY)
set_text(bridge, [
    ([("다음 단계   ", 13, AMBER, True),
      ("뉴스 논조·유튜브 댓글은 '텍스트' — 모델에 넣으려면 숫자로 바꿔야 함", 13, WHITE, True)],
     PP_ALIGN.LEFT, 2),
    ([("→ FinBERT로 댓글 감성을 점수화 (다음 슬라이드)", 12.5, "D9E2F3", False)],
     PP_ALIGN.LEFT, 0),
], anchor=MSO_ANCHOR.MIDDLE, ml=14)

# ------------------------------------------------------------
# 6. 신규 슬라이드 B — 기존 방법과 비교 (벤치마크, S22 뒤)
# ------------------------------------------------------------
sB = prs.slides.add_slide(blank)
for ph in list(sB.placeholders):
    remove_shape(ph)
for nm in ["Rectangle 1", "Rectangle 2", "TextBox 3", "TextBox 4"]:
    clone_shape(find(s10, nm), sB)
tsh = find(sB, "Rectangle 1")
p0 = tsh.text_frame.paragraphs[0]
for r in list(p0.runs):
    r._r.getparent().remove(r._r)
r = p0.add_run(); r.text = "  기존 방법과 비교 — 동일 데이터로 정면 평가"
r.font.size = Pt(20); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = C(WHITE)
sub = find(sB, "TextBox 3")
p0 = sub.text_frame.paragraphs[0]
for r in list(p0.runs):
    r._r.getparent().remove(r._r)
r = p0.add_run()
r.text = "기존 리셀 연구의 모델(선형회귀·랜덤포레스트)을 우리와 같은 데이터에 적용해 정량 비교"
r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = C(GRAY)

sB.shapes.add_picture(BENCH_PNG, Inches(0.35), Inches(1.20), width=Inches(7.55))
glance = box(sB, 0.35, 6.00, 7.55, 1.00, fill=NAVY)
set_text(glance, [
    ([("한눈에 보기   ", 12.5, AMBER, True),
      ("기존 선형회귀는 AUC 0.60~0.71에 그침,", 12.5, WHITE, False)], PP_ALIGN.LEFT, 1),
    ([("우리 모델은 0.84~0.96 — 머신러닝 + 미디어 채널 접근이 리셀 예측에 효과적",
       12.5, WHITE, False)], PP_ALIGN.LEFT, 0),
], anchor=MSO_ANCHOR.MIDDLE, ml=12)

bx, bw = 8.20, 4.80
tb = textbox(sB, bx, 1.20, bw, 0.40)
set_text(tb, [([("■ 비교한 방법 (출처)", 14, NAVY, True)], PP_ALIGN.LEFT, 0)])
tb = textbox(sB, bx, 1.62, bw, 1.95)
set_text(tb, [
    ([("선형회귀 ", 11.5, INK, True), ("[리셀 · Raditya 2021] ", 9, BLUE, True),
      ("— 직선으로 관계를 긋는 전통 통계", 10, GRAY, False)], PP_ALIGN.LEFT, 3),
    ([("랜덤포레스트 ", 11.5, INK, True), ("[리셀 · Raditya 2021] ", 9, BLUE, True),
      ("— 여러 결정트리를 합친 머신러닝", 10, GRAY, False)], PP_ALIGN.LEFT, 3),
    ([("우리 모델 ", 11.5, ORANGE, True), ("[주식→리셀] ", 9, ORANGE, True),
      ("— 주식에서 검증된 미디어 예측법(검색·감성)을 리셀에 적용", 10, GRAY, False)],
     PP_ALIGN.LEFT, 3),
    ([("※ Raditya 2021은 선형·RF 2개 모델만 비교 — 우리는 XGBoost·미디어 채널을 추가 확장",
       8.5, MGRAY, False)], PP_ALIGN.LEFT, 0),
])
tb = textbox(sB, bx, 3.62, bw, 0.40)
set_text(tb, [([("■ AUC가 뭔가요?", 14, NAVY, True)], PP_ALIGN.LEFT, 0)])
tb = textbox(sB, bx, 4.02, bw, 0.80)
set_text(tb, [
    ([("가격이 ", 11, INK, False), ('"오를지·내릴지" 방향', 11, INK, True),
      ("을 맞히는 정확도.", 11, INK, False)], PP_ALIGN.LEFT, 1),
    ([("0.5 = 찍기(동전) · 1.0 = 완벽 · 0.9↑ = 매우 잘 맞힘", 11, INK, False)],
     PP_ALIGN.LEFT, 0),
])
res = box(sB, bx, 4.95, bw, 2.05, fill=LORANGE, line=ORANGE, line_w=1.0)
set_text(res, [
    ([("■ 이 비교가 보여주는 것", 13, ORANGE, True)], PP_ALIGN.LEFT, 4),
    ([("· 트리 머신러닝이 전통 선형회귀보다 리셀 방향 예측에 효과적 (AUC +0.24~+0.36)",
       10.5, INK, False)], PP_ALIGN.LEFT, 2),
    ([("· 우리 모델은 랜덤포레스트와 대등 + 어떤 미디어가 선행하는지 해석(SHAP) 가능",
       10.5, INK, False)], PP_ALIGN.LEFT, 2),
    ([("· 동일 데이터·동일 검증 조건 → 객관적 평가 기준 확보", 10.5, INK, True)],
     PP_ALIGN.LEFT, 0),
], anchor=MSO_ANCHOR.MIDDLE, ml=10)
tb = textbox(sB, 0.35, 7.08, 12.60, 0.35)
set_text(tb, [([("※ 모든 모델 동일 데이터·동일 검증(시간순 교차검증 5-fold) 평균 AUC — 표·그래프 수치는 results/model_comparison_auc.csv 실측값",
                 9, GRAY, False)], PP_ALIGN.LEFT, 0)])

# ------------------------------------------------------------
# 7. 신규 슬라이드 위치 이동 (A -> index 12, B -> index 23)
# ------------------------------------------------------------
xml_slides = prs.slides._sldIdLst
ids = list(xml_slides)
elA, elB = ids[-2], ids[-1]
xml_slides.remove(elA); xml_slides.insert(12, elA)   # 파이프라인(12번째) 뒤
ids = list(xml_slides)
xml_slides.remove(elB); xml_slides.insert(23, elB)   # '왜 채널 선별...'(23번째) 뒤

# ------------------------------------------------------------
# 8. S17 — Granger 수치 풀이 보강
# ------------------------------------------------------------
cap = find(s17, "TextBox 6")
cap.height = Inches(1.35)
set_text(cap, [
    ([("그래프 해석  ", 9, BLUE, True),
      ("자산×채널 15쌍의 Granger F값 막대그래프 — 막대가 길수록 선행성 강함.", 9, GRAY, False)],
     PP_ALIGN.LEFT, 2),
    ([("F값", 9, INK, True),
      (" = 미디어의 과거값을 추가했을 때 가격 예측 오차가 줄어드는 정도 — 클수록 '먼저 움직이는' 신호가 강함.",
       9, GRAY, False)], PP_ALIGN.LEFT, 2),
    ([("p값", 9, INK, True),
      (" = 이 결과가 우연일 확률 — 0.05 미만이면 유의로 판정. 예: 스니커즈 뉴스감성 p=0.001은 우연일 확률 0.1%.",
       9, GRAY, False)], PP_ALIGN.LEFT, 0),
])

# ------------------------------------------------------------
# 9. S18 — AUC·ΔAUC 풀이 보강
# ------------------------------------------------------------
cap = find(s18, "TextBox 22")
cap.height = Inches(0.78)
tfp = cap.text_frame.add_paragraph()
tfp.space_before = Pt(2)
r = tfp.add_run()
r.text = ("AUC = 가격이 오를지·내릴지 방향을 맞히는 정확도(0.5=동전 던지기 · 1.0=완벽). "
          "ΔAUC −0.005는 0.5%p 차이 — 사실상 동일한 성능.")
r.font.size = Pt(9); r.font.name = FONT; r.font.color.rgb = C(GRAY)

# ------------------------------------------------------------
# 9b. S18 — 하단 결론 배너가 슬라이드 밖으로 잘리는 문제 수정
# ------------------------------------------------------------
for nm in ["TextBox 5", "TextBox 9"]:
    sh = find(s18, nm)
    if sh is not None:
        remove_shape(sh)
concl = box(s18, 0.30, 6.62, 12.73, 0.62, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
set_text(concl, [
    ([("결론   ", 12.5, AMBER, True),
      ("Granger 유의 채널만 쓴 모델 B(스니커즈 CH1+CH3+CH4 · 카드 CH1)는 전채널 모델 A보다 AUC가 높지 않음 "
       "→ 채널 선별이 예측력 향상으로 이어지지 않음", 12.5, WHITE, True)], PP_ALIGN.CENTER, 0),
], anchor=MSO_ANCHOR.MIDDLE)

# ------------------------------------------------------------
# 10. S21 — DM p값 풀이 보강
# ------------------------------------------------------------
cap = find(s21, "TextBox 999")
if cap is not None:
    tfp = cap.text_frame.add_paragraph()
    r = tfp.add_run()
    r.text = "DM p값이 클수록 '두 모델의 차이는 우연 범위'라는 뜻 — p=0.678이면 차이 없음에 매우 가까움."
    r.font.size = Pt(9); r.font.name = FONT; r.font.color.rgb = C(GRAY)

# ------------------------------------------------------------
# 11. S25 — 결론에서 Jaccard 제거
# ------------------------------------------------------------
sub = find(s25, "TextBox 3")
p0 = sub.text_frame.paragraphs[0]
for r in list(p0.runs):
    r._r.getparent().remove(r._r)
r = p0.add_run()
r.text = "Granger → 유의 채널 집합 비교 → XGBoost+DM, 세 단계 분석의 핵심 결과"
r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = C(GRAY)

rq2 = find(s25, "TextBox 15")
set_text(rq2, [
    ([("스니커즈 {뉴스감성·유튜브뷰·검색량} · 카드 {검색량} · 레고 {없음}", 13, INK, True)],
     PP_ALIGN.LEFT, 1),
    ([("— 유의 채널 구성이 자산마다 다름", 12, GRAY, False)], PP_ALIGN.LEFT, 0),
])

# ------------------------------------------------------------
# 12. 전역 치환 — '수치화'->'점수화', 'IRF'->'기존 방법 비교'
# ------------------------------------------------------------
for slide in prs.slides:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if "수치화" in r.text:
                    r.text = r.text.replace("수치화", "점수화")
                if "IRF" in r.text:
                    r.text = r.text.replace("IRF", "기존 방법 비교")

# ------------------------------------------------------------
# 13. 페이지 번호 재계산 (콘텐츠 슬라이드만, 'n / N' 통일)
# ------------------------------------------------------------
pat = re.compile(r"^\s*\d{1,2}(\s*/\s*\d{1,2})?\s*$")
numbered = []
for slide in prs.slides:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.left is None or sh.top is None:
            continue
        if sh.left < Inches(11.0) or sh.top < Inches(6.8):
            continue
        if pat.match(sh.text_frame.text.strip()):
            numbered.append(sh)
            break
total = len(numbered)
for i, sh in enumerate(numbered, 1):
    p = sh.text_frame.paragraphs[0]
    if p.runs:
        for r in list(p.runs)[1:]:
            r._r.getparent().remove(r._r)
        p.runs[0].text = "%d / %d" % (i, total)
    else:
        r = p.add_run(); r.text = "%d / %d" % (i, total)
        r.font.size = Pt(9); r.font.name = FONT; r.font.color.rgb = C(GRAY)

# ------------------------------------------------------------
# 검증 출력
# ------------------------------------------------------------
prs.save(DST)
print("saved:", DST, "| slides:", len(prs.slides._sldIdLst), "| numbered:", total)

chk = Presentation(DST)
bad = []
for i, slide in enumerate(chk.slides, 1):
    for sh in slide.shapes:
        texts = []
        if sh.has_text_frame:
            texts.append(sh.text_frame.text)
        if sh.has_table:
            texts += [c.text for row in sh.table.rows for c in row.cells]
        for t in texts:
            for kw in ["Jaccard", "jaccard", "BH", "Benjamini", "수치화", "IRF"]:
                if kw in t:
                    bad.append((i, kw, t[:60]))
print("금지어 검사:", "통과" if not bad else bad)
