"""
1) References 슬라이드(25)의 '데이터' 섹션에 실제 수집 플랫폼 5건 추가
   (StockX·PriceCharting·BrickRanker·Google Trends API·YouTube Data API v3)
   - FinBERT는 이미 [14] Araci(2019) 로 머신러닝·NLP 섹션에 인용되어 있어 중복 추가하지 않음
2) Granger 결과 슬라이드(10) 하단에 'RQ1·RQ2 지지' 결론 배너 추가
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PATH = "최종_버전_수정.pptx"

NAVY = RGBColor(0x1F, 0x2D, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DATA_PREFIX_COLOR = RGBColor(0x37, 0x47, 0x51)
DESC_COLOR = RGBColor(0x1A, 0x1A, 0x1A)
SEP_COLOR = RGBColor(0xCC, 0xCC, 0xCC)
DETAIL_COLOR = RGBColor(0x77, 0x77, 0x77)

NEW_DATA_REFS = [
    ("[44] ", "StockX (2025)", "Sneaker resale marketplace data, stockx.com"),
    ("[45] ", "PriceCharting (2025)", "Trading card price guide, pricecharting.com"),
    ("[46] ", "BrickRanker (2025)", "LEGO secondary-market index, brickranker.com"),
    ("[47] ", "Google Trends API (pytrends)", "trends.google.com"),
    ("[48] ", "YouTube Data API v3", "developers.google.com/youtube"),
]


def shape_by_id(slide, shape_id):
    for shp in slide.shapes:
        if shp.shape_id == shape_id:
            return shp
    raise KeyError(shape_id)


def add_ref_paragraph(text_frame, num, source_desc, detail):
    p = text_frame.add_paragraph()
    r1 = p.add_run()
    r1.text = num
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = DATA_PREFIX_COLOR
    r1.font.name = "Calibri"

    r2 = p.add_run()
    r2.text = source_desc
    r2.font.size = Pt(10)
    r2.font.bold = False
    r2.font.color.rgb = DESC_COLOR
    r2.font.name = "Calibri"

    r3 = p.add_run()
    r3.text = "  ·  "
    r3.font.size = Pt(10)
    r3.font.color.rgb = SEP_COLOR
    r3.font.name = "Calibri"

    r4 = p.add_run()
    r4.text = detail
    r4.font.size = Pt(10)
    r4.font.bold = False
    r4.font.color.rgb = DETAIL_COLOR
    r4.font.name = "Calibri"


def replace_textframe_keep_style(shape, new_text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    ref = p.runs[0]
    size, bold, color, name = ref.font.size, ref.font.bold, ref.font.color.rgb, ref.font.name
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = new_text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_textbox(slide, left, top, width, height, text, size, color, bold=False):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    box.fill.background()
    box.line.fill.background()
    return box


def main():
    prs = Presentation(PATH)

    # ---------- 1) References: 데이터 출처 5건 추가 ----------
    s25 = prs.slides[25]
    data_box = shape_by_id(s25, 31)   # '[43] Leetaru & Schrodt ...'
    for num, desc, detail in NEW_DATA_REFS:
        add_ref_paragraph(data_box.text_frame, num, desc, detail)
    # 박스 높이를 6줄 분량으로 보정 (auto_size가 PPT에서 다시 맞춰주지만 보험으로 갱신)
    line_h = data_box.height
    data_box.height = Emu(int(line_h) * 6)

    footer = shape_by_id(s25, 33)
    replace_textframe_keep_style(
        footer, "48 references  ·  리셀 시장 미디어 선행 지표 연구  ·  draft.tex 인용 + 데이터 출처 기준"
    )

    # ---------- 2) Granger 결과 슬라이드: RQ1·RQ2 지지 결론 배너 ----------
    s10 = prs.slides[10]
    bar_left, bar_top, bar_w, bar_h = Emu(182880), Emu(6300000), Emu(11201400), Emu(420000)
    add_rect(s10, bar_left, bar_top, bar_w, bar_h, NAVY)
    add_textbox(
        s10,
        Emu(int(bar_left) + 120000), bar_top, Emu(int(bar_w) - 240000), bar_h,
        "결론 RQ1·RQ2  |  BH 보정 후에도 스니커즈←CH3 뉴스감성, 카드←CH1 Google Trends 유의 → 미디어의 가격 선행성 부분 지지(RQ1)\n"
        "자산마다 유의한 채널 구성이 다름 (스니커즈 3개·카드 1개·레고 0개) → 채널 구성의 자산별 이질성 확인(RQ2)",
        10, WHITE, bold=True,
    )

    prs.save(PATH)
    print("saved:", PATH)


if __name__ == "__main__":
    main()
