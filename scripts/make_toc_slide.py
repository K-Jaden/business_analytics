"""최종_버전_수정.pptx 와 동일한 톤(네이비 헤더 + 블루 포인트)으로
'목차(Contents)' 슬라이드 한 장을 별도 파일로 생성한다.

원본 발표(19매)는 목차 슬라이드가 없어 신규 제작 — 6개 섹션으로 재구성.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1F, 0x2D, 0x4E)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ROW_ALT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)

SLIDE_W = Emu(12188825)
SLIDE_H = Emu(6858000)

SECTIONS = [
    ("01", "Introduction · 연구 배경과 질문",
     "선행 연구의 공백 · 연구 문제(RQ1–RQ3) · 연구 대상(15개 아이템) · 5개 미디어 채널"),
    ("02", "Methodology · 분석 설계",
     "ADF 정상성 검정 · Granger 인과 분석 설계 · XGBoost 모델 A/B/C 설계"),
    ("03", "RQ1 & RQ2 · Granger 인과 검정 결과",
     "자산×채널 15회 검정 (Benjamini-Hochberg 보정) · Jaccard 유사도로 본 채널 구성 차이"),
    ("04", "RQ3 · 예측 성능과 SHAP·DM 검정",
     "모델 A/B/C 성능 비교 · SHAP 채널 기여도 분석 · Diebold-Mariano 검정"),
    ("05", "강건성 검증과 연구 제한점",
     "충격반응함수(IRF) 분석 · 데이터·방법론 한계"),
    ("06", "결론과 시사점",
     "연구 문제별 결론 · 학술적·실무적 기여 · 참고 문헌"),
]


def add_textbox(slide, left, top, width, height, text, size, color,
                bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    box.line.fill.background()
    return box


def add_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 헤더 바
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Emu(658368), NAVY)
    add_textbox(slide, Emu(182880), Emu(0), Emu(8000000), Emu(658368),
                "목차", 22, WHITE, bold=True)
    add_rect(slide, Emu(0), Emu(658368), SLIDE_W, Emu(22860), BLUE)
    add_textbox(slide, Emu(182880), Emu(700000), Emu(11000000), Emu(280000),
                "발표 흐름 한눈에 보기 — 19매 슬라이드를 6개 섹션으로 구성",
                10, GRAY)

    # 섹션 행
    margin = Emu(500000)
    row_w = SLIDE_W - margin - margin
    row_h = Emu(760000)
    gap = Emu(130000)
    start_y = 1180000

    for i, (num, title, desc) in enumerate(SECTIONS):
        y = Emu(start_y + i * (760000 + 130000))
        bg = WHITE if i % 2 == 0 else ROW_ALT
        add_rect(slide, margin, y, row_w, row_h, bg)
        add_rect(slide, margin, y, Emu(60000), row_h, NAVY)

        add_textbox(slide, Emu(620000), y, Emu(750000), row_h,
                    num, 22, NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, Emu(1450000), Emu(int(y) + 180000), Emu(12000), Emu(400000), LIGHT_GRAY)

        add_textbox(slide, Emu(1620000), Emu(int(y) + 140000), Emu(9000000), Emu(320000),
                    title, 15, NAVY, bold=True)
        add_textbox(slide, Emu(1620000), Emu(int(y) + 460000), Emu(9000000), Emu(260000),
                    desc, 11, GRAY)

    out_path = "slide_toc_목차.pptx"
    prs.save(out_path)
    print(f"saved -> {out_path}")


if __name__ == "__main__":
    build()
