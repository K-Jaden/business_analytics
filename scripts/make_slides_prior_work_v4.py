#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
3장 구성 (v4) — 교집합 프레이밍
 ① 선행연구 3흐름 (미디어→주식 칸에 유재필·김문선 2023 추가)
 ② 방법론 교집합 비교 (유재필 + Olaniyan/MDPI vs 본 연구, XGBoost 이유)
 ③ 4축 차별점 (v3 7페이지 스타일)
검증된 실제 논문만 인용.
Usage : python scripts/make_slides_prior_work_v4.py
Output: slides_prior_work_v4.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slides_prior_work_v4.pptx")

NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE= RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
LGREEN= RGBColor(0xE7, 0xF2, 0xE9)
RED   = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK   = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 21, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 1 — 선행연구 3흐름 (미디어→주식 칸에 유재필 추가)
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "기존 연구 & 한계 — 세 연구 흐름의 빈 교점",
       [("미디어→주식 · 리셀 예측 · 멀티채널 — 셋 다 ", 13, False, DGRAY),
        ("‘리셀 × 다채널 × 선행성 비교’", 13, True, ORANGE),
        ("의 교점은 비워 둠", 13, False, DGRAY)])

C_T, C_H, C_W, c_gap = 1.45, 3.95, 4.13, 0.22
xs = [0.30, 0.30 + C_W + c_gap, 0.30 + 2 * (C_W + c_gap)]
streams = [
    ("미디어 → 주식 선행성",
     [("Tetlock (2007)", "뉴스 비관론 → S&P500 하락"),
      ("Da et al. (2011)", "Google 검색량 → 주가 선행"),
      ("유재필·김문선 (2023)", "유튜브 조회·검색 → 주가 선행 (Granger)")],
     "전통 금융자산(주식)에만 적용"),
    ("리셀 상품 가격 예측",
     [("Campello et al. (2021)", "StockX 스니커즈 가격 예측"),
      ("Dobrynskaya & Kishilova (2023)", "레고를 대체자산으로 분석")],
     "미디어 채널 신호 미포함 · 단일 자산"),
    ("멀티채널 예측 방법론",
     [("Bollen et al. (2011)", "Twitter 감성 → DJIA 선행"),
      ("Zhang et al. (2018)", "뉴스+감성+가격 결합 예측")],
     "주식 단일 · 채널 비교 없음 · Granger 미적용"),
]
for x, (sname, papers, limit) in zip(xs, streams):
    box(s1, x, C_T, C_W, C_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s1, x, C_T, C_W, 0.46, fill=BLUE, round_=True)
    text(s1, x, C_T + 0.03, C_W, 0.40, [[(sname, 12.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    pruns = []
    for (auth, desc) in papers:
        pruns.append([(auth, 11, True, NAVY)])
        pruns.append([("   " + desc, 10.5, False, INK)])
    text(s1, x + 0.12, C_T + 0.60, C_W - 0.24, 2.3, pruns, space=5)
    lb_t = C_T + C_H - 0.78
    box(s1, x + 0.12, lb_t, C_W - 0.24, 0.66, fill=RGBColor(0xFC, 0xF0, 0xE6),
        line=ORANGE, line_w=1.0, round_=True)
    text(s1, x + 0.18, lb_t + 0.02, C_W - 0.36, 0.62,
         [[("⚠ 한계  ", 10.5, True, ORANGE), (limit, 10.5, False, INK)]],
         anchor=MSO_ANCHOR.MIDDLE, space=1)

KEY_T = C_T + C_H + 0.18
box(s1, 0.30, KEY_T, 12.73, 0.74, fill=GREEN, round_=True)
text(s1, 0.45, KEY_T + 0.03, 12.4, 0.68,
     [[("본 연구  ", 13, True, WHITE),
       ("스니커즈·카드·레고 3자산 × 5채널 × Granger × SHAP+DM", 12, True, WHITE),
       (" — 세 흐름의 빈 교점을 메움", 12, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — 방법론 교집합 비교
# ════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
header(s2, "방법론 교집합 — 본 연구는 세 요소를 모두 결합",
       [("미디어 채널 · Granger 선별 · 예측 검증 — ", 13, False, DGRAY),
        ("선행연구는 일부만, 본 연구는 전부", 13, True, BLUE)])

def detail_card(x, t, w, h, tag, title, auth, lines):
    box(s2, x, t, w, h, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s2, x, t, w, 0.40, fill=BLUE, round_=True)
    text(s2, x, t + 0.02, w, 0.36, [[(tag, 11, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s2, x + 0.06, t + 0.44, w - 0.12, 0.74,
         [[(title, 11.5, True, NAVY)], [(auth, 9, False, MGRAY)]], space=2)
    body = [[("•  " + ln, 10, False, INK)] for ln in lines]
    text(s2, x + 0.10, t + 1.18, w - 0.20, h - 1.24, body, space=5)

cw2 = 6.16
cardx = [0.30, 0.30 + cw2 + 0.31]
ct, ch = 1.40, 1.92
detail_card(cardx[0], ct, cw2, ch, "국내 · Granger · 2023",
            "유튜브 기반 데이터 분석을 통한 주가 선행성 분석",
            "유재필·김문선 (2023), 정보화연구",
            ["채널 = 조회수·좋아요·댓글·검색빈도 (채널 단위)",
             "방법 = YouTube Data API V3 (상위 5채널·일별) → Granger",
             "결과 = 유튜브가 ETF(주식) 시장 선행 (예측모델 없음)"])
detail_card(cardx[1], ct, cw2, ch, "Granger + 예측 · Electronics 2024",
            "Granger 인과 + LSTM 기반 금융 예측",
            "Olaniyan 외 (2024), Electronics·MDPI",
            ["채널 = 가격 내부 변수(OHLCV·재무비율) · 미디어 없음",
             "방법 = Granger+상관분석 변수 선별 → 베이지안 최적화 LSTM",
             "결과 = 종가 예측, 기존 기법 대비 오차 감소"])

cx = [0.30, 1.95, 5.30, 8.30]
cw = [1.65, 3.35, 3.00, 4.73]
chead = ["구분", "유재필·김문선(2023)", "Olaniyan 외(2024)", "본 연구"]
tt, hh, rh = ct + ch + 0.16, 0.42, 0.52
for x, w, htxt in zip(cx, cw, chead):
    fill = GREEN if htxt == "본 연구" else NAVY
    box(s2, x, tt, w, hh, fill=fill, line=WHITE, line_w=1.0)
    text(s2, x, tt, w, hh, [[(htxt, 11, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("미디어 채널", "유튜브 (채널 단위)", "없음 (가격 변수)", "5채널 (검색·뉴스량·뉴스감성·조회수·댓글)"),
    ("Granger 선별", "사용 (선행성 검정)", "사용 (변수 선별)", "사용 (채널 선별)"),
    ("예측 모델", "없음 (선행성만)", "LSTM (딥러닝)", "XGBoost (소표본)"),
    ("분석 대상", "ETF (주식)", "중국 주식", "리셀 3종 (스니커즈·카드·레고)"),
]
for i, row in enumerate(rows):
    y = tt + hh + i * rh
    for j, (x, w) in enumerate(zip(cx, cw)):
        if j == 0:
            fill = RGBColor(0xEC, 0xEF, 0xF3); bold = True; col = NAVY; al = PP_ALIGN.CENTER; sz = 10.5
        elif j == 3:
            fill = LGREEN; bold = True; col = INK; al = PP_ALIGN.LEFT; sz = 10
        else:
            fill = WHITE if i % 2 == 0 else RGBColor(0xF7, 0xF8, 0xFA); bold = False; col = INK; al = PP_ALIGN.LEFT; sz = 10
        box(s2, x, y, w, rh, fill=fill, line=RGBColor(0xD5, 0xDA, 0xE0), line_w=0.75)
        text(s2, x + (0 if j == 0 else 0.06), y, w - (0 if j == 0 else 0.10), rh,
             [[(row[j], sz, bold, col)]], align=al, anchor=MSO_ANCHOR.MIDDLE)

note_t = tt + hh + 4 * rh + 0.14
box(s2, 0.30, note_t, 12.73, 0.62, fill=GREEN, round_=True)
text(s2, 0.45, note_t + 0.02, 12.4, 0.58,
     [[("본 연구만 셋을 모두 결합  ", 12.5, True, WHITE),
       ("— N=48 소표본이라 딥러닝(LSTM) 대신 과적합에 강한 ", 11, False, WHITE),
       ("XGBoost", 11, True, WHITE), (" 채택", 11, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 3 — 4축 차별점 (v3 7페이지 스타일)
# ════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
header(s3, "기존 연구와의 차별점 — 4가지 축",
       [("무엇이 새로운가 — ", 13, False, DGRAY),
        ("자산 · 채널 수 · 채널 선정 · 분석 방향", 13, True, BLUE),
        (" 네 축에서 차별화", 13, False, DGRAY)])

dx = [0.30, 2.85, 7.95]
dw = [2.55, 5.10, 5.08]
dhead = ["구분", "기존 연구", "본 연구 (제안)"]
dt, dh_h, drow_h = 1.60, 0.55, 1.10
for x, w, htxt in zip(dx, dw, dhead):
    fill = GREEN if htxt.startswith("본 연구") else NAVY
    box(s3, x, dt, w, dh_h, fill=fill, line=WHITE, line_w=1.0)
    text(s3, x, dt, w, dh_h, [[(htxt, 13, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

drows = [
    ("분석 자산", "주식 단일 자산", "리셀 수집품 3종 동시 비교 (스니커즈·카드·레고)"),
    ("채널 수",   "단일·소수 채널 (트위터 or 구글)", "다중 5채널 (보도량·감성 분리)"),
    ("채널 선정", "채널 사전 고정 후 그대로 사용 (자산별 비교 없음)", "5채널은 사전 지정 — 단, 자산별 선행 채널을 Granger로 비교·선별"),
    ("분석 방향", "‘선행하는가 Yes/No’", "XGBoost·SHAP로 기여 정량화 + 예측력 검증(DM)"),
]
for i, (a, b, c) in enumerate(drows):
    y = dt + dh_h + i * drow_h
    for j, (x, w, val) in enumerate(zip(dx, dw, (a, b, c))):
        if j == 0:
            fill = RGBColor(0xEC, 0xEF, 0xF3); bold = True; col = NAVY; al = PP_ALIGN.CENTER
        elif j == 2:
            fill = LGREEN; bold = True; col = INK; al = PP_ALIGN.LEFT
        else:
            fill = WHITE if i % 2 == 0 else RGBColor(0xF7, 0xF8, 0xFA); bold = False; col = DGRAY; al = PP_ALIGN.LEFT
        box(s3, x, y, w, drow_h, fill=fill, line=RGBColor(0xD5, 0xDA, 0xE0), line_w=0.75)
        text(s3, x + (0 if j == 0 else 0.12), y, w - (0 if j == 0 else 0.18), drow_h,
             [[(val, 11.5, bold, col)]], align=al, anchor=MSO_ANCHOR.MIDDLE)

ky3 = dt + dh_h + 4 * drow_h + 0.20
box(s3, 0.30, ky3, 12.73, 0.70, fill=GREEN, round_=True)
text(s3, 0.45, ky3 + 0.03, 12.4, 0.64,
     [[("본 연구 = ", 13, True, WHITE),
       ("다채널 × 리셀 × Granger 선별 × 예측 검증", 13, True, WHITE),
       (" — 선행연구가 비워 둔 교점을 메우는 최초 연구", 12, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
