# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu

p = Presentation(r"발표자료_최종_피드백반영_v2.pptx")
targets = [int(x) for x in sys.argv[1].split(",")]
for i, s in enumerate(p.slides, 1):
    if i not in targets:
        continue
    print(f"=== SLIDE {i} ===")
    for sh in s.shapes:
        l = Emu(sh.left).inches if sh.left is not None else -1
        t = Emu(sh.top).inches if sh.top is not None else -1
        w = Emu(sh.width).inches if sh.width is not None else -1
        h = Emu(sh.height).inches if sh.height is not None else -1
        txt = ""
        if sh.has_text_frame:
            txt = sh.text_frame.text.replace("\n", " / ")[:80]
        print(f"  {sh.name:28s} L{l:5.2f} T{t:5.2f} W{w:5.2f} H{h:5.2f} | {txt}")
    # font check on first text shape
print("layouts:", len(p.slide_layouts))
for j, ly in enumerate(p.slide_layouts):
    print(" layout", j, ly.name, "placeholders:", len(ly.placeholders))
