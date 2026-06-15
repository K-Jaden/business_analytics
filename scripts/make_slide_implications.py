#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
시사점 & 향후 연구 — 직관적 레이아웃 (1장)
좌: 시사점 — 대상별(투자자/운영자/연구자) 카드 3개, 아이콘 배지
우: 향후 연구 — 세로 로드맵 4단계 (번호 원 + 연결선)
원문 출처: 발표자료_최종_피드백반영_v3.pptx SLIDE28 "시사점 & 향후 연구" (쉬운 말로 재구성)
Usage : python scripts/make_slide_implications.py
Output: slide_implications_future.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide_implications_future.pptx")

NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE= RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
TEAL  = RGBColor(0x2A, 0x9D, 0x8F)
LGREEN= RGBColor(0xE7, 0xF2, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0xAA, 0xB1, 0xBC)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK   = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False, oval=False):
    shape_type = MSO_SHAPE.OVAL if oval else (MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE)
    shp = sld.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 21, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "시사점 & 향후 연구 — 그래서 무엇이 달라지는가",
       [("탐색적 발견을 ", 13, False, DGRAY), ("실무에서 참고할 점", 13, True, BLUE),
        ("과 ", 13, False, DGRAY), ("다음 단계 연구 방향", 13, True, ORANGE),
        ("으로 정리", 13, False, DGRAY)])

MARGIN = 0.30
GAP = 0.31
LEFT_X, LEFT_W = MARGIN, 6.16
RIGHT_X, RIGHT_W = MARGIN + LEFT_W + GAP, 13.33 - 2 * MARGIN - LEFT_W - GAP

LBL_T, LBL_H = 1.15, 0.46
BODY_T = LBL_T + LBL_H + 0.14   # 1.75
BANNER_H = 0.50
BODY_H = 7.5 - BODY_T - BANNER_H - 0.10   # 5.0

# ── 좌측 라벨 ────────────────────────────────────────────────────────────
box(s1, LEFT_X, LBL_T, LEFT_W, LBL_H, fill=BLUE, round_=True)
text(s1, LEFT_X, LBL_T, LEFT_W, LBL_H,
     [[("이 연구는 누구에게 도움이 될까?  ", 13, True, WHITE), ("(시사점)", 11, False, WHITE)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── 우측 라벨 ────────────────────────────────────────────────────────────
box(s1, RIGHT_X, LBL_T, RIGHT_W, LBL_H, fill=ORANGE, round_=True)
text(s1, RIGHT_X, LBL_T, RIGHT_W, LBL_H,
     [[("다음엔 무엇을 더 봐야 할까?  ", 13, True, WHITE), ("(향후 연구)", 11, False, WHITE)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════════
# 좌측 — 시사점 카드 3개 (대상별 배지)
# ════════════════════════════════════════════════════════════════════════
CARD_GAP = 0.16
CARD_H = (BODY_H - 1 * CARD_GAP) / 2   # ≈2.42

stakeholders = [
    (ORANGE, "운영자", "리셀 플랫폼을 운영한다면",
     ["모든 미디어 데이터를 다 모을 필요는 없음",
      "자산마다 효과 있는 채널 몇 개만 골라 운영해도 예측력은 비슷 → 데이터 수집·운영 비용 절감",
      "어떤 채널이 효과적인지는 자산마다 다름 — 스니커즈는 뉴스감성·검색량·조회수, 카드는 검색량 중심"]),
    (GREEN, "연구자", "학술 연구를 한다면",
     ["이번에 쓴 방법(공개데이터 수집 → 먼저 움직이는 신호 찾기 → 머신러닝 예측)은",
      "명품·시계 등 다른 수집품 시장에도 그대로 적용해볼 수 있음",
      "Granger·SHAP·DM처럼 이미 검증된 통계·머신러닝 기법을 그대로 재사용 가능"]),
]

for i, (badge_c, badge_txt, title_txt, lines) in enumerate(stakeholders):
    y = BODY_T + i * (CARD_H + CARD_GAP)
    box(s1, LEFT_X, y, LEFT_W, CARD_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    bd = 0.85
    box(s1, LEFT_X + 0.18, y + (CARD_H - bd) / 2, bd, bd, fill=badge_c, oval=True)
    text(s1, LEFT_X + 0.18, y + (CARD_H - bd) / 2, bd, bd, [[(badge_txt, 12.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s1, LEFT_X + 1.20, y + 0.16, LEFT_W - 1.34, 0.40, [[(title_txt, 15, True, NAVY)]])
    body = [[("•  " + ln, 12, False, INK)] for ln in lines]
    text(s1, LEFT_X + 1.20, y + 0.66, LEFT_W - 1.34, CARD_H - 0.72, body, space=8)

# ════════════════════════════════════════════════════════════════════════
# 우측 — 향후 연구 세로 로드맵 (1→4, 번호 원 + 연결선)
# ════════════════════════════════════════════════════════════════════════
ITEM_H = BODY_H / 3   # ≈1.667
CIRC_D = 0.56
CIRC_CX = RIGHT_X + 0.30   # 원 중심 x

# 연결선 (첫 원 중심 ~ 마지막 원 중심)
line_top = BODY_T + ITEM_H / 2
line_bot = BODY_T + 2 * ITEM_H + ITEM_H / 2
box(s1, CIRC_CX - 0.015, line_top, 0.03, line_bot - line_top, fill=MGRAY)

roadmap = [
    ("기간 늘리고 상품 수 늘리기",
     ["지금은 4년(48개월)·15개 상품 → 5년 이상, 더 많은 상품으로 검증하면", "결과를 더 믿을 수 있게 됨"]),
    ("뉴스 감성 분석 더 정교하게",
     ["지금은 범용 자동분석 도구 사용 → 금융 전문 AI로 바꾸면", "신호의 품질이 좋아질 것"]),
    ("월 단위 → 주 단위로 더 세밀하게",
     ["한정판 출시 같은 특별한 사건의 효과를 따로 떼어서", "분석할 수 있음"]),
]

for i, (title_txt, lines) in enumerate(roadmap):
    y = BODY_T + i * ITEM_H
    cy = y + (ITEM_H - CIRC_D) / 2
    box(s1, CIRC_CX - CIRC_D / 2, cy, CIRC_D, CIRC_D, fill=ORANGE, oval=True)
    text(s1, CIRC_CX - CIRC_D / 2, cy, CIRC_D, CIRC_D, [[(str(i + 1), 13, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx = CIRC_CX + CIRC_D / 2 + 0.20
    box(s1, tx, y + 0.04, RIGHT_X + RIGHT_W - tx, ITEM_H - 0.08, fill=LGRAY,
        line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    text(s1, tx + 0.14, y + 0.10, RIGHT_X + RIGHT_W - tx - 0.24, 0.32, [[(title_txt, 12, True, NAVY)]])
    body = [[("•  " + ln, 10, False, INK)] for ln in lines]
    text(s1, tx + 0.14, y + 0.46, RIGHT_X + RIGHT_W - tx - 0.24, ITEM_H - 0.50, body, space=2)

# ════════════════════════════════════════════════════════════════════════
# 하단 요약 배너
# ════════════════════════════════════════════════════════════════════════
BAN_T = BODY_T + BODY_H + 0.10
box(s1, MARGIN, BAN_T, 13.33 - 2 * MARGIN, BANNER_H, fill=GREEN, round_=True)
text(s1, MARGIN + 0.15, BAN_T, 13.33 - 2 * MARGIN - 0.3, BANNER_H,
     [[("한 줄 요약  ", 13, True, WHITE),
       ("시사점은 ", 12, False, WHITE), ("'운영·연구에 바로 적용 가능한 결과'", 12, True, WHITE),
       ("로, 향후 연구는 ", 12, False, WHITE), ("'더 확실하게 만드는 다음 단계'", 12, True, WHITE),
       ("로 — 탐색적 발견을 실전·후속 연구로 잇는 다리", 12, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
