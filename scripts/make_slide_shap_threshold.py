#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
'평균 중요도 임계값' = 기준 점선  — 개념 설명 중심 슬라이드 (1장)
핵심 비유: 반 평균 점수. 채널 5개 SHAP의 평균에 점선을 긋고, 그 위(초록)만 선택.
데이터: results/model_c_channel_selection.csv, results/model_c_threshold_results.csv
Usage : python scripts/make_slide_shap_threshold.py
Output: slide_shap_threshold.pptx
"""
import os
import io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(BASE, "slide_shap_threshold_v2.pptx")

for f in ["Malgun Gothic", "맑은 고딕"]:
    if any(f.lower() in fn.name.lower() for fn in fm.fontManager.ttflist):
        plt.rcParams["font.family"] = f
        break
plt.rcParams["axes.unicode_minus"] = False

NAVY = RGBColor(0x1F, 0x2D, 0x4E)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK = RGBColor(0x0A, 0x0A, 0x0A)

CH = ["CH1", "CH2", "CH3", "CH4", "CH5"]
SHAP = {
    "스니커즈": {"CH1": 0.3081, "CH2": 0.0, "CH3": 0.0, "CH4": 0.1306, "CH5": 0.2236},
    "카드":     {"CH1": 0.0588, "CH2": 0.0, "CH3": 0.0, "CH4": 0.3334, "CH5": 0.1145},
    "레고":     {"CH1": 0.0778, "CH2": 0.0562, "CH3": 0.0045, "CH4": 0.2058, "CH5": 0.0925},
}
SELECTED = {"스니커즈": ["CH1", "CH5"], "카드": ["CH4", "CH5"], "레고": ["CH4", "CH5"]}
DM_P = {"스니커즈": 0.678, "카드": 0.254, "레고": 0.646}
ASSETS = ["스니커즈", "카드", "레고"]

# ════════════════════════════════════════════════════════════════════════
# 그래프 — 21쪽 스타일(채널별 색 + ★) + 기준 점선(평균)
# ════════════════════════════════════════════════════════════════════════
CH_COLOR = {"CH1": "#1F8FE0", "CH2": "#E08A20", "CH3": "#9AA0A6",
            "CH4": "#E8B62C", "CH5": "#C0392B"}
CH_NAME = {"CH1": "GT", "CH2": "뉴스량", "CH3": "뉴스감성", "CH4": "YT조회", "CH5": "YT감성"}

fig, axes = plt.subplots(1, 3, figsize=(8.6, 2.95), dpi=200)
for ax, asset in zip(axes, ASSETS):
    d = SHAP[asset]
    vals = [d[c] for c in CH]
    mean_imp = sum(vals) / len(CH)
    sel_set = set(SELECTED[asset])
    xpos = np.arange(len(CH))
    for xi, c in zip(xpos, CH):
        sel = c in sel_set
        ax.bar(xi, d[c], width=0.70, color=CH_COLOR[c],
               alpha=1.0 if sel else 0.30, edgecolor="white", linewidth=0.6, zorder=3)
        ax.text(xi, d[c] + max(vals) * 0.03, f"{d[c]:.2f}", ha="center", va="bottom",
                fontsize=7.5, fontweight="bold" if sel else "normal",
                color="#222222" if sel else "#9AA0A6")
        if sel:
            ax.text(xi, d[c] + max(vals) * 0.13, "★", ha="center", va="bottom",
                    fontsize=10, color="#E8730F", zorder=4)
    # 기준 점선 (평균) — 강조. 라벨은 중앙(CH3, 값≈0) 위 빈 공간에 배치해 막대 값과 겹침 방지
    ax.axhline(mean_imp, color="#E8730F", linestyle="--", linewidth=1.8, zorder=5)
    ax.text(2, mean_imp + max(vals) * 0.04, f"기준선 {mean_imp:.2f}",
            color="#E8730F", fontsize=7.5, fontweight="bold", ha="center", va="bottom", zorder=6)
    ax.set_xticks(xpos)
    ax.set_xticklabels(CH, fontsize=8.5)
    ax.set_title(f"{asset}  →  {'+'.join(SELECTED[asset])}",
                 fontsize=10.5, fontweight="bold", pad=12, color="#1F2D4E")
    ax.set_ylim(0, max(vals) * 1.34)
    ax.set_xlim(-0.7, len(CH) - 0.3)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(axis="y", labelsize=7)
fig.suptitle("채널별 평균|SHAP|  —  ★ = 기준 점선(채널 평균) 위 → 선택 · 흐린 막대 = 탈락",
             fontsize=10.5, fontweight="bold", y=1.06, color="#1F2D4E")
plt.tight_layout()
fig_buf = io.BytesIO()
fig.savefig(fig_buf, dpi=200, bbox_inches="tight", facecolor="white")
plt.close(fig)

# ════════════════════════════════════════════════════════════════════════
# 슬라이드
# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 21, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "채널 선택 기준 = '평균 중요도 임계값' (기준 점선)",
       [("채널 5개 SHAP의 ", 13, False, DGRAY), ("평균에 점선을 긋고, 그 위만 선택", 13, True, ORANGE),
        (" — 반 평균 점수와 같은 원리", 13, False, DGRAY)])

MARGIN = 0.30

# ── 좌측: 개념 설명 (반 평균 비유) ────────────────────────────────────────
EXP_T, EXP_W, EXP_H = 1.30, 4.30, 3.82
box(s1, MARGIN, EXP_T, EXP_W, EXP_H, fill=RGBColor(0xEA, 0xF1, 0xFA), line=BLUE, line_w=1.0, round_=True)
text(s1, MARGIN + 0.18, EXP_T + 0.12, EXP_W - 0.36, 0.40,
     [[("'평균 중요도 임계값'이란?", 14, True, BLUE)]])
text(s1, MARGIN + 0.18, EXP_T + 0.60, EXP_W - 0.36, EXP_H - 0.72,
     [[("반 평균 점수와 똑같은 원리", 12, True, ORANGE)]],
     space=4)
text(s1, MARGIN + 0.18, EXP_T + 0.98, EXP_W - 0.36, EXP_H - 1.10,
     [[("①  채널 5개의 SHAP(기여도) 점수를", 11.5, False, INK)],
      [("     모두 더해 5로 나눔 → ", 11.5, False, INK), ("평균", 11.5, True, GREEN)],
      [("", 5, False, INK)],
      [("②  그 평균값에 ", 11.5, False, INK), ("기준 점선", 11.5, True, ORANGE), ("을 그음", 11.5, False, INK)],
      [("", 5, False, INK)],
      [("③  점선보다 ", 11.5, False, INK), ("위 = 선택", 11.5, True, GREEN),
       (" / ", 11.5, False, INK), ("아래 = 탈락", 11.5, True, RGBColor(0x99, 0x99, 0x99))],
      [("", 7, False, INK)],
      [("핵심: 내가 '2개'라고 정하는 게 아니라", 11, True, NAVY)],
      [("점선(평균)이 개수를 자동으로 정함", 11, True, NAVY)],
      [("→ 임의성 없는 객관적 기준", 11, False, DGRAY)]],
     space=3)

# ── 우측: 그래프 (기준 점선 주인공) ───────────────────────────────────────
FIG_X = MARGIN + EXP_W + 0.24
FIG_W = 13.33 - MARGIN - FIG_X
FIG_T = 1.30
FIG_BOX_H = 3.82
box(s1, FIG_X, FIG_T, FIG_W, FIG_BOX_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
from PIL import Image as _PILImage
fig_buf.seek(0)
_w_px, _h_px = _PILImage.open(fig_buf).size
ratio = _h_px / _w_px
pic_w = FIG_W - 0.26
pic_h = pic_w * ratio
if pic_h > FIG_BOX_H - 0.22:
    pic_h = FIG_BOX_H - 0.22
    pic_w = pic_h / ratio
fig_buf.seek(0)
s1.shapes.add_picture(fig_buf, Inches(FIG_X + (FIG_W - pic_w) / 2), Inches(FIG_T + 0.12),
                      width=Inches(pic_w))

# ── 하단: 결과 + DM 동등 ──────────────────────────────────────────────────
RES_T = EXP_T + EXP_H + 0.18
RES_H = 7.5 - RES_T - 0.62
box(s1, MARGIN, RES_T, 13.33 - 2 * MARGIN, RES_H, fill=RGBColor(0xE7, 0xF2, 0xE9),
    line=GREEN, line_w=1.0, round_=True)
text(s1, MARGIN + 0.18, RES_T + 0.08, 13.33 - 2 * MARGIN - 0.36, 0.32,
     [[("기준선을 적용한 결과", 12.5, True, GREEN)]])
text(s1, MARGIN + 0.18, RES_T + 0.42, 13.33 - 2 * MARGIN - 0.36, RES_H - 0.50,
     [[("•  세 자산 모두 ", 11.5, False, INK), ("기준선 위 채널이 정확히 2개씩", 11.5, True, GREEN),
       (" 선택됨  —  스니커즈 CH1+CH5 · 카드 CH4+CH5 · 레고 CH4+CH5", 11.5, False, INK)],
      [("•  선택 모델(2채널)을 전채널(5채널) 모델과 DM 검정 비교 → ", 11.5, False, INK),
       ("셋 다 동등", 11.5, True, GREEN),
       (f"  (스니커즈 p={DM_P['스니커즈']:.3f} · 카드 p={DM_P['카드']:.3f} · 레고 p={DM_P['레고']:.3f})", 11.5, False, INK)],
      [("※ 방법 출처: ", 9.5, False, DGRAY),
       ("scikit-learn SelectFromModel 기본 기준(threshold='mean')", 9.5, True, BLUE),
       ("  ·  Wang, Liang, Hancock & Khoshgoftaar (2024), Journal of Big Data 11(1):44", 9.5, False, DGRAY)]],
     space=6)

# ── 결론 배너 ─────────────────────────────────────────────────────────────
BAN_T = 7.5 - 0.52
box(s1, MARGIN, BAN_T, 13.33 - 2 * MARGIN, 0.40, fill=NAVY)
text(s1, MARGIN + 0.15, BAN_T, 13.33 - 2 * MARGIN - 0.3, 0.40,
     [[("결론  ", 12.5, True, WHITE),
       ("'평균(기준 점선)보다 기여가 큰 채널만 남긴다'", 11.5, True, WHITE),
       ("는 객관적 규칙 → 자산별 2채널로도 전채널과 동등 → 간결성 입증", 11.5, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
