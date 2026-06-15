#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RQ3 관련 슬라이드 2장 재작성 (v2 — 쉬운 말로)
 ① "왜 채널 선별이 예측력 향상으로 이어지지 않는가?" — Olaniyan(2024) 비교
    전문용어(OHLC+Volume, 자기상관 등) 없이 일상어로 재작성
 ② "연구 제한점" — 6개 한계를 쉬운 말로 재작성

데이터 출처: results/ablation_results.csv (실측값, 재학습 없음)
Usage : python scripts/make_slides_rq3_revised.py
Output: slides_rq3_revised.pptx (2 slides)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slides_rq3_revised_v3.pptx")

NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE= RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
LGREEN= RGBColor(0xE7, 0xF2, 0xE9)
LORANGE=RGBColor(0xFC, 0xF0, 0xE6)
RED   = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK   = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 20, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 1 — 왜 채널 선별이 예측력 향상으로 이어지지 않는가? (쉬운 말로)
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "왜 채널을 골라도 예측력이 좋아지지 않을까? — 비교 논문과 다른 이유",
       [("그랜저 검정으로 ", 13, False, DGRAY), ("'무엇을 고르느냐'", 13, True, BLUE),
        ("에 따라 결과가 달라짐 — 본 연구 데이터로 직접 확인", 13, False, DGRAY)])

# 상단 긴장(질문) 박스
TQ_T = 1.00
box(s1, 0.30, TQ_T, 12.73, 0.96, fill=LORANGE, line=ORANGE, line_w=1.5, round_=True)
text(s1, 0.45, TQ_T + 0.06, 12.4, 0.86,
     [[("비교 논문 (Olaniyan, 2024)  ", 12, True, ORANGE),
       ("그랜저 검정으로 가격 데이터(시가·고가·저가·종가·거래량) 중 중요한 것만 골라 모델에 넣음 → 예측 오차 감소", 11.5, False, INK)],
      [("본 연구  ", 12, True, ORANGE),
       ("그랜저 검정으로 미디어 신호(검색량·뉴스·유튜브 등) 중 먼저 움직이는 것만 골라 모델에 넣음 → 예측 정확도 변화 없음", 11.5, False, INK)],
      [("→  같은 ", 11.5, False, INK), ("'그랜저로 골라서 넣기'", 11.5, True, ORANGE),
       ("인데 결과가 다른 이유는 — ", 11.5, False, INK), ("고른 대상 자체가 다르기 때문", 11.5, True, ORANGE)]],
     anchor=MSO_ANCHOR.MIDDLE, space=2)

# 2개 카드: 무엇을 골랐는가
CARD_T = TQ_T + 0.96 + 0.16
CARD_H = 2.15
CARD_W = 6.16
cardx = [0.30, 0.30 + CARD_W + 0.31]

def compare_card(x, tag, title, lines, tag_fill):
    box(s1, x, CARD_T, CARD_W, CARD_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s1, x, CARD_T, CARD_W, 0.42, fill=tag_fill, round_=True)
    text(s1, x, CARD_T + 0.02, CARD_W, 0.38, [[(tag, 12, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s1, x + 0.10, CARD_T + 0.50, CARD_W - 0.20, 0.36, [[(title, 12, True, NAVY)]])
    body = [[("•  " + ln, 11, False, INK)] for ln in lines]
    text(s1, x + 0.12, CARD_T + 0.92, CARD_W - 0.24, CARD_H - 1.0, body, space=7)

compare_card(cardx[0], "비교 논문이 고른 것", "가격·거래량 데이터 자체 (시가·고가·저가·종가·거래량)",
             ["가격은 원래부터 예측의 핵심 재료 — 어제 가격으로 오늘 가격을 짐작 가능",
              "그랜저 검정은 '이 핵심 재료 중 어떤 걸 쓸지'만 정리한 것",
              "→ 핵심 정보는 그대로 모델에 남아서 예측이 좋아짐"],
             tag_fill=BLUE)
compare_card(cardx[1], "본 연구가 고른 것", "검색량·뉴스·유튜브 같은 미디어 신호 (CH1~5)",
             ["가격의 최근 흐름(상승/하락 추세) 정보는 모든 모델에 이미 들어있음",
              "미디어 신호는 그 위에 '추가'로 얹는 보조 정보일 뿐",
              "그랜저 검정은 '어떤 신호가 가격보다 먼저 움직이는가'를 찾은 것"],
             tag_fill=ORANGE)

# 근거 표 (ablation_results.csv 실측값, 쉬운 말로)
TB_T = CARD_T + CARD_H + 0.18
box(s1, 0.30, TB_T, 12.73, 0.40, fill=NAVY, round_=True)
text(s1, 0.45, TB_T, 12.4, 0.40,
     [[("참고  ", 11.5, True, ORANGE),
       ("3개 자산 평균 AUC — 표본이 적어(48개월) 나타난 패턴일 수 있어, 추후 더 많은 자산·기간으로 검증 필요", 11, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

rows3 = [
    ("가격 흐름만 사용", "미디어 신호 없이, 가격의 최근 흐름만으로 예측", "0.84 ~ 0.96"),
    ("가격 흐름 + 미디어 5개 전부", "미디어 신호(검색·뉴스·유튜브 등)를 다 추가", "0.84 ~ 0.96  (위와 거의 동일)"),
    ("미디어 신호만 사용", "가격의 최근 흐름 없이, 미디어 신호만으로 예측", "0.47 ~ 0.53  (= 동전 던지기 수준)"),
]
rh = 0.50
cxw = [(0.30, 2.60), (2.90, 6.20), (9.10, 3.93)]
for i, (a, b, c) in enumerate(rows3):
    y = TB_T + 0.40 + i * rh
    for j, (val, (x, w)) in enumerate(zip((a, b, c), cxw)):
        if j == 0:
            fill = RGBColor(0xEC, 0xEF, 0xF3); bold = True; col = NAVY; al = PP_ALIGN.CENTER
        elif j == 2 and i == 2:
            fill = LORANGE; bold = True; col = RED; al = PP_ALIGN.LEFT
        elif j == 2 and i == 1:
            fill = LGREEN; bold = True; col = GREEN; al = PP_ALIGN.LEFT
        else:
            fill = WHITE if i % 2 == 0 else RGBColor(0xF7, 0xF8, 0xFA); bold = False; col = INK; al = PP_ALIGN.LEFT
        box(s1, x, y, w, rh, fill=fill, line=RGBColor(0xD5, 0xDA, 0xE0), line_w=0.75)
        text(s1, x + (0 if j == 0 else 0.12), y, w - (0 if j == 0 else 0.18), rh,
             [[(val, 10.5, bold, col)]], align=al, anchor=MSO_ANCHOR.MIDDLE)

# 하단 결론 박스
KY_T = TB_T + 0.40 + 3 * rh + 0.16
box(s1, 0.30, KY_T, 12.73, 0.80, fill=GREEN, round_=True)
text(s1, 0.45, KY_T + 0.05, 12.4, 0.72,
     [[("결론  ", 13, True, WHITE),
       ("가격의 최근 흐름이 예측력의 거의 전부 — 미디어 신호는 ", 12, False, WHITE),
       ("'먼저 움직인다(선행성)'는 맞지만 예측에 추가 도움은 거의 없음", 12, True, WHITE)],
      [("비교 논문은 '가격 데이터 중 핵심을 정리'한 것이고, 본 연구는 '보조 신호가 먼저 움직이는지'를 확인한 것 — 서로 다른 것을 비교했으니 결과가 다른 건 당연함", 11, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE, space=2)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — 연구 제한점 (쉬운 말로)
# ════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
header(s2, "연구의 한계 — 솔직하게 짚어두는 6가지",
       [("과장된 결론을 막기 위해 한계를 투명하게 밝힘 — ", 13, False, DGRAY),
        ("탐색적 발견으로 한정", 13, True, BLUE)])

GRID_T = 1.15
GAP_X, GAP_Y = 0.20, 0.18
MARGIN = 0.30
COL_W = (13.33 - 2 * MARGIN - 2 * GAP_X) / 3
ROW_H = 2.55
xs = [MARGIN, MARGIN + COL_W + GAP_X, MARGIN + 2 * (COL_W + GAP_X)]
ys = [GRID_T, GRID_T + ROW_H + GAP_Y]


def limit_box(x, t, w, h, num, title, desc, fix):
    box(s2, x, t, w, h, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s2, x, t, w, 0.42, fill=BLUE, round_=True)
    text(s2, x, t + 0.02, w, 0.38, [[(f"{num}.  {title}", 11.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s2, x + 0.12, t + 0.50, w - 0.24, h - 1.00,
         [[("•  " + ln, 10, False, INK)] for ln in desc], space=4)
    box(s2, x + 0.12, t + h - 0.50, w - 0.24, 0.42, fill=LORANGE, line=ORANGE, line_w=0.75, round_=True)
    text(s2, x + 0.18, t + h - 0.48, w - 0.36, 0.38, [[("→  " + fix, 9.5, True, ORANGE)]],
         anchor=MSO_ANCHOR.MIDDLE)


limit_box(xs[0], ys[0], COL_W, ROW_H, 1, "표본이 적음 (48개월)",
          ["자산마다 4년치(48개월) 데이터만 사용",
           "이 정도로는 '확실한 증거'라고 말하기엔 부족함"],
          "단정적 결론 대신 '탐색적 발견'으로 제시")

limit_box(xs[1], ys[0], COL_W, ROW_H, 2, "스니커즈 검색량 데이터의 불안정함",
          ["스니커즈의 구글 검색량 데이터는 시간이 지나도",
           "패턴이 들쭉날쭉함 (경계선 수준)"],
          "해석 시 주의가 필요하다고 명시")

limit_box(xs[2], ys[0], COL_W, ROW_H, 3, "뉴스 감성 분석의 정밀도",
          ["원래는 금융 전문 AI로 뉴스 감성을 분석할 계획이었으나,",
           "범용 자동분석 도구로 대체 → 금융 특화도가 다소 낮음"],
          "한계로 명시 + 추가 검증(민감도 분석) 수행")

limit_box(xs[0], ys[1], COL_W, ROW_H, 4, "레고 일부 모델의 검색량 부족",
          ["레고 중 포르쉐·부가티 모델은",
           "구글 검색량이 0인 달이 많아 신호가 약함"],
          "이 모델들에는 강한 결론을 내리지 않음")

limit_box(xs[1], ys[1], COL_W, ROW_H, 5, "유튜브만 분석함",
          ["미디어 채널 중 유튜브만 분석",
           "틱톡·트위터(X) 같은 다른 플랫폼은 포함 안 됨"],
          "향후 연구에서 다른 플랫폼으로 확장 가능")

limit_box(xs[2], ys[1], COL_W, ROW_H, 6, "데이터 표준화 방식의 한계",
          ["머신러닝에 넣기 전 데이터를 표준화할 때",
           "전체 기간을 한 번에 사용함 (정확히는 학습 구간만 써야 함)"],
          "표본이 적어 영향은 작지만 한계로 명시")

ky2 = ys[1] + ROW_H + 0.18
box(s2, 0.30, ky2, 12.73, 0.66, fill=GREEN, round_=True)
text(s2, 0.45, ky2 + 0.02, 12.4, 0.62,
     [[("그래도  ", 13, True, WHITE),
       ("자산별·기간별로 나눠 다시 검증해봐도 같은 결과 → ", 12, False, WHITE),
       ("'확정된 결론'이 아닌 '앞으로 더 검증해볼 가설'로 제시", 12, True, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
