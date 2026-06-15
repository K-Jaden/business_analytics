#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DM 검정 전체 결과 그래프
Usage : python scripts/make_dm_graph.py
Output: slide_dm_results.pptx
"""
import os, io
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide_dm_results.pptx")

# ── 데이터 로드 ───────────────────────────────────────────────────────────
abl = pd.read_csv(os.path.join(BASE, "results", "ablation_results.csv"))
mc  = pd.read_csv(os.path.join(BASE, "results", "model_c_results.csv"))

ASSETS = ["sneakers", "cards", "lego"]
ASSET_KR = {"sneakers": "스니커즈", "cards": "카드", "lego": "레고"}

# 모델별 데이터 조립 (Channels-only 제외 — 스케일 극단)
MODELS = [
    ("Baseline",      "가격래그만\n(채널 없음)"),
    ("C",             "C\n(SHAP선별)"),
    ("A-dropGranger", "N-A\n(Granger제거)"),
    ("CH1-only",      "CH1\n단독"),
    ("CH2-only",      "CH2\n단독"),
    ("CH3-only",      "CH3\n단독"),
    ("CH4-only",      "CH4\n단독"),
    ("CH5-only",      "CH5\n단독"),
]

# asset → model → (dm_stat, dm_p, vs_A)
DATA = {}
for asset in ASSETS:
    DATA[asset] = {}
    sub = abl[abl["asset_type"] == asset]
    for mkey, _ in MODELS:
        if mkey == "C":
            row = mc[mc["asset_type"] == asset].iloc[0]
            dm_stat = float(row["dm_stat"])
            dm_p    = float(row["dm_p"])
            vs_A    = row["verdict"]
        else:
            row = sub[sub["model"] == mkey]
            if row.empty:
                dm_stat, dm_p, vs_A = 0.0, 1.0, "no diff"
            else:
                dm_stat = float(row["dm_stat"].iloc[0])
                dm_p    = float(row["dm_p"].iloc[0])
                vs_A    = str(row["vs_A"].iloc[0])
        DATA[asset][mkey] = (dm_stat, dm_p, vs_A)

# Channels-only DM stats (주석용)
CO_STATS = {a: abl[(abl["asset_type"]==a) & (abl["model"]=="Channels-only")]["dm_stat"].iloc[0]
            for a in ASSETS}

# ── 그래프 생성 ───────────────────────────────────────────────────────────
plt.rcParams["font.family"] = ["Malgun Gothic", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False

fig, axes = plt.subplots(1, 3, figsize=(15.5, 5.2), sharey=False)
fig.patch.set_facecolor("white")

# 색상: DM > 0 → A가 더 좋음(빨강), DM < 0 → 해당 모델이 더 좋음(파랑)
C_A_BETTER  = "#D94F4F"   # 빨강 (A 우수)
C_MOD_BETTER= "#2E6DA4"   # 파랑 (모델 우수)
C_NODIFF    = "#AAAAAA"   # 회색 (차이없음)
SIG_ALPHA   = 1.0         # 유의 (p<0.05)
NS_ALPHA    = 0.42        # 비유의

for ax, asset in zip(axes, ASSETS):
    mkeys  = [m for m, _ in MODELS]
    labels = [lbl for _, lbl in MODELS]
    stats  = [DATA[asset][m][0] for m in mkeys]
    pvals  = [DATA[asset][m][1] for m in mkeys]
    vs_as  = [DATA[asset][m][2] for m in mkeys]

    xs = np.arange(len(mkeys))

    for xi, (stat, pv, vsa, lbl) in enumerate(zip(stats, pvals, vs_as, labels)):
        sig = (pv < 0.05)
        alpha = SIG_ALPHA if sig else NS_ALPHA

        if stat > 0:
            col = C_A_BETTER
        elif stat < 0:
            col = C_MOD_BETTER
        else:
            col = C_NODIFF

        bar = ax.bar(xi, stat, 0.62, color=col, alpha=alpha,
                     edgecolor="white", linewidth=0.5, zorder=3)

        # 수치 레이블
        offset = 0.05 if stat >= 0 else -0.05
        va = "bottom" if stat >= 0 else "top"
        fw = "bold" if sig else "normal"
        ax.text(xi, stat + offset, f"{stat:.2f}",
                ha="center", va=va, fontsize=8.5,
                fontweight=fw, color="#1A1A1A")

        # 유의 별표
        if sig:
            star_y = stat + (0.25 if stat >= 0 else -0.25)
            ax.text(xi, star_y, "★", ha="center", va="center",
                    fontsize=9, color=col)

    # ±1.96 참조선
    ax.axhline( 1.96, color=C_A_BETTER,   linestyle="--", linewidth=1.3,
                alpha=0.6, zorder=4, label="±1.96 (유의 경계)")
    ax.axhline(-1.96, color=C_MOD_BETTER, linestyle="--", linewidth=1.3,
                alpha=0.6, zorder=4)
    ax.axhline(0, color="#333333", linewidth=1.0, zorder=4)

    # Channels-only 주석
    co = CO_STATS[asset]
    ax.annotate(f"채널만(CH-only)\nDM={co:.1f} → A 압도",
                xy=(0.98, 0.97), xycoords="axes fraction",
                ha="right", va="top", fontsize=7.5,
                color="#888888", style="italic",
                bbox=dict(boxstyle="round,pad=0.3", fc="white",
                          ec="#CCCCCC", lw=0.8, alpha=0.9))

    ymax = max(3.2, max(abs(s) for s in stats) * 1.35 + 0.5)
    ax.set_ylim(-ymax, ymax)
    ax.set_xlim(-0.6, len(mkeys) - 0.4)
    ax.set_xticks(xs)
    ax.set_xticklabels(labels, fontsize=8.5)
    ax.set_title(ASSET_KR[asset], fontsize=14, fontweight="bold", pad=8)
    ax.set_ylabel("DM 통계량" if asset == "sneakers" else "", fontsize=10)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.35, zorder=0)
    ax.set_facecolor("#F8F9FB")

    # 방향 라벨
    ax.text(-0.55, ymax * 0.85, "A가\n더 나쁨\n↑",
            ha="left", va="top", fontsize=7.5,
            color=C_A_BETTER, alpha=0.7, fontweight="bold")
    ax.text(-0.55, -ymax * 0.85, "A가\n더 좋음\n↓",
            ha="left", va="bottom", fontsize=7.5,
            color=C_MOD_BETTER, alpha=0.7, fontweight="bold")

# 범례
legend_handles = [
    mpatches.Patch(facecolor=C_MOD_BETTER, alpha=SIG_ALPHA,
                   label="모델이 A보다 유의미하게 우수 (p<0.05)"),
    mpatches.Patch(facecolor=C_A_BETTER,   alpha=SIG_ALPHA,
                   label="A가 모델보다 유의미하게 우수 (p<0.05)"),
    mpatches.Patch(facecolor=C_MOD_BETTER, alpha=NS_ALPHA,
                   label="비유의 — 모델 방향 (p≥0.05)"),
    mpatches.Patch(facecolor=C_A_BETTER,   alpha=NS_ALPHA,
                   label="비유의 — A 방향 (p≥0.05)"),
]
fig.legend(handles=legend_handles, loc="lower center", ncol=4,
           fontsize=9.5, bbox_to_anchor=(0.5, -0.01),
           framealpha=0.96, edgecolor="#cccccc")

fig.suptitle(
    "DM 검정 전체 결과 — 각 모델 vs Model A  (DM > 0: A 우수 / DM < 0: 해당 모델 우수 / 점선: ±1.96 유의 경계  |  ★ = p<0.05)",
    fontsize=11.5, fontweight="bold", y=1.01)

fig.text(0.5, -0.06,
    "비교 대상: 가격래그만(Baseline), Model C(SHAP선별), N-A(Granger유의채널제거), CH1~5단독  |  "
    "채널만(Channels-only)은 DM≥7.0으로 스케일 극단 → 별도 주석 처리",
    ha="center", fontsize=8.5, color="#888888", style="italic")

plt.tight_layout(rect=[0, 0.10, 1, 1])

fig_buf = io.BytesIO()
fig.savefig(fig_buf, dpi=165, bbox_inches="tight", facecolor="white")
plt.close(fig)

# ── PPT 생성 ─────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld = prs.slides.add_slide(prs.slide_layouts[6])

hbar = sld.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(0.72))
hbar.fill.solid(); hbar.fill.fore_color.rgb = NAVY; hbar.line.fill.background()
tf = hbar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  DM 검정 전체 결과 — 각 모델 vs Model A (전채널)"
r.font.bold=True; r.font.size=Pt(20); r.font.color.rgb=WHITE; r.font.name="Calibri"

sep = sld.shapes.add_textbox(Inches(0), Inches(0.72), Inches(13.33), Inches(0.025))
sep.fill.solid(); sep.fill.fore_color.rgb = BLUE; sep.line.fill.background()

sub = sld.shapes.add_textbox(Inches(0.2), Inches(0.755), Inches(13.0), Inches(0.30))
sub.fill.background(); sub.line.fill.background()
tf2 = sub.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run()
r2.text = ("DM > 0 → A(전채널)이 더 좋음  |  DM < 0 → 해당 모델이 더 좋음  |  "
           "점선 ±1.96 = 95% 유의 경계  |  ★ = p<0.05  |  진한색 = 유의, 연한색 = 비유의")
r2.font.size=Pt(8.5); r2.font.color.rgb=DGRAY; r2.font.italic=True; r2.font.name="Calibri"

fig_buf.seek(0)
sld.shapes.add_picture(fig_buf,
    Inches(0.10), Inches(1.05),
    Inches(13.13), Inches(6.33))

prs.save(OUT)
print(f"[OK] {OUT}")
