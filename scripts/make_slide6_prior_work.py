#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
슬라이드 6 재작성 — 기존 연구 & 한계 (선행연구 계보 + 연구 공백 + 본 연구 위치)
저작권 그래프 없이 텍스트/박스 기반. 단일 슬라이드 파일.
Usage : python scripts/make_slide6_prior_work.py
Output: slide6_prior_work_new.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide6_prior_work_new.pptx")

NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE= RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK   = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld = prs.slides.add_slide(prs.slide_layouts[6])
S = sld.shapes


def box(l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = S.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=2, fill=None):
    """runs: list of paragraphs; each paragraph = list of (txt, size, bold, color)."""
    tb = S.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


# ── 헤더 ──────────────────────────────────────────────────────────────────
box(0, 0, 13.33, 0.82, fill=NAVY)
text(0.25, 0.06, 12.8, 0.70,
     [[("기존 연구 & 한계 — 미디어 선행지표 연구의 계보", 23, True, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)
box(0, 0.82, 13.33, 0.045, fill=BLUE)
text(0.25, 0.92, 12.8, 0.40,
     [[("선행연구는 ‘감성 → Granger 선별 → 예측’ 패러다임을 ", 13, False, DGRAY),
       ("주식", 13, True, BLUE),
       ("에 적용해 성공 — 그러나 ", 13, False, DGRAY),
       ("리셀 등 대안자산의 다채널 선행성은 미답", 13, True, ORANGE)]],
     anchor=MSO_ANCHOR.MIDDLE)

# ── 3개 연구 카드 ─────────────────────────────────────────────────────────
CARD_T = 1.52
CARD_H = 3.35
CARD_W = 4.13
gap = 0.22
xs = [0.30, 0.30 + CARD_W + gap, 0.30 + 2 * (CARD_W + gap)]

cards = [
    {
        "tag": "전설적 시초 · 2011",
        "title": "Twitter mood predicts the stock market",
        "auth": "Bollen, Mao & Zeng (2011)",
        "lines": [
            "수천만 트윗을 6개 감정으로 분류",
            "→ 다우존스(DJIA)와 Granger 인과검정",
            "‘차분함(Calm)’이 주가를 2~4일 선행",
            "신경망 결합 시 방향예측 정확도 86.7%",
        ],
    },
    {
        "tag": "국내 적용 · 2019",
        "title": "뉴스·SNS 감성분석 + Granger 주가 예측",
        "auth": "국내 학술지 (2019)",
        "lines": [
            "뉴스·블로그·트위터 감성지수 추출",
            "→ Granger로 유의 채널만 선별",
            "→ LSTM 입력으로 예측 정확도 향상",
            "‘선별 후 딥러닝’ 파이프라인 정착",
        ],
    },
    {
        "tag": "최신 트렌드 · 2024",
        "title": "CausalStock: News-driven Causal Discovery",
        "auth": "OpenReview (2024)",
        "lines": [
            "LLM으로 뉴스·SNS 노이즈 제거",
            "→ 시차 기반 인과그래프 자동 탐색",
            "→ 다수 종목 움직임 동시 예측",
            "Granger를 넘어 종단간(end-to-end) 인과",
        ],
    },
]

for x, c in zip(xs, cards):
    box(x, CARD_T, CARD_W, CARD_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    # 태그 띠
    box(x, CARD_T, CARD_W, 0.42, fill=BLUE, round_=True)
    text(x, CARD_T + 0.02, CARD_W, 0.38, [[(c["tag"], 11.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 제목 + 저자
    text(x + 0.05, CARD_T + 0.48, CARD_W - 0.1, 0.95,
         [[(c["title"], 12.5, True, NAVY)],
          [(c["auth"], 10, False, RGBColor(0x77, 0x77, 0x77))]],
         space=3)
    # 본문 라인
    body = [[("•  " + ln, 11, False, INK)] for ln in c["lines"]]
    text(x + 0.08, CARD_T + 1.42, CARD_W - 0.16, CARD_H - 1.5, body, space=6)

# ── 공통 한계(연구 공백) 띠 ────────────────────────────────────────────────
GAP_T = CARD_T + CARD_H + 0.18
box(0.30, GAP_T, 12.73, 0.78, fill=RGBColor(0xFC, 0xF0, 0xE6),
    line=ORANGE, line_w=1.5, round_=True)
text(0.45, GAP_T + 0.04, 12.4, 0.70,
     [[("공통 한계 (연구 공백)   ", 12.5, True, ORANGE),
       ("① 모두 ", 11.5, False, INK), ("효율적 시장인 주식", 11.5, True, INK),
       ("에 한정   ② ", 11.5, False, INK), ("단일·소수 채널", 11.5, True, INK),
       (" 위주   ③ ", 11.5, False, INK),
       ("리셀 등 대안자산의 다채널 선행성은 검증된 바 없음", 11.5, True, INK)]],
     anchor=MSO_ANCHOR.MIDDLE, space=2)

# ── 본 연구 위치(핵심 메시지) ──────────────────────────────────────────────
KEY_T = GAP_T + 0.92
box(0.30, KEY_T, 12.73, 0.70, fill=GREEN, round_=True)
text(0.45, KEY_T + 0.03, 12.4, 0.64,
     [[("본 연구 ", 13, True, WHITE),
       ("= 동일 방법론(감성→Granger 선별→예측)을 ", 12.5, False, WHITE),
       ("리셀 3종(스니커즈·카드·레고)", 12.5, True, WHITE),
       ("에 ", 12.5, False, WHITE),
       ("5개 채널로 최초 적용", 12.5, True, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}")
