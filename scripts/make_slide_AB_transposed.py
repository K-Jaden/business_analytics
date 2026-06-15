#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
스크린샷 2026-06-07 002201 표 — 행/열 전치 버전
원본: 행=자산(스니커즈/카드/레고), 열=모델A AUC / 모델B AUC / ΔAUC / 해석
전치: 행=모델A AUC / 모델B AUC / ΔAUC / 해석, 열=자산

Usage : python scripts/make_slide_AB_transposed.py
Output: slide_AB_transposed.pptx
"""
import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide_AB_transposed.pptx")

# ── 색상 ──────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x33, 0x33, 0x33)

# 자산별 배경색 (스크린샷과 동일)
SNK_BG = RGBColor(0xD0, 0xE4, 0xF7)   # 연파랑 (스니커즈)
CRD_BG = RGBColor(0xD4, 0xED, 0xDA)   # 연초록 (카드)
LGO_BG = RGBColor(0xFD, 0xF3, 0xCF)   # 연노랑 (레고)
HDR_BG = RGBColor(0xD8, 0xDF, 0xEE)   # 열 헤더 배경

# ΔAUC 음수 = 빨강, 0 = 회색
RED  = RGBColor(0xC0, 0x00, 0x00)
DGRAY = RGBColor(0x66, 0x66, 0x66)

# ── 원본 데이터 (스크린샷 2026-06-07 002201 그대로) ──────────────────────
# 원본: 행=자산, 열=[모델A AUC, 모델B AUC, ΔAUC, 해석]
original = {
    "header_cols": ["모델 A AUC\n(전채널 5개)", "모델 B AUC\n(Granger 선별)", "ΔAUC\n(B−A)", "해석"],
    "rows": [
        {
            "asset":  "스니커즈",
            "A_AUC":  "0.841",
            "B_AUC":  "0.836",
            "delta":  "−0.005",
            "interp": "Granger 유의 채널(CH3)만\n투입해도 성능 소폭 하락만",
        },
        {
            "asset":  "카드",
            "A_AUC":  "0.859",
            "B_AUC":  "0.855",
            "delta":  "−0.004",
            "interp": "Granger 유의 채널(CH1)만\n투입 → A와 거의 동등",
        },
        {
            "asset":  "레고",
            "A_AUC":  "0.960",
            "B_AUC":  "0.960",
            "delta":  "0.000",
            "interp": "Granger 유의 채널 없음\n→ A와 동일 (채널 제외 불가)",
        },
    ],
}

ASSETS      = [r["asset"]  for r in original["rows"]]
A_AUCS      = [r["A_AUC"]  for r in original["rows"]]
B_AUCS      = [r["B_AUC"]  for r in original["rows"]]
DELTAS      = [r["delta"]  for r in original["rows"]]
INTERPS     = [r["interp"] for r in original["rows"]]
ASSET_BGS   = [SNK_BG, CRD_BG, LGO_BG]

# 전치 테이블: 5행 × 4열 (헤더+3자산)
# transposed[0] = 헤더행  ["", 스니커즈, 카드, 레고]
# transposed[1] = 모델A AUC 행
# transposed[2] = 모델B AUC 행
# transposed[3] = ΔAUC 행
# transposed[4] = 해석 행

row_hdrs  = original["header_cols"]   # 4개
row_datas = [A_AUCS, B_AUCS, DELTAS, INTERPS]

# ── 유틸 ──────────────────────────────────────────────────────────────────
def set_cell_bg(cell, rgb):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    sf   = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(sf,  qn("a:srgbClr"))
    srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")

def fill_cell(cell, text, bold=False, size=13, color=GRAY,
              align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.bold  = bold
    r.font.size  = Pt(size)
    r.font.name  = "Calibri"
    r.font.color.rgb = color

def set_border_white(tbl, n_rows, n_cols):
    for ri in range(n_rows):
        for ci in range(n_cols):
            tc   = tbl.cell(ri, ci)._tc
            tcPr = tc.get_or_add_tcPr()
            for tag in ["a:lnL","a:lnR","a:lnT","a:lnB"]:
                for old in tcPr.findall(qn(tag)):
                    tcPr.remove(old)
                ln   = etree.SubElement(tcPr, qn(tag))
                ln.set("w","9525"); ln.set("cap","flat"); ln.set("cmpd","sng")
                sf   = etree.SubElement(ln, qn("a:solidFill"))
                srgb = etree.SubElement(sf, qn("a:srgbClr"))
                srgb.set("val","FFFFFF")

# ── PPT 생성 ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld = prs.slides.add_slide(prs.slide_layouts[6])

# 헤더 바
hbar = sld.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(0.68))
hbar.fill.solid(); hbar.fill.fore_color.rgb = NAVY
hbar.line.fill.background()
tf = hbar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  RQ3 — Model A vs Model B AUC 비교 (행/열 전치)"
r.font.bold=True; r.font.size=Pt(21); r.font.color.rgb=WHITE; r.font.name="Calibri"

sep = sld.shapes.add_textbox(Inches(0), Inches(0.68), Inches(13.33), Inches(0.022))
sep.fill.solid(); sep.fill.fore_color.rgb = BLUE; sep.line.fill.background()

sub = sld.shapes.add_textbox(Inches(0.2), Inches(0.72), Inches(13.0), Inches(0.28))
sub.fill.background(); sub.line.fill.background()
tf2 = sub.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run()
r2.text = ("A=전채널(CH1~5+가격래그)  |  B=Granger raw p<0.05 유의 채널  "
           "|  스니커즈B: CH3 / 카드B: CH1 / 레고B: 없음 (→ B=A)")
r2.font.size=Pt(9); r2.font.color.rgb=RGBColor(0x55,0x55,0x55)
r2.font.italic=True; r2.font.name="Calibri"

# ── 테이블: 5행 × 4열 ─────────────────────────────────────────────────────
# 행: 헤더, 모델A AUC, 모델B AUC, ΔAUC, 해석
# 열: [행이름], 스니커즈, 카드, 레고
N_ROWS = 5
N_COLS = 4

col_w = [2.90, 3.15, 3.15, 3.15]   # 합 = 12.35
row_h = [0.72, 1.12, 1.12, 1.12, 1.88]   # 합 = 5.96

tbl_shape = sld.shapes.add_table(
    N_ROWS, N_COLS,
    Inches(0.49), Inches(1.06),
    sum(Inches(w) for w in col_w),
    sum(Inches(h) for h in row_h))
tbl = tbl_shape.table

for ci, w in enumerate(col_w):
    tbl.columns[ci].width = Inches(w)
for ri, h in enumerate(row_h):
    tbl.rows[ri].height = Inches(h)

# ── 행 0: 헤더 (열 이름 = 자산) ──────────────────────────────────────────
# cell(0,0): 빈 좌상단 (짙은 배경)
set_cell_bg(tbl.cell(0, 0), NAVY)
fill_cell(tbl.cell(0, 0), "", bold=True, size=12, color=WHITE)

for ci, (asset, bg) in enumerate(zip(ASSETS, ASSET_BGS), start=1):
    cell = tbl.cell(0, ci)
    set_cell_bg(cell, NAVY)
    fill_cell(cell, asset, bold=True, size=15, color=WHITE)

# ── 행 1~4: 데이터 ───────────────────────────────────────────────────────
for ri, (row_hdr, row_data) in enumerate(zip(row_hdrs, row_datas), start=1):
    # 첫 열: 행 이름
    cell0 = tbl.cell(ri, 0)
    set_cell_bg(cell0, HDR_BG)
    fill_cell(cell0, row_hdr, bold=True, size=12, color=NAVY)

    # 데이터 열
    for ci, (val, bg) in enumerate(zip(row_data, ASSET_BGS), start=1):
        cell = tbl.cell(ri, ci)
        set_cell_bg(cell, bg)

        # ΔAUC 행은 색으로 강조
        if row_hdr == "ΔAUC\n(B−A)":
            col = RED if val.startswith("−") or val.startswith("-") else DGRAY
            fill_cell(cell, val, bold=True, size=16, color=col)
        elif row_hdr in ["모델 A AUC\n(전채널 5개)", "모델 B AUC\n(Granger 선별)"]:
            fill_cell(cell, val, bold=True, size=17, color=GRAY)
        else:
            fill_cell(cell, val, bold=False, size=11, color=GRAY)

set_border_white(tbl, N_ROWS, N_COLS)

# ── 하단 주석 ─────────────────────────────────────────────────────────────
note = sld.shapes.add_textbox(Inches(0.49), Inches(7.14),
                               Inches(12.35), Inches(0.28))
note.fill.background(); note.line.fill.background()
tf3 = note.text_frame; p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.LEFT
r3 = p3.add_run()
r3.text = ("평가지표: AUC-ROC  |  교차검증: TimeSeriesSplit(n_splits=5)  "
           "|  레고 B=A: Granger 유의 채널 없음, 채널 제외 불가")
r3.font.size=Pt(8.5); r3.font.color.rgb=RGBColor(0x77,0x77,0x77)
r3.font.italic=True; r3.font.name="Calibri"

prs.save(OUT)
print(f"[OK] {OUT}")
