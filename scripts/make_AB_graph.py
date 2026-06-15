#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Model A vs B 단독 비교 그래프
Usage : python scripts/make_AB_graph.py
Output: slide_AB_comparison.pptx
"""
import os, io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide_AB_comparison.pptx")

# ── 데이터 ────────────────────────────────────────────────────────────────
# B = Granger raw p<0.05 유의 채널
#   스니커즈: CH1+CH3+CH4 (3채널) — 근사
#   카드:     CH1 (정확)
#   레고:     유의 채널 없음 = A
ASSETS   = ["스니커즈", "카드", "레고"]
A_VALS   = [0.8413, 0.8587, 0.9601]
B_VALS   = [0.8461, 0.8612, 0.9601]
B_CH     = ["CH1+CH3+CH4 †", "CH1", "없음 (= A)"]
DM_LABEL = ["2/5 A 우수\n3/5 차이없음", "1/5 A 우수\n4/5 차이없음", "B = A\n(비교불가)"]

C_A  = "#1F2D4E"   # 네이비   (A 채움)
C_B  = "#E87722"   # 오렌지   (B 점선)

# ── 그래프 생성 ───────────────────────────────────────────────────────────
plt.rcParams["font.family"] = ["Malgun Gothic", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False

fig, ax = plt.subplots(figsize=(12.5, 5.4))
fig.patch.set_facecolor("white")
ax.set_facecolor("#F8F9FB")

xs  = np.array([0, 1.8, 3.6])   # 자산 그룹 위치
bw  = 0.55                        # 막대 너비
off = 0.30                        # A/B 간격

# ── A 막대 (채움) ─────────────────────────────────────────────────────────
bars_a = ax.bar(xs - off, A_VALS, bw,
                color=C_A, edgecolor="white", linewidth=0.6,
                label="Model A  전채널 (CH1~5)", zorder=3)

# ── B 막대 (채움) ────────────────────────────────────────────────────────
bars_b = ax.bar(xs + off, B_VALS, bw,
                color=C_B, edgecolor="white", linewidth=0.6,
                label="Model B  Granger 유의 채널 (raw p<0.05)", zorder=3)

# ── 수치 레이블 ───────────────────────────────────────────────────────────
for bar, val in zip(bars_a, A_VALS):
    ax.text(bar.get_x() + bar.get_width()/2, val + 0.0008,
            f"{val:.4f}", ha="center", va="bottom",
            fontsize=10.5, fontweight="bold", color=C_A)

for bar, val, lego_same in zip(bars_b, B_VALS, [False, False, True]):
    label = f"{val:.4f}" + (" (= A)" if lego_same else "")
    ax.text(bar.get_x() + bar.get_width()/2, val + 0.0008,
            label, ha="center", va="bottom",
            fontsize=10.5, fontweight="bold", color=C_B)

# ── Δ AUC 화살표 + 텍스트 ────────────────────────────────────────────────
YMAX = 0.975
for xi, (av, bv) in enumerate(zip(A_VALS, B_VALS)):
    xc = xs[xi]
    delta = bv - av
    if abs(delta) > 0.0001:
        sign = "+" if delta > 0 else ""
        ax.annotate("",
            xy=(xc + off, bv + 0.003),
            xytext=(xc - off, av + 0.003),
            arrowprops=dict(arrowstyle="<->", color="#666666",
                            lw=1.3, connectionstyle="arc3,rad=0.0"))
        ax.text(xc, max(av, bv) + 0.006,
                f"Δ {sign}{delta:.4f}",
                ha="center", va="bottom", fontsize=9.5,
                color="#444444", fontweight="bold")


# ── B 채널 구성 표시 ──────────────────────────────────────────────────────
CH_Y = [0.838, 0.857, 0.958]
for xi, ch in enumerate(B_CH):
    ax.text(xs[xi] + off, CH_Y[xi], ch,
            ha="center", va="bottom", fontsize=8,
            color=C_B, style="italic")

# ── 축 설정 ───────────────────────────────────────────────────────────────
ax.set_xlim(-0.65, 4.25)
ax.set_ylim(0.815, YMAX)
ax.set_xticks(xs)
ax.set_xticklabels(ASSETS, fontsize=14, fontweight="bold")
ax.set_ylabel("AUC-ROC", fontsize=12)
ax.spines["top"].set_visible(False)
ax.spines["right"].set_visible(False)
ax.yaxis.grid(True, linestyle="--", alpha=0.4, zorder=0)
ax.tick_params(axis="y", labelsize=10)

# 레고 브레이크 구분 표시 (y축 값이 많이 다르므로)
ax.annotate("※ 레고는\n절대값 다름\n(~0.96)",
            xy=(3.6, 0.963), xytext=(4.05, 0.940),
            fontsize=8, color="#888888", style="italic",
            arrowprops=dict(arrowstyle="->", color="#BBBBBB", lw=1.0))

# ── 범례 ─────────────────────────────────────────────────────────────────
legend_handles = [
    mpatches.Patch(facecolor=C_A, edgecolor="white",
                   label="A  전채널 (CH1~5 + 가격래그)"),
    mpatches.Patch(facecolor=C_B, edgecolor="white",
                   label="B  Granger 유의 채널 (raw p<0.05)\n     스니커즈: CH1+CH3+CH4 / 카드: CH1 / 레고: 없음"),
]
ax.legend(handles=legend_handles, loc="upper left",
          fontsize=10, framealpha=0.96, edgecolor="#cccccc",
          bbox_to_anchor=(0.0, 1.0))

ax.set_title("Model A (전채널) vs Model B (Granger 유의 채널) — AUC-ROC 비교",
             fontsize=14, fontweight="bold", pad=12)

# 주석
fig.text(0.5, -0.01,
    "† 스니커즈 B: CH1+CH3+CH4 조합 AUC 미산출 → 개별 AUC 평균 근사 (0.8461)",
    ha="center", fontsize=8.5, color="#888888", style="italic")

plt.tight_layout()

fig_buf = io.BytesIO()
fig.savefig(fig_buf, dpi=170, bbox_inches="tight", facecolor="white")
plt.close(fig)

# ── PPT 생성 ─────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld = prs.slides.add_slide(prs.slide_layouts[6])

# 헤더
hbar = sld.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(0.72))
hbar.fill.solid(); hbar.fill.fore_color.rgb = NAVY
hbar.line.fill.background()
tf = hbar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  RQ3 — Model A vs Model B AUC 비교 (Granger raw p<0.05 채널 선별)"
r.font.bold=True; r.font.size=Pt(19); r.font.color.rgb=WHITE; r.font.name="Calibri"

sep = sld.shapes.add_textbox(Inches(0), Inches(0.72), Inches(13.33), Inches(0.025))
sep.fill.solid(); sep.fill.fore_color.rgb = BLUE; sep.line.fill.background()

sub = sld.shapes.add_textbox(Inches(0.2), Inches(0.755), Inches(13.0), Inches(0.30))
sub.fill.background(); sub.line.fill.background()
tf2 = sub.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run()
r2.text = ("A=전채널(CH1~5+가격래그)  |  B=Granger 유의 채널(스니커즈:CH1+CH3+CH4, 카드:CH1, 레고:없음→B=A)  "
           "|  A=네이비 / B=오렌지  |  Δ = B − A AUC")
r2.font.size=Pt(8.5); r2.font.color.rgb=DGRAY; r2.font.italic=True; r2.font.name="Calibri"

# 그래프 삽입 (전체 슬라이드 활용)
fig_buf.seek(0)
sld.shapes.add_picture(fig_buf,
    Inches(0.13), Inches(1.05),
    Inches(13.07), Inches(6.30))

prs.save(OUT)
print(f"[OK] {OUT}")
