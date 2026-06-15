import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

prs = Presentation('발표자료_최종_피드백반영_v2.pptx')

def set_run_text(para, new_text):
    if not para.runs:
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ''

def delete_shape(slide, shape_name):
    for shape in slide.shapes:
        if shape.name == shape_name:
            sp = shape._element
            sp.getparent().remove(sp)
            return True
    return False

def add_caption(slide, text, left_in, top_in, width_in, font_size=8.5):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    return tb

# =====================================================================
# SLIDE 18: Jaccard 설명 박스 추가
# 이미지: top=2.21, h=3.07 → bottom=5.28
# 오른쪽 패널: 6.8in~13in
# Rounded Rectangle 8 (Jaccard 유사도 박스) 수정
# =====================================================================
slide18 = prs.slides[17]
for shape in slide18.shapes:
    if shape.name == 'Rounded Rectangle 8':
        # Jaccard 결과값에 정의 추가
        paras = shape.text_frame.paragraphs
        set_run_text(paras[0], 'Jaccard 유사도란?')
        if len(paras) > 1:
            set_run_text(paras[1], '두 집합의 교집합 ÷ 합집합 (범위 0~1)')
        # 추가 줄이 없으면 새 paragraph 추가하기 어려우니 caption box로 추가
        break

# Jaccard 결과값 캡션 박스 별도 추가 (이미지 아래)
add_caption(slide18,
    '※ Jaccard = 교집합 채널 수 ÷ 합집합 채널 수  |  '
    '스니커즈∩카드={CH1} → 1÷3=0.333  |  레고∩기타=∅ → 0÷N=0.000',
    0.33, 5.6, 6.0, font_size=8.5)
print('Slide 18 Jaccard 설명 추가 완료')

# =====================================================================
# SLIDE 19: 캡션 위치 수정
# 이미지: top=1.62, h=3.94 → bottom=5.56in
# 기존 TextBox 999(top=4.30) 삭제 후 5.6in에 재배치
# =====================================================================
slide19 = prs.slides[18]
delete_shape(slide19, 'TextBox 999')

# "모델 A vs 모델 B AUC 비교" 텍스트(top=4.43, 이미지 안)도 이미지 아래로
for shape in slide19.shapes:
    if shape.name == 'TextBox 4':  # 그래프 제목 "모델 A vs 모델 B AUC 비교"
        shape.top = Inches(5.57)  # 이미지 아래로 이동

add_caption(slide19,
    '그림. Model A(전채널 5개) vs Model B(Granger 유의 채널)의 자산별 AUC-ROC 비교 막대그래프. '
    '각 막대: TimeSeriesSplit(5-fold) 교차검증 평균 AUC. 세 자산 모두 두 모델 간 차이 0.01 내외.',
    0.1, 5.57, 8.0, font_size=8.5)
print('Slide 19 캡션 이동 완료')

# =====================================================================
# SLIDE 20: 캡션 위치 확인
# 이미지: top=1.16, h=3.01 → bottom=4.17in
# SHAP 설명 텍스트들: top=4.31~ (이미지 아래) → 문제 없음
# TextBox 19(top=7.15): 너무 아래 → 5.5로 이동
# =====================================================================
slide20 = prs.slides[19]
for shape in slide20.shapes:
    if shape.name == 'TextBox 19':
        shape.top = Inches(4.20)  # 이미지 바로 아래
        shape.left = Inches(0.25)
        shape.width = Inches(12.5)
        break
print('Slide 20 캡션 이동 완료')

# =====================================================================
# SLIDE 21: 캡션 위치 수정
# 이미지: top=1.47, h=4.99 → bottom=6.46in
# 기존 TextBox 999(top=5.20, 이미지 안) 삭제
# 결론 TextBox 8(top=6.73)이 있어 캡션을 subtitle 위치에 병기
# =====================================================================
slide21 = prs.slides[20]
delete_shape(slide21, 'TextBox 999')

# 부제목(TextBox 3)에 캡션 내용을 간결하게 병기
for shape in slide21.shapes:
    if shape.name == 'TextBox 3':
        set_run_text(shape.text_frame.paragraphs[0],
            'A=전채널  |  B=Granger 유의 채널  |  C=SHAP 상위 2채널  |  N-A=유의 채널 제거  |  '
            '그림: 각 모델×자산별 AUC-ROC (5-fold 평균) — 모든 구성에서 최대 0.01 내외 차이')
        break
print('Slide 21 캡션 처리 완료')

# =====================================================================
# SLIDE 22: 캡션 위치 수정
# 이미지: top=1.21, h=4.05 → bottom=5.26in
# 기존 TextBox 999(top=5.10, 이미지 안) 삭제 후 5.3에 재배치
# =====================================================================
slide22 = prs.slides[21]
delete_shape(slide22, 'TextBox 999')

add_caption(slide22,
    '그림. DM(Diebold-Mariano) 검정: Model A(전채널)와 Model C(SHAP 상위 2채널)의 '
    '월별 예측오차 시계열을 비교하여 성능 차이가 통계적으로 유의한지 검정. '
    '3자산 모두 p>0.50 → H0 기각 실패(두 모델 간 유의미한 성능 차이 없음).',
    2.1, 5.35, 8.1, font_size=8.5)
print('Slide 22 캡션 이동 완료')

# =====================================================================
# SLIDE 23: 부제목 수정 (성능 차이 없음 → 왜 향상이 안 됐는가 맥락으로)
# =====================================================================
slide23 = prs.slides[22]
for shape in slide23.shapes:
    if shape.name == 'TextBox 3':
        set_run_text(shape.text_frame.paragraphs[0],
            'Granger 선별(B)·SHAP 선별(C) 모두 전채널(A)보다 예측력이 높지 않음 — '
            '세 가지 메커니즘으로 설명')
        break
print('Slide 23 수정 완료')

prs.save('발표자료_최종_피드백반영_v2.pptx')
print('\n저장 완료: 발표자료_최종_피드백반영_v2.pptx')
