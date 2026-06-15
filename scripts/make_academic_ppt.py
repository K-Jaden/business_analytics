#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
학술 발표 PPT v3 — 리셀 시장 미디어 선행 지표 연구
Usage : python scripts/make_academic_ppt.py
Output: 수업발표자료_v2.pptx  (15 slides)
"""
import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE  = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG   = os.path.join(BASE, "paper", "fig")
RES   = os.path.join(BASE, "results", "figures")
OUT   = os.path.join(BASE, "ppt_v3.pptx")
TOTAL = 16

# ─── 색상 팔레트 ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x2D, 0x4E)
BLUE   = RGBColor(0x2E, 0x6D, 0xA4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x0A, 0x0A, 0x0A)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
ALTROW = RGBColor(0xF5, 0xF5, 0xFB)
SIG    = RGBColor(0xC8, 0xE6, 0xC9)
SIG_HL = RGBColor(0xD4, 0xED, 0xDA)
MARG   = RGBColor(0xFF, 0xF1, 0xC0)
NS     = RGBColor(0xF5, 0xF5, 0xF5)
STAT   = RGBColor(0xCE, 0xEF, 0xD3)
DIFF   = RGBColor(0xFF, 0xE0, 0xB2)
SNK_BG = RGBColor(0xE3, 0xF2, 0xFD)
CRD_BG = RGBColor(0xE8, 0xF5, 0xE9)
LGO_BG = RGBColor(0xFF, 0xF9, 0xC4)
RAW_SG = RGBColor(0xFD, 0xF3, 0xD8)
CARD1  = RGBColor(0x2E, 0x6D, 0xA4)
CARD2  = RGBColor(0x27, 0x7A, 0x4A)
CARD3  = RGBColor(0x6C, 0x3D, 0x91)
BORDER = RGBColor(0xBB, 0xBB, 0xBB)
EXPBG  = RGBColor(0xF0, 0xF4, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════
# 헬퍼 함수
# ═══════════════════════════════════════════════════════════════════════════

def _set_cell_bg(cell, rgb):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    for c in list(tcPr):
        if c.tag.split('}')[-1] in ('solidFill','gradFill','noFill','blipFill','pattFill'):
            tcPr.remove(c)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    clr = etree.SubElement(sf, qn('a:srgbClr'))
    clr.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')
    tcPr.remove(sf); tcPr.insert(0, sf)


def add_box(slide, l, t, w, h, fill=None, border=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border:
        tb.line.color.rgb = BORDER; tb.line.width = Pt(0.75)
    else:
        tb.line.fill.background()
    return tb


def add_text(slide, text, l, t, w, h,
             bold=False, italic=False, size=12, color=BLACK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.bold = bold; r.font.italic = italic
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font
    return tb


def add_header(slide, title, sub=None):
    bar = add_box(slide, 0, 0, 13.33, 0.72, fill=NAVY)
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"  {title}"
    r.font.bold = True; r.font.size = Pt(22)
    r.font.color.rgb = WHITE; r.font.name = "Calibri"
    add_box(slide, 0, 0.72, 13.33, 0.025, fill=BLUE)
    if sub:
        add_text(slide, sub, 0.2, 0.755, 13.0, 0.33,
                 size=10, color=DGRAY, italic=True)


def add_pgnum(slide, n):
    add_text(slide, f"{n} / {TOTAL}", 12.45, 7.17, 0.85, 0.3,
             size=9, color=DGRAY, align=PP_ALIGN.RIGHT)


def add_img(slide, path, l, t, w, h, caption="", cap_sz=9):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        add_box(slide, l, t, w, h, fill=LGRAY)
        add_text(slide, f"[{os.path.basename(path)}]",
                 l+0.1, t+h/2-0.15, w-0.2, 0.35,
                 size=9.5, color=DGRAY, align=PP_ALIGN.CENTER)
    if caption:
        add_text(slide, f"주. {caption}", l, t+h+0.05, w, 0.45,
                 size=cap_sz, color=DGRAY, italic=True)


def build_table(slide, rows_data, col_w, l, t, h,
                hdr_bg=NAVY, hdr_fg=WHITE, hdr_sz=11, body_sz=10,
                row_bgs=None, center_cols=None):
    nr = len(rows_data); nc = len(col_w)
    center_cols = center_cols or set()
    tbl = slide.shapes.add_table(
        nr, nc, Inches(l), Inches(t), Inches(sum(col_w)), Inches(h)
    ).table
    for i, w in enumerate(col_w): tbl.columns[i].width = Inches(w)
    rh = Inches(h / nr)
    for i in range(nr): tbl.rows[i].height = rh
    for ri, row in enumerate(rows_data):
        is_hdr = (ri == 0)
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (is_hdr or ci in center_cols) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(txt)
            r.font.name = "Calibri"
            r.font.size = Pt(hdr_sz if is_hdr else body_sz)
            r.font.bold = is_hdr
            r.font.color.rgb = hdr_fg if is_hdr else BLACK
            if is_hdr:
                _set_cell_bg(cell, hdr_bg)
            elif row_bgs and (ri-1) < len(row_bgs) and row_bgs[ri-1]:
                _set_cell_bg(cell, row_bgs[ri-1])
            elif ri % 2 == 0:
                _set_cell_bg(cell, ALTROW)
            else:
                _set_cell_bg(cell, WHITE)


def add_card(slide, l, t, w, h, title, body, limitation, title_bg=NAVY):
    HDR_H = 0.54; SEP_H = 0.02; LIM_H = 1.15; PAD = 0.14
    add_box(slide, l, t, w, h, fill=WHITE, border=True)
    hdr = add_box(slide, l, t, w, HDR_H, fill=title_bg)
    tf = hdr.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.bold=True; r.font.size=Pt(12); r.font.color.rgb=WHITE; r.font.name="Calibri"
    body_t = t + HDR_H + 0.10
    body_h = h - HDR_H - 0.10 - SEP_H - 0.07 - LIM_H - 0.08
    add_text(slide, body, l+PAD, body_t, w-2*PAD, body_h, size=10, color=BLACK)
    sep_t = t + h - LIM_H - SEP_H - 0.07
    add_box(slide, l+PAD, sep_t, w-2*PAD, SEP_H, fill=RGBColor(0xCC,0xCC,0xCC))
    lim_t = sep_t + SEP_H + 0.07
    add_text(slide, "한계: " + limitation, l+PAD, lim_t, w-2*PAD, LIM_H,
             size=9.5, color=DGRAY, italic=True)


def expbox(slide, l, t, w, h, text, title=None):
    add_box(slide, l, t, w, h, fill=EXPBG, border=True)
    if title:
        add_text(slide, title, l+0.12, t+0.08, w-0.24, 0.38,
                 bold=True, size=12, color=NAVY)
        add_text(slide, text, l+0.12, t+0.50, w-0.24, h-0.62,
                 size=11, color=BLACK)
    else:
        add_text(slide, text, l+0.12, t+0.10, w-0.24, h-0.20,
                 size=11, color=BLACK)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 1 — 제목
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_box(sld, 0, 0, 13.33, 7.5, fill=NAVY)
add_box(sld, 0, 5.1, 13.33, 0.08, fill=BLUE)

add_text(sld, "리셀 시장에서 미디어 채널은\n가격을 선행하는가?",
         0.7, 0.85, 11.93, 2.4,
         bold=True, size=34, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sld, "스니커즈  ·  트레이딩 카드  ·  레고 — 다중 자산 Granger 인과 분석 및 XGBoost 검증",
         0.7, 3.30, 11.93, 0.7,
         size=18, color=RGBColor(0x9B,0xC4,0xE2), align=PP_ALIGN.CENTER)
add_text(sld, "2022.01 – 2025.12  |  48개월 × 15개 아이템 × 5개 미디어 채널",
         0.7, 4.10, 11.93, 0.55,
         size=14, color=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)
add_text(sld, "K-Jaden  |  파이낸스학과  |  2026.06",
         0.7, 6.65, 11.93, 0.45,
         size=12, color=RGBColor(0x77,0x99,0xBB), align=PP_ALIGN.CENTER)
add_pgnum(sld, 1)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — 선행 연구 갭 (3열 카드)
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "선행 연구의 공백",
    sub="세 연구 흐름이 각자 발전해 왔지만, 그 교점 — '리셀 시장 × 다채널 × 선행성 비교' — 은 아직 아무도 다루지 않았다")

CW = 4.10; CG = 0.245
CL = [0.25, 0.25+CW+CG, 0.25+2*(CW+CG)]
CT = 1.15; CH = 4.88

add_card(sld, CL[0], CT, CW, CH,
    title="미디어 채널 -> 가격 선행성",
    body=(
        "Tetlock (2007)  Journal of Finance\n"
        "  WSJ 뉴스 비관론 지수 -> S&P 500 수익률\n"
        "  뉴스 감성이 단기 주가를 예측함을 실증\n\n"
        "Da, Engelberg & Gao (2011)  J. Finance\n"
        "  Google 검색량(SVI) -> 개별 주가\n"
        "  Granger 인과 검정으로 선행성 공식 확인\n"
        "  검색량 급증 -> 미래 주가 상승 선행"
    ),
    limitation="전통 금융 자산(주식)에만 적용.\n리셀·대체 자산 사례 전무.",
    title_bg=CARD1)

add_card(sld, CL[1], CT, CW, CH,
    title="리셀 상품 가격 예측",
    body=(
        "Campello et al. (2021)  Procedia CS\n"
        "  StockX 스니커즈 + 출시 정보로 예측\n"
        "  RF·XGBoost 적용 -> 소셜 미디어 변수\n"
        "  추가 필요성을 제안에 그침\n\n"
        "Dobrynskaya & Kishilova (2023)\n"
        "  레고 세트 한정판 가격 상승 메커니즘\n"
        "  대체 자산으로서 레고의 특성 규명"
    ),
    limitation="미디어 채널 신호 미포함.\n단일 자산 클래스만 분석.",
    title_bg=CARD2)

add_card(sld, CL[2], CT, CW, CH,
    title="멀티채널 예측 방법론",
    body=(
        "Bollen et al. (2011)  J. Comp. Sci.\n"
        "  Twitter 감성 6차원 -> DJIA 지수\n"
        "  Granger + SVM 검증\n"
        "  Twitter가 주가를 이틀 선행함\n\n"
        "Zhang et al. (2018)  Knowledge-Based Sys.\n"
        "  뉴스 이벤트 + 감성 + 가격 정량 결합\n"
        "  멀티소스 ML로 주가 방향 예측"
    ),
    limitation="주식 단일 자산. 자산 간 채널 구성\n비교 없음. Granger 검정 미적용.",
    title_bg=CARD3)

banner_t = CT + CH + 0.16
add_box(sld, 0.25, banner_t, 12.83, 0.72, fill=NAVY)
add_text(sld,
    "이 연구 (2026):  스니커즈 · 카드 · 레고  3개 리셀 자산에 걸쳐  5개 미디어 채널 x Granger 인과 검정 x SHAP + DM 검증  — 세 흐름의 교점을 최초로 분석",
    0.38, banner_t+0.14, 12.57, 0.48,
    size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_pgnum(sld, 2)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 3 — 연구 문제 (RQ 1-3)
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "연구 문제 (RQ 1 - 3)",
    sub="시간적 선행성 -> 채널 이질성 -> 예측 간결성 — 세 단계로 리셀 시장 미디어 효과를 해부한다")

rq_data = [
    ["RQ", "핵심 질문", "사용한 검정", "측정 지표", "귀무가설"],
    ["RQ1",
     "5개 미디어 채널 중\n어떤 채널이 리셀 가격을\nGranger-cause 하는가?",
     "이변량 Granger 인과 검정\n-> 3자산 x 5채널 = 15회\nBH-FDR 보정 (q=0.05)",
     "F 통계량\n원본 p값\nBH 보정 p값",
     "H0: 채널이 가격을\n선행하지 않는다\n-> 기각 시 선행성 확인"],
    ["RQ2",
     "유의한 선행 채널의\n구성이 자산마다\n다른가?",
     "RQ1 유의 채널 집합 비교\n-> Jaccard 유사도 계산\n(3자산 쌍별 비교)",
     "Jaccard 지수\n(0 = 완전 다름\n1 = 완전 같음)",
     "H0: 모든 자산에서\n같은 채널이 선행\n-> 기각 시 자산별 채널 특이성"],
    ["RQ3",
     "Granger·SHAP 기반\n선별 채널 모델이\n전채널 모델과 동등한가?",
     "XGBoost 이진 분류\nModel A (전채널) vs Model B (선별)\nDiebold-Mariano 검정",
     "AUC-ROC\nDelta AUC\nDM p값",
     "H0: 예측 정확도 동일\n-> 기각 실패 시 간결성 확인\n(적은 채널로 같은 성능)"],
]
rq_bgs = [SNK_BG, CRD_BG, LGO_BG]
build_table(sld, rq_data,
    col_w=[0.75, 3.3, 3.1, 2.1, 3.0],
    l=0.25, t=1.12, h=5.75,
    hdr_bg=NAVY, hdr_sz=11, body_sz=10,
    row_bgs=rq_bgs, center_cols={0})
add_pgnum(sld, 3)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 4 — 연구 대상: 15개 아이템
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "연구 대상: 15개 리셀 아이템  (2022.01 - 2025.12, N = 48개월/아이템)",
    sub="스니커즈(StockX US Sz10 평균가) · 트레이딩 카드(PriceCharting PSA 10) · 레고(BrickRanker 새 봉인)")

items_data = [
    ["아이템 ID", "아이템 명", "자산", "출처", "가격 지표", "커버리지"],
    ["sneakers_jordan1", "Jordan 1 Retro High OG Bordeaux",    "스니커즈", "StockX",        "US Sz10 평균가", "48/48"],
    ["sneakers_panda",   "Nike Dunk Low Panda",                 "스니커즈", "StockX",        "US Sz10 평균가", "48/48"],
    ["sneakers_yeezy",   "Yeezy 350 V2 Zebra",                  "스니커즈", "StockX",        "US Sz10 평균가", "48/48 *"],
    ["sneakers_travis",  "Travis Scott Jordan 1",               "스니커즈", "StockX",        "US Sz10 평균가", "48/48"],
    ["sneakers_nb550",   "New Balance 550 White Green",         "스니커즈", "StockX",        "US Sz10 평균가", "48/48"],
    ["cards_charizard1", "Charizard VMAX Shining Fates SV107",  "카드",     "PriceCharting", "PSA 10 가격",    "48/48"],
    ["cards_umbreon",    "Umbreon VMAX Alt Art",                 "카드",     "PriceCharting", "PSA 10 가격",    "48/48"],
    ["cards_rayquaza",   "Rayquaza VMAX Alt Art",               "카드",     "PriceCharting", "PSA 10 가격",    "48/48"],
    ["cards_pikachu",    "Pikachu VMAX",                        "카드",     "PriceCharting", "PSA 10 가격",    "48/48"],
    ["cards_charizard2", "Charizard GX Hidden Fates",           "카드",     "PriceCharting", "PSA 10 가격",    "48/48"],
    ["lego_falcon",      "Millennium Falcon 75192",             "레고",     "BrickRanker",   "새 봉인 가격",   "48/48"],
    ["lego_hogwarts",    "Hogwarts Castle 71043",               "레고",     "BrickRanker",   "새 봉인 가격",   "48/48"],
    ["lego_titanic",     "Titanic 10294",                       "레고",     "BrickRanker",   "새 봉인 가격",   "48/48"],
    ["lego_porsche",     "Porsche 911 42096",                   "레고",     "BrickRanker",   "새 봉인 가격",   "48/48 +"],
    ["lego_bugatti",     "Bugatti Chiron 42083",                "레고",     "BrickRanker",   "새 봉인 가격",   "48/48 +"],
]
item_bgs = [SNK_BG]*5 + [CRD_BG]*5 + [LGO_BG]*5
build_table(sld, items_data,
    col_w=[2.3, 2.85, 1.25, 1.65, 2.0, 1.2],
    l=0.25, t=1.12, h=5.88,
    hdr_bg=NAVY, hdr_sz=10, body_sz=8.5,
    row_bgs=item_bgs, center_cols={0, 2, 3, 5})
add_text(sld,
    "* Yeezy: StockX 샘플링 아티팩트로 4개월 결측 -> 선형보간 처리    "
    "+ Porsche(21/48)·Bugatti(25/48): Google Trends 0주 비율 높음 — 제한점으로 명시",
    0.25, 7.1, 13.0, 0.36, size=8.5, color=DGRAY, italic=True)
add_pgnum(sld, 4)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 5 — 5개 미디어 채널
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "5개 미디어 채널 (CH1 - CH5)",
    sub="각 채널은 소비자·미디어 행동의 다른 측면을 포착 — Granger 투입 전 반드시 변환 적용")

ch_data = [
    ["ID", "채널명", "원본 데이터 출처", "변환 방식", "스케일", "개념적 프록시"],
    ["CH1", "Google Trends",
     "pytrends API (월평균, geo='')",
     "직접 사용 (변환 없음)",
     "0 - 100",
     "소비자 검색 관심도\n제품 인지도 대리 지표"],
    ["CH2", "뉴스 보도량",
     "GDELT timelinevolnorm\n(일별 수집 -> 월평균)",
     "아이템 내 최댓값으로\n0-1 정규화",
     "0 - 1",
     "미디어 노출 강도\n언론 보도 빈도"],
    ["CH3", "뉴스 감성",
     "GDELT timelinetone\n(일별 수집 -> 월평균)",
     "/10 정규화 -> -1~+1",
     "-1 ~ +1",
     "미디어 논조·감성\n긍정/부정 여론"],
    ["CH4", "유튜브 조회수",
     "YouTube Data API v3\n(resell+review 키워드, 월별)",
     "log1p 변환\n-> 아이템 내 z-score",
     "z-score",
     "정보 탐색 의도\n리뷰·리셀 관심도"],
    ["CH5", "유튜브 댓글 감성",
     "FinBERT (금융 특화 BERT)\n상위 15 영상 x 30댓글",
     "P(pos) - P(neg) 월평균",
     "-1 ~ +1",
     "소매 투자자 감성\n커뮤니티 온도"],
]
ch_bgs = [
    RGBColor(0xE3,0xF2,0xFD),
    RGBColor(0xF3,0xE5,0xF5),
    RGBColor(0xE8,0xF5,0xE9),
    RGBColor(0xFF,0xF9,0xC4),
    RGBColor(0xFB,0xE9,0xE7),
]
build_table(sld, ch_data,
    col_w=[0.7, 1.85, 2.6, 2.5, 1.2, 2.7],
    l=0.25, t=1.12, h=5.2,
    hdr_bg=NAVY, hdr_sz=11, body_sz=10,
    row_bgs=ch_bgs, center_cols={0, 4})
add_text(sld,
    "제한점  |  CH3: GDELT tone은 자동 추출 방식 (FinBERT 대비 도메인 특화도 낮음)  "
    "|  CH4: 'resell+review' 키워드 필터로 일반 hype 과소 계상 가능  |  두 항목 모두 논문 §5.3 명시",
    0.25, 6.45, 13.0, 0.52, size=9.5, color=DGRAY, italic=True)
add_pgnum(sld, 5)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 6 — ADF 전처리 (쉽게 설명)
# 좌: 개념 설명 / 우: 결과 요약 표
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "분석 전처리: ADF 단위근 검정으로 정상성 확보",
    sub="Granger 검정 투입 전 모든 시계열이 '정상(stationary)' 이어야 함 — 비정상 시계열 사용 시 허구 인과 위험")

# 좌측: 개념 설명 박스
expbox(sld, 0.25, 1.12, 4.5, 5.85,
    title="왜 정상성 검정이 필요한가?",
    text=(
        "Granger 검정의 전제 조건:\n"
        "분석하는 시계열이 '정상(I(0))'\n"
        "이어야 통계 검정이 유효함\n\n"
        "정상 시계열이란?\n"
        "-> 시간이 지나도 평균·분산 일정\n"
        "-> 추세나 계절성이 없는 상태\n\n"
        "비정상이면 어떤 문제?\n"
        "-> 관계없는 시계열끼리도\n"
        "   '유의하다'고 잘못 판단\n"
        "   (허구 상관: spurious correlation)\n\n"
        "검정 기준 (ADF 단위근 검정):\n"
        "  p < 0.05 -> 정상 (그대로 사용)\n"
        "  p >= 0.05 -> 비정상 -> 1차 차분\n"
        "  dx_t = x_t - x_(t-1) -> 재검정\n\n"
        "최종 결과:\n"
        "  정상 I(0): 9개 시계열\n"
        "  차분 후 I(1): 9개 시계열\n"
        "  예외: 스니커즈 CH1\n"
        "   차분 후에도 p=0.120 (경계)"
    ))

# 우측: 결과 표
adf_data = [
    ["자산", "시계열", "ADF p값", "정상?", "조치"],
    ["스니커즈", "mean_price",  "0.232", "비정상", "차분 (p=0.000)"],
    ["스니커즈", "score_ch1",   "0.926", "비정상", "차분 (p=0.120) *"],
    ["스니커즈", "score_ch2",   "0.000", "정상",   "—"],
    ["스니커즈", "score_ch3",   "0.007", "정상",   "—"],
    ["스니커즈", "score_ch4",   "0.002", "정상",   "—"],
    ["스니커즈", "score_ch5",   "0.000", "정상",   "—"],
    ["카드",     "mean_price",  "0.267", "비정상", "차분 (p=0.007)"],
    ["카드",     "score_ch1",   "0.216", "비정상", "차분 (p=0.000)"],
    ["카드",     "score_ch2",   "0.000", "정상",   "—"],
    ["카드",     "score_ch3",   "0.364", "비정상", "차분 (p=0.000)"],
    ["카드",     "score_ch4",   "0.305", "비정상", "차분 (p=0.000)"],
    ["카드",     "score_ch5",   "0.002", "정상",   "—"],
    ["레고",     "mean_price",  "0.113", "비정상", "차분 (p=0.000)"],
    ["레고",     "score_ch1",   "0.693", "비정상", "차분 (p=0.000)"],
    ["레고",     "score_ch2",   "0.000", "정상",   "—"],
    ["레고",     "score_ch3",   "0.000", "정상",   "—"],
    ["레고",     "score_ch4",   "0.000", "정상",   "—"],
    ["레고",     "score_ch5",   "0.035", "정상",   "—"],
]
adf_bgs = [STAT if r[3]=="정상" else DIFF for r in adf_data[1:]]
build_table(sld, adf_data,
    col_w=[1.5, 1.65, 1.35, 1.2, 2.55],
    l=4.95, t=1.12, h=5.85,
    hdr_bg=NAVY, hdr_sz=10, body_sz=9,
    row_bgs=adf_bgs, center_cols={0, 2, 3, 4})
add_pgnum(sld, 6)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 7 — RQ1: Granger 검정 결과
# 좌: 결과 표 / 우: F통계량 그래프
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "RQ1: Granger 인과 검정 — 전체 15회 결과",
    sub="이변량 VAR | 래그 수 AIC로 선택(1-4) | Benjamini-Hochberg FDR 보정 적용 (q=0.05, m=15)")

granger_data = [
    ["자산", "채널", "래그", "F값", "원본 p", "BH p", "최종 판정"],
    ["스니커즈", "CH3 뉴스감성",  "1", "13.385", "0.001", "0.011", "유의 (BH 통과)"],
    ["카드",     "CH1 Google",    "1",  "8.417", "0.006", "0.044", "유의 (BH 통과)"],
    ["스니커즈", "CH4 유튜브뷰",  "1",  "4.729", "0.035", "0.153", "원본 유의 / BH 탈락"],
    ["스니커즈", "CH1 Google",    "2",  "3.473", "0.041", "0.153", "경계 / BH 탈락"],
    ["레고",     "CH1 Google",    "1",  "3.922", "0.054", "0.162", "경계 p<0.10"],
    ["스니커즈", "CH2 뉴스량",    "1",  "1.551", "0.220",  "—",   "비유의"],
    ["스니커즈", "CH5 YT감성",    "1",  "0.091", "0.764",  "—",   "비유의"],
    ["카드",     "CH2 뉴스량",    "1",  "0.728", "0.398",  "—",   "비유의"],
    ["카드",     "CH3 뉴스감성",  "1",  "1.409", "0.242",  "—",   "비유의"],
    ["카드",     "CH4 유튜브뷰",  "1",  "0.247", "0.622",  "—",   "비유의"],
    ["카드",     "CH5 YT감성",    "1",  "0.614", "0.438",  "—",   "비유의"],
    ["레고",     "CH2 뉴스량",    "1",  "0.318", "0.576",  "—",   "비유의"],
    ["레고",     "CH3 뉴스감성",  "1",  "0.401", "0.530",  "—",   "비유의"],
    ["레고",     "CH4 유튜브뷰",  "1",  "1.833", "0.183",  "—",   "비유의"],
    ["레고",     "CH5 YT감성",    "1",  "0.173", "0.680",  "—",   "비유의"],
]
def _gbg(r):
    v=r[6]
    if "BH 통과" in v: return SIG
    if "BH 탈락" in v: return RAW_SG
    if "경계" in v:    return MARG
    return NS
granger_bgs = [_gbg(r) for r in granger_data[1:]]
build_table(sld, granger_data,
    col_w=[1.55, 1.6, 0.5, 0.95, 0.95, 0.85, 2.1],
    l=0.25, t=1.12, h=5.62,
    hdr_bg=NAVY, hdr_sz=10, body_sz=9,
    row_bgs=granger_bgs, center_cols={0,2,3,4,5,6})

add_img(sld, os.path.join(FIG, "fig_granger.png"),
        l=8.75, t=1.12, w=4.35, h=4.58,
        caption="F통계량 히트맵. 셀 색이 진할수록 F값 높음.\n굵은 테두리 = BH 보정 통과 (2개)")

add_text(sld,
    "결론  |  BH 보정 통과: 2/15  (스니커즈<-CH3 뉴스감성, 카드<-CH1 Google Trends)"
    "  |  원본만 유의: 2/15  |  비유의: 11/15",
    0.25, 6.85, 13.0, 0.4, size=10, color=NAVY, bold=True)
add_pgnum(sld, 7)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 8 — RQ2: 채널 이질성 분석
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "RQ2: 자산별 채널 구성이 다른가? — Jaccard 유사도",
    sub="RQ1 유의 채널 집합을 자산 쌍별로 비교 | 채널 구성이 달라야 '자산별 맞춤 미디어 전략'이 필요함")

expbox(sld, 0.25, 1.12, 5.5, 2.05,
    title="이 분석이 왜 필요한가?",
    text=(
        "만약 모든 자산에서 동일한 채널이 선행한다면\n"
        "-> 채널 선택 불필요, 5개 모두 모니터링하면 그만\n\n"
        "하지만 자산마다 다른 채널이 선행한다면\n"
        "-> 스니커즈 투자자 != 카드 투자자 != 레고 투자자\n"
        "-> 자산 특이적(asset-specific) 정보 채널 존재\n"
        "-> 맞춤형 미디어 전략이 의미 있음"
    ))

sets_data = [
    ["자산", "유의 채널 (원본 p<0.05)", "BH 보정 유의 채널"],
    ["스니커즈", "CH1·CH3·CH4  (3개)", "CH3 뉴스감성  (1개)"],
    ["카드",     "CH1  (1개)",           "CH1 Google Trends  (1개)"],
    ["레고",     "—  (없음)",            "—  (없음)"],
]
build_table(sld, sets_data,
    col_w=[1.8, 3.2, 3.2],
    l=5.95, t=1.12, h=1.52,
    hdr_bg=NAVY, hdr_sz=10, body_sz=10,
    row_bgs=[SNK_BG, CRD_BG, LGO_BG],
    center_cols={0,1,2})

jac_data = [
    ["Jaccard", "스니커즈", "카드",    "레고"],
    ["스니커즈",  "1.000",   "0.333",   "0.000"],
    ["카드",      "0.333",   "1.000",   "0.000"],
    ["레고",      "0.000",   "0.000",   "1.000"],
]
build_table(sld, jac_data,
    col_w=[1.65, 1.3, 1.3, 1.3],
    l=5.95, t=2.76, h=1.3,
    hdr_bg=NAVY, hdr_sz=10, body_sz=10,
    row_bgs=[SNK_BG, CRD_BG, LGO_BG],
    center_cols={0,1,2,3})

add_img(sld, os.path.join(FIG, "fig_jaccard.png"),
        l=0.25, t=4.17, w=12.83, h=2.55,
        caption=(
            "그림 설명: 좌측 = 자산 쌍별 Jaccard 유사도 막대 그래프 (높을수록 채널 구성이 유사)"
            "  |  우측 = 유의 채널 벤다이어그램 (원의 교집합이 공통 채널)"
            "  -> 스니커즈·카드는 CH1을 공유(J=0.333), 레고와는 교집합 없음(J=0.000)"
        ))
add_text(sld,
    "결론  |  J(스니커즈, 카드) = 0.333 — 부분 겹침(CH1만 공유)  |  J(*, 레고) = 0.000 — 레고는 완전히 다른 채널 구조  ->  RQ2 지지: 자산별 채널 구성 이질성 확인",
    0.25, 7.1, 13.0, 0.35, size=10, color=NAVY, bold=True)
add_pgnum(sld, 8)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 9 — RQ3: 모델 설계 (A + B)
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "RQ3: XGBoost 모델 설계 — 모델 A vs 모델 B",
    sub="전채널(모델A) -> SHAP으로 채널 기여도 분석 -> 상위 채널만 선별(모델B) -> DM 검정으로 성능 비교")

flow_txt = (
    "분석 순서\n\n"
    "STEP 1  모델 A 학습 (전채널 5개)\n"
    "   -> CH1~CH5 + 통제 피처 전부 투입\n"
    "   -> AUC 측정: 베이스라인 확보\n\n"
    "STEP 2  모델 B 학습 (Granger 선별)\n"
    "   -> Granger 인과 유의 채널만 투입\n"
    "   -> A 대비 성능 비교\n"
    "   -> '선행 채널 = 예측 채널?' 검증\n\n"
    "STEP 3  SHAP 분석 (모델 A 기반)\n"
    "   -> 각 채널의 실제 예측 기여도 측정\n"
    "   -> Granger 순위와 차이 있는지 확인\n\n"
    "STEP 4  모델 C 학습 (SHAP 선별)\n"
    "   -> SHAP 기여 상위 2채널만 투입\n"
    "   -> DM 검정으로 A vs C 최종 비교\n"
    "   -> 간결성 확보 여부 판단"
)
expbox(sld, 0.25, 1.12, 4.8, 5.85, title=None, text=flow_txt)

model_data = [
    ["구성 요소", "모델 A  (전채널)", "모델 B  (Granger 선별)", "모델 C  (SHAP 선별)"],
    ["투입 채널",
     "통제 피처\n+ CH1~CH5 전부\n(5채널)",
     "통제 피처\n+ Granger 유의 채널만",
     "통제 피처\n+ SHAP 상위\n2채널만"],
    ["자산별\n선별 채널",
     "모두 동일\n(5채널 전부)",
     "스니커즈: CH3\n카드: CH1\n레고: (없음)",
     "스니커즈: CH1+CH5\n카드: CH4+CH1\n레고: CH5+CH4"],
    ["통제 피처",
     "price_vs_ma3\nprice_chg_lag1~lag3",
     "동일", "동일"],
    ["교차검증",
     "TimeSeriesSplit\n(n_splits=5)",
     "동일", "동일"],
    ["비교 방식",
     "기준 모델\n(베이스라인)",
     "A 대비 AUC\n직접 비교",
     "A 대비\nDiebold-Mariano 검정"],
]
build_table(sld, model_data,
    col_w=[1.9, 2.1, 2.2, 2.2],
    l=5.25, t=1.12, h=5.85,
    hdr_bg=NAVY, hdr_sz=10.5, body_sz=9.5,
    row_bgs=[NS]*5, center_cols={0})
add_pgnum(sld, 9)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 10 — 모델 A & B 결과 + SHAP 분석
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "모델 A & B 결과 — 그리고 SHAP이 밝힌 것",
    sub="모델 A(전채널) vs 모델 B(Granger 선별 채널) AUC 비교 -> SHAP으로 실제 예측 기여 채널 파악 -> 모델 C 설계 근거")

# ── 상단 좌: 모델 A vs B AUC 비교 ──────────────────────────────────────
add_text(sld, "모델 A vs 모델 B  AUC 비교",
    0.25, 1.12, 6.2, 0.38, bold=True, size=12, color=NAVY)

ab_data = [
    ["자산",     "모델 A AUC\n(전채널 5개)", "모델 B AUC\n(Granger 선별)", "ΔAUC\n(B-A)", "해석"],
    ["스니커즈", "0.841",  "0.836",  "-0.005",
     "Granger 유의 채널(CH3)만\n투입해도 성능 소폭 하락만"],
    ["카드",     "0.859",  "0.855",  "-0.004",
     "Granger 유의 채널(CH1)만\n투입 → A와 거의 동등"],
    ["레고",     "0.960",  "0.960",   "0.000",
     "Granger 유의 채널 없음\n→ A와 동일 (채널 제외 불가)"],
]
ab_bgs = [SNK_BG, CRD_BG, LGO_BG]
build_table(sld, ab_data,
    col_w=[1.55, 1.65, 1.65, 1.1, 3.8],
    l=0.25, t=1.55, h=2.15,
    hdr_bg=NAVY, hdr_sz=10, body_sz=9.5,
    row_bgs=ab_bgs, center_cols={0,1,2,3})

# ── 상단 우: SHAP 핵심 발견 ──────────────────────────────────────────────
add_text(sld, "SHAP 분석 — 실제 예측 기여 채널 (모델 A 기반)",
    6.55, 1.12, 6.55, 0.38, bold=True, size=12, color=NAVY)

shap_box_data = [
    ("스니커즈", SNK_BG, CARD1,
     "SHAP 1위: CH1 Google (0.125)\nSHAP 2위: CH5 YT감성 (0.124)\nSHAP 3위: CH4 유튜브뷰 (0.080)\n\n"
     "주목: CH3 뉴스감성은 Granger 1위지만\nSHAP에서는 4위 (0.011)\n"
     "→ 선행성과 예측 기여도가 다름"),
    ("카드", CRD_BG, CARD2,
     "SHAP 1위: CH4 유튜브뷰 (0.091)\nSHAP 2위: CH1 Google (0.048)\nSHAP 3위: CH2 뉴스량 (0.031)\n\n"
     "CH1이 Granger 1위 & SHAP 2위\n→ 두 검정 일관성 확인\n"
     "→ CH4가 예측엔 더 강한 기여"),
    ("레고", LGO_BG, CARD3,
     "SHAP 1위: CH5 YT감성 (0.070)\nSHAP 2위: CH4 유튜브뷰 (0.036)\nSHAP 3위: CH1 Google (0.021)\n\n"
     "Granger 유의 채널 없었지만\nSHAP에서는 CH5, CH4 기여\n→ 선행성은 없으나 동시적 연관성 존재"),
]
bx_w = 2.1; bx_t = 1.55
for i, (label, bg, hdr_c, txt) in enumerate(shap_box_data):
    bx_l = 6.55 + i * (bx_w + 0.12)
    add_box(sld, bx_l, bx_t, bx_w, 2.15, fill=bg, border=True)
    add_text(sld, label, bx_l+0.08, bx_t+0.06, bx_w-0.16, 0.32,
             bold=True, size=11, color=hdr_c)
    add_text(sld, txt, bx_l+0.08, bx_t+0.40, bx_w-0.16, 1.68,
             size=8.5, color=BLACK)

# ── 하단 구분선 + SHAP 그래프 ────────────────────────────────────────────
add_box(sld, 0.25, 3.82, 12.83, 0.03, fill=BLUE)
add_text(sld, "SHAP 상위 2채널 -> 모델 C 구성 근거 (상세는 다음 슬라이드)",
    0.25, 3.90, 12.83, 0.32, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
add_img(sld, os.path.join(FIG, "fig_shap.png"),
        l=0.25, t=4.28, w=12.83, h=2.85,
        caption="자산별 채널 평균 |SHAP| 막대 그래프. 높을수록 XGBoost 예측에 실제로 기여한 채널. "
                "스니커즈는 CH1·CH5, 카드는 CH4·CH1, 레고는 CH5·CH4가 상위권 -> Model C 채널 선택 근거.")
add_pgnum(sld, 10)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 11 — 모델 C 결과 + DM 검정
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "모델 C 결과 & DM 검정 — SHAP 선별 채널로도 충분한가?",
    sub="SHAP 상위 2채널로 구성한 모델 C와 전채널 모델 A를 Diebold-Mariano 검정으로 통계적 비교")

# DM 검정 설명 박스 (좌)
expbox(sld, 0.25, 1.12, 3.5, 2.3,
    title="DM 검정이란?",
    text=(
        "두 예측 모델의 정확도가\n"
        "통계적으로 동일한지 검정\n\n"
        "H0: 두 모델의 예측 오차\n"
        "     크기가 같다\n\n"
        "p > 0.05 -> H0 기각 실패\n"
        "-> 더 간결한 모델(C) 선택\n"
        "    가능 (정확도 희생 없음)"
    ))

# DM 결과 표 (우)
dm_data = [
    ["자산",     "AUC 모델 A\n(전채널)",  "AUC 모델 C\n(SHAP 선별)",
     "선별 채널", "DeltaAUC",  "DM 통계량",  "p값",  "판정"],
    ["스니커즈", "0.8413", "0.8401",
     "CH1 + CH5", "-0.001", "-0.417", "0.678", "H0 기각 실패\n-> 동등"],
    ["카드",     "0.8587", "0.8590",
     "CH4 + CH1", "+0.000", "-0.670", "0.504", "H0 기각 실패\n-> 동등"],
    ["레고",     "0.9601", "0.9626",
     "CH5 + CH4", "+0.003", "+0.460", "0.646", "H0 기각 실패\n-> 동등"],
]
dm_bgs = [SIG_HL, SIG_HL, SIG_HL]
build_table(sld, dm_data,
    col_w=[1.45, 1.7, 1.7, 1.55, 1.05, 1.45, 0.9, 1.75],
    l=3.95, t=1.12, h=2.3,
    hdr_bg=NAVY, hdr_sz=10, body_sz=9.5,
    row_bgs=dm_bgs, center_cols={0,1,2,3,4,5,6})

add_img(sld, os.path.join(FIG, "fig_dm.png"),
        l=0.25, t=3.55, w=12.83, h=3.0,
        caption=(
            "DM 통계량 95% 신뢰구간 (점선: +-1.96). 세 자산 모두 |DM| < 1.96 구간 내 -> "
            "모델 A(전채널)와 모델 C(SHAP 2채널) 간 예측 정확도 차이 통계적으로 유의하지 않음."
        ))
add_text(sld,
    "결론 RQ3  |  SHAP 상위 2채널만으로 구성한 모델 C가 전채널 모델 A와 동등한 정확도 달성 (DM p > 0.50, 3자산 전부)  ->  간결성(parsimony) 확인",
    0.25, 6.88, 13.0, 0.38, size=10.5, color=NAVY, bold=True)
add_pgnum(sld, 11)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 12 — 왜 모델 B/C 결과가 이렇게 나왔나? (해석)
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "왜 모델 B·C가 모델 A와 동등한 결과를 냈는가?",
    sub="Granger 선별(B)과 SHAP 선별(C) 모두 전채널(A)과 성능 차이 없음 — 세 가지 메커니즘으로 설명")

# 박스 1: 채널 간 중복 정보
expbox(sld, 0.25, 1.12, 4.15, 2.6,
    title="① 채널 간 정보 중복",
    text=(
        "CH1(Google)과 CH4(유튜브뷰)는\n"
        "모두 '소비자 관심도'를 측정\n"
        "-> 높은 상관 -> 한 채널로도\n"
        "   충분한 정보 커버 가능\n\n"
        "XGBoost는 이미 중복 채널에\n"
        "낮은 가중치를 자동 부여\n"
        "-> 채널 추가해도 성능 미변화"
    ))

# 박스 2: Granger vs SHAP 불일치 원인
expbox(sld, 4.57, 1.12, 4.15, 2.6,
    title="② Granger ≠ SHAP 순위 불일치",
    text=(
        "Granger(선형 VAR)와\nXGBoost SHAP(비선형)은\n다른 것을 측정\n\n"
        "스니커즈 CH3:\n"
        "  Granger F=13.4 (1위)\n"
        "  SHAP 0.011 (4위)\n"
        "-> 선형 선행성 강하지만\n"
        "   비선형 예측 기여는 낮음\n"
        "-> 두 지표는 상호 보완적"
    ))

# 박스 3: 소표본 한계
expbox(sld, 8.89, 1.12, 4.19, 2.6,
    title="③ 소표본(N=48) 효과",
    text=(
        "N=48개월의 한계:\n"
        "채널 간 미세한 성능 차이를\n검출할 통계적 검정력 부족\n\n"
        "DM 검정 p값이 모두 0.50 이상\n"
        "-> '차이 없음' 결론이지만\n"
        "   '차이가 없다'는 증명도 아님\n\n"
        "N=200+이면 미세 차이도\n통계적으로 탐지 가능"
    ))

# 하단: 자산별 해석
add_box(sld, 0.25, 3.88, 12.83, 0.03, fill=BLUE)
add_text(sld, "자산별 해석",
    0.25, 3.97, 12.83, 0.38, bold=True, size=13, color=NAVY)

interp_data = [
    ["자산", "모델 B 결과 해석\n(Granger 선별)", "모델 C 결과 해석\n(SHAP 선별)", "핵심 인사이트"],
    ["스니커즈",
     "CH3만 투입 -> AUC 소폭 하락 (-0.005)\n나머지 4채널도 잔여 정보 보유",
     "CH1+CH5 투입 -> AUC 거의 동일 (-0.001)\nSHAP이 실제 기여 채널 더 정확히 파악",
     "Granger 선행 채널(CH3)과\n예측 기여 채널(CH1·CH5)이 다름\n-> 선행성 != 예측력"],
    ["카드",
     "CH1만 투입 -> AUC 소폭 하락 (-0.004)\nGranger 유의 채널 1개로 거의 유지",
     "CH4+CH1 투입 -> AUC 오히려 소폭 상승\nSHAP이 CH4의 잠재력 발견",
     "Granger·SHAP 모두 CH1 중요\n-> 가장 일관성 있는 자산\n두 방법론 일치"],
    ["레고",
     "Granger 유의 채널 없음\n-> B = A 동일 (비교 불가)",
     "CH5+CH4 투입 -> AUC 소폭 상승 (+0.003)\n선행성 없어도 예측 기여 채널 존재",
     "선행성(Granger) = 없음\n예측 기여(SHAP) = 존재\n-> 동시적 연관성만 있음"],
]
interp_bgs = [SNK_BG, CRD_BG, LGO_BG]
build_table(sld, interp_data,
    col_w=[1.5, 3.45, 3.45, 3.85],
    l=0.25, t=4.40, h=2.87,
    hdr_bg=NAVY, hdr_sz=10, body_sz=9,
    row_bgs=interp_bgs, center_cols={0})
add_pgnum(sld, 12)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 13 — IRF 강건성 분석
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "추가 강건성 검증: 충격반응함수 (IRF) 분석",
    sub="Granger 검정이 '유의한지 여부'만 알려주는 한계를 보완 — 채널 충격의 방향·크기·지속 기간까지 추적")

expbox(sld, 0.25, 1.12, 4.0, 2.65,
    title="왜 IRF를 추가했나?",
    text=(
        "Granger 검정의 한계:\n"
        "'선행한다 / 아니다'만 판단\n"
        "-> 효과의 방향·크기·지속 기간 불명\n\n"
        "IRF(충격반응함수)로 보완:\n"
        "채널 값이 갑자기 1sigma 증가 시\n"
        "-> 가격이 이후 6개월 어떻게 반응?\n"
        "-> 방향(위/아래), 크기, 소멸 시점 파악"
    ))

expbox(sld, 4.45, 1.12, 4.1, 2.65,
    title="분석 대상 및 설정",
    text=(
        "대상: BH 보정 유의 쌍 2개\n"
        "  ① 스니커즈 <- CH3 뉴스감성\n"
        "  ② 카드 <- CH1 Google Trends\n\n"
        "설정:\n"
        "  반응 기간: 6개월 추적\n"
        "  95% 신뢰구간:\n"
        "  부트스트랩 1,000회 반복\n"
        "  해석: 0 포함 안 되면 유의"
    ))

expbox(sld, 8.75, 1.12, 4.33, 2.65,
    title="경제적 해석 기준",
    text=(
        "CH3 뉴스감성 -> 스니커즈 가격:\n"
        "  부정 뉴스 충격 -> 가격 하락?\n"
        "  (소비자 심리 위축 -> 수요 감소)\n\n"
        "CH1 Google Trends -> 카드 가격:\n"
        "  검색량 급증 충격 -> 가격 상승?\n"
        "  (관심 증가 -> 수요 증가 선행)"
    ))

add_img(sld, os.path.join(RES, "irf_sneakers_ch3.png"),
        l=0.25, t=3.9, w=6.25, h=2.95,
        caption="스니커즈 <- CH3 뉴스감성 충격반응. 래그 1-2에서 음의 반응 확인, 래그 4에서 소멸. "
                "부정 뉴스 -> 단기 가격 하락 선행을 IRF로 재확인.")
add_img(sld, os.path.join(RES, "irf_cards_ch1.png"),
        l=6.82, t=3.9, w=6.26, h=2.95,
        caption="카드 <- CH1 Google Trends 충격반응. 래그 1에서 양의 반응, 래그 3에서 감쇠. "
                "검색량 증가 -> 카드 가격 상승 선행을 IRF로 재확인.")
add_pgnum(sld, 13)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 13 — 제한점
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "연구 제한점",
    sub="방법론적 한계 및 데이터 품질 이슈 — 논문 §5.3에 전부 명시하여 과도한 일반화 방지")

lim_data = [
    ["#", "제한점", "구체적 내용", "심각도", "대응"],
    ["L1", "소표본 문제",
     "자산별 N=48개월 — 복잡한 VAR 구조 검정하기엔 관측치 부족. 검정력(power) 낮을 수 있음.",
     "높음", "원본 p + BH 보정 p 병기. 검정력 명시."],
    ["L2", "스니커즈 CH1\n약한 정상성",
     "1차 차분 후에도 ADF p=0.120 — 경계적 비정상 상태. Granger 결과 해석 시 주의 필요.",
     "중간", "해당 결과 해석 시 주의 문구 명시."],
    ["L3", "CH3 도메인\n특화도 한계",
     "GDELT tone은 기사 전반 자동 추출. FinBERT(금융 특화) 대비 도메인 정밀도 낮음.",
     "중간", "CH3 민감도 분석 진행. 논문 §5.3 명시."],
    ["L4", "CH4 키워드\n편향",
     "'resell+review' 키워드 필터 -> 일반 hype 콘텐츠 과소 계상 가능성.",
     "중간", "전 아이템 동일 키워드 프로토콜 적용."],
    ["L5", "레고 GT 0주\n비율 높음",
     "Porsche 21/48, Bugatti 25/48 구간에서 Google Trends 값 = 0.",
     "중간", "0비율 테이블 제시. 해당 아이템 강한 결론 배제."],
    ["L6", "단일 플랫폼\n편향",
     "YouTube만 수집. TikTok·Reddit·Twitter/X 제외 -> CH4·CH5 대표성 한계.",
     "낮음", "향후 연구 기회로 명시."],
    ["L7", "z-score\n데이터 누출",
     "전체 패널로 StandardScaler 피팅 후 TimeSeriesSplit 적용 — 엄밀한 fold별 fit 아님.",
     "낮음", "N=48 특성상 실질 영향 미미. §5.3 명시."],
]
def _lbg(r):
    s=r[3]
    if s=="높음": return RGBColor(0xFF,0xCC,0xCC)
    if s=="중간": return MARG
    return NS
lim_bgs = [_lbg(r) for r in lim_data[1:]]
build_table(sld, lim_data,
    col_w=[0.45, 1.9, 4.5, 1.1, 4.75],
    l=0.25, t=1.12, h=5.92,
    hdr_bg=NAVY, hdr_sz=11, body_sz=9.5,
    row_bgs=lim_bgs, center_cols={0,3})
add_pgnum(sld, 14)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 14 — 결론
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "결론: 연구 문제별 답변",
    sub="Granger 인과 검정 -> Jaccard 채널 비교 -> XGBoost + DM 검정 — 세 단계 분석의 핵심 결과")

concl_data = [
    ["RQ", "연구 문제 (요약)", "실증 결과", "핵심 근거"],
    ["RQ1",
     "미디어 채널이 리셀\n가격을 선행하는가?",
     "예 — BH 보정 기준 2/15 유의\n\n"
     "  스니커즈 <- CH3 뉴스감성\n"
     "  (F=13.4,  p=0.001,  BH p=0.011)\n\n"
     "  카드 <- CH1 Google Trends\n"
     "  (F=8.4,   p=0.006,  BH p=0.044)",
     "Granger 인과 검정\nBH FDR 보정\n(q=0.05, m=15)"],
    ["RQ2",
     "선행 채널 구성이\n자산마다 다른가?",
     "예 — 자산별 채널 구성 이질성 확인\n\n"
     "  J(스니커즈, 카드) = 0.333\n"
     "  (CH1만 공유, 부분적 겹침)\n\n"
     "  J(*, 레고) = 0.000\n"
     "  (레고: 유의 채널 없음)",
     "Jaccard 유사도\n채널 집합 비교"],
    ["RQ3",
     "선별 채널 모델이\n전채널 모델과\n동등한가?",
     "예 — 3자산 전부 DM H0 기각 실패\n\n"
     "  스니커즈: DM p=0.678\n"
     "  카드:     DM p=0.504\n"
     "  레고:     DM p=0.646\n\n"
     "  2채널만으로도 5채널 동등 성능",
     "Diebold-Mariano\n검정\nAUC-ROC 비교"],
]
build_table(sld, concl_data,
    col_w=[0.75, 2.5, 6.0, 2.6],
    l=0.25, t=1.12, h=5.6,
    hdr_bg=NAVY, hdr_sz=11, body_sz=10.5,
    row_bgs=[SNK_BG, CRD_BG, LGO_BG], center_cols={0,1,3})
add_pgnum(sld, 15)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 15 — 학술적 기여 & 마무리
# ═══════════════════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
add_header(sld, "학술적 기여 및 실무적 시사점",
    sub="세 연구 흐름의 교점을 최초로 분석 — 리셀 시장 미디어 모니터링 전략의 이론적 근거 제공")

contrib_data = [
    ["기여 유형", "내용"],
    ["새로운 연구 도메인",
     "리셀 시장(스니커즈·카드·레고)에 다채널 Granger 인과 분석을 최초로 적용."
     " 기존 연구는 전통 금융 자산이나 단일 리셀 자산에 국한되어 있었음."],
    ["다자산 채널 비교",
     "단일 채널이 모든 자산을 선행하지 않음을 실증."
     " 스니커즈는 뉴스감성(CH3)이, 카드는 검색량(CH1)이 핵심 — 자산 특이적 채널 존재 확인."],
    ["방법론 삼각 검증",
     "Granger 인과(시간적 선행성) + SHAP(모델 내 채널 기여도) + DM 검정(예측 비교)"
     " 세 방법을 결합하여 단일 방법의 한계를 상호 보완."],
    ["실무적 시사점",
     "스니커즈 투자자: 뉴스 논조(CH3) 모니터링 우선"
     " | 카드 투자자: Google 검색 트렌드(CH1) 추적"
     " | 레고: 현재 유효한 미디어 선행 지표 없음 -> 가격 자체 모멘텀에 의존."],
    ["향후 연구 방향",
     "TikTok·Reddit 등 플랫폼 확장 | 월별 -> 주별 데이터로 고빈도 분석"
     " | 경매·한정판 드롭 이벤트 효과 통제 | 장기 패널(4년 이상) 구축."],
]
build_table(sld, contrib_data,
    col_w=[2.5, 10.2],
    l=0.25, t=1.12, h=4.8,
    hdr_bg=BLUE, hdr_sz=11, body_sz=10.5,
    center_cols={0})

add_box(sld, 0.25, 6.1, 12.83, 0.88, fill=NAVY)
add_text(sld, "감사합니다  —  질문 환영합니다",
    0.25, 6.22, 12.83, 0.65,
    size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_pgnum(sld, 16)


# ═══════════════════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"[OK] Saved -> {len(prs.slides)} slides")
