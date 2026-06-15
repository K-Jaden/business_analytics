import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

prs = Presentation('발표자료_최종_피드백반영.pptx')

def set_run_text(para, new_text):
    """Replace paragraph text in run[0], clear all other runs."""
    if not para.runs:
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ''

def set_shape_para(shape, para_idx, new_text):
    """Helper: set paragraph para_idx in shape's text frame."""
    if shape.has_text_frame:
        paras = shape.text_frame.paragraphs
        if para_idx < len(paras):
            set_run_text(paras[para_idx], new_text)

def add_caption_box(slide, text, left_in, top_in, width_in, font_size=8.5):
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.65)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    return txBox

# ===================== SLIDE 10: 채널 전처리 =====================
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if shape.shape_type == 19:  # TABLE
        table = shape.table
        set_run_text(table.cell(0, 3).text_frame.paragraphs[0], '전처리 → 최종 스케일 (모두 z-score)')
        set_run_text(table.cell(1, 3).text_frame.paragraphs[0], '월평균 → z-score 표준화')
        set_run_text(table.cell(2, 3).text_frame.paragraphs[0], 'GDELT → z-score 표준화')
        set_run_text(table.cell(3, 3).text_frame.paragraphs[0], 'GDELT tone → z-score 표준화')
        set_run_text(table.cell(4, 3).text_frame.paragraphs[0], 'log1p → z-score 표준화')
        set_run_text(table.cell(5, 3).text_frame.paragraphs[0], 'FinBERT → z-score 표준화')

    if shape.name == 'Rounded Rectangle 7':  # ADF 정상성 박스
        paras = shape.text_frame.paragraphs
        set_run_text(paras[0], '정상성(ADF) 검정 사용 이유')
        if len(paras) > 1:
            set_run_text(paras[1], '정상성 검정: 시계열의 평균·분산이 시간에 따라 일정한지 검증')
        if len(paras) > 2:
            set_run_text(paras[2], '비정상 시계열은 허구 인과 위험 → Granger의 전제조건')
print('Slide 10 완료')

# ===================== SLIDE 17: Granger 결과 =====================
slide17 = prs.slides[16]
for shape in slide17.shapes:
    if shape.name == 'TextBox 6':
        set_shape_para(shape, 0,
            '자산×채널 15쌍의 Granger F값 히트맵. '
            'F값: 미디어 변수를 추가했을 때 가격 예측 오차가 줄어드는 정도를 나타내는 통계량 — '
            'F가 크고 p<0.05이면 해당 채널이 가격에 선행함. '
            '스니커즈 CH3(F=13.4)·카드 CH1(F=8.4)이 가장 강한 선행성, 나머지 채널은 F<5, p>0.05.'
        )
print('Slide 17 완료')

# ===================== SLIDE 18: Jaccard =====================
slide18 = prs.slides[17]
for shape in slide18.shapes:
    if shape.name == 'TextBox 3':
        set_shape_para(shape, 0,
            '유의 채널 집합이 자산마다 다른가? — '
            'Jaccard 유사도: 두 집합의 교집합÷합집합 (0=완전이질, 1=완전동일)으로 '
            '자산 간 채널 구성 이질성을 정량화하는 지표'
        )
    if shape.name == 'TextBox 6':
        set_shape_para(shape, 0,
            '자산별 유의 채널 집합 비교. 원이 겹칠수록 두 자산의 유의 채널이 동일함. '
            '스니커즈-카드는 CH1(Google Trends)을 공유(Jaccard=0.333), '
            '레고는 유의 채널이 없어 나머지 쌍 모두 Jaccard=0.000.'
        )
print('Slide 18 완료')

# ===================== SLIDE 19: 모델 A & B =====================
slide19 = prs.slides[18]
for shape in slide19.shapes:
    if shape.name == 'TextBox 20':
        set_shape_para(shape, 0, '12 / 19')
    if shape.name == 'TextBox 9':
        set_shape_para(shape, 0,
            '결론: Granger 유의 채널만 쓴 Model B(스니커즈 CH1+CH3+CH4, 카드 CH1)가 '
            '전채널 Model A보다 AUC가 높지 않음 → 채널 선별이 예측력 향상으로 이어지지 않음.'
        )
add_caption_box(
    slide19,
    'Model A(전채널 5개)와 Model B(Granger 유의 채널)의 자산별 AUC-ROC 비교 막대 그래프. '
    '각 막대는 TimeSeriesSplit(5-fold) 교차검증 평균 AUC. '
    '세 자산 모두 두 모델 간 차이 0.01 내외 — 채널 선별이 예측력을 높이지 못함.',
    0, 4.3, 8.0, font_size=8.5
)
print('Slide 19 완료')

# ===================== SLIDE 20: SHAP =====================
slide20 = prs.slides[19]
for shape in slide20.shapes:
    if shape.name == 'TextBox 20':
        set_shape_para(shape, 0, '13 / 19')
    if shape.name == 'TextBox 19':
        set_shape_para(shape, 0,
            '주. 자산별 채널 평균 |SHAP| 막대 그래프. 높을수록 XGBoost 예측에 실제로 기여한 채널. '
            '스니커즈: CH1(Google)·CH5(YT감성) 상위, 카드: CH4(YT조회수)·CH1 상위, 레고: CH5·CH4 상위. '
            'Granger 1위인 스니커즈 CH3(F=13.4)가 SHAP에서는 4위(0.011) — 선행성과 예측 기여도가 일치하지 않음.'
        )
print('Slide 20 완료')

# ===================== SLIDE 21: 전체 모델 비교 =====================
slide21 = prs.slides[20]
for shape in slide21.shapes:
    if shape.name == 'PageNum':
        set_shape_para(shape, 0, '14 / 19')
add_caption_box(
    slide21,
    '모델 A(전채널)·B(Granger 선별)·C(SHAP 상위 2채널)·N-A(유의 채널 제거)·채널 단독(CH1~CH5)의 '
    '자산별 AUC-ROC 비교. 모든 모델 구성에서 AUC 차이 최대 0.01 내외 → '
    '채널 수를 줄이거나 늘려도 예측 성능이 실질적으로 변하지 않음.',
    0, 5.2, 12.0, font_size=8.5
)
print('Slide 21 완료')

# ===================== SLIDE 22: DM 검정 =====================
slide22 = prs.slides[21]
for shape in slide22.shapes:
    if shape.name == 'TextBox 28':
        set_shape_para(shape, 0, 'Granger 유의 채널이 전채널보다 예측력이 높은가? — DM 검정')
    if shape.name == 'TextBox 30':
        set_shape_para(shape, 0,
            "Diebold-Mariano 검정으로 'Model A(전채널) vs Model C(SHAP 선별 2채널)' "
            "예측오차 차이를 통계적으로 검증"
        )
    if shape.name in ('TextBox 33', 'TextBox 35'):
        set_shape_para(shape, 0,
            '결론 RQ3  |  3자산 모두 DM p > 0.50 → '
            'Granger 유의 채널(Model C)이 전채널(A)보다 우수하다는 가설 기각 실패 → '
            '채널 선별의 예측력 향상 효과 확인되지 않음'
        )
add_caption_box(
    slide22,
    'Diebold-Mariano(DM) 검정 결과. Model A(전채널)와 Model C(SHAP 상위 2채널)의 월별 예측오차 시계열 비교. '
    '세 자산 모두 DM p>0.50 → 두 모델의 예측오차 차이가 통계적으로 유의하지 않음 → H0 기각 실패.',
    2, 5.1, 8.0, font_size=8.5
)
print('Slide 22 완료')

# ===================== SLIDE 23: 왜 선별이 효과 없나 =====================
slide23 = prs.slides[22]
for shape in slide23.shapes:
    if shape.name == 'TextBox 16':
        set_shape_para(shape, 0, '16 / 19')
    if shape.name == 'TextBox 1':
        set_shape_para(shape, 0, '왜 채널 선별이 예측력 향상으로 이어지지 않는가?')
print('Slide 23 완료')

# ===================== SLIDE 24: 제한점 =====================
slide24 = prs.slides[23]
for shape in slide24.shapes:
    if shape.name == 'TextBox 4':
        paras = shape.text_frame.paragraphs
        if paras and paras[0].runs and paras[0].runs[0].text.strip() == '14':
            set_run_text(paras[0], '17')
print('Slide 24 완료')

# ===================== SLIDE 26: 결론 =====================
slide26 = prs.slides[25]
for shape in slide26.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        t = para.text.strip()
        if t == '15':
            set_run_text(para, '18')
        elif t == '선별 모델이 전채널과 동등한가?':
            set_run_text(para, 'Granger 선별 채널이 전채널보다 예측력이 높은가?')
        elif '2채널로 5채널 동등' in t:
            set_run_text(para, 'DM p>0.50 (스 0.678·카 0.504·레 0.646) — 채널 선별 예측력 향상 미확인')
        elif t == '동등(간결성)':
            set_run_text(para, '미지지')
print('Slide 26 완료')

# ===================== SLIDE 27: 시사점 =====================
slide27 = prs.slides[26]
for shape in slide27.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        if para.text.strip() == '16':
            set_run_text(para, '19')
            break
print('Slide 27 완료')

# ===================== SLIDE 28: 참고문헌 =====================
slide28 = prs.slides[27]
for shape in slide28.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        if para.text.strip() == '17':
            set_run_text(para, '20')
            break
print('Slide 28 완료')

prs.save('발표자료_최종_피드백반영_v2.pptx')
print('\n저장 완료: 발표자료_최종_피드백반영_v2.pptx')
