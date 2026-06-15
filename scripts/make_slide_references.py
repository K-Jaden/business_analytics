#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
참고문헌 슬라이드 — slide_references_v2.pptx 2페이지 스타일 기준
(얇은 다크 탑바 / 대형 오픈 제목 / 회색 구분선 / 미니멀 레이아웃)
Usage : python scripts/make_slide_references.py
Output: slide_references_v3.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide_references_v5.pptx")

# ── 색상 ─────────────────────────────────────────────────────────────────────
DARK   = RGBColor(0x1A, 0x1A, 0x1A)   # 탑 스트립
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SEP_C  = RGBColor(0xE0, 0xE0, 0xE0)   # 수평/수직 구분선
TITLE  = RGBColor(0x1A, 0x1A, 0x1A)   # 제목 텍스트
SUB_C  = RGBColor(0x88, 0x88, 0x88)   # 부제목
NUM_C  = RGBColor(0x2E, 0x6D, 0xA4)   # ref 번호 (파랑)
AUT_C  = RGBColor(0x1A, 0x1A, 0x1A)   # 저자
SRC_C  = RGBColor(0x77, 0x77, 0x77)   # 출처 (이탤릭)
DOT_C  = RGBColor(0xCC, 0xCC, 0xCC)   # 구분점

# 카테고리 언더라인 색 (카테고리별 고유)
C_STAT = RGBColor(0x1F, 0x4E, 0x79)
C_ML   = RGBColor(0x14, 0x5A, 0x32)
C_ALT  = RGBColor(0x4A, 0x23, 0x5A)
C_MED  = RGBColor(0x78, 0x28, 0x00)
C_DIG  = RGBColor(0x0B, 0x3D, 0x66)
C_DAT  = RGBColor(0x37, 0x47, 0x51)

# ── 3열 콘텐츠 ────────────────────────────────────────────────────────────────
COLS = [
    [   # 열 1
        {"cat": "계량경제 · 통계",   "uc": C_STAT, "refs": [
            ("[1]",  "Granger (1969)",             "Econometrica 37(3)"),
            ("[2]",  "Dickey & Fuller (1979)",      "JASA 74"),
            ("[3]",  "Schwarz (1978)",              "Ann. Stat. 6(2)"),
            ("[4]",  "Diebold & Mariano (1995)",    "JBES 13(3)"),
            ("[5]",  "Benjamini & Hochberg (1995)", "JRSS-B 57(1)"),
            ("[6]",  "West (1996)",                 "Econometrica 64(5)"),
            ("[7]",  "Hansen (2005)",               "JBES 23(4)"),
            ("[8]",  "Hansen et al. (2011)",        "Econometrica 79(2)"),
            ("[9]",  "Holtz-Eakin et al. (1988)",  "Econometrica 56(6)"),
            ("[10]", "Dumitrescu & Hurlin (2012)", "Econ. Modelling 29(4)"),
            ("[11]", "Welch & Goyal (2008)",        "RFS 21(4)"),
        ]},
        {"cat": "머신러닝 · NLP",    "uc": C_ML, "refs": [
            ("[12]", "Chen & Guestrin (2016)", "KDD 2016"),
            ("[13]", "Lundberg & Lee (2017)",  "NeurIPS 30"),
            ("[14]", "Araci (2019) FinBERT",   "arXiv:1908.10063"),
            ("[15]", "Devlin et al. (2019)",    "NAACL-HLT 2019"),
        ]},
    ],
    [   # 열 2
        {"cat": "머신러닝 · NLP (cont.)", "uc": C_ML, "refs": [
            ("[16]", "Gu et al. (2020)",      "RFS 33(5)"),
            ("[17]", "Akiba et al. (2019)",   "KDD 2019"),
        ]},
        {"cat": "대안 자산 · 수집품 시장", "uc": C_ALT, "refs": [
            ("[18]", "Burton & Jacobsen (1999)",       "JEP 13(4)"),
            ("[19]", "Mei & Moses (2002)",             "AER 92(5)"),
            ("[20]", "Renneboog & Spaenjers (2013)",   "Mgmt. Sci. 59(1)"),
            ("[21]", "Dimson et al. (2015)",           "JFE 118(2)"),
            ("[22]", "Dobrynskaya & Kishilova (2022)", "RIBF 59"),
            ("[23]", "Raditya et al. (2021)",          "Procedia CS 179"),
        ]},
        {"cat": "미디어 · 감성 → 자산", "uc": C_MED, "refs": [
            ("[24]", "Tetlock (2007)",              "JF 62(3)"),
            ("[25]", "Tetlock et al. (2008)",        "JF 63(3)"),
            ("[26]", "Antweiler & Frank (2004)",     "JF 59(3)"),
            ("[27]", "Engelberg & Parsons (2011)",   "JF 66(1)"),
            ("[28]", "Loughran & McDonald (2011)",   "JF 66(1)"),
            ("[29]", "Loughran & McDonald (2016)",   "JAR 54(4)"),
        ]},
    ],
    [   # 열 3
        {"cat": "미디어 · 감성 → 자산 (cont.)", "uc": C_MED, "refs": [
            ("[30]", "García (2013)",          "JF 68(3)"),
            ("[31]", "Kearney & Liu (2014)",   "IRFA 33"),
            ("[32]", "Heston & Sinha (2017)",  "FAJ 73(3)"),
            ("[33]", "Bollen et al. (2011)",   "J. Comput. Sci. 2(1)"),
            ("[34]", "Sprenger et al. (2014)", "EFM 20(5)"),
        ]},
        {"cat": "검색 · 디지털 흔적",  "uc": C_DIG, "refs": [
            ("[35]", "Da et al. (2011)",               "JF 66(5)"),
            ("[36]", "Preis et al. (2013)",            "Sci. Rep. 3"),
            ("[37]", "Joseph et al. (2011)",           "IJF 27(4)"),
            ("[38]", "Vlastakis & Markellos (2012)",   "JBF 36(6)"),
            ("[39]", "Vozlyublennaia (2014)",           "JBF 41"),
            ("[40]", "Bordino et al. (2012)",           "PLoS ONE 7(7)"),
            ("[41]", "Moat et al. (2013)",              "Sci. Rep. 3"),
            ("[42]", "Curme et al. (2014)",             "PNAS 111(32)"),
        ]},
        {"cat": "데이터",              "uc": C_DAT, "refs": [
            ("[43]", "Leetaru & Schrodt (2013) GDELT", "ISA Annual Conv."),
        ]},
    ],
]

# ── 레이아웃 상수 ─────────────────────────────────────────────────────────────
L_MARGIN = 0.73     # 좌 여백 (슬라이드 2 기준)
R_MARGIN = 0.60     # 우 여백
COL_GAP  = 0.30     # 열 간 간격
NCOLS    = 3
COL_W    = (13.33 - L_MARGIN - R_MARGIN - (NCOLS-1)*COL_GAP) / NCOLS  # ≈ 3.80"
COL_X    = [L_MARGIN + i*(COL_W + COL_GAP) for i in range(NCOLS)]

TOP      = 1.06     # 콘텐츠 시작 y
CAT_LBL  = 0.245    # 카테고리 레이블 텍스트박스 높이
CAT_UL   = 0.022    # 카테고리 언더라인 높이
REF_H    = 0.210    # 참고문헌 한 줄 높이
CAT_GAP  = 0.140    # 카테고리 간 여백

# ── PPT 생성 ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld = prs.slides.add_slide(prs.slide_layouts[6])

# ① 얇은 다크 탑 스트립 (슬라이드 2 스타일 그대로)
strip = sld.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(0.055))
strip.fill.solid(); strip.fill.fore_color.rgb = DARK
strip.line.fill.background()

# ② 제목 (오픈 스타일, 배경 없음)
title_box = sld.shapes.add_textbox(
    Inches(L_MARGIN), Inches(0.09),
    Inches(13.33 - L_MARGIN - R_MARGIN), Inches(0.60))
title_box.fill.background(); title_box.line.fill.background()
tf_t = title_box.text_frame; tf_t.word_wrap = False
p_t = tf_t.paragraphs[0]; p_t.alignment = PP_ALIGN.LEFT

r_t1 = p_t.add_run()
r_t1.text = "참고 문헌"
r_t1.font.bold = True; r_t1.font.size = Pt(22)
r_t1.font.color.rgb = TITLE; r_t1.font.name = "Calibri"

r_t2 = p_t.add_run()
r_t2.text = "  References"
r_t2.font.bold = False; r_t2.font.size = Pt(18)
r_t2.font.color.rgb = RGBColor(0x55, 0x55, 0x55); r_t2.font.name = "Calibri"

# 부제목
sub_box = sld.shapes.add_textbox(
    Inches(L_MARGIN), Inches(0.64),
    Inches(12.0), Inches(0.28))
sub_box.fill.background(); sub_box.line.fill.background()
tf_s = sub_box.text_frame; p_s = tf_s.paragraphs[0]; p_s.alignment = PP_ALIGN.LEFT
r_s = p_s.add_run()
r_s.text = "draft.tex 실제 인용 43건  ·  Granger 인과 / XGBoost+SHAP / DM 검정 / 대안자산 / 미디어·감성 / 디지털 흔적"
r_s.font.size = Pt(10.5); r_s.font.color.rgb = SUB_C; r_s.font.name = "Calibri"

# ③ 수평 구분선
h_sep = sld.shapes.add_textbox(
    Inches(L_MARGIN), Inches(0.92),
    Inches(13.33 - L_MARGIN - R_MARGIN + 0.3), Inches(0.015))
h_sep.fill.solid(); h_sep.fill.fore_color.rgb = SEP_C
h_sep.line.fill.background()

# ④ 수직 구분선 (열 사이)
for i in range(1, NCOLS):
    vx = L_MARGIN + i*(COL_W + COL_GAP) - COL_GAP/2 - 0.007
    v_sep = sld.shapes.add_textbox(
        Inches(vx), Inches(TOP - 0.04),
        Inches(0.014), Inches(7.5 - TOP - 0.10))
    v_sep.fill.solid(); v_sep.fill.fore_color.rgb = SEP_C
    v_sep.line.fill.background()

# ── 카테고리 블록 그리기 ──────────────────────────────────────────────────────
def draw_block(slide, col_x, y, col_w, cat_name, ul_color, refs):
    """카테고리 라벨(텍스트+언더라인) + 참고문헌 텍스트박스 그리기"""
    # 카테고리 라벨
    lbl = slide.shapes.add_textbox(
        Inches(col_x), Inches(y),
        Inches(col_w), Inches(CAT_LBL))
    lbl.fill.background(); lbl.line.fill.background()
    tf = lbl.text_frame; tf.word_wrap = False
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0); tf.margin_left = Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = cat_name.upper()
    r.font.bold = True; r.font.size = Pt(11)
    r.font.color.rgb = ul_color; r.font.name = "Calibri"

    # 카테고리 언더라인
    ul = slide.shapes.add_textbox(
        Inches(col_x), Inches(y + CAT_LBL - 0.03),
        Inches(col_w), Inches(CAT_UL))
    ul.fill.solid(); ul.fill.fore_color.rgb = ul_color
    ul.line.fill.background()

    y_refs = y + CAT_LBL + 0.01

    # 참고문헌 텍스트박스
    refs_h = len(refs) * REF_H + 0.02
    rb = slide.shapes.add_textbox(
        Inches(col_x), Inches(y_refs),
        Inches(col_w), Inches(refs_h))
    rb.fill.background(); rb.line.fill.background()
    tf2 = rb.text_frame; tf2.word_wrap = False
    tf2.margin_top = Pt(0); tf2.margin_bottom = Pt(0); tf2.margin_left = Pt(1)

    for i, (num, author, src) in enumerate(refs):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.space_before = Pt(0); p2.space_after = Pt(0)

        rn = p2.add_run()        # [번호] — 카테고리색 볼드
        rn.text = num + " "
        rn.font.bold = True; rn.font.size = Pt(10)
        rn.font.color.rgb = ul_color; rn.font.name = "Calibri"

        ra = p2.add_run()        # 저자·연도
        ra.text = author
        ra.font.bold = False; ra.font.size = Pt(10)
        ra.font.color.rgb = AUT_C; ra.font.name = "Calibri"

        rd = p2.add_run()        # 구분점
        rd.text = "  ·  "
        rd.font.size = Pt(10); rd.font.color.rgb = DOT_C; rd.font.name = "Calibri"

        rs = p2.add_run()        # 출처 (이탤릭)
        rs.text = src
        rs.font.bold = False; rs.font.size = Pt(10); rs.font.italic = True
        rs.font.color.rgb = SRC_C; rs.font.name = "Calibri"

    return y_refs + refs_h + 0.01

# ── 열별 렌더링 ───────────────────────────────────────────────────────────────
for ci, col_blocks in enumerate(COLS):
    y_cur = TOP
    for bi, blk in enumerate(col_blocks):
        y_cur = draw_block(
            sld, COL_X[ci], y_cur, COL_W,
            blk["cat"], blk["uc"], blk["refs"])
        if bi < len(col_blocks) - 1:
            y_cur += CAT_GAP

# ⑤ 하단 수평선 + 푸터
ft_sep = sld.shapes.add_textbox(
    Inches(L_MARGIN), Inches(7.26),
    Inches(13.33 - L_MARGIN - R_MARGIN + 0.3), Inches(0.014))
ft_sep.fill.solid(); ft_sep.fill.fore_color.rgb = SEP_C
ft_sep.line.fill.background()

ft = sld.shapes.add_textbox(
    Inches(L_MARGIN), Inches(7.30),
    Inches(12.0), Inches(0.18))
ft.fill.background(); ft.line.fill.background()
tf_f = ft.text_frame; p_f = tf_f.paragraphs[0]; p_f.alignment = PP_ALIGN.LEFT
rf = p_f.add_run()
rf.text = "43 references  ·  리셀 시장 미디어 선행 지표 연구  ·  draft.tex 인용 기준"
rf.font.size = Pt(8); rf.font.color.rgb = SUB_C; rf.font.name = "Calibri"

prs.save(OUT)
print(f"[OK] {OUT}")
