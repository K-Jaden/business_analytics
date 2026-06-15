#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2장 구성 — ① 채널 선정 근거(앞)  ② 방법론 계보 & 연구 공백(뒤)
저작권 그래프 없이 텍스트/박스 기반. 진짜 검증된 논문만 인용.
Usage : python scripts/make_slides_prior_work2.py
Output: slides_prior_work_v2.pptx (2 slides)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slides_prior_work_v2.pptx")

NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE= RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0x77, 0x77, 0x77)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
TEAL  = RGBColor(0x2A, 0x9D, 0x8F)
INK   = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 22, True, WHITE)]],
         anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 1 — 채널 선정 근거
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "채널 선정 근거 — 5개 미디어 채널은 무엇을 포착하는가",
       [("가격에 ", 13, False, DGRAY),
        ("선행할 수 있는 소비자·시장 신호", 13, True, BLUE),
        ("를 단일 지표가 아닌 5개 층위로 측정", 13, False, DGRAY)])

rows = [
    ("CH1", "Google Trends", "소비자 검색 관심도",
     "관심 급증은 구매 수요·가격 상승에 선행할 수 있음"),
    ("CH2", "뉴스 보도량 (GDELT)", "미디어 노출·주목도",
     "보도량 증가는 대중 관심을 환기해 수요를 자극"),
    ("CH3", "뉴스 감성 (GDELT)", "보도 논조(긍정/부정)",
     "긍정 보도는 매수 심리·가격 기대를 형성"),
    ("CH4", "유튜브 조회수", "정보 탐색 의도(구매 전 리뷰)",
     "유재필·김문선(2023): 유튜브 조회수의 시장 선행성을 Granger로 확인"),
    ("CH5", "유튜브 댓글 감성 (FinBERT)", "커뮤니티 심리·반응",
     "댓글 긍/부정은 잠재 수요자의 태도를 반영"),
]
row_t = 1.50
row_h = 0.82
row_g = 0.10
col = [0.30, 1.55, 4.55, 7.30]  # x positions: tag, name, measures, logic
for i, (tag, name, meas, logic) in enumerate(rows):
    y = row_t + i * (row_h + row_g)
    box(s1, 0.30, y, 12.73, row_h, fill=LGRAY,
        line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s1, 0.42, y + 0.16, 0.95, row_h - 0.32, fill=BLUE, round_=True)
    text(s1, 0.42, y + 0.16, 0.95, row_h - 0.32, [[(tag, 14, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s1, 1.55, y, 3.0, row_h, [[(name, 12.5, True, NAVY)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s1, 4.60, y, 2.7, row_h, [[(meas, 11.5, False, DGRAY)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s1, 7.35, y, 5.55, row_h, [[(logic, 11.5, False, INK)]],
         anchor=MSO_ANCHOR.MIDDLE)

ky = row_t + 5 * (row_h + row_g) + 0.04
box(s1, 0.30, ky, 12.73, 0.62, fill=GREEN, round_=True)
text(s1, 0.45, ky + 0.02, 12.4, 0.58,
     [[("핵심  ", 13, True, WHITE),
       ("관심(검색)·노출(보도량)·태도(논조)·탐색(조회수)·심리(댓글) — "
        "미디어 신호를 5개 층위로 다면 포착", 12, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — 방법론 계보 & 연구 공백
# ════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
header(s2, "방법론 계보 — ‘감성 → Granger 선별 → 예측’의 흐름",
       [("선행연구는 이 패러다임을 ", 13, False, DGRAY),
        ("주식", 13, True, BLUE),
        ("에 적용 — 그러나 ", 13, False, DGRAY),
        ("리셀 등 대안자산의 다채널 선행성은 미답", 13, True, ORANGE)])

CARD_T, CARD_H, CARD_W, gap = 1.52, 3.35, 4.13, 0.22
xs = [0.30, 0.30 + CARD_W + gap, 0.30 + 2 * (CARD_W + gap)]
cards = [
    {"tag": "전설적 시초 · 2011",
     "title": "Twitter mood predicts the stock market",
     "auth": "Bollen, Mao & Zeng (2011), J. Comput. Sci.",
     "lines": ["트윗 수천만 건을 6개 기분 차원으로 분류",
               "→ 다우존스와 Granger 인과검정",
               "‘차분함(Calm)’이 주가를 2~4일 선행",
               "신경망 결합 시 방향예측 정확도 86.7%"]},
    {"tag": "국내 · Granger · 2023",
     "title": "유튜브 기반 데이터 분석을 통한 주가 선행성 분석",
     "auth": "유재필·김문선 (2023), 정보화연구",
     "lines": ["유튜브 채널 데이터 크롤링 + ETF 주가",
               "→ Granger 인과분석 수행",
               "유튜브 조회수가 금융시장을 선행함을 확인",
               "검색·조회 급증이 거래량의 선행지표"]},
    {"tag": "최신 트렌드 · NeurIPS 2024",
     "title": "CausalStock: News-driven Causal Discovery",
     "auth": "Li et al. (2024), NeurIPS",
     "lines": ["LLM으로 뉴스 노이즈 제거(Denoised Encoder)",
               "→ 시차 기반 인과그래프 자동 탐색",
               "美·中·日·英 6개 데이터셋서 최고 성능",
               "Granger를 넘어 종단간(end-to-end) 인과"]},
]
for x, c in zip(xs, cards):
    box(s2, x, CARD_T, CARD_W, CARD_H, fill=LGRAY,
        line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s2, x, CARD_T, CARD_W, 0.42, fill=BLUE, round_=True)
    text(s2, x, CARD_T + 0.02, CARD_W, 0.38, [[(c["tag"], 11.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s2, x + 0.05, CARD_T + 0.48, CARD_W - 0.1, 0.95,
         [[(c["title"], 12, True, NAVY)], [(c["auth"], 9.5, False, MGRAY)]], space=3)
    body = [[("•  " + ln, 10.5, False, INK)] for ln in c["lines"]]
    text(s2, x + 0.08, CARD_T + 1.42, CARD_W - 0.16, CARD_H - 1.5, body, space=6)

GAP_T = CARD_T + CARD_H + 0.18
box(s2, 0.30, GAP_T, 12.73, 0.78, fill=RGBColor(0xFC, 0xF0, 0xE6),
    line=ORANGE, line_w=1.5, round_=True)
text(s2, 0.45, GAP_T + 0.04, 12.4, 0.70,
     [[("공통 한계 (연구 공백)   ", 12.5, True, ORANGE),
       ("① 모두 ", 11.5, False, INK), ("효율적 시장인 주식", 11.5, True, INK),
       ("에 한정   ② ", 11.5, False, INK), ("단일·소수 채널", 11.5, True, INK),
       (" 위주   ③ ", 11.5, False, INK),
       ("리셀 등 대안자산의 다채널 선행성은 검증된 바 없음", 11.5, True, INK)]],
     anchor=MSO_ANCHOR.MIDDLE)

KEY_T = GAP_T + 0.92
box(s2, 0.30, KEY_T, 12.73, 0.70, fill=GREEN, round_=True)
text(s2, 0.45, KEY_T + 0.03, 12.4, 0.64,
     [[("본 연구 ", 13, True, WHITE),
       ("= 동일 방법론(감성→Granger 선별→예측)을 ", 12.5, False, WHITE),
       ("리셀 3종(스니커즈·카드·레고)", 12.5, True, WHITE),
       ("에 ", 12.5, False, WHITE),
       ("5개 채널로 최초 적용", 12.5, True, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
