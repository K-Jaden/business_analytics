#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Slide 9 모델 설계 표 — 행/열 전치 버전 (make_academic_ppt.py 원본 불변)
원본: 행=구성요소(6개), 열=모델A/B/C
전치: 행=모델A/B/C, 열=구성요소(6개)

Usage : python scripts/make_slide9_transposed.py
Output: slide9_transposed.pptx
"""
import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slide9_transposed.pptx")

# ── 색상 ──────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x44, 0x44, 0x44)
BG1   = RGBColor(0xF4, 0xF6, 0xF9)   # 모델A/C 배경
BG2   = RGBColor(0xFF, 0xFF, 0xFF)   # 모델B 배경

# ── 원본 데이터 (make_academic_ppt.py 슬라이드 9 그대로) ──────────────────
# original[0] = 헤더행, original[1..6] = 구성요소행
original = [
    ["구성 요소",       "모델 A  (전채널)",                     "모델 B  (Granger 선별)",               "모델 C  (SHAP 선별)"],
    ["투입 채널",       "통제 피처\n+ CH1~CH5 전부\n(5채널)",   "통제 피처\n+ Granger 유의 채널만",     "통제 피처\n+ SHAP 상위\n2채널만"],
    ["자산별\n선별 채널","모두 동일\n(5채널 전부)",              "스니커즈: CH3\n카드: CH1\n레고: (없음)","스니커즈: CH1+CH5\n카드: CH4+CH1\n레고: CH5+CH4"],
    ["통제 피처",       "price_vs_ma3\nprice_chg_lag1~lag3",   "동일",                                  "동일"],
    ["교차검증",        "TimeSeriesSplit\n(n_splits=5)",        "동일",                                  "동일"],
    ["비교 방식",       "기준 모델\n(베이스라인)",               "A 대비 AUC\n직접 비교",                 "A 대비\nDiebold-Mariano 검정"],
]

# ── 전치 ─────────────────────────────────────────────────────────────────
# transposed[0] = 헤더행 ["모델", 구성요소1..6]
# transposed[1..3] = 모델A, B, C 행
transposed = []

# 헤더행: ["모델"] + [구성요소 이름들]
hdr = ["모델"] + [row[0] for row in original[1:]]
transposed.append(hdr)

# 각 모델 행
for col_idx, model_name in enumerate(["모델 A  (전채널)", "모델 B  (Granger 선별)", "모델 C  (SHAP 선별)"], start=1):
    row = [model_name] + [original[ri][col_idx] for ri in range(1, len(original))]
    transposed.append(row)

# ── 유틸 ──────────────────────────────────────────────────────────────────
def set_cell_bg(cell, rgb):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    sf   = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(sf,  qn("a:srgbClr"))
    srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")

def set_cell_text(cell, text, bold=False, size=10, color=GRAY,
                  align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.bold  = bold
    r.font.size  = Pt(size)
    r.font.name  = "Calibri"
    r.font.color.rgb = color

# ── PPT 생성 ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld = prs.slides.add_slide(prs.slide_layouts[6])

# 헤더 바
hbar = sld.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(0.72))
hbar.fill.solid(); hbar.fill.fore_color.rgb = NAVY
hbar.line.fill.background()
tf = hbar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  RQ3 — 모델 설계 요약 (A / B / C)"
r.font.bold = True; r.font.size = Pt(22); r.font.color.rgb = WHITE; r.font.name = "Calibri"

# 구분선
sep = sld.shapes.add_textbox(Inches(0), Inches(0.72), Inches(13.33), Inches(0.025))
sep.fill.solid(); sep.fill.fore_color.rgb = BLUE; sep.line.fill.background()

# 부제목
sub = sld.shapes.add_textbox(Inches(0.2), Inches(0.755), Inches(13.0), Inches(0.28))
sub.fill.background(); sub.line.fill.background()
tf2 = sub.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run()
r2.text = ("평가지표: AUC-ROC  |  교차검증: TimeSeriesSplit(n_splits=5)  "
           "|  통제 피처 공통: price_vs_ma3, price_chg_lag1~lag3  |  모델: XGBClassifier")
r2.font.size = Pt(9); r2.font.color.rgb = RGBColor(0x55,0x55,0x55)
r2.font.italic = True; r2.font.name = "Calibri"

# ── 테이블: 4행 × 7열 ────────────────────────────────────────────────────
n_rows = 4    # 헤더 + A + B + C
n_cols = 7    # 모델명 + 구성요소 6개

# 열 너비 (합 = 12.5인치)
# [모델명, 투입채널, 자산별선별채널, 통제피처, 교차검증, 비교방식]
col_w = [1.55, 1.90, 2.65, 1.90, 1.55, 2.00, 0.95]   # 7열: 모델 + 6 구성요소
# 아, 열이 7개: [모델, 투입채널, 자산별선별채널, 통제피처, 교차검증, 비교방식] + 모델명 열
# 실제 transposed[0] = ["모델", "투입 채널", "자산별\n선별 채널", "통제 피처", "교차검증", "비교 방식"] — 6개
# 그러므로 n_cols=6
n_cols = 6
col_w  = [1.65, 2.05, 2.75, 2.00, 1.55, 2.35]   # 합 = 12.35

tbl_l = Inches(0.49)
tbl_t = Inches(1.10)
tbl_w = sum(Inches(w) for w in col_w)
tbl_h = Inches(6.10)

tbl_shape = sld.shapes.add_table(n_rows, n_cols, tbl_l, tbl_t, tbl_w, tbl_h)
tbl = tbl_shape.table

# 열 너비
for ci, w in enumerate(col_w):
    tbl.columns[ci].width = Inches(w)

# 행 높이
row_h = [0.58, 1.84, 1.84, 1.84]
for ri, h in enumerate(row_h):
    tbl.rows[ri].height = Inches(h)

# ── 헤더 행 채우기 ─────────────────────────────────────────────────────────
for ci, hdr_txt in enumerate(transposed[0]):
    cell = tbl.cell(0, ci)
    set_cell_bg(cell, NAVY)
    set_cell_text(cell, hdr_txt, bold=True, size=11, color=WHITE,
                  align=PP_ALIGN.CENTER)

# ── 데이터 행 채우기 (모델A=row1, B=row2, C=row3) ─────────────────────────
row_bgs   = [BG1, BG2, BG1]
model_colors = [
    RGBColor(0x1F, 0x2D, 0x4E),   # A — 네이비
    RGBColor(0xE8, 0x77, 0x22),   # B — 오렌지
    RGBColor(0x2E, 0x6D, 0xA4),   # C — 블루
]

for ri, (data_row, bg, mcol) in enumerate(zip(transposed[1:], row_bgs, model_colors), start=1):
    for ci, val in enumerate(data_row):
        cell = tbl.cell(ri, ci)
        if ci == 0:
            # 모델명 열 — 해당 모델 색
            set_cell_bg(cell, mcol)
            set_cell_text(cell, val, bold=True, size=11, color=WHITE,
                          align=PP_ALIGN.CENTER)
        else:
            set_cell_bg(cell, bg)
            # 자산별선별채널(ci==2) — 왼쪽 정렬
            al = PP_ALIGN.LEFT if ci == 2 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bold=False, size=10.5, color=GRAY, align=al)

# ── 테이블 테두리 흰색 ────────────────────────────────────────────────────
for ri in range(n_rows):
    for ci in range(n_cols):
        tc   = tbl.cell(ri, ci)._tc
        tcPr = tc.get_or_add_tcPr()
        for tag in ["a:lnL","a:lnR","a:lnT","a:lnB"]:
            for old in tcPr.findall(qn(tag)):
                tcPr.remove(old)
            ln   = etree.SubElement(tcPr, qn(tag))
            ln.set("w","6350"); ln.set("cap","flat"); ln.set("cmpd","sng")
            sf   = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(sf, qn("a:srgbClr"))
            srgb.set("val","FFFFFF")

# ── 색상 범례 ─────────────────────────────────────────────────────────────
legend = sld.shapes.add_textbox(Inches(0.49), Inches(7.18), Inches(12.35), Inches(0.25))
legend.fill.solid(); legend.fill.fore_color.rgb = RGBColor(0xEF,0xF2,0xF7)
legend.line.color.rgb = RGBColor(0xBB,0xBB,0xBB)
tf3 = legend.text_frame; p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.LEFT
r3 = p3.add_run()
r3.text = ("■ 네이비 = 모델 A  ■ 오렌지 = 모델 B  ■ 블루 = 모델 C  |  "
           "B 스니커즈: CH3만 유의 (BH 보정 기준)  |  레고 B: 유의 채널 없음 → B = A")
r3.font.size = Pt(8); r3.font.color.rgb = GRAY; r3.font.italic = True; r3.font.name = "Calibri"

prs.save(OUT)
print(f"[OK] {OUT}")
