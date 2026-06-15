# -*- coding: utf-8 -*-
import sys
from pptx import Presentation

path = sys.argv[1]
p = Presentation(path)
for i, s in enumerate(p.slides, 1):
    print(f"=== SLIDE {i} ===")
    for sh in s.shapes:
        if sh.has_text_frame:
            txt = " | ".join(
                r.text
                for para in sh.text_frame.paragraphs
                for r in para.runs
                if r.text.strip()
            )
            if txt.strip():
                print(f"  [TXT:{sh.name}] {txt[:400]}")
        elif sh.shape_type == 13:
            print(f"  [PIC] {sh.name}")
        if sh.has_table:
            for row in sh.table.rows:
                print("  [TBL] " + " | ".join(c.text.replace(chr(10), " / ") for c in row.cells))
