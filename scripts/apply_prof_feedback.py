# -*- coding: utf-8 -*-
"""
교수님 피드백 반영 — 최종_버전_수정.pptx -> 최종_버전_수정_교수님피드백.pptx
원본은 손대지 않고 항상 새 파일로 저장 (idempotent: 매 실행마다 원본에서 복사).

피드백 요약(물리 슬라이드 번호 기준):
  - 전체: 글씨 줄이고 핵심 강조, 검정마다 '왜 이 검정을 썼는지' 명시, 그래프 의미 캡션
  - S4  : 기존 연구 -> 한계 강조, 본문 디테일 축소
  - S5  : 연구문제 표 -> RQ1->RQ2->RQ3 인포그래픽 (검정 사용 이유 포함)
  - S8+9: 채널 전처리 + 정상성 한 장으로 병합 (스케일 다른 이유 + z-score 통일, ADF 요약). S9 삭제
  - S13 : 모델설계 텍스트 -> ㅁ->ㅁ->ㅁ 플로우 도형
  - S19 : DM 설명 축소
  - S21 : IRF 설명 축소
  - 페이지 번호 재계산, 새 파일 저장
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = "최종_버전_수정.pptx"
DST = "최종_버전_수정_교수님피드백.pptx"

# 팔레트
NAVY   = "1F2D4E"
BLUE   = "2E6DA4"
GREEN  = "277A4A"
PURPLE = "6C3D91"
GRAY_T = "555555"
INK    = "0A0A0A"
LBLUE  = "EAF1F8"   # 연한 블루 카드
LGREEN = "E6F2EC"
LPUR   = "EFE9F4"
LRED   = "F8D7DA"   # 한계 강조 박스
DRED   = "B23B3B"
LGRAY  = "F2F3F5"
WHITE  = "FFFFFF"
FONT   = "Calibri"

# ---------------- helpers ----------------
def C(s):
    return RGBColor.from_string(s)

def find(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

def remove_shape(sh):
    sh._element.getparent().remove(sh._element)

def no_shadow(shp):
    try:
        shp.shadow.inherit = False
    except Exception:
        pass

def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = C(line); shp.line.width = Pt(line_w)
    no_shadow(shp)
    # 둥근 모서리 살짝
    return shp

def set_text(shp, paras, anchor=MSO_ANCHOR.TOP, wrap=True,
             ml=5, mr=5, mt=3, mb=3):
    """paras: list of paragraphs; each paragraph = (runs, align, space_after)
       runs = list of (text, size, color, bold)"""
    tf = shp.text_frame
    tf.clear()  # 기존 내용 제거 (기존 도형 편집 시 중복 방지)
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(ml); tf.margin_right = Pt(mr)
    tf.margin_top = Pt(mt); tf.margin_bottom = Pt(mb)
    first = True
    for runs, align, space in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if space is not None:
            p.space_after = Pt(space)
        p.space_before = Pt(0)
        for text, size, color, bold in runs:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = C(color)
    return shp

def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    return tb

def add_arrow(slide, l, t, w, h, fill=BLUE):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = C(fill)
    a.line.fill.background()
    no_shadow(a)
    return a

def resize_pic(slide, name, new_h, top=None, left=None, center=False):
    pic = find(slide, name)
    if pic is None:
        return None
    ratio = pic.width / pic.height
    pic.height = Inches(new_h)
    pic.width = Inches(new_h * ratio)
    if top is not None: pic.top = Inches(top)
    if left is not None: pic.left = Inches(left)
    if center:
        pic.left = Inches((13.333 - new_h * ratio) / 2.0)
    return pic

def caption(slide, l, t, w, text):
    """그래프 의미 캡션 (회색 작은 글씨, 왼쪽에 파란 라벨)"""
    tb = textbox(slide, l, t, w, 0.5)
    set_text(tb, [
        ([("그래프 해석  ", 9, BLUE, True), (text, 9, GRAY_T, False)], PP_ALIGN.LEFT, 0),
    ], anchor=MSO_ANCHOR.TOP)
    return tb

# ================= main =================
import shutil
shutil.copyfile(SRC, DST)
prs = Presentation(DST)
S = prs.slides

# 슬라이드 객체 (물리 번호 -1)
s4  = S[3]
s5  = S[4]
s8  = S[7]
s9  = S[8]
s11 = S[10]
s13 = S[12]
s14 = S[13]
s15 = S[14]
s16 = S[15]
s19 = S[18]
s21 = S[20]

# -------- 공통: 제목줄 부제 교체 헬퍼 --------
def set_subtitle(slide, name, text):
    sh = find(slide, name)
    if sh and sh.has_text_frame:
        sh.text_frame.paragraphs[0].clear() if False else None
        # 첫 문단 첫 런만 교체
        p = sh.text_frame.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        r = p.add_run(); r.text = text
        r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = C(GRAY_T)

# ============================================================
# S4 — 기존 연구 한계 강조 + 본문 축소
# ============================================================
def edit_s4():
    # 본문 축소
    bodies = {
        "TextBox 6":  "Tetlock (2007) · 뉴스 비관론 → S&P500 수익률 예측\n"
                      "Da, Engelberg & Gao (2011) · Google 검색량(SVI) → 주가 선행 (Granger 확인)",
        "TextBox 11": "Campello et al. (2021) · StockX 스니커즈 가격 예측\n(소셜미디어 변수 추가 필요성 제안에 그침)\n"
                      "Dobrynskaya & Kishilova (2023) · 레고의 대체자산 특성 규명",
        "TextBox 16": "Bollen et al. (2011) · Twitter 감성 → DJIA 2일 선행 (Granger+SVM)\n"
                      "Zhang et al. (2018) · 뉴스+감성+가격 결합 멀티소스 예측",
    }
    for nm, txt in bodies.items():
        sh = find(s4, nm)
        paras = []
        for line in txt.split("\n"):
            paras.append(([(line, 12, INK, False)], PP_ALIGN.LEFT, 3))
        set_text(sh, paras)

    # 한계 박스 강조 (빨간 박스 + 굵게)
    lims = {
        "TextBox 8":  "전통 금융 자산(주식)에만 적용",
        "TextBox 13": "미디어 채널 신호 미포함 · 단일 자산만 분석",
        "TextBox 18": "주식 단일 자산 · 자산 간 채널 비교 없음 · Granger 미적용",
    }
    # 컬럼 x 좌표
    xmap = {"TextBox 8": 0.4, "TextBox 13": 4.7, "TextBox 18": 9.1}
    for nm, txt in lims.items():
        sh = find(s4, nm)
        x = xmap[nm]
        # 위치/크기 재설정 (본문 아래 빈 공간 채우기)
        sh.left = Inches(x); sh.top = Inches(4.55)
        sh.width = Inches(3.8); sh.height = Inches(1.45)
        sh.fill.solid(); sh.fill.fore_color.rgb = C(LRED)
        sh.line.color.rgb = C(DRED); sh.line.width = Pt(1.25)
        no_shadow(sh)
        set_text(sh, [
            ([("⚠ 한계", 13, DRED, True)], PP_ALIGN.LEFT, 4),
            ([(txt, 12.5, INK, True)], PP_ALIGN.LEFT, 0),
        ], anchor=MSO_ANCHOR.MIDDLE, mt=5, mb=5)
    # 부제 강조
    set_subtitle(s4, "TextBox 3",
                 "세 연구 흐름의 공통 한계 → '리셀 시장 × 다채널 × 선행성 비교'는 누구도 다루지 않았다")

# ============================================================
# S5 — RQ 표 -> 인포그래픽 (RQ1 -> RQ2 -> RQ3)
# ============================================================
def edit_s5():
    tbl = find(s5, "Table 4")
    if tbl: remove_shape(tbl)
    set_subtitle(s5, "TextBox 3",
                 "시간적 선행성 → 채널 이질성 → 예측 간결성 · 세 단계로 리셀 시장 미디어 효과를 해부")

    cards = [
        dict(x=0.30, head=BLUE,  lhead=LBLUE,  badge="RQ1",
             stage="시간적 선행성",
             q="어떤 미디어 채널이\n리셀 가격을 선행하는가?",
             test="Granger 인과 검정",
             why="과거 채널값이 미래 가격을\n예측하는지(시간 선행) 판별",
             metric="F 통계량 · p값 · 래그"),
        dict(x=4.62, head=GREEN, lhead=LGREEN, badge="RQ2",
             stage="채널 이질성",
             q="선행 채널 구성이\n자산마다 다른가?",
             test="자산별 유의 채널 집합 비교",
             why="자산마다 유의 채널이 다른지\n직접 비교해 특이성 확인",
             metric="유의 채널 집합 일치 여부"),
        dict(x=8.94, head=PURPLE, lhead=LPUR, badge="RQ3",
             stage="예측 간결성",
             q="선별 채널 모델이\n전채널 모델과 동등한가?",
             test="XGBoost A/B/C · DM 검정",
             why="적은 채널로도 같은 성능인지\n예측오차를 통계 검정",
             metric="AUC · ΔAUC · DM p값"),
    ]
    cw = 3.55
    ctop = 1.55
    chh = 5.05
    for c in cards:
        x = c["x"]
        # 카드 배경
        card = box(s5, x, ctop, cw, chh, fill=WHITE, line="DDDDDD", line_w=1.0)
        # 헤더 띠
        hd = box(s5, x, ctop, cw, 0.95, fill=c["head"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_text(hd, [
            ([(c["badge"] + "   ", 22, WHITE, True), (c["stage"], 15, WHITE, True)], PP_ALIGN.CENTER, 0),
        ], anchor=MSO_ANCHOR.MIDDLE)
        # 질문
        q = textbox(s5, x+0.18, ctop+1.08, cw-0.36, 0.95)
        set_text(q, [([(line, 13.5, INK, True)], PP_ALIGN.CENTER, 2) for line in c["q"].split("\n")],
                 anchor=MSO_ANCHOR.MIDDLE)
        # 검정 박스
        tb = box(s5, x+0.18, ctop+2.05, cw-0.36, 0.55, fill=c["lhead"], line=c["head"], line_w=1.0)
        set_text(tb, [([("검정  ", 10.5, c["head"], True), (c["test"], 11.5, INK, True)], PP_ALIGN.CENTER, 0)],
                 anchor=MSO_ANCHOR.MIDDLE)
        # 왜 이 검정
        wy = textbox(s5, x+0.18, ctop+2.70, cw-0.36, 0.95)
        runs = [([("왜 이 검정?", 10.5, c["head"], True)], PP_ALIGN.LEFT, 2)]
        for line in c["why"].split("\n"):
            runs.append(([(line, 11, GRAY_T, False)], PP_ALIGN.LEFT, 0))
        set_text(wy, runs, anchor=MSO_ANCHOR.TOP)
        # 측정 지표 (하단 띠)
        mt = box(s5, x+0.0, ctop+chh-0.62, cw, 0.62, fill=LGRAY, shape=MSO_SHAPE.RECTANGLE)
        set_text(mt, [([("측정  ", 10, c["head"], True), (c["metric"], 11, INK, True)], PP_ALIGN.CENTER, 0)],
                 anchor=MSO_ANCHOR.MIDDLE)
    # 화살표 (카드 사이)
    add_arrow(s5, 4.10, ctop+2.0, 0.50, 0.9, fill=BLUE)
    add_arrow(s5, 8.42, ctop+2.0, 0.50, 0.9, fill=GREEN)

# ============================================================
# S8 (+9 병합) — 채널 전처리 & 정상성 한 장
# ============================================================
def edit_s8_merge():
    # 제목/부제
    title = find(s8, "TextBox 1")
    if title:
        p = title.text_frame.paragraphs[0]
        for r in list(p.runs): r._r.getparent().remove(r._r)
        r = p.add_run(); r.text = "  전처리: 5개 미디어 채널 변환 및 정상성 확보"
        r.font.size = Pt(20); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = C(WHITE)
    set_subtitle(s8, "TextBox 3",
                 "채널마다 측정 단위가 달라 스케일이 다름 → 비교·모델 투입을 위해 전처리로 통일")

    # 기존 표 컬럼 축소 (개념적 프록시 열 제거, 스케일 = 전처리 결과로 표기)
    tbl = find(s8, "Table 4").table
    # 표를 새 내용으로 덮어쓰기: ID | 채널명 | 측정 대상 | 전처리(최종 스케일)
    new_rows = [
        ["ID", "채널명", "측정 대상", "전처리 → 최종 스케일"],
        ["CH1", "Google Trends", "소비자 검색 관심도", "월평균 → 0–100"],
        ["CH2", "뉴스 보도량", "미디어 노출 강도", "GDELT 월평균 → 0–1 정규화"],
        ["CH3", "뉴스 감성", "미디어 논조", "GDELT tone → −1~+1"],
        ["CH4", "유튜브 조회수", "정보 탐색 의도", "log1p → z-score"],
        ["CH5", "유튜브 댓글 감성", "소매 투자자 감성", "FinBERT P(pos)−P(neg) → −1~+1"],
    ]
    # 표 행 수가 6이므로 그대로 채움 (열은 6 -> 앞 4열만 사용, 뒤 2열은 비우고 폭 축소 어려우니 재구성)
    # 간단히: 기존 표 제거하고 새 표 생성
    remove_shape(find(s8, "Table 4"))
    rows, cols = 6, 4
    gt = s8.shapes.add_table(rows, cols, Inches(0.55), Inches(1.35), Inches(7.6), Inches(3.4)).table
    gt.columns[0].width = Inches(0.7)
    gt.columns[1].width = Inches(1.9)
    gt.columns[2].width = Inches(2.2)
    gt.columns[3].width = Inches(2.8)
    for ri in range(rows):
        for ci in range(cols):
            cell = gt.cell(ri, ci)
            cell.margin_left = Pt(4); cell.margin_right = Pt(4)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci >= 1 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = new_rows[ri][ci]
            r.font.name = FONT
            if ri == 0:
                r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = C(WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = C(NAVY)
            else:
                r.font.size = Pt(10.5); r.font.bold = (ci == 0)
                r.font.color.rgb = C(INK)
                cell.fill.solid(); cell.fill.fore_color.rgb = C(WHITE if ri % 2 else LGRAY)

    # 오른쪽: 스케일 통일 박스
    b1 = box(s8, 8.4, 1.35, 4.55, 1.55, fill=LBLUE, line=BLUE, line_w=1.0)
    set_text(b1, [
        ([("왜 채널마다 스케일이 다른가?", 12, BLUE, True)], PP_ALIGN.LEFT, 3),
        ([("각 채널은 검색량·보도량·감성 등 ", 11, INK, False),
          ("서로 다른 단위", 11, INK, True), ("를 측정.", 11, INK, False)], PP_ALIGN.LEFT, 2),
        ([("→ 모델(XGBoost) 투입 직전 ", 11, INK, False),
          ("StandardScaler로 z-score 통일", 11, BLUE, True)], PP_ALIGN.LEFT, 0),
    ], anchor=MSO_ANCHOR.MIDDLE)

    # 오른쪽: 정상성(ADF) 요약 박스 (왜 + 결과)
    b2 = box(s8, 8.4, 3.05, 4.55, 1.70, fill=LGREEN, line=GREEN, line_w=1.0)
    set_text(b2, [
        ([("왜 정상성(ADF) 검정?", 12, GREEN, True)], PP_ALIGN.LEFT, 3),
        ([("비정상 시계열은 ", 11, INK, False), ("허구 인과", 11, DRED, True),
          (" 위험 → Granger 전제조건", 11, INK, False)], PP_ALIGN.LEFT, 3),
        ([("ADF 결과: ", 11, INK, True), ("정상 9개 그대로 · 비정상 9개 1차 차분 후 사용", 11, INK, False)],
         PP_ALIGN.LEFT, 2),
        ([("예외: 스니커즈 CH1 차분 후도 p=0.120 (경계, 한계 명시)", 10, GRAY_T, False)], PP_ALIGN.LEFT, 0),
    ], anchor=MSO_ANCHOR.MIDDLE)

    # 하단 제한점 한 줄 유지 (TextBox 5) — 위치 조정
    note = find(s8, "TextBox 5")
    if note:
        note.top = Inches(5.0); note.left = Inches(0.55); note.width = Inches(7.6)

# ============================================================
# S13 — 모델설계 텍스트 -> ㅁ->ㅁ->ㅁ 플로우
# ============================================================
def edit_s13():
    # 제거: 표, 기존 그림, 부제 재설정
    for nm in ["Table 6", "그림 13"]:
        sh = find(s13, nm)
        if sh: remove_shape(sh)
    set_subtitle(s13, "TextBox 3",
                 "전채널(A) → SHAP으로 기여 채널 분석 → 상위만 선별(B·C) → DM 검정으로 성능 비교")

    flow_top = 1.45
    fh = 1.55
    # 4단계 박스
    steps = [
        dict(x=0.30, w=2.75, fill=NAVY,   t="① Model A\n전채널",
             sub="CH1~CH5 전부\n+ 가격 통제피처"),
        dict(x=3.55, w=2.75, fill=BLUE,   t="② SHAP 분석",
             sub="각 채널의 실제\n예측 기여도 측정"),
        dict(x=6.80, w=2.75, fill=GREEN,  t="③ 채널 선별\nModel B · C",
             sub="B: Granger 유의\nC: SHAP 상위 2채널"),
        dict(x=10.05, w=2.95, fill=PURPLE, t="④ DM 검정",
             sub="A vs 선별 모델\n예측오차 동등성"),
    ]
    for st in steps:
        b = box(s13, st["x"], flow_top, st["w"], fh, fill=st["fill"])
        set_text(b, [
            ([(line, 14, WHITE, True)], PP_ALIGN.CENTER, 2) for line in st["t"].split("\n")
        ] + [([(line, 10.5, WHITE, False)], PP_ALIGN.CENTER, 0) for line in st["sub"].split("\n")],
                 anchor=MSO_ANCHOR.MIDDLE)
    # 화살표
    for ax in [3.07, 6.32, 9.57]:
        add_arrow(s13, ax, flow_top+0.5, 0.46, 0.55, fill="9AA5B1")

    # 하단: A/B/C 채널 구성 미니 카드 (핵심 텍스트만)
    mini_top = 3.55
    minis = [
        dict(x=0.30, w=4.05, fill=LGRAY, head=NAVY,  name="Model A · 전채널",
             body="CH1+CH2+CH3+CH4+CH5\n(기준 베이스라인)"),
        dict(x=4.62, w=4.05, fill=LGREEN, head=GREEN, name="Model B · Granger 선별",
             body="스니커즈 CH3 · 카드 CH1 · 레고 —"),
        dict(x=8.94, w=4.05, fill=LPUR, head=PURPLE, name="Model C · SHAP 상위 2",
             body="스니커즈 CH1+CH5 · 카드 CH4+CH1\n레고 CH5+CH4"),
    ]
    for m in minis:
        b = box(s13, m["x"], mini_top, m["w"], 1.55, fill=m["fill"], line="DDDDDD", line_w=1.0)
        runs = [([(m["name"], 12.5, m["head"], True)], PP_ALIGN.LEFT, 4)]
        for line in m["body"].split("\n"):
            runs.append(([(line, 11.5, INK, False)], PP_ALIGN.LEFT, 1))
        set_text(b, runs, anchor=MSO_ANCHOR.MIDDLE, ml=10)

    # 공통 설정 한 줄
    cm = textbox(s13, 0.30, 5.30, 12.7, 0.4)
    set_text(cm, [([("공통  ", 11, BLUE, True),
                    ("통제피처 price_vs_ma3 · price_chg_lag1~3   |   검증 TimeSeriesSplit(5)   |   평가 AUC-ROC",
                     11, GRAY_T, False)], PP_ALIGN.LEFT, 0)])

# ============================================================
# 그래프 캡션 추가 (S11, S14, S16, S19)
# ============================================================
def add_captions():
    caption(s11, 0.2, 6.30, 4.9,
            "자산×채널 15회 Granger F값. 막대가 높을수록(스니커즈 CH3·카드 CH1) 가격 선행성이 강함.")
    # S14 AUC 비교 그래프
    caption(s14, 0.1, 5.55, 8.1,
            "자산별 Model A(전채널) vs B(선별) AUC. 막대 높이 차이가 거의 없어 두 모델 성능이 동등.")
    # S16 전체 모델 비교 그래프 (그래프 축소해 캡션 자리 확보)
    pic = resize_pic(s16, "Picture 7", 4.55, top=1.5, center=True)
    cl = Emu(pic.left).inches if pic else 0.9
    cw = Emu(pic.width).inches if pic else 11.5
    caption(s16, cl, 6.15, cw,
            "A/B/C/단독채널 모델의 AUC. 어느 구성이든 차이 최대 0.01 내외 → 채널 줄여도 성능 유지.")

# ============================================================
# S19 — DM 설명 축소
# ============================================================
def edit_s19():
    # 왼쪽 설명 박스들 제거(중복 다수) 후 핵심 1개로 대체
    for nm in ["TextBox 17","TextBox 18","TextBox 19","TextBox 20","TextBox 21",
               "TextBox 22","TextBox 23","TextBox 24","TextBox 25","TextBox 26","TextBox 27"]:
        sh = find(s19, nm)
        if sh: remove_shape(sh)
    # 중복 결론 박스 제거 (TextBox 34/35 - 둘 중 하나)
    for nm in ["TextBox 34","TextBox 35"]:
        sh = find(s19, nm)
        if sh: remove_shape(sh)
    # 핵심 'why DM' 박스 신설 (왼쪽)
    b = box(s19, 0.30, 1.15, 4.40, 3.85, fill=LBLUE, line=BLUE, line_w=1.0)
    set_text(b, [
        ([("왜 DM 검정인가?", 14, BLUE, True)], PP_ALIGN.LEFT, 6),
        ([("AUC 숫자(0.841 vs 0.840)만으론", 12, INK, False)], PP_ALIGN.LEFT, 1),
        ([("차이가 ", 12, INK, False), ("우연인지 실제인지", 12, DRED, True),
          (" 알 수 없다.", 12, INK, False)], PP_ALIGN.LEFT, 8),
        ([("DM 검정이 하는 일", 12.5, BLUE, True)], PP_ALIGN.LEFT, 2),
        ([("두 모델이 매달 낸 예측오차를", 12, INK, False)], PP_ALIGN.LEFT, 1),
        ([("시계열로 통째 비교", 12, INK, True), (" → 그 차이가", 12, INK, False)], PP_ALIGN.LEFT, 1),
        ([("우연 범위인지 통계로 판정.", 12, INK, False)], PP_ALIGN.LEFT, 8),
        ([("→ '채널 줄여도 성능 유지'를", 12, GREEN, True)], PP_ALIGN.LEFT, 1),
        ([("   통계적으로 입증 (간결성)", 12, GREEN, True)], PP_ALIGN.LEFT, 0),
    ], anchor=MSO_ANCHOR.TOP, mt=8, ml=10)
    # DM 그래프 축소 + 캡션 (표와 겹치지 않게)
    pic = resize_pic(s19, "Picture 15", 3.45, top=1.05, left=4.9)
    cw = Emu(pic.width).inches if pic else 7.0
    caption(s19, 4.9, 4.62, cw,
            "Model A vs C의 월별 예측오차 시계열. 두 곡선이 겹쳐 차이가 통계적으로 유의하지 않음(동등).")

# ============================================================
# S21 — IRF 설명 축소
# ============================================================
def edit_s21():
    # 박스1: 왜 IRF
    sh = find(s21, "TextBox 6")
    set_text(sh, [
        ([("Granger는 ", 13, INK, False), ("'선행 유무'만", 13, DRED, True), (" 판단.", 13, INK, False)], PP_ALIGN.LEFT, 4),
        ([("IRF는 충격의 ", 13, INK, False), ("방향·크기·지속기간", 13, BLUE, True),
          ("까지 추적.", 13, INK, False)], PP_ALIGN.LEFT, 4),
        ([("채널값 1σ 충격 → 이후 6개월", 12, GRAY_T, False)], PP_ALIGN.LEFT, 1),
        ([("가격 반응을 관측.", 12, GRAY_T, False)], PP_ALIGN.LEFT, 0),
    ], anchor=MSO_ANCHOR.TOP)
    # 박스2: 분석 대상/설정
    sh = find(s21, "TextBox 9")
    set_text(sh, [
        ([("분석 대상 (최고 유의 채널)", 12.5, BLUE, True)], PP_ALIGN.LEFT, 3),
        ([("① 스니커즈 ← CH3 뉴스감성", 12.5, INK, False)], PP_ALIGN.LEFT, 1),
        ([("② 카드 ← CH1 Google Trends", 12.5, INK, False)], PP_ALIGN.LEFT, 6),
        ([("설정: 6개월 추적 · 부트스트랩 1,000회", 11.5, GRAY_T, False)], PP_ALIGN.LEFT, 1),
        ([("95% CI가 0 미포함 → 유의", 11.5, GRAY_T, False)], PP_ALIGN.LEFT, 0),
    ], anchor=MSO_ANCHOR.TOP)
    # 박스3: 기대 해석
    sh = find(s21, "TextBox 12")
    set_text(sh, [
        ([("CH3↓(부정 뉴스) → 스니커즈 가격 하락?", 12.5, INK, False)], PP_ALIGN.LEFT, 4),
        ([("CH1↑(검색 급증) → 카드 가격 상승?", 12.5, INK, False)], PP_ALIGN.LEFT, 6),
        ([("→ 방향·크기를 IRF로 실증 확인", 12, GREEN, True)], PP_ALIGN.LEFT, 0),
    ], anchor=MSO_ANCHOR.TOP)

# ---------------- 실행 ----------------
edit_s4()
edit_s5()
edit_s8_merge()
edit_s13()
add_captions()
edit_s19()
edit_s21()

# ---------------- S9 삭제 ----------------
xml_slides = prs.slides._sldIdLst
slide_list = list(xml_slides)
xml_slides.remove(slide_list[8])   # 물리 슬라이드 9 (index 8)

# ---------------- 페이지 번호 재계산 ----------------
import re
pat = re.compile(r'^\s*\d+\s*/\s*\d+\s*$')
numbered = []
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.has_text_frame and pat.match(sh.text_frame.text.strip()):
            numbered.append(sh)
            break
total = len(numbered)
for i, sh in enumerate(numbered, 1):
    p = sh.text_frame.paragraphs[0]
    # 기존 서식 유지: 첫 런 텍스트만 교체
    if p.runs:
        for r in list(p.runs)[1:]:
            r._r.getparent().remove(r._r)
        p.runs[0].text = "%d / %d" % (i, total)
    else:
        r = p.add_run(); r.text = "%d / %d" % (i, total)
        r.font.size = Pt(9); r.font.name = FONT

prs.save(DST)
print("saved:", DST, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst), "| numbered:", total)
