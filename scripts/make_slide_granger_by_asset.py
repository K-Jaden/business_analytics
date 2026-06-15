#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RQ1·RQ2 Granger 인과검정 결과 — 자산별(스니커즈→카드→레고) 정리 (1장)
원문: 발표자료_최종_피드백반영_v3.pptx SLIDE19 (이미지 기반 슬라이드, 사진으로 직접 확인한 값 사용)
추가: F/p 유의성 판단 기준 설명, 자산별 그룹 막대그래프
Usage : python scripts/make_slide_granger_by_asset.py
Output: slide_granger_by_asset.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(BASE, "slide_granger_by_asset_v3.pptx")
FIG = os.path.join(BASE, "results", "figures", "granger_by_asset.png")

NAVY = RGBColor(0x1F, 0x2D, 0x4E)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
LGREEN = RGBColor(0xE7, 0xF2, 0xE9)
LORANGE = RGBColor(0xFC, 0xF0, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0xAA, 0xB1, 0xBC)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 21, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "RQ1·RQ2 Granger 인과검정 결과 — 자산별 정리(스니커즈·카드·레고)",
       [("자산별로 묶어 다시 정리 + ", 13, False, DGRAY), ("유의성 판단 기준(F·p)", 13, True, ORANGE),
        (" 추가", 13, False, DGRAY)])

MARGIN = 0.30

# ════════════════════════════════════════════════════════════════════════
# 좌측 — 자산별 그룹 막대그래프
# ════════════════════════════════════════════════════════════════════════
CHART_T, CHART_W, CHART_H = 1.36, 6.55, 3.34
box(s1, MARGIN, CHART_T, CHART_W, CHART_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)

from PIL import Image as _PILImage
_fig_w_px, _fig_h_px = _PILImage.open(FIG).size
_fig_ratio = _fig_h_px / _fig_w_px
_pic_h = CHART_H - 0.12
_pic_w = _pic_h / _fig_ratio
pic = s1.shapes.add_picture(FIG, Inches(MARGIN + (CHART_W - _pic_w) / 2), Inches(CHART_T + 0.06),
                             height=Inches(_pic_h))

# ════════════════════════════════════════════════════════════════════════
# 우측 — 자산별 정렬 표 (스니커즈 → 카드 → 레고, F값 내림차순)
# ════════════════════════════════════════════════════════════════════════
TB_X = MARGIN + CHART_W + 0.20
TB_W = 13.33 - MARGIN - TB_X
TB_T = 1.36

rows = [
    ("스니커즈", "CH3 뉴스감성",      "1", "13.38", "0.0007", "유의"),
    ("스니커즈", "CH4 YT조회수",      "1", "4.73",  "0.035",  "유의"),
    ("스니커즈", "CH1 Google Trends", "2", "3.47",  "0.041",  "유의"),
    ("스니커즈", "CH2 뉴스량",        "1", "1.55",  "0.220",  "비유의"),
    ("스니커즈", "CH5 YT댓글감성",    "1", "0.09",  "0.764",  "비유의"),
    ("카드",     "CH1 Google Trends", "1", "8.42",  "0.006",  "유의"),
    ("카드",     "CH3 뉴스감성",      "1", "1.41",  "0.242",  "비유의"),
    ("카드",     "CH2 뉴스량",        "1", "0.73",  "0.398",  "비유의"),
    ("카드",     "CH5 YT댓글감성",    "1", "0.61",  "0.438",  "비유의"),
    ("카드",     "CH4 YT조회수",      "1", "0.25",  "0.622",  "비유의"),
    ("레고",     "CH1 Google Trends", "1", "3.92",  "0.054",  "비유의"),
    ("레고",     "CH4 YT조회수",      "1", "1.83",  "0.183",  "비유의"),
    ("레고",     "CH3 뉴스감성",      "1", "0.40",  "0.530",  "비유의"),
    ("레고",     "CH2 뉴스량",        "1", "0.32",  "0.576",  "비유의"),
    ("레고",     "CH5 YT댓글감성",    "1", "0.17",  "0.680",  "비유의"),
]

VERDICT_COLOR = {"유의": GREEN, "비유의": MGRAY}
ASSET_BG = {"스니커즈": LGREEN, "카드": RGBColor(0xEA, 0xF1, 0xFA), "레고": LORANGE}

col_w = [1.05, 2.10, 0.55, 0.70, 0.85, 0.85]
col_x = [TB_X]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)
hdr_h = 0.34
row_h = (CHART_H - hdr_h) / len(rows)

headers = ["자산", "채널", "래그", "F", "p", "판정"]
for x, w, htxt in zip(col_x, col_w, headers):
    box(s1, x, TB_T, w, hdr_h, fill=NAVY, line=WHITE, line_w=0.75)
    text(s1, x, TB_T, w, hdr_h, [[(htxt, 10.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prev_asset = None
for i, (asset, ch, lag, F, p, verdict) in enumerate(rows):
    y = TB_T + hdr_h + i * row_h
    bg = ASSET_BG[asset]
    cells = [asset if asset != prev_asset else "", ch, lag, F, p, verdict]
    prev_asset = asset
    for j, (x, w, val) in enumerate(zip(col_x, col_w, cells)):
        fill = bg
        col = INK; bold = False; sz = 9.5
        if j == 0 and val:
            bold = True; col = NAVY; sz = 10
        if j == 5:
            fill = VERDICT_COLOR[verdict]; col = WHITE; bold = True; sz = 9.5
        box(s1, x, y, w, row_h, fill=fill, line=WHITE, line_w=0.5)
        text(s1, x, y, w, row_h, [[(val, sz, bold, col)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════════
# 하단 — 유의성 판단 기준 설명
# ════════════════════════════════════════════════════════════════════════
CRIT_T = CHART_T + CHART_H + 0.16
CRIT_H = 1.55
box(s1, MARGIN, CRIT_T, 13.33 - 2 * MARGIN, CRIT_H, fill=RGBColor(0xEA, 0xF1, 0xFA),
    line=BLUE, line_w=1.0, round_=True)
text(s1, MARGIN + 0.18, CRIT_T + 0.08, 13.33 - 2 * MARGIN - 0.36, 0.36,
     [[("F값·p값, 어느 기준으로 '유의'를 판단하나?", 13, True, BLUE)]])
text(s1, MARGIN + 0.18, CRIT_T + 0.46, 13.33 - 2 * MARGIN - 0.36, CRIT_H - 0.54,
     [[("•  F값에는 '이 값 이상이면 유의'라는 고정된 기준이 없음 — 임계값이 표본 수(n=48)·래그(lag)마다 달라지기 때문", 11, False, INK)],
      [("    (lag=1이면 F ", 11, False, INK), ("≥ 약 4.06", 11, True, BLUE),
       (", lag=2이면 F ", 11, False, INK), ("≥ 약 3.21", 11, True, BLUE),
       ("일 때 p<0.05 — 같은 'p<0.05'도 lag에 따라 F 값 자체는 다름)", 11, False, INK)],
      [("•  그래서 그래프는 F가 아닌 ", 11, False, INK), ("p값이 작을수록 막대가 높아지도록 환산", 11, True, GREEN),
       ("해 표시 — 모든 채널에 동일한 기준선(p=0.05) 하나로 색상을 그대로 판단 가능", 11, False, INK)],
      [("•  본 연구의 최종 판정 기준은 ", 11, False, INK), ("p < 0.05", 11, True, GREEN),
       (" 단 하나 — 기준선을 넘으면 '유의', 못 넘으면 '비유의'. F·lag 원값은 우측 표 참고", 11, False, INK)],
      [("•  다중비교(BH) 보정을 적용하면 '유의' 4개 중 2개(스니커즈 CH3, 카드 CH1)만 남음 — 보수적 해석 시 참고", 11, False, INK)]],
     space=3)

# ════════════════════════════════════════════════════════════════════════
# 최하단 — RQ1 / RQ2 결론
# ════════════════════════════════════════════════════════════════════════
RQ_T = CRIT_T + CRIT_H + 0.14
RQ_H = 7.5 - RQ_T - 0.15
RQ_W = (13.33 - 2 * MARGIN - 0.20) / 2

box(s1, MARGIN, RQ_T, RQ_W, RQ_H, fill=GREEN, round_=True)
text(s1, MARGIN + 0.16, RQ_T, RQ_W - 0.32, RQ_H,
     [[("결론 RQ1  ", 12.5, True, WHITE)],
      [("미디어가 가격 선행성 부분 지지 — 스니커즈 3개(CH3·CH4·CH1)·카드 1개(CH1) 채널이 유의(p<0.05)",
        11.5, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE, space=3)

box(s1, MARGIN + RQ_W + 0.20, RQ_T, RQ_W, RQ_H, fill=GREEN, round_=True)
text(s1, MARGIN + RQ_W + 0.36, RQ_T, RQ_W - 0.32, RQ_H,
     [[("결론 RQ2  ", 12.5, True, WHITE)],
      [("자산 유형마다 선행하는 채널 구성이 다름 → 레고는 유의 채널 없음, 스니커즈·카드는 서로 다른 채널 조합 → RQ2 지지",
        11.5, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE, space=3)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
