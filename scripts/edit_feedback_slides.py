"""최종_버전_수정.pptx 개별 슬라이드 텍스트만 수정 (전체 재생성 없이).

피드백 반영 내용:
1) RQ2: Jaccard 유사도 설명/수치 제거 -> "Granger 결과에서 자산별로 유의 채널이
   다르게 나왔다" 는 직접 비교로 단순화 (slide 1, 4, 23)
2) ADF(slide 8): '왜 정상성 검정이 필요한가' 이론 설명 축약
   (정상 시계열이란?/비정상이면 어떤 문제? 절을 한 줄로 압축)
3) DM(slide 18): 좌측 설명 패널 3블록 본문을 더 간결하게 압축
   (헤더/레이아웃은 그대로 두고 본문만 축약 -> 레이아웃 깨짐 방지)

기존 런(run)의 폰트(size/bold/color/name)를 그대로 유지한 채 텍스트만 교체한다.
"""
from pptx import Presentation

PATH = "최종_버전_수정.pptx"


def replace_textframe(shape, new_text):
    """첫 문단의 첫 런 서식을 유지한 채 본문 텍스트를 통째로 교체."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    ref = p.runs[0]
    size, bold, color, name = ref.font.size, ref.font.bold, ref.font.color.rgb, ref.font.name
    # 여분 문단 제거
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    # 기존 런 제거
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = new_text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def shape_by_id(slide, shape_id):
    for shp in slide.shapes:
        if shp.shape_id == shape_id:
            return shp
    raise KeyError(shape_id)


def main():
    prs = Presentation(PATH)

    # ---------- 1) RQ2: Jaccard 제거 ----------
    # slide 1: 목차 섹션 설명에서 'Jaccard 유사도로 본 채널 구성 차이' 표현 제거
    s1 = prs.slides[1]
    replace_textframe(
        shape_by_id(s1, 23),
        "자산×채널 15회 검정 (Benjamini-Hochberg 보정) · 자산별 유의 채널 구성 비교",
    )

    # slide 4: 연구문제 표 RQ2 행 — 검정 방법 / 측정 지표
    s4 = prs.slides[4]
    table4 = next(shp for shp in s4.shapes if shp.has_table).table
    replace_textframe(table4.cell(2, 2),
                      "RQ1에서 자산별로\n유의했던 채널 집합을\n직접 비교")
    replace_textframe(table4.cell(2, 3),
                      "유의 채널 집합의\n자산 간 일치 여부")

    # ---------- 2) ADF 슬라이드 축약 ----------
    s8 = prs.slides[8]
    adf_box = shape_by_id(s8, 7)
    replace_textframe(adf_box,
        "Granger 검정의 전제 조건:\n"
        "시계열이 정상(I(0))이어야\n"
        "통계 검정이 유효함\n"
        "-> 비정상 시계열 사용 시\n"
        "   허구 상관(spurious correlation) 위험\n"
        "\n"
        "검정 기준 (ADF 단위근 검정):\n"
        "  p < 0.05 -> 정상 (그대로 사용)\n"
        "  p >= 0.05 -> 비정상 -> 1차 차분 후 재검정\n"
        "\n"
        "최종 결과:\n"
        "  정상 I(0) 9개  /  차분 후 I(1) 9개\n"
        "  예외: 스니커즈 CH1 (차분 후도 p=0.120, 경계)"
    )

    # ---------- 3) DM 슬라이드 좌측 설명 패널 본문 축약 ----------
    s18 = prs.slides[18]
    replace_textframe(shape_by_id(s18, 22),
        'A AUC = 0.841, C AUC = 0.840\n'
        '→ 차이가 우연인지 실제 차이인지\n'
        '   숫자만으로는 알 수 없음'
    )
    replace_textframe(shape_by_id(s18, 25),
        '예측 오차 시계열을 직접 비교해\n'
        '통계적 유의성을 검정\n'
        '(H₀: 두 모델의 오차가 같다)\n'
        '\n'
        'p > 0.05 → H₀ 기각 실패\n'
        '→ "통계적으로 동등함"을 공식 확인'
    )
    replace_textframe(shape_by_id(s18, 28),
        'RQ3 핵심 주장 —\n'
        '"5개 채널 대신 2개만 써도 충분하다"\n'
        '\n'
        '→ DM 검정으로 그 근거를\n'
        '   통계적으로 확보 (parsimony)'
    )

    # ---------- 1) RQ2: 결론 슬라이드 표에서도 Jaccard 수치 -> Granger 직접비교로 ----------
    s23 = prs.slides[23]
    table23 = next(shp for shp in s23.shapes if shp.has_table).table
    replace_textframe(table23.cell(2, 2),
        "  스니커즈: CH1·CH3·CH4 유의\n"
        "  카드: CH1만 유의\n"
        "  레고: 유의 채널 없음\n"
        "\n"
        "  -> 자산마다 선행 채널 구성이 다름"
    )
    replace_textframe(table23.cell(2, 3),
        "RQ1 Granger 결과의\n유의 채널 집합 비교"
    )

    prs.save(PATH)
    print("saved:", PATH)


if __name__ == "__main__":
    main()
