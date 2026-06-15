#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2장 구성 (v3)
 ① 기존 연구 & 한계 — 연구의 계보 (3카드: Bollen 2011 · 유재필2023 · CausalStock 2024)
 ② 핵심 선행연구 2편 (Bollen 제외) + 하단 비교표(사용 채널 / 방법론 / 본 연구와 차이)
검증된 실제 논문만 인용.
Usage : python scripts/make_slides_prior_work_v3.py
Output: slides_prior_work_v3.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(BASE, "slides_prior_work_v3.pptx")

NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
ORANGE= RGBColor(0xE8, 0x77, 0x22)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
LGREEN= RGBColor(0xE7, 0xF2, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0x77, 0x77, 0x77)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
INK   = RGBColor(0x0A, 0x0A, 0x0A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def box(sld, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = sld.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sld, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def header(sld, title, subtitle_runs):
    box(sld, 0, 0, 13.33, 0.82, fill=NAVY)
    text(sld, 0.25, 0.06, 12.8, 0.70, [[(title, 22, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    box(sld, 0, 0.82, 13.33, 0.045, fill=BLUE)
    text(sld, 0.25, 0.92, 12.8, 0.40, [subtitle_runs], anchor=MSO_ANCHOR.MIDDLE)


def paper_card(sld, x, t, w, h, tag, title, auth, lines, tag_fill=BLUE, ts=12, ls=10.5):
    box(sld, x, t, w, h, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(sld, x, t, w, 0.42, fill=tag_fill, round_=True)
    text(sld, x, t + 0.02, w, 0.38, [[(tag, 11.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(sld, x + 0.06, t + 0.48, w - 0.12, 0.92,
         [[(title, ts, True, NAVY)], [(auth, 9.5, False, MGRAY)]], space=3)
    body = [[("•  " + ln, ls, False, INK)] for ln in lines]
    text(sld, x + 0.10, t + 1.40, w - 0.20, h - 1.48, body, space=6)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 1 — 기존 연구 & 한계 (세 연구 흐름, v2 6페이지 논문 반영)
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "기존 연구 & 한계 — 세 연구 흐름의 빈 교점",
       [("미디어→주식 · 리셀 예측 · 멀티채널 — 셋 다 ", 13, False, DGRAY),
        ("‘리셀 × 다채널 × 선행성 비교’", 13, True, ORANGE),
        ("의 교점은 비워 둠", 13, False, DGRAY)])

C_T, C_H, C_W, c_gap = 1.45, 3.95, 4.13, 0.22
xs = [0.30, 0.30 + C_W + c_gap, 0.30 + 2 * (C_W + c_gap)]
streams = [
    ("미디어 → 주식 선행성",
     [("Tetlock (2007)", "뉴스 비관론 → S&P500 하락"),
      ("Da et al. (2011)", "Google 검색량 → 주가 선행")],
     "전통 금융자산(주식)에만 적용"),
    ("리셀 상품 가격 예측",
     [("Campello et al. (2021)", "StockX 스니커즈 가격 예측"),
      ("Dobrynskaya & Kishilova (2023)", "레고를 대체자산으로 분석")],
     "미디어 채널 신호 미포함 · 단일 자산"),
    ("멀티채널 예측 방법론",
     [("Bollen et al. (2011)", "Twitter 감성 → DJIA 선행"),
      ("Zhang et al. (2018)", "뉴스+감성+가격 결합 예측")],
     "주식 단일 · 채널 비교 없음 · Granger 미적용"),
]
for x, (sname, papers, limit) in zip(xs, streams):
    box(s1, x, C_T, C_W, C_H, fill=LGRAY, line=RGBColor(0xCF, 0xD6, 0xDF), line_w=1.0, round_=True)
    box(s1, x, C_T, C_W, 0.46, fill=BLUE, round_=True)
    text(s1, x, C_T + 0.03, C_W, 0.40, [[(sname, 12.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    pruns = []
    for (auth, desc) in papers:
        pruns.append([(auth, 11.5, True, NAVY)])
        pruns.append([("   " + desc, 11, False, INK)])
    text(s1, x + 0.12, C_T + 0.62, C_W - 0.24, 2.0, pruns, space=6)
    lb_t = C_T + C_H - 0.92
    box(s1, x + 0.12, lb_t, C_W - 0.24, 0.80, fill=RGBColor(0xFC, 0xF0, 0xE6),
        line=ORANGE, line_w=1.0, round_=True)
    text(s1, x + 0.18, lb_t + 0.02, C_W - 0.36, 0.76,
         [[("⚠ 한계  ", 10.5, True, ORANGE), (limit, 10.5, False, INK)]],
         anchor=MSO_ANCHOR.MIDDLE, space=1)

KEY_T = C_T + C_H + 0.18
box(s1, 0.30, KEY_T, 12.73, 0.74, fill=GREEN, round_=True)
text(s1, 0.45, KEY_T + 0.03, 12.4, 0.68,
     [[("본 연구  ", 13, True, WHITE),
       ("스니커즈·카드·레고 3자산 × 5채널 × Granger × SHAP+DM", 12, True, WHITE),
       (" — 세 흐름의 빈 교점을 메움", 12, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — 핵심 선행연구 2편 + 비교표
# ════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
header(s2, "본 연구와 가장 맞닿은 선행연구 2편 — 직접 비교",
       [("사용 ", 13, False, DGRAY), ("채널·방법론·분석 대상", 13, True, BLUE),
        ("을 본 연구와 나란히 비교", 13, False, DGRAY)])

c_t, c_h, c_w = 1.45, 2.60, 6.16
xs2 = [0.30, 0.30 + c_w + 0.31]
paper_card(s2, xs2[0], c_t, c_w, c_h, "국내 · Granger · 2023",
           "유튜브 기반 데이터 분석을 통한 주가 선행성 분석",
           "유재필·김문선 (2023), 정보화연구",
           ["채널 = 유튜브 조회수·검색빈도 (2개 지표)",
            "학습법 = Granger 인과분석만 (예측모델 없음)",
            "결과 = 유튜브가 ETF(주식) 시장을 선행함을 통계 확인"], ts=12.5, ls=11)
paper_card(s2, xs2[1], c_t, c_w, c_h, "Granger + 예측 · Electronics 2024",
           "Granger 인과 + LSTM 기반 금융 예측",
           "Intelligent Financial Forecasting (Electronics·MDPI, 2024)",
           ["채널 = 금융·거시 시계열 변수 (가격·지표 등)",
            "변수 선별 = Granger 인과 + 상관분석",
            "학습법 = 베이지안 최적화 + LSTM (딥러닝)",
            "본 연구와 동일한 ‘Granger 선별 → 예측’ 구조"], ts=12.5, ls=11)

# ── 하단 비교표 ──────────────────────────────────────────────────────────
tb_t = c_t + c_h + 0.20
cols_x = [0.30, 1.95, 5.55, 9.15]
cols_w = [1.65, 3.60, 3.60, 3.88]
col_hdr = ["구분", "유재필·김문선(2023)", "Granger+LSTM (2024)", "본 연구"]
rows = [
    ("사용 채널", "유튜브 조회수·검색빈도 (2지표)", "금융·거시 시계열 (미디어 아님)",
     "5개 채널 (검색·뉴스량·뉴스감성·조회수·댓글)"),
    ("방법론", "Granger 인과분석", "Granger 선별 + LSTM",
     "Granger + XGBoost · SHAP · DM"),
    ("분석 대상", "ETF (주식)", "시장지수·자산가격 (주식)",
     "리셀 3종 (스니커즈·카드·레고)"),
]
hdr_h, row_h = 0.46, 0.58

# 헤더 행
for j, (cx, cw, htxt) in enumerate(zip(cols_x, cols_w, col_hdr)):
    fill = GREEN if j == 3 else NAVY
    box(s2, cx, tb_t, cw, hdr_h, fill=fill, line=WHITE, line_w=1.0)
    text(s2, cx, tb_t, cw, hdr_h, [[(htxt, 11.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 데이터 행
for i, row in enumerate(rows):
    y = tb_t + hdr_h + i * row_h
    for j, (cx, cw) in enumerate(zip(cols_x, cols_w)):
        if j == 0:
            cell_fill = RGBColor(0xEC, 0xEF, 0xF3); bold = True; col = NAVY; sz = 11
        elif j == 3:
            cell_fill = LGREEN; bold = True; col = INK; sz = 10.5
        else:
            cell_fill = WHITE if i % 2 == 0 else RGBColor(0xF7, 0xF8, 0xFA)
            bold = False; col = INK; sz = 10.5
        box(s2, cx, y, cw, row_h, fill=cell_fill, line=RGBColor(0xD5, 0xDA, 0xE0), line_w=0.75)
        text(s2, cx + 0.02, y, cw - 0.04, row_h, [[(row[j], sz, bold, col)]],
             align=PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)

# ── 하단 핵심 요약 ───────────────────────────────────────────────────────
ky = tb_t + hdr_h + 3 * row_h + 0.16
box(s2, 0.30, ky, 12.73, 0.66, fill=GREEN, round_=True)
text(s2, 0.45, ky + 0.02, 12.4, 0.62,
     [[("차이 요약  ", 13, True, WHITE),
       ("선행연구는 같은 'Granger 선별→예측'을 ", 12, False, WHITE),
       ("주식·금융 도메인", 12, True, WHITE),
       ("에 적용 — 본 연구는 ", 12, False, WHITE),
       ("‘미디어 다채널 × 리셀’", 12, True, WHITE),
       ("로 확장한 최초 검증", 12, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"[OK] {OUT}  (slides={len(prs.slides._sldIdLst)})")
