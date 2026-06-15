#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
슬라이드 — A/B/C/N-A + 채널단독 전체 비교 (전부 단색 채움 버전)
make_slide11.py 와 동일한 데이터·표 구성이지만, 그래프의 B/C/N-A 점선-빈칸
막대와 CH1~5 파스텔(alpha) 막대를 모두 단색 채움 막대로 변경.
Usage : python scripts/make_slide_model_comparison_filled.py
Output: slide_model_comparison_filled.pptx
"""
import os, io
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RES  = os.path.join(BASE, "results")
OUT  = os.path.join(BASE, "slide_model_comparison_filled.pptx")

# ── 데이터 로드 ───────────────────────────────────────────────────────────
abl = pd.read_csv(os.path.join(RES, "ablation_results.csv"))
mc  = pd.read_csv(os.path.join(RES, "model_c_results.csv"))

abl_auc = abl.pivot_table(index="model", columns="asset_type", values="auc")
abl_vsA = abl.pivot_table(index="model", columns="asset_type",
                          values="vs_A", aggfunc="first")

def get(m, a):
    try:   return float(abl_auc.loc[m, a])
    except: return None

def vsa(m, a):
    try:
        v = str(abl_vsA.loc[m, a])
        return {"no diff": "차이없음", "A better": "A 우수"}.get(
            v, "ablation 우수" if "ablation" in v else v)
    except: return "—"

ASSETS = ["sneakers", "cards", "lego"]
A_auc  = {a: get("A", a)             for a in ASSETS}
C_auc  = {r["asset_type"]: r["auc_C"] for _, r in mc.iterrows()}
C_vsa  = {r["asset_type"]: ("차이없음" if "no diff" in r["verdict"] else "A 우수")
          for _, r in mc.iterrows()}
NA_auc = {a: get("A-dropGranger", a)  for a in ASSETS}
NA_vsa = {a: vsa("A-dropGranger", a)  for a in ASSETS}

# B (Granger raw p<0.05 유의 채널):
#   sneakers → CH1+CH3+CH4 조합: ablation 미보유 → 3채널 평균 근사
#   cards    → CH1-only (exact)
#   lego     → 유의 채널 없음 = A
B_auc = {
    "sneakers": round((get("CH1-only","sneakers") +
                       get("CH3-only","sneakers") +
                       get("CH4-only","sneakers")) / 3, 4),  # ≈ 0.8461
    "cards":    get("CH1-only","cards"),                     # 0.8612 exact
    "lego":     get("A","lego"),                             # = A (no channels)
}

CH_models = ["CH1-only","CH2-only","CH3-only","CH4-only","CH5-only"]

# ── 색상 정의 (전부 단색 채움) ──────────────────────────────────────────────
C_A   = "#1F2D4E"
C_B   = "#E87722"
C_C   = "#2E6DA4"
C_NA  = "#7B3FA0"
C_CH  = ["#1F8FE0","#3FA84A","#E08A20","#D9B500","#D9534F"]
ASSET_TITLE = {"sneakers": "스니커즈", "cards": "카드", "lego": "레고"}

# 세 패널 공통 y축 (무작위 0.5 ~ 완벽 1.0) — 실제 AUC 차이(0.00x)가
# 시각적으로도 "차이 없음"으로 보이도록 자산별로 다른 축을 쓰지 않음
Y_LO, Y_HI = 0.5, 1.0
YLIM = {a: (Y_LO, Y_HI) for a in ASSETS}

# ── 그래프 생성 ───────────────────────────────────────────────────────────
plt.rcParams["font.family"] = ["Malgun Gothic", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False

fig, axes = plt.subplots(1, 3, figsize=(15.5, 4.8), sharey=True)
fig.patch.set_facecolor("white")

def draw_bar(ax, x, height, color, width=0.52):
    """모든 막대를 단색으로 채움 (점선/빈칸/파스텔 alpha 제거)"""
    ax.bar(x, height, width, color=color,
           edgecolor="white", linewidth=0.5, zorder=3)

def annotate_bar(ax, x, y, txt, color, fontsize=7.0):
    ax.text(x, y + (ax.get_ylim()[1] - ax.get_ylim()[0]) * 0.007,
            txt, ha="center", va="bottom", fontsize=fontsize,
            color=color, fontweight="bold", rotation=0)

for ax, asset in zip(axes, ASSETS):
    ax.set_facecolor("#F9F9F9")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4, zorder=0)

    if asset == "lego":
        # A=B=N-A 병합 (레고 Granger 유의 없음)
        bars_info = [
            (0,  A_auc[asset], C_A, "A=B=N-A"),
            (1,  C_auc[asset], C_C, "C"),
        ]
        for xi, h, col, lbl in bars_info:
            draw_bar(ax, xi, h, col)
        # 채널 단독
        for j, (ch, col) in enumerate(zip(CH_models, C_CH)):
            xi = 2 + j
            h  = get(ch, asset)
            draw_bar(ax, xi, h, col)
        # 수치 표시
        ymin, ymax = YLIM[asset]
        ax.set_ylim(ymin, ymax)
        ax.set_xticks(range(7))
        ax.set_xticklabels(
            ["A=B=N-A","C","CH1","CH2","CH3","CH4","CH5"],
            fontsize=8.5, fontweight="bold")
        for xi, h, col, _ in bars_info:
            annotate_bar(ax, xi, h, f"{h:.4f}", col)
        for j, (ch, col) in enumerate(zip(CH_models, C_CH)):
            h = get(ch, asset)
            annotate_bar(ax, 2+j, h, f"{h:.4f}", "#555555", fontsize=6.5)
        # A 기준선
        ax.axhline(A_auc[asset], color=C_A, linestyle="--",
                   linewidth=0.9, alpha=0.5, zorder=4)
    else:
        # sneakers / cards: A/B/C/N-A/CH1-5 모두 단색 채움
        main_bars = [
            (0, A_auc[asset],  C_A,  "A"),
            (1, B_auc[asset],  C_B,  "B≈" if asset=="sneakers" else "B"),
            (2, C_auc[asset],  C_C,  "C"),
            (3, NA_auc[asset], C_NA, "N-A"),
        ]
        for xi, h, col, lbl in main_bars:
            draw_bar(ax, xi, h, col)

        for j, (ch, col) in enumerate(zip(CH_models, C_CH)):
            xi = 4 + j
            h  = get(ch, asset)
            draw_bar(ax, xi, h, col)

        ymin, ymax = YLIM[asset]
        ax.set_ylim(ymin, ymax)
        ax.set_xticks(range(9))
        ax.set_xticklabels(
            ["A","B≈" if asset=="sneakers" else "B","C","N-A",
             "CH1","CH2","CH3","CH4","CH5"],
            fontsize=8.5, fontweight="bold")

        for xi, h, col, _ in main_bars:
            annotate_bar(ax, xi, h, f"{h:.4f}", col)
        for j, (ch, col) in enumerate(zip(CH_models, C_CH)):
            h = get(ch, asset)
            annotate_bar(ax, 4+j, h, f"{h:.4f}", "#555555", fontsize=6.5)

        # A 기준선
        ax.axhline(A_auc[asset], color=C_A, linestyle="--",
                   linewidth=0.9, alpha=0.5, zorder=4)
        # 구분선 (주요 모델 / 채널단독)
        ax.axvline(3.5, color="#BBBBBB", linewidth=0.9, linestyle=":", zorder=4)
        ax.annotate("← 주요 모델", xy=(3.4, ymin + (ymax-ymin)*0.03),
                    ha="right", fontsize=7, color="#999999")
        ax.annotate("채널 단독 →", xy=(3.6, ymin + (ymax-ymin)*0.03),
                    ha="left",  fontsize=7, color="#999999")

    # 무작위 추측 기준선 (0.5) — 모든 막대가 이 기준선 위에서 얼마나
    # "거의 동일한 높이"인지를 공통 축에서 직접 비교할 수 있게 함
    ax.axhline(0.5, color="#C0392B", linestyle="--", linewidth=1.0, alpha=0.7, zorder=2)
    if asset == "sneakers":
        ax.annotate("무작위 수준 0.5", xy=(0.05, 0.508), fontsize=7.5,
                    color="#C0392B", va="bottom")

    ax.set_title(ASSET_TITLE[asset], fontsize=13, fontweight="bold", pad=6)
    ax.set_ylabel("AUC-ROC" if asset == "sneakers" else "", fontsize=9)

# 범례 (모두 단색 채움 패치)
legend_handles = [
    mpatches.Patch(facecolor=C_A,  edgecolor="white", label="A  전채널"),
    mpatches.Patch(facecolor=C_B,  edgecolor="white", label="B  Granger 유의 채널"),
    mpatches.Patch(facecolor=C_C,  edgecolor="white", label="C  SHAP 선별 채널"),
    mpatches.Patch(facecolor=C_NA, edgecolor="white", label="N-A  Granger 유의 채널 제거"),
] + [mpatches.Patch(facecolor=c, edgecolor="white", label=f"CH{i+1} 단독")
     for i, c in enumerate(C_CH)]

fig.legend(handles=legend_handles, loc="lower center", ncol=4,
           fontsize=8.5, bbox_to_anchor=(0.5, -0.02),
           framealpha=0.97, edgecolor="#cccccc")
fig.suptitle("모델별 AUC-ROC 비교  (공통 축 0.5~1.0 · 전체 막대 단색 채움)",
             fontsize=12, fontweight="bold", y=1.01)
fig.text(0.5, -0.10, "† B(스니커즈): CH1+CH3+CH4 조합 AUC 미산출 → 3채널 개별 AUC 평균 근사치 (0.8461)",
         ha="center", fontsize=8, color="#888888", style="italic")

plt.tight_layout(rect=[0, 0.11, 1, 1])
fig_buf = io.BytesIO()
fig.savefig(fig_buf, dpi=160, bbox_inches="tight", facecolor="white")
plt.close(fig)

# ── PPT 헬퍼 ─────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1F, 0x2D, 0x4E)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DGRAY = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0xF2, 0xF2, 0xF2)
BORD  = RGBColor(0xBB, 0xBB, 0xBB)
BG = {
    "A":    RGBColor(0xD9,0xDE,0xF5),
    "B":    RGBColor(0xFF,0xE8,0xCC),
    "C":    RGBColor(0xCC,0xE5,0xFF),
    "N-A":  RGBColor(0xEA,0xD5,0xFF),
    "CH":   RGBColor(0xF5,0xF5,0xF5),
    "hdr":  NAVY,
}
G_ND  = RGBColor(0xD4,0xED,0xDA)
G_AB  = RGBColor(0xFF,0xCC,0xCC)
G_ABL = RGBColor(0xCC,0xE5,0xFF)

def _set_bg(cell, rgb):
    tc = cell._tc; p = tc.get_or_add_tcPr()
    for c in list(p):
        if c.tag.split("}")[-1] in ("solidFill","gradFill","noFill","blipFill","pattFill"):
            p.remove(c)
    sf  = etree.SubElement(p, qn("a:solidFill"))
    clr = etree.SubElement(sf, qn("a:srgbClr"))
    clr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
    p.remove(sf); p.insert(0, sf)

def add_box(sld, l, t, w, h, fill=None):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else:    tb.fill.background()
    tb.line.fill.background()
    return tb

def add_text(sld, text, l, t, w, h, bold=False, italic=False,
             size=11, color=BLACK, align=PP_ALIGN.LEFT):
    tb = add_box(sld, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.bold=bold; r.font.italic=italic
    r.font.size=Pt(size); r.font.color.rgb=color; r.font.name="Calibri"
    return tb

def verdict_bg(v):
    if v == "차이없음":      return G_ND
    if v == "A 우수":        return G_AB
    if "ablation" in v:     return G_ABL
    return LGRAY

# ── PPT 생성 ─────────────────────────────────────────────────────────────
prs  = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sld  = prs.slides.add_slide(prs.slide_layouts[6])

# 헤더 (0 ~ 0.745)
hbar = add_box(sld, 0, 0, 13.33, 0.72, fill=NAVY)
tf = hbar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  RQ3 전체 모델 성능 비교 — A / B / C / N-A / 채널단독 (CH1~CH5) (단색 채움 · 공통 축)"
r.font.bold=True; r.font.size=Pt(19); r.font.color.rgb=WHITE; r.font.name="Calibri"
add_box(sld, 0, 0.72, 13.33, 0.025, fill=BLUE)
add_text(sld,
    "A=전채널  |  B=Granger 유의 채널(raw p<0.05: 스니커즈 CH1+CH3+CH4, 카드 CH1)  "
    "|  C=SHAP 상위 채널  |  N-A=Granger 유의 채널 제거  |  레고: B=N-A=A  "
    "—  공통 축(0.5~1.0)에서 보면 실제 차이(0.00x)는 시각적으로도 거의 없음",
    0.2, 0.755, 13.0, 0.30, size=8.5, color=DGRAY, italic=True)

# ── AUC 표 (t=1.05 ~ 3.45, h=2.40) ──────────────────────────────────────
TBL_T = 1.05; TBL_H = 2.42

def fmt(auc, verdict, is_ref=False, note=""):
    if auc is None: return "—"
    v = "(기준)" if is_ref else f"({verdict})"
    n = f" {note}" if note else ""
    return f"{auc:.4f}{n}\n{v}"

rows = [
    # header
    ["모델 (구성 채널)", "스니커즈\nAUC (vs A)", "카드\nAUC (vs A)", "레고\nAUC (vs A)"],
    # A
    ["A  전채널 CH1~5 + 가격래그",
     fmt(A_auc["sneakers"],  "기준", True),
     fmt(A_auc["cards"],     "기준", True),
     fmt(A_auc["lego"],      "기준", True)],
    # B
    ["B  Granger 유의 채널 (raw p<0.05)\n  스니커즈: CH1+CH3+CH4 / 카드: CH1 / 레고: —",
     fmt(B_auc["sneakers"],  "차이없음", note="†근사"),
     fmt(B_auc["cards"],     "차이없음"),
     "0.9601\n(= A, 유의채널없음)"],
    # C
    ["C  SHAP 상위 2채널\n  스니커즈: CH1+CH5 / 카드: CH4+CH1 / 레고: CH5+CH4",
     fmt(C_auc["sneakers"],  C_vsa["sneakers"]),
     fmt(C_auc["cards"],     C_vsa["cards"]),
     fmt(C_auc["lego"],      C_vsa["lego"])],
    # N-A
    ["N-A  Granger 유의 채널 제거\n  스니커즈: CH2+CH5 / 카드: CH2~5 / 레고: = A",
     fmt(NA_auc["sneakers"], NA_vsa["sneakers"]),
     fmt(NA_auc["cards"],    NA_vsa["cards"]),
     "0.9601\n(= A, 유의채널없음)"],
    # CH1-5
    ["CH1 단독  Google Trends + 가격래그",
     fmt(get("CH1-only","sneakers"), vsa("CH1-only","sneakers")),
     fmt(get("CH1-only","cards"),    vsa("CH1-only","cards")),
     fmt(get("CH1-only","lego"),     vsa("CH1-only","lego"))],
    ["CH2 단독  뉴스 보도량 + 가격래그",
     fmt(get("CH2-only","sneakers"), vsa("CH2-only","sneakers")),
     fmt(get("CH2-only","cards"),    vsa("CH2-only","cards")),
     fmt(get("CH2-only","lego"),     vsa("CH2-only","lego"))],
    ["CH3 단독  뉴스 감성 + 가격래그",
     fmt(get("CH3-only","sneakers"), vsa("CH3-only","sneakers")),
     fmt(get("CH3-only","cards"),    vsa("CH3-only","cards")),
     fmt(get("CH3-only","lego"),     vsa("CH3-only","lego"))],
    ["CH4 단독  YT 조회수 + 가격래그",
     fmt(get("CH4-only","sneakers"), vsa("CH4-only","sneakers")),
     fmt(get("CH4-only","cards"),    vsa("CH4-only","cards")),
     fmt(get("CH4-only","lego"),     vsa("CH4-only","lego"))],
    ["CH5 단독  YT 댓글 감성 + 가격래그",
     fmt(get("CH5-only","sneakers"), vsa("CH5-only","sneakers")),
     fmt(get("CH5-only","cards"),    vsa("CH5-only","cards")),
     fmt(get("CH5-only","lego"),     vsa("CH5-only","lego"))],
]

ROW_MODEL = [None,"A","B","C","N-A","CH","CH","CH","CH","CH"]
CELL_BG_COL = {"A": BG["A"], "B": BG["B"], "C": BG["C"], "N-A": BG["N-A"], "CH": BG["CH"]}
CELL_BG_VERDICT = lambda cell_txt: (
    G_ND  if "차이없음" in cell_txt else
    G_AB  if "A 우수"   in cell_txt else
    G_ABL if "ablation" in cell_txt or "우수" in cell_txt else
    LGRAY
)

nr, nc = len(rows), 4
col_w  = [4.30, 2.92, 2.92, 2.92]
tbl_l  = (13.33 - sum(col_w)) / 2   # ≈ 0.135

tbl = sld.shapes.add_table(
    nr, nc, Inches(tbl_l), Inches(TBL_T),
    Inches(sum(col_w)), Inches(TBL_H)).table

for i, w in enumerate(col_w): tbl.columns[i].width = Inches(w)
rh = Inches(TBL_H / nr)
for i in range(nr): tbl.rows[i].height = rh

for ri, row in enumerate(rows):
    is_hdr = (ri == 0)
    mk = ROW_MODEL[ri] if ri < len(ROW_MODEL) else "CH"
    for ci, txt in enumerate(row):
        cell = tbl.cell(ri, ci)
        tf   = cell.text_frame; tf.word_wrap = True
        p    = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if (is_hdr or ci > 0) else PP_ALIGN.LEFT
        r = p.add_run(); r.text = str(txt)
        r.font.name = "Calibri"
        r.font.size = Pt(8.5 if is_hdr else 8)
        r.font.bold = is_hdr or (ci == 0 and not is_hdr)
        r.font.color.rgb = WHITE if is_hdr else BLACK
        if is_hdr:
            _set_bg(cell, NAVY)
        elif ci == 0:
            _set_bg(cell, CELL_BG_COL.get(mk, LGRAY))
        else:
            _set_bg(cell, CELL_BG_VERDICT(txt))

# ── 구분선 ────────────────────────────────────────────────────────────────
SEP_T = TBL_T + TBL_H + 0.06
add_box(sld, 0.13, SEP_T, 13.07, 0.022, fill=BLUE)
add_text(sld, "† 스니커즈 B: CH1+CH3+CH4 조합 AUC 미산출 → 개별 AUC 평균 근사치 (0.8461)   "
              "| 셀 색상: 연초록=차이없음 / 연빨강=A 우수 / 연파랑=ablation 우수",
         0.13, SEP_T + 0.03, 13.07, 0.25,
         size=7.5, color=DGRAY, italic=True)

# ── 그래프 삽입 ───────────────────────────────────────────────────────────
GRF_T = SEP_T + 0.30
GRF_H = 7.36 - GRF_T
fig_buf.seek(0)
sld.shapes.add_picture(fig_buf,
    Inches(0.13), Inches(GRF_T),
    Inches(13.07), Inches(GRF_H))

# ── 저장 ─────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"[OK] {OUT}")
print(f"  Table : {TBL_T:.2f} ~ {TBL_T+TBL_H:.2f}")
print(f"  Graph : {GRF_T:.2f} ~ {GRF_T+GRF_H:.2f}  (no overlap)")
